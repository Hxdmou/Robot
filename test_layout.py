# -*- coding: utf-8 -*-
"""精确测试10pt字号在不同设置下的实际行高"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# 1英寸 = 914400 EMU
# 1pt = 12700 EMU
print("=== 单位换算 ===")
print(f"1pt = {Pt(1)} EMU = {Pt(1)/914400:.4f} 英寸")
print(f"10pt = {Pt(10)} EMU = {Pt(10)/914400:.4f} 英寸")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 测试：20条10pt内容，行距1.0，需要多少高度？
test_items = [f'【测试标签{i:02d}】这是第{i+1}条测试内容，用来精确测量实际占用高度' for i in range(20)]

# 用不同行距测试
for ls_idx, line_spacing in enumerate([0.95, 1.0, 1.05, 1.1]):
    y_start = Inches(0.5 + ls_idx * 1.8)
    box = slide.shapes.add_textbox(Inches(0.5), y_start, Inches(6), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    
    for i, item in enumerate(test_items[:10]):  # 先测10条
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = '▸ ' + item
        run.font.size = Pt(10)
        run.font.name = '微软雅黑'

prs.save('F:\\个人作品\\具身智能\\test_layout.pptx')
print("\n测试PPT已保存到 test_layout.pptx")
print("请手动打开查看，确认10条/20条内容需要的高度")
print()
print("=== 理论计算 ===")
print("假设每行纯文字高度约为 10pt * 1.2 (ascender/descender) = 12pt = 0.1667英寸")
print("10条：10 * 0.1667 = 1.67英寸")
print("20条：20 * 0.1667 = 3.33英寸")
print()
print("内容页可用高度（从标题栏下到页脚）约：7.5 - 0.82 - 0.24 = 6.44英寸")
print("  - 上部两栏区域：约 4.2英寸，足够放10条+标题")
print("  - 下部过程区域：约 6.44 - 4.2 - 0.05 = 2.19英寸，足够放3条详细过程")
print()
print("细节页可用高度约：6.44英寸")
print("  - 20条需要约3.33英寸，完全足够！")
