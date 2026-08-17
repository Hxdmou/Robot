# -*- coding: utf-8 -*-
"""
全页面布局测试脚本 v2
核心修正：细节页去掉圆角卡片，直接在深蓝背景上放文字，最大化宽高！
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import re, os

# ========== 颜色定义 ==========
DARK_BLUE = RGBColor(0x0A, 0x16, 0x2F)
MID_BLUE = RGBColor(0x10, 0x25, 0x48)
ACCENT_BLUE = RGBColor(0x1E, 0x5F, 0xA8)
GOLD = RGBColor(0xD4, 0xA5, 0x37)
LGRAY = RGBColor(0xC8, 0xD0, 0xDC)
MGRAY = RGBColor(0x80, 0x88, 0x98)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = 13.333
SLIDE_H = 7.5

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fc=None, lc=None, lw=1):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.line.fill.background()
    if fc:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fc
    else:
        shp.fill.background()
    return shp

def rrect(slide, x, y, w, h, fc=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.05
    shp.line.fill.background()
    if fc:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fc
    else:
        shp.fill.background()
    return shp

def tb(slide, x, y, w, h, text, sz=10, b=False, c=LGRAY, al=PP_ALIGN.LEFT, an=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0); tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    tf.vertical_anchor = an
    p = tf.paragraphs[0]
    p.alignment = al
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = b
    run.font.color.rgb = c
    run.font.name = '微软雅黑'
    return txBox

def add_bullets(tf, items, start_idx=0, sz=10, color=LGRAY, line_spacing=1.1, space_after=0):
    for i, item in enumerate(items):
        idx = start_idx + i
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        run_prefix = p.add_run()
        run_prefix.text = '▸ '
        run_prefix.font.size = Pt(sz)
        run_prefix.font.color.rgb = ACCENT_BLUE
        run_prefix.font.name = '微软雅黑'
        parts = re.split(r'(【[^】]+】)', item)
        for part in parts:
            run = p.add_run()
            run.text = part
            run.font.size = Pt(sz)
            run.font.name = '微软雅黑'
            if part.startswith('【') and part.endswith('】'):
                run.font.color.rgb = GOLD
                run.font.bold = True
            else:
                run.font.color.rgb = color

def add_page_header_compact(slide, part_num, title):
    """超紧凑页眉：高度0.4英寸，最大化内容区空间"""
    HEADER_Y = 0.4
    rect(slide, 0, 0, SLIDE_W, 0.04, fc=GOLD)
    rrect(slide, 0.12, 0.1, 0.75, 0.25, fc=GOLD)
    tb(slide, 0.12, 0.1, 0.75, 0.25, part_num, sz=8, b=True, c=DARK_BLUE, al=PP_ALIGN.CENTER, an=MSO_ANCHOR.MIDDLE)
    tb(slide, 0.95, 0.08, 10, 0.3, title, sz=14, b=True, c=WHITE, an=MSO_ANCHOR.MIDDLE)
    rect(slide, 0.1, HEADER_Y, SLIDE_W - 0.2, 0.015, fc=RGBColor(0x20, 0x35, 0x60))
    return HEADER_Y + 0.03  # 返回内容起始Y=0.43

def add_page_tag(slide, tag_text, tag_color):
    rrect(slide, SLIDE_W - 1.5, 0.14, 1.2, 0.26, fc=tag_color)
    tb(slide, SLIDE_W - 1.5, 0.15, 1.2, 0.24, tag_text, sz=8, b=True, c=WHITE, al=PP_ALIGN.CENTER)

# ========== 测试内容（最长真实内容） ==========
content_left = [
'【市场规模】2026年全球人形机器人市场规模突破1200亿元人民币，中国市场占比超40%达480亿元，同比增长128%；全球出货量预计达18万台，其中中国市场出货8万台，占比44%',
'【产业定位】人形机器人是继智能手机、新能源汽车之后下一代通用智能终端，被国家列为未来产业重点培育方向，纳入十四五、十五五规划战略性新兴产业',
'【量产里程碑】2026年是人形机器人量产元年：特斯拉Optimus Gen3年产能目标10万台，优必选Walker X2年产能5万台，波士顿动力Atlas年产能1万台',
'【成本下降】人形机器人整机成本从2022年200万元/台快速下降到2026年20-30万元/台，预计2028年降至10-15万元，2030年进一步降至5-8万元',
'【核心技术突破】关节减速器国产化率从2023年15%提升到2026年65%；伺服电机国产化率达70%；六维力传感器国产化率突破50%',
'【政策支持】工信部发布《人形机器人创新发展指导意见》提出2025年实现技术突破、2027年形成完整产业生态',
'【产业链完善】已形成从核心零部件→本体制造→系统集成→场景应用完整产业链，四大产业集群快速成型',
'【资本热度】2026年H1人形机器人领域融资额超380亿元，单笔融资额平均超5亿元，多家企业估值超百亿元',
'【应用场景拓展】已从工业制造→商业服务→家庭服务→特种作业多场景落地',
'【安徽布局】安徽将人形机器人列为十大新兴产业重点方向，芜湖埃夫特、合肥欣奕华、蚌埠传感器形成配套'
]
content_right = [
'【特斯拉Optimus Gen3】2026量产版采用纯电驱动+谐波减速器+28个自由度，身高1.73米体重57kg，负载20kg，行走速度7km/h',
'【优必选Walker X2】中国首个万台级量产人形机器人，身高1.55米体重45kg，负载15kg，行走速度5km/h，比亚迪工厂订单超3万台',
'【波士顿动力Atlas电动版】放弃液压改为纯电驱动，身高1.75米体重75kg，可完成跑酷、后空翻等高难度动作',
'【小米CyberOne二代】身高1.70米体重52kg，搭载小米MiLM具身大模型，语音交互能力突出，已进入小米工厂测试',
'【小鹏Iron】搭载小鹏XNGP同源智驾技术，擅长自主导航和避障，物流场景表现突出，首批交付小鹏汽车工厂',
'【傅利叶GR-2】康复人形机器人领先者，身高1.65米体重50kg，已进入200+医院辅助患者康复训练',
'【宇树H1】四足机器人龙头宇树科技首款人形机器人，身高1.80米体重47kg，奔跑速度15km/h全球最快',
'【达闼机器人】云端大脑架构，5G远程操控，商用服务场景部署量国内领先',
'【智元远征A2】前华为天才少年创业项目，具身智能算法领先，通用操作能力强',
'【安徽本地企业】芜湖埃夫特工业机器人本体产能全国前三，合肥欣奕华人形机器人样机下线，蚌埠中国传感谷提供核心传感器'
]
content_process = [
'【研发阶段（2022-2025）】技术验证期：波士顿动力Atlas实现后空翻但液压驱动成本高昂；特斯拉Optimus Gen1亮相仅能简单行走；国内企业完成原型机验证，单台成本超百万元',
'【量产元年（2026年）】规模化落地期：特斯拉/优必选等头部企业实现万台级量产交付，核心零部件国产化率超80%，成本快速下降至15-25万元区间，工业场景率先规模化部署',
'【规模爆发（2027-2028）】成本下探至10-15万元，应用场景从工业向商业服务/家庭陪护渗透，年销量突破20万台，产业生态逐步成熟',
'【普及阶段（2029-2030）】成本进一步降至5-8万元接近家用轿车价格，具身智能大模型成熟，年销量突破100万台，中国成为全球最大生产应用市场',
'【产业价值】带动核心零部件/AI算法/系统集成/场景应用全产业链发展，预计2030年带动相关产业规模超万亿元，创造百万级就业岗位'
]
detail_items_long = [
'【华为Mate 80 Pro】麒麟9020（7nm），8核CPU+12核GPU+24核NPU（120TOPS）；6.8英寸2K OLED 1-120Hz LTPO，3000nit亮度；后置5000万超光变主摄+4000万超广角+800万潜望长焦；5000mAh硅碳负极电池，100W有线+80W无线快充；鸿蒙OS 5.0支持天通卫星通话+北斗卫星消息；IP68防水；存储：8+256G 6499/12+256G 6999/12+512G 7999/16+1T 8999元',
'【华为Pura 80 Ultra】麒麟9020；6.9英寸2K LTPO OLED；一英寸超感光主摄（f/1.2-f/4.0可变光圈）+4000万超广角+2亿像素潜望长焦（200倍数码变焦）+TOF；5500mAh电池100W有线+80W无线；XMAGE影像系统+AI大模型计算摄影；IP68；纳米微晶陶瓷机身；存储：12+256G 8999/16+512G 9999/16+1T 11999元；2026年7月发布',
'【华为MateBook X Pro 2026】英特尔酷睿Ultra 9 288V + 昇腾310B独立NPU（128TOPS AI算力）；32GB LPDDR5X内存+2TB PCIe 4.0 SSD；14.2英寸3.1K柔性OLED原色屏，120Hz刷新率，100% DCI-P3色域；70Wh电池140W快充；1.15kg重量；鸿蒙OS for PC支持多屏协同/AI摘要/AI修图/AI会议纪要；存储：16+512G 9999/32+1T 12999/32+2T 15999元',
'【苹果iPhone 18 Pro Max】A19 Pro芯片（3nm第二代），8核CPU+8核GPU+16核NPU（150TOPS）；6.9英寸Super Retina XDR OLED，1-120Hz ProMotion，3500nit峰值亮度；后置4800万主摄（传感器位移防抖）+4800万超广角+1200万5倍潜望长焦；4800mAh电池45W有线+25W MagSafe无线；iOS 20全功能支持Apple Intelligence；钛金属中框IP69防水；存储：256G 9999/512G 11499/1T 13499元；2026年9月发布',
'【苹果MacBook Pro 16 M5】Apple M5 Max芯片：16核CPU（12性能+4能效）+40核GPU+32核NPU（180TOPS）；最高192GB统一内存+最高8TB SSD；16.2英寸Liquid Retina XDR显示屏，3456×2234分辨率，120Hz ProMotion，1600nit HDR峰值亮度；100Wh电池21小时视频播放；140W MagSafe快充；macOS Sequoia深度整合Apple Intelligence；重量2.1kg；存储：36+512G 19999/64+1T 24999/96+2T 29999/192+8T 45999元',
'【苹果Vision Pro 2】M5芯片+R2芯片；Micro-OLED双眼4K分辨率（单眼2300万像素），120Hz刷新率，120度视场角；重量450g（比一代减轻30%）；眼动追踪+手势追踪+空间音频；visionOS 3支持空间视频拍摄/沉浸式办公/空间游戏；外接电池续航4小时；存储：256G 14999元；2026年6月WWDC发布',
'【小米16 Ultra】高通骁龙8 Gen4（3nm）+澎湃P2快充芯片+澎湃G2电池管理芯片；6.9英寸2K LTPO AMOLED华星C9屏，1-144Hz自适应刷新率，4000nit峰值亮度，1920Hz PWM调光；后置徕卡四摄：一英寸LYT-900主摄（f/1.4-f/4.0可变光圈）+4000万超广角+2亿像素5倍潜望长焦+5000万像素10倍超长焦；6000mAh硅碳负极电池120W有线+50W无线；澎湃OS 3.0支持双向卫星通信；IP68防水陶瓷机身；存储：12+256G 5999/16+512G 6499/16+1T 7499/24+2T 8499元',
'【小米16 Pro】骁龙8 Gen4；6.7英寸2K 144Hz LTPO AMOLED；后置5000万LYT-800主摄+5000万超广角+5000万3倍潜望长焦；5500mAh电池120W有线+50W无线快充；澎湃OS 3.0；IP68防水金属中框玻璃后盖；存储：12+256G 4999/12+512G 5499/16+1T 6299元',
'【小米RedmiBook Pro 16 2026】英特尔酷睿Ultra 7 268V + 小米自研NPU（80TOPS AI算力）；32GB LPDDR5X内存+1TB PCIe 4.0 SSD（可扩展）；16英寸3.2K 120Hz IPS屏，100% sRGB色域，500nit亮度；80Wh电池100W快充；1.8kg重量；澎湃OS for Windows支持AI字幕/AI会议/AI画图/AI写作；存储：16+512G 4999/32+1T 5999元；2026年3月发布',
'【小米SU7 Ultra】三电机四驱系统：前220kW+后350kW+后350kW，系统总功率960kW（1306马力），系统总扭矩1680N·m；0-100km/h加速1.98秒，0-200km/h加速5.8秒，最高车速350km/h；宁德时代麒麟II电池130kWh，CLTC续航800km，800V高压平台10分钟补能400km；小米全栈自研智驾系统双Orin-X芯片（508TOPS）+激光雷达+11摄像头+12超声波+5毫米波雷达；21英寸轮毂碳陶瓷刹车；空气悬架+CDC电磁减振；售价52.99万元；2026年3月发布已交付',
'【AI功能共性】三大品牌旗舰全部支持：自然语言语音助手（连续对话/复杂指令/多轮交互）、AI拍照修图（AI消除/AI扩图/AI增强/AI夜景）、AI会议（实时录音转写/智能摘要/待办提取/多语种翻译）、AI写作（文案/邮件/报告/代码生成）、AI搜索（自然语言搜索/信息整合总结）、AI安全（隐私计算/端侧处理不上云）',
'【端侧大模型】华为盘古大模型端侧版（13B参数）、苹果Apple Intelligence（云端30B参数+端侧3B混合推理）、小米MiLM大模型端侧版（7B/13B），推理延迟<1秒，隐私敏感数据全程端侧处理不上传',
'【生态互联】华为鸿蒙分布式软总线支持手机/PC/平板/手表/车机/智能家居无缝流转接续；苹果Continuity支持iPhone/Mac/iPad/Watch/Vision Pro全生态设备无缝协同接力；小米澎湃OS HyperConnect支持全生态AIoT设备互联互通智能联动',
'【通信技术】华为独家支持天通一号卫星通话+北斗卫星消息+星闪NearLink；苹果支持卫星SOS紧急求助+卫星消息共享；小米支持双向卫星通信（天通+北斗）+5.5G NTN；三家旗舰均标配Wi-Fi 7/蓝牙5.4/NFC/红外遥控',
'【材料工艺】华为采用昆仑玻璃2代+玄武架构+纳米微晶陶瓷；苹果采用Grade 5钛金属中框+超瓷晶面板；小米采用龙晶玻璃+航空铝中框+陶瓷后盖；抗跌落抗刮擦耐用性较上一代提升2-3倍',
'【市场份额】2026年H1中国智能手机市场：华为28%份额重回第一，苹果18%第二，小米15%第三，三家合计占据61%市场份额；AI PC市场联想/华为/苹果/小米位列前四',
'【技术趋势1】端侧NPU算力快速提升：2026年旗舰手机NPU算力达100-150TOPS，AI PC NPU算力达80-180TOPS，可本地运行7B-13B参数大模型',
'【技术趋势2】多模态AI成为消费电子标配：支持文字/图像/视频/语音/3D空间等多模态信息理解、生成、编辑能力',
'【技术趋势3】AI代理（Agent）深度整合系统：手机/PC上的AI助理能够自主理解指令并完成订票/订餐/安排日程/处理邮件等复杂任务',
'【未来方向】2027-2030年AI终端形态向AR智能眼镜/脑机接口/智能家居服务机器人延伸，全场景个人AI助理成为现实'
]

# ========== 开始测试 ==========
prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)

# 参数：超紧凑页眉，内容区从0.43开始，页脚在7.42，可用高度=6.97英寸
# 细节页：无卡片，直接放文字，文本框高度≈6.92英寸足够放45行6.88英寸
# 内容页：上部2.9英寸（左右栏各10条18行≈2.75英寸）+ 间隙0.06 + 下部4.01英寸（写更长内容填满）
MARGIN = 0.12
CONTENT_START = 0.43  # 超紧凑页眉下
FOOTER_Y = 7.42
CONTENT_END = FOOTER_Y - 0.02
DETAIL_MARGIN_X = 0.15  # 细节页左右边距更小，宽度更大

print(f'测试参数：')
print(f'  CONTENT_START={CONTENT_START}, CONTENT_END={CONTENT_END}')
print(f'  内容区可用高度={CONTENT_END - CONTENT_START:.2f}英寸')

# ========== 测试页1：内容页（单数页） ==========
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide1, DARK_BLUE)
add_page_header_compact(slide1, 'PART 12', '消费电子：华为/苹果/小米三强争霸')
add_page_tag(slide1, '内容描述', ACCENT_BLUE)

col_gap = 0.15
col_w = (SLIDE_W - 2*MARGIN - col_gap) / 2
upper_h = 2.9  # 左右栏10条约2.75英寸，2.9足够
gap = 0.06
lower_h = CONTENT_END - CONTENT_START - upper_h - gap  # ≈4.01英寸
upper_y = CONTENT_START
left_x = MARGIN
right_x = MARGIN + col_w + col_gap
bot_y = upper_y + upper_h + gap

print(f'  内容页：上部{upper_h}英寸，下部{lower_h:.2f}英寸')

# 左栏
rrect(slide1, left_x, upper_y, col_w, upper_h, fc=MID_BLUE)
rect(slide1, left_x, upper_y, 0.05, upper_h, fc=ACCENT_BLUE)
box_l = slide1.shapes.add_textbox(Inches(left_x + 0.08), Inches(upper_y + 0.05), Inches(col_w - 0.14), Inches(upper_h - 0.1))
tf_l = box_l.text_frame; tf_l.word_wrap = True
tf_l.margin_left = Pt(3); tf_l.margin_right = Pt(3); tf_l.margin_top = Pt(0); tf_l.margin_bottom = Pt(0)
p_tl = tf_l.paragraphs[0]; p_tl.line_spacing = 1.1; p_tl.space_after = Pt(0)
r_tl = p_tl.add_run(); r_tl.text = '▎核心内容'
r_tl.font.size = Pt(10); r_tl.font.bold = True; r_tl.font.color.rgb = GOLD; r_tl.font.name = '微软雅黑'
add_bullets(tf_l, content_left, start_idx=1, sz=10, line_spacing=1.1, space_after=0)

# 右栏
rrect(slide1, right_x, upper_y, col_w, upper_h, fc=MID_BLUE)
rect(slide1, right_x, upper_y, 0.05, upper_h, fc=ACCENT_BLUE)
box_r = slide1.shapes.add_textbox(Inches(right_x + 0.08), Inches(upper_y + 0.05), Inches(col_w - 0.14), Inches(upper_h - 0.1))
tf_r = box_r.text_frame; tf_r.word_wrap = True
tf_r.margin_left = Pt(3); tf_r.margin_right = Pt(3); tf_r.margin_top = Pt(0); tf_r.margin_bottom = Pt(0)
p_tr = tf_r.paragraphs[0]; p_tr.line_spacing = 1.1; p_tr.space_after = Pt(0)
r_tr = p_tr.add_run(); r_tr.text = '▎代表动态'
r_tr.font.size = Pt(10); r_tr.font.bold = True; r_tr.font.color.rgb = GOLD; r_tr.font.name = '微软雅黑'
add_bullets(tf_r, content_right, start_idx=1, sz=10, line_spacing=1.1, space_after=0)

# 下部通栏
rrect(slide1, MARGIN, bot_y, SLIDE_W - 2*MARGIN, lower_h, fc=MID_BLUE)
rect(slide1, MARGIN, bot_y, SLIDE_W - 2*MARGIN, 0.03, fc=GOLD)
box_b = slide1.shapes.add_textbox(Inches(MARGIN + 0.08), Inches(bot_y + 0.05), Inches(SLIDE_W - 2*MARGIN - 0.16), Inches(lower_h - 0.1))
tf_b = box_b.text_frame; tf_b.word_wrap = True
tf_b.margin_left = Pt(3); tf_b.margin_right = Pt(3); tf_b.margin_top = Pt(0); tf_b.margin_bottom = Pt(0)
p_tb = tf_b.paragraphs[0]; p_tb.line_spacing = 1.1; p_tb.space_after = Pt(1)
r_tb = p_tb.add_run(); r_tb.text = '▎消费电子产业发展过程阐述'
r_tb.font.size = Pt(10); r_tb.font.bold = True; r_tb.font.color.rgb = GOLD; r_tb.font.name = '微软雅黑'
add_bullets(tf_b, content_process, start_idx=1, sz=10, line_spacing=1.1, space_after=3)

# 页脚
tb(slide1, 0, FOOTER_Y, SLIDE_W, 0.12, '具身智能&AI产业最新进展 · 2026年8月15日', sz=7, c=MGRAY, al=PP_ALIGN.CENTER)

# ========== 测试页2：细节页（双数页）—— 去掉卡片！直接在背景放文字！ ==========
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide2, DARK_BLUE)
add_page_header_compact(slide2, 'PART 12', '消费电子：华为/苹果/小米三强争霸')
add_page_tag(slide2, '细节描述', GOLD)

detail_h = CONTENT_END - CONTENT_START
print(f'  细节页（无卡片）：内容高度={detail_h:.2f}英寸，文本框宽度={SLIDE_W - 2*DETAIL_MARGIN_X:.2f}英寸')

# 左侧金色竖条装饰
rect(slide2, DETAIL_MARGIN_X, CONTENT_START, 0.04, detail_h, fc=GOLD)
# 文本框直接放，宽度更大，折行更少
box = slide2.shapes.add_textbox(Inches(DETAIL_MARGIN_X + 0.1), Inches(CONTENT_START + 0.02), Inches(SLIDE_W - 2*DETAIL_MARGIN_X - 0.15), Inches(detail_h - 0.05))
tf = box.text_frame; tf.word_wrap = True
tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
p_title = tf.paragraphs[0]; p_title.line_spacing = 1.1; p_title.space_after = Pt(2)
run_t = p_title.add_run(); run_t.text = '▎旗舰参数 · 价格版本 · AI功能 · 技术亮点'
run_t.font.size = Pt(11); run_t.font.bold = True; run_t.font.color.rgb = GOLD; run_t.font.name = '微软雅黑'
add_bullets(tf, detail_items_long, start_idx=1, sz=10, line_spacing=1.1, space_after=0)

# 页脚
tb(slide2, 0, FOOTER_Y, SLIDE_W, 0.12, '具身智能&AI产业最新进展 · 2026年8月15日 · 最新行情/研发/成果', sz=7, c=MGRAY, al=PP_ALIGN.CENTER)

# 保存并自动打开
out_path = r'F:\个人作品\具身智能\test_layout_v3.pptx'
prs.save(out_path)
print(f'\n✅ 测试PPT已生成并自动打开：{out_path}')
print('核心改进：细节页去掉圆角卡片，直接在深蓝背景放文字，宽度更大折行更少！')
os.startfile(out_path)
