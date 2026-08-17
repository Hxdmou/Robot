# -*- coding: utf-8 -*-
"""
布局调试测试脚本 - 用消费电子最长内容测试，加边界线标记
红色：卡片边界
蓝色：文本框边界
黄色：内容边界
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import re

# 颜色定义
DARK_BLUE = RGBColor(0x0A, 0x16, 0x2F)
MID_BLUE = RGBColor(0x10, 0x25, 0x48)
ACCENT_BLUE = RGBColor(0x1E, 0x5F, 0xA8)
GOLD = RGBColor(0xD4, 0xA5, 0x37)
LGRAY = RGBColor(0xC8, 0xD0, 0xDC)
RED = RGBColor(0xFF, 0x00, 0x00)
BLUE = RGBColor(0x00, 0xFF, 0xFF)
YELLOW = RGBColor(0xFF, 0xFF, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = 13.333
SLIDE_H = 7.5

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def debug_rect(slide, x, y, w, h, color, label=''):
    """调试用：画带边框的矩形，标记边界"""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.background()
    shp.line.color.rgb = color
    shp.line.width = Pt(1)
    if label:
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.size = Pt(6)
        run.font.color.rgb = color
        run.font.name = '微软雅黑'

def add_bullets(tf, items, start_idx=0, sz=10, color=LGRAY, line_spacing=1.1, space_after=0):
    for i, item in enumerate(items):
        idx = start_idx + i
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        run_prefix = p.add_run()
        run_prefix.text = '▸ '
        run_prefix.font.size = Pt(sz)
        run_prefix.font.color.rgb = ACCENT_BLUE
        run_prefix.font.name = '微软雅黑'
        parts = re.split(r'(【[^】]+】)', item)
        for part in parts:
            run = p.add_run()
            run.text = part
            run.font.size = Pt(sz)
            run.font.name = '微软雅黑'
            if part.startswith('【') and part.endswith('】'):
                run.font.color.rgb = GOLD
                run.font.bold = True
            else:
                run.font.color.rgb = color

# 消费电子最长的20条细节内容（真实数据）
detail_items_long = [
'【华为Mate 80 Pro】麒麟9020（7nm），8核CPU+12核GPU+24核NPU（120TOPS）；6.8英寸2K OLED 1-120Hz LTPO，3000nit亮度；后置5000万超光变主摄+4000万超广角+800万潜望长焦；5000mAh硅碳负极电池，100W有线+80W无线快充；鸿蒙OS 5.0支持天通卫星通话+北斗卫星消息；IP68防水；存储：8+256G 6499/12+256G 6999/12+512G 7999/16+1T 8999元',
'【华为Pura 80 Ultra】麒麟9020；6.9英寸2K LTPO OLED；一英寸超感光主摄（f/1.2-f/4.0可变光圈）+4000万超广角+2亿像素潜望长焦（200倍数码变焦）+TOF；5500mAh电池100W有线+80W无线；XMAGE影像系统+AI大模型计算摄影；IP68；纳米微晶陶瓷机身；存储：12+256G 8999/16+512G 9999/16+1T 11999元；2026年7月发布',
'【华为MateBook X Pro 2026】英特尔酷睿Ultra 9 288V + 昇腾310B独立NPU（128TOPS AI算力）；32GB LPDDR5X内存+2TB PCIe 4.0 SSD；14.2英寸3.1K柔性OLED原色屏，120Hz刷新率，100% DCI-P3色域；70Wh电池140W快充；1.15kg重量；鸿蒙OS for PC支持多屏协同/AI摘要/AI修图/AI会议纪要；存储：16+512G 9999/32+1T 12999/32+2T 15999元',
'【苹果iPhone 18 Pro Max】A19 Pro芯片（3nm第二代），8核CPU+8核GPU+16核NPU（150TOPS）；6.9英寸Super Retina XDR OLED，1-120Hz ProMotion，3500nit峰值亮度；后置4800万主摄（传感器位移防抖）+4800万超广角+1200万5倍潜望长焦；4800mAh电池45W有线+25W MagSafe无线；iOS 20全功能支持Apple Intelligence；钛金属中框IP69防水；存储：256G 9999/512G 11499/1T 13499元；2026年9月发布',
'【苹果MacBook Pro 16 M5】Apple M5 Max芯片：16核CPU（12性能+4能效）+40核GPU+32核NPU（180TOPS）；最高192GB统一内存+最高8TB SSD；16.2英寸Liquid Retina XDR显示屏，3456×2234分辨率，120Hz ProMotion，1600nit HDR峰值亮度；100Wh电池21小时视频播放；140W MagSafe快充；macOS Sequoia深度整合Apple Intelligence；重量2.1kg；存储：36+512G 19999/64+1T 24999/96+2T 29999/192+8T 45999元',
'【苹果Vision Pro 2】M5芯片+R2芯片；Micro-OLED双眼4K分辨率（单眼2300万像素），120Hz刷新率，120度视场角；重量450g（比一代减轻30%）；眼动追踪+手势追踪+空间音频；visionOS 3支持空间视频拍摄/沉浸式办公/空间游戏；外接电池续航4小时；存储：256G 14999元；2026年6月WWDC发布',
'【小米16 Ultra】高通骁龙8 Gen4（3nm）+澎湃P2快充芯片+澎湃G2电池管理芯片；6.9英寸2K LTPO AMOLED华星C9屏，1-144Hz自适应刷新率，4000nit峰值亮度，1920Hz PWM调光；后置徕卡四摄：一英寸LYT-900主摄（f/1.4-f/4.0可变光圈）+4000万超广角+2亿像素5倍潜望长焦+5000万像素10倍超长焦；6000mAh硅碳负极电池120W有线+50W无线；澎湃OS 3.0支持双向卫星通信；IP68防水陶瓷机身；存储：12+256G 5999/16+512G 6499/16+1T 7499/24+2T 8499元',
'【小米16 Pro】骁龙8 Gen4；6.7英寸2K 144Hz LTPO AMOLED；后置5000万LYT-800主摄+5000万超广角+5000万3倍潜望长焦；5500mAh电池120W有线+50W无线快充；澎湃OS 3.0；IP68防水金属中框玻璃后盖；存储：12+256G 4999/12+512G 5499/16+1T 6299元',
'【小米RedmiBook Pro 16 2026】英特尔酷睿Ultra 7 268V + 小米自研NPU（80TOPS AI算力）；32GB LPDDR5X内存+1TB PCIe 4.0 SSD（可扩展）；16英寸3.2K 120Hz IPS屏，100% sRGB色域，500nit亮度；80Wh电池100W快充；1.8kg重量；澎湃OS for Windows支持AI字幕/AI会议/AI画图/AI写作；存储：16+512G 4999/32+1T 5999元；2026年3月发布',
'【小米SU7 Ultra】三电机四驱系统：前220kW+后350kW+后350kW，系统总功率960kW（1306马力），系统总扭矩1680N·m；0-100km/h加速1.98秒，0-200km/h加速5.8秒，最高车速350km/h；宁德时代麒麟II电池130kWh，CLTC续航800km，800V高压平台10分钟补能400km；小米全栈自研智驾系统双Orin-X芯片（508TOPS）+激光雷达+11摄像头+12超声波+5毫米波雷达；21英寸轮毂碳陶瓷刹车；空气悬架+CDC电磁减振；售价52.99万元；2026年3月发布已交付',
'【AI功能共性】三大品牌旗舰全部支持：自然语言语音助手（连续对话/复杂指令/多轮交互）、AI拍照修图（AI消除/AI扩图/AI增强/AI夜景）、AI会议（实时录音转写/智能摘要/待办提取/多语种翻译）、AI写作（文案/邮件/报告/代码生成）、AI搜索（自然语言搜索/信息整合总结）、AI安全（隐私计算/端侧处理不上云）',
'【端侧大模型】华为盘古大模型端侧版（13B参数）、苹果Apple Intelligence（云端30B参数+端侧3B混合推理）、小米MiLM大模型端侧版（7B/13B），推理延迟<1秒，隐私敏感数据全程端侧处理不上传',
'【生态互联】华为鸿蒙分布式软总线支持手机/PC/平板/手表/车机/智能家居无缝流转接续；苹果Continuity支持iPhone/Mac/iPad/Watch/Vision Pro全生态设备无缝协同接力；小米澎湃OS HyperConnect支持全生态AIoT设备互联互通智能联动',
'【通信技术】华为独家支持天通一号卫星通话+北斗卫星消息+星闪NearLink；苹果支持卫星SOS紧急求助+卫星消息共享；小米支持双向卫星通信（天通+北斗）+5.5G NTN；三家旗舰均标配Wi-Fi 7/蓝牙5.4/NFC/红外遥控',
'【材料工艺】华为采用昆仑玻璃2代+玄武架构+纳米微晶陶瓷；苹果采用Grade 5钛金属中框+超瓷晶面板；小米采用龙晶玻璃+航空铝中框+陶瓷后盖；抗跌落抗刮擦耐用性较上一代提升2-3倍',
'【市场份额】2026年H1中国智能手机市场：华为28%份额重回第一，苹果18%第二，小米15%第三，三家合计占据61%市场份额；AI PC市场联想/华为/苹果/小米位列前四',
'【技术趋势1】端侧NPU算力快速提升：2026年旗舰手机NPU算力达100-150TOPS，AI PC NPU算力达80-180TOPS，可本地运行7B-13B参数大模型',
'【技术趋势2】多模态AI成为消费电子标配：支持文字/图像/视频/语音/3D空间等多模态信息理解、生成、编辑能力',
'【技术趋势3】AI代理（Agent）深度整合系统：手机/PC上的AI助理能够自主理解指令并完成订票/订餐/安排日程/处理邮件等复杂任务',
'【未来方向】2027-2030年AI终端形态向AR智能眼镜/脑机接口/智能家居服务机器人延伸，全场景个人AI助理成为现实'
]

prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)

# ========== 测试页：细节页（通栏20条） ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide, DARK_BLUE)

# 测试不同布局参数
# 参数1：内容区起始位置 y_start，结束位置 y_end
test_cases = [
    # (y_start, y_end, 说明)
    (0.7, 7.35, 'y0.7-y7.35'),
]

for case_idx, (y_start, y_end, label) in enumerate(test_cases):
    content_h = y_end - y_start
    margin = 0.15
    x_start = margin
    x_end = SLIDE_W - margin
    w = x_end - x_start
    
    # 画卡片边界（红色）
    debug_rect(slide, x_start, y_start, w, content_h, RED, f'卡片: {label} h={content_h:.2f}in')
    
    # 标题区域高度
    title_h = 0.25
    # 文本框：从标题结束位置开始
    text_x = x_start + 0.1
    text_y = y_start + title_h
    text_w = w - 0.2
    text_h = content_h - title_h - 0.05
    
    # 画文本框边界（蓝色）
    debug_rect(slide, text_x, text_y, text_w, text_h, BLUE, f'文本框: h={text_h:.2f}in')
    
    # 添加文本
    box = slide.shapes.add_textbox(Inches(text_x), Inches(text_y), Inches(text_w), Inches(text_h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    # 标题
    p_title = tf.paragraphs[0]
    p_title.line_spacing = 1.1
    p_title.space_after = Pt(0)
    run_t = p_title.add_run()
    run_t.text = '▎旗舰参数 · 价格版本 · AI功能 · 技术亮点'
    run_t.font.size = Pt(11)
    run_t.font.bold = True
    run_t.font.color.rgb = GOLD
    run_t.font.name = '微软雅黑'
    add_bullets(tf, detail_items_long, start_idx=1, sz=10, line_spacing=1.1, space_after=0)

# 保存
prs.save(r'F:\个人作品\具身智能\test_layout_debug.pptx')
print('调试PPT已生成：test_layout_debug.pptx')
print('红色框：卡片边界')
print('蓝色框：文本框边界')
print('请打开查看：20条最长内容在文本框内是溢出还是有空隙')
