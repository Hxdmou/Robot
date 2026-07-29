
# ============================================================================
# 免责声明与AI使用规范
# ============================================================================
# 本文件仅供技术研究与学习交流使用，不得用于任何非法用途。
#
# AI使用规范：
#   1. 使用本文件相关内容时须遵守所在地法律法规及伦理准则
#   2. 不得用于侵犯他人合法权益、危害网络安全、破坏公共秩序的活动
#   3. 涉及自动化决策的场景须确保人工复核机制与可解释性
#   4. 处理个人信息时须符合数据保护相关法规要求
#
# 风险提示：
#   本文件内容按"现状"提供，不保证绝对准确无误。
#   使用者须自行评估风险，因使用本文件导致的任何损失由使用者承担。
# ============================================================================

from pptx import Presentation

ppt_path = EXTERNAL_PROJECT_DIR  # was: A2A_PROTOCOL_AI_AGENT_V14_完美清晰版.pptx
prs = Presentation(ppt_path)

print("=" * 60)
print("PPT表格结构诊断")
print("=" * 60)

for slide_idx, slide in enumerate(prs.slides):
    print(f"\n📊 第 {slide_idx + 1} 张幻灯片:")
    for shape_idx, shape in enumerate(slide.shapes):
        if shape.has_table:
            table = shape.table
            rows = len(table.rows)
            cols = len(table.columns)
            print(f"  表格 {shape_idx + 1}: {rows} 行 x {cols} 列")
            print(f"  表头内容: ", end="")
            header = []
            for col in range(min(cols, 5)):
                cell = table.cell(0, col)
                text = cell.text.strip()[:10]
                header.append(text)
            print(header)
            print(f"  第二行内容: ", end="")
            if rows > 1:
                row2 = []
                for col in range(min(cols, 5)):
                    cell = table.cell(1, col)
                    text = cell.text.strip()[:10]
                    row2.append(f"'{text}'" if text else "'(空)'")
                print(row2)