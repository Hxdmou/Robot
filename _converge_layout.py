# -*- coding: utf-8 -*-
'''V3.26收敛校准：COM只读测量BoundHeight，python-pptx写spcPts段间距（绝不COM写SpaceAfter）
流程：每轮 ①python-pptx读当前sa ②COM只读测B ③计算新sa ④python-pptx应用+必要时裁剪 ⑤保存
对双版本执行，迭代到0溢出0空隙'''
import sys
sys.stdout.reconfigure(encoding='utf-8')
import time, json
import win32com.client
from pptx import Presentation
from pptx.util import Pt

FILES = [
    r"F:\个人作品\具身智能\具身智能AI产业最新进展_20260817_商务汇报_无水印_v29.pptx",
    r"F:\个人作品\具身智能\具身智能AI产业最新进展_20260817_商务汇报_水印版_v29.pptx",
]
OV_TOL = 2.0   # 溢出容差
GP_TOL = 6.0   # 空隙容差
MAX_ROUNDS = 6

def com_measure(fp):
    """COM只读测量：返回 {(slide_idx, sig): {B,H,n,W}} sig=首段前15字"""
    App = win32com.client.DispatchEx("PowerPoint.Application")
    App.Visible = True
    time.sleep(1)
    P = App.Presentations.Open(fp, ReadOnly=True)
    time.sleep(2)
    result = {}
    for si in range(1, P.Slides.Count + 1):
        try:
            P.Slides(si).Select()
        except Exception:
            pass
        time.sleep(0.05)
        for shape in P.Slides(si).Shapes:
            if not shape.HasTextFrame:
                continue
            try:
                tr = shape.TextFrame.TextRange
                if len(tr.Text.strip()) == 0 or tr.Paragraphs().Count < 2:
                    continue
                time.sleep(0.12)
                vals = sorted([tr.BoundHeight for _ in range(5)])
                B = vals[2]
                n_par = tr.Paragraphs().Count
                sig = tr.Paragraphs(1).Text[:15].replace('\r', '')
                # 找最长段（裁剪目标）
                worst_idx, worst_len = 1, 0
                for i in range(1, n_par + 1):
                    L = len(tr.Paragraphs(i).Text)
                    if L > worst_len:
                        worst_len, worst_idx = L, i
                result[(si, sig)] = {'B': B, 'H': shape.Height, 'n': n_par, 'W': shape.Width,
                                     'worst_idx': worst_idx - 1, 'worst_len': worst_len}
            except Exception:
                pass
    P.Close()
    App.Quit()
    return result

def pptx_apply(fp, adjustments):
    """python-pptx应用新sa（写spcPts）+裁剪。adjustments: {(slide_idx,sig): {'sa':x,'trim':[(par_idx,chars)]}}"""
    prs = Presentation(fp)
    applied = 0
    for si, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            paras = tf.paragraphs
            if len(paras) < 2:
                continue
            sig = paras[0].text[:15]
            key = (si, sig)
            if key not in adjustments:
                continue
            adj = adjustments[key]
            # 应用段间距：前n-1段设sa，末段设0（BoundHeight不含末段段距，避免尾部空隙）
            if 'sa' in adj:
                sa = adj['sa']
                for i, p in enumerate(paras):
                    if i < len(paras) - 1:
                        p.space_after = Pt(sa)
                    else:
                        p.space_after = Pt(0)
            # 裁剪指定段落（从段尾删除chars个字符）
            if adj.get('trim'):
                for pidx, chars in adj['trim']:
                    if pidx < len(paras):
                        p = paras[pidx]
                        full = ''.join(r.text for r in p.runs)
                        if chars < len(full) - 42:
                            remain = chars
                            for r in reversed(p.runs):
                                if remain <= 0:
                                    break
                                if len(r.text) <= remain:
                                    remain -= len(r.text)
                                    r.text = ''
                                else:
                                    r.text = r.text[:-remain]
                                    remain = 0
            applied += 1
    prs.save(fp)
    return applied

def process_file(fp):
    for rnd in range(1, MAX_ROUNDS + 1):
        meas = com_measure(fp)
        # 读当前sa
        prs = Presentation(fp)
        cur_sa = {}
        for si, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                paras = shape.text_frame.paragraphs
                if len(paras) < 2:
                    continue
                sig = paras[0].text[:15]
                try:
                    sa = paras[1].space_after.pt if paras[1].space_after else 0.0
                except Exception:
                    sa = 0.0
                cur_sa[(si, sig)] = sa
        adjustments = {}
        ov = gp = ok = 0
        for key, m in meas.items():
            B, H, n, W = m['B'], m['H'], m['n'], m['W']
            sa_cur = cur_sa.get(key, 0.0)
            # 自然高度 = 当前B - 当前段距贡献
            B0 = B - sa_cur * (n - 1)
            if B > H + OV_TOL:
                ov += 1
                if B0 > H - 1 and m.get('worst_len', 0) > 45:
                    # 内容本身超高：裁剪最长段，sa先归0（下轮再填充）
                    overflow_lines = (B0 - H + 2) / 11.0
                    chars_per_line = max(int(W / 10.6), 20)
                    cut = int(overflow_lines * chars_per_line) + 8
                    adjustments[key] = {'sa': 0.0, 'trim': [(m['worst_idx'], cut)]}
                else:
                    # 内容不超高、只是段距过大：精确计算目标sa（绝不粗暴归0，避免振荡）
                    sa_target = max(0.0, min((H - 1 - B0) / (n - 1), 40.0))
                    adjustments[key] = {'sa': sa_target}
            elif H - B > GP_TOL:
                gp += 1
                # 空隙：增大sa
                new_sa = sa_cur + (H - 1 - B) / (n - 1)
                new_sa = max(0.0, min(new_sa, 40.0))
                adjustments[key] = {'sa': new_sa}
            else:
                ok += 1
        print(f"  轮{rnd}: 溢出{ov} 空隙{gp} 达标{ok} 调整{len(adjustments)}")
        if ov == 0 and gp == 0:
            print(f"  ✓ 第{rnd}轮收敛：0溢出0空隙")
            return True
        pptx_apply(fp, adjustments)
        time.sleep(1)
    return False

for fp in FILES:
    print('=' * 70)
    print('收敛校准:', fp.split('\\')[-1])
    done = process_file(fp)
    if not done:
        print('  ⚠ 达到最大轮数仍未完全收敛')
print('V3.26收敛校准全部完成')
