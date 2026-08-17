
# 精确测量每个文本框在10pt微软雅黑下，1.1行距，space_after=3pt时能放多少字
import win32com.client.dynamic
import time

# 先生成一个测试PPT，测量最大容量
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = 13.333
SLIDE_H = 7.5
CONTENT_X = 0.35
CONTENT_W = SLIDE_W - 2*CONTENT_X  # = 12.633in
col_gap = 0.12
col_w = (CONTENT_W - col_gap)/2  # = (12.633-0.12)/2 = 6.2565in
CONTENT_UPPER_H = 3.25  # 上部卡片高度
CONTENT_LOWER_H = 3.62  # 下部卡片高度
DETAIL_CONTENT_H = 7.0  # detail页高度
DETAIL_MARGIN_X = 0.35

# 测试左栏文本框容量
slide = prs.slides.add_slide(prs.slide_layouts[6])
test_box = slide.shapes.add_textbox(Inches(CONTENT_X + 0.05), Inches(0.5), Inches(col_w - 0.08), Inches(CONTENT_UPPER_H - 0.04))
tf = test_box.text_frame
tf.word_wrap = True
tf.margin_left = Pt(0); tf.margin_right = Pt(0); tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
tf.vertical_anchor = MSO_ANCHOR.TOP
p = tf.paragraphs[0]
p.line_spacing = Pt(11)
p.space_after = Pt(0)
run = p.add_run()
run.text = '▎核心内容'
run.font.size = Pt(10); run.font.name = '微软雅黑'

# 逐步增加文字直到溢出
def measure_capacity(tf, items, space_after_pt, box_h_in):
    from pptx.util import Inches, Pt
    # 添加条目直到溢出
    prev_len = 0
    for i in range(50):
        if len(tf.paragraphs) > 1:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.line_spacing = Pt(11)
        p.space_after = Pt(space_after_pt)
        # 用"测"字测试，中文宽度一致
        test_text = f"{i+1}. " + "测" * 80
        r = p.add_run()
        r.text = test_text
        r.font.size = Pt(10); r.font.name = '微软雅黑'
        
        # 临时保存打开检查
        prs.save('F:\\个人作品\\具身智能\\test_capacity.pptx')
        time.sleep(0.5)
        
        import win32com.client.dynamic
        app = win32com.client.dynamic.Dispatch("PowerPoint.Application")
        pres = app.Presentations.Open('F:\\个人作品\\具身智能\\test_capacity.pptx', ReadOnly=True)
        time.sleep(0.5)
        s = pres.Slides(1)
        has_overflow = False
        text_h_pt = 0
        for shape_idx in range(1, s.Shapes.Count+1):
            sh = s.Shapes(shape_idx)
            if sh.HasTextFrame:
                try:
                    has_overflow = sh.TextFrame.HasOverflowText
                    text_h_pt = sh.TextFrame.TextRange.BoundHeight
                except:
                    pass
        pres.Close()
        app.Quit()
        
        box_h_pt = box_h_in * 72 - 0.04*72  # 减去边距
        gap = box_h_pt - text_h_pt
        print(f"  {i+1}条: 文本高度={text_h_pt:.1f}pt, 框高度={box_h_pt:.1f}pt, gap={gap:.1f}pt, 溢出={has_overflow}")
        if has_overflow or gap < 5:
            print(f"  → 最大容量: {i}条 (每条约{len(test_text)-3}字)")
            return i
    return 0

print("测量内容页左栏容量:")
left_cap = measure_capacity(tf, [], 2, CONTENT_UPPER_H)
print()
prs.save('F:\\个人作品\\具身智能\\test_capacity.pptx')
