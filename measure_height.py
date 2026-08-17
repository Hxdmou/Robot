# -*- coding: utf-8 -*-
"""精确测量不同条目数需要的实际高度"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import re

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 测试内容：每条长度适中（约60-70字）
test_items = [
    '【测试标签01】这是第一条测试内容，用来精确测量10pt字号1.1行距下的实际占用高度，确保计算准确',
    '【测试标签02】第二条内容包含足够文字用于折行测试，因为PPT自动折行会影响总高度需要考虑',
    '【测试标签03】第三条内容保持类似长度，让我们能够准确计算20条内容总共需要多少空间才够',
    '【测试标签04】第四条内容同样是标准长度，我们将测量10条单栏和20条通栏的实际高度',
    '【测试标签05】第五条内容继续填充，每条都包含中文标点和标签，可以模拟真实的使用场景',
    '【测试标签06】第六条内容保持一致长度，这样测量出来的结果才具有代表性和参考价值',
    '【测试标签07】第七条内容还是一样的标准长度，我们将通过多个测试找到最优的布局参数',
    '【测试标签08】第八条内容继续，接下来我们测量左栏10条右栏10条以及下部3条的高度分配',
    '【测试标签09】第九条内容，测量完成后我们将根据结果精确设置每个文本框的高度值',
    '【测试标签10】第十条内容，单栏10条测量完成后再测试通栏20条内容需要的准确高度',
    '【测试标签11】第十一条开始是通栏内容，通栏宽度更宽折行更少高度会更矮一些需要区分',
    '【测试标签12】第十二条内容，通栏宽度约12.9英寸而双栏宽度约6.3英寸折行情况不同',
    '【测试标签13】第十三条内容，我们将用测量得到的精确数值来设置布局确保完全填满',
    '【测试标签14】第十四条内容，精确设置高度后既不会溢出也不会留下任何空隙完美贴合',
    '【测试标签15】第十五条内容，这次是最后一次机会我们必须确保结果完美符合用户要求',
    '【测试标签16】第十六条内容，内容完全铺满背景图案不浪费任何空间对观众表示尊重',
    '【测试标签17】第十七条内容，重点标签金色加粗高亮方便观众一眼就能识别重点内容',
    '【测试标签18】第十八条内容，所有页面格式统一10pt字号1.1行距边距设置合理舒适',
    '【测试标签19】第十九条内容，封面目录内容页细节页封底页每一页都经过精心设计布局',
    '【测试标签20】第二十条内容，总共47页完整包含22个模块双版本同步生成带牢固防盗窃水印',
]

def add_bullets_measured(tf, items, sz=10, line_spacing=1.1, space_after=2):
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        run_prefix = p.add_run()
        run_prefix.text = '▸ '
        run_prefix.font.size = Pt(sz)
        run_prefix.font.name = '微软雅黑'
        parts = re.split(r'(【[^】]+】)', item)
        for part in parts:
            run = p.add_run()
            run.text = part
            run.font.size = Pt(sz)
            run.font.name = '微软雅黑'

# 测试1：双栏模式（每栏宽约6.3英寸），10条内容
col_w = 6.3
box1 = slide.shapes.add_textbox(Inches(0.2), Inches(0.3), Inches(col_w), Inches(10))
tf1 = box1.text_frame
tf1.word_wrap = True
tf1.margin_left = Pt(5); tf1.margin_right = Pt(5); tf1.margin_top = Pt(2); tf1.margin_bottom = Pt(2)
add_bullets_measured(tf1, test_items[:10])
print(f"双栏10条内容 - 文本框高度设为10英寸，自动计算实际高度")

# 测试2：通栏模式（宽约12.9英寸），20条内容
full_w = 12.933
box2 = slide.shapes.add_textbox(Inches(0.2), Inches(0.3), Inches(full_w), Inches(10))
tf2 = box2.text_frame
tf2.word_wrap = True
tf2.margin_left = Pt(5); tf2.margin_right = Pt(5); tf2.margin_top = Pt(2); tf2.margin_bottom = Pt(2)
add_bullets_measured(tf2, test_items[:20])
print(f"通栏20条内容 - 文本框高度设为10英寸，自动计算实际高度")

prs.save('F:\\个人作品\\具身智能\\measure_result.pptx')
print("测量PPT已保存到 measure_result.pptx")
print("请打开measure_result.pptx，查看两个文本框的实际高度：")
print("1. 左上方双栏10条（约6.3英寸宽）实际高度是多少？")
print("2. 下方通栏20条（约12.9英寸宽）实际高度是多少？")
