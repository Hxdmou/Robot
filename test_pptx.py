
# ============================================================================
# 免责声明与AI使用规范
# ============================================================================
# 本文件仅供技术研究与学习交流使用，不得用于任何非法用途。
#
# AI使用规范：
#   1. 使用本文件相关内容时须遵守所在地法律法规及伦理准则
#   2. 不得用于侵犯他人合法权益、危害网络安全、破坏公共秩序的活动
#   3. 涉及自动化决策的场景须确保人工复核机制与可解释性
#   4. 处理个人信息时须符合数据保护相关法规要求
#
# 风险提示：
#   本文件内容按"现状"提供，不保证绝对准确无误。
#   使用者须自行评估风险，因使用本文件导致的任何损失由使用者承担。
# ============================================================================

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 创建一个简单的PPTX来测试表格显示
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# 添加一个空白幻灯片
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)

# 添加标题
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(12), Inches(0.8))
title_frame = title_box.text_frame
p = title_frame.add_paragraph()
p.text = "测试章节标题"
p.font.size = Pt(24)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# 创建一个简单的表格
table_data = [
    ['列1', '列2', '列3', '列4'],
    ['数据1', '数据2', '数据3', '数据4'],
    ['数据5', '数据6', '数据7', '数据8'],
]

num_rows = len(table_data)
num_cols = len(table_data[0])
left = Inches(0.4)
top = Inches(2.0)
width = Inches(14.8)
height = Inches(4.2)
table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

for i, row in enumerate(table_data):
    for j, cell in enumerate(row):
        table.cell(i, j).text = cell
        paragraph = table.cell(i, j).text_frame.paragraphs[0]
        paragraph.font.size = Pt(10)
        paragraph.alignment = PP_ALIGN.CENTER
        if i == 0:
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.bold = True
            table.cell(i, j).fill.solid()
            table.cell(i, j).fill.fore_color.rgb = RGBColor(0, 100, 200)

# 保存
prs.save(EXTERNAL_PROJECT_DIR  # was: test.pptx)
print('测试PPTX已生成')

# 验证生成的PPTX
prs2 = Presentation(EXTERNAL_PROJECT_DIR  # was: test.pptx)
for slide in prs2.slides:
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            print('幻灯片内容:', shape.text[:50])