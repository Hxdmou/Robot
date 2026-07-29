
# ============================================================================
# 免责声明与AI使用规范
# ============================================================================
# 本文件仅供技术研究与学习交流使用，不得用于任何非法用途。
#
# AI使用规范：
#   1. 使用本文件相关内容时须遵守所在地法律法规及伦理准则
#   2. 不得用于侵犯他人合法权益、危害网络安全、破坏公共秩序的活动
#   3. 涉及自动化决策的场景须确保人工复核机制与可解释性
#   4. 处理个人信息时须符合数据保护相关法规要求
#
# 风险提示：
#   本文件内容按"现状"提供，不保证绝对准确无误。
#   使用者须自行评估风险，因使用本文件导致的任何损失由使用者承担。
# ============================================================================

from pptx import Presentation

def verify_fix():
    print("=" * 120)
    print("✅ 验证修复结果")
    print("=" * 120)
    
    ppt_path = EXTERNAL_PROJECT_DIR  # was: A2A_ENTERPRISE_PPT_V12_完整版.pptx
    
    prs = Presentation(ppt_path)
    
    # 检查项目3（幻灯片3）
    if len(prs.slides) > 2:
        slide = prs.slides[2]
        print(f"\n📄 幻灯片3 (项目3):")
        
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                print(f"   表格: {len(table.rows)}行 × {len(table.columns)}列")
                
                # 打印列标题
                headers = []
                for j in range(len(table.columns)):
                    headers.append(table.cell(0, j).text.strip())
                print(f"   列标题: {headers}")
                
                # 打印所有行
                print(f"\n   完整内容:")
                for i in range(len(table.rows)):
                    row_content = []
                    for j in range(len(table.columns)):
                        row_content.append(table.cell(i, j).text.strip())
                    print(f"   {row_content}")
                
                # 检查是否包含L10边缘层
                has_l10 = any("边缘层" in table.cell(i, 0).text for i in range(len(table.rows)))
                print(f"\n   是否包含边缘层(L10): {'✅ 是' if has_l10 else '❌ 否'}")
                
                break

if __name__ == "__main__":
    verify_fix()