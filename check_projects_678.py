
# ============================================================================
# 免责声明与AI使用规范
# ============================================================================
# 本文件仅供技术研究与学习交流使用，不得用于任何非法用途。
#
# AI使用规范：
#   1. 使用本文件相关内容时须遵守所在地法律法规及伦理准则
#   2. 不得用于侵犯他人合法权益、危害网络安全、破坏公共秩序的活动
#   3. 涉及自动化决策的场景须确保人工复核机制与可解释性
#   4. 处理个人信息时须符合数据保护相关法规要求
#
# 风险提示：
#   本文件内容按"现状"提供，不保证绝对准确无误。
#   使用者须自行评估风险，因使用本文件导致的任何损失由使用者承担。
# ============================================================================

from pptx import Presentation

def check_projects():
    ppt_path = EXTERNAL_PROJECT_DIR  # was: A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx
    prs = Presentation(ppt_path)
    
    print("=" * 120)
    print("📌 项目6、7、8当前内容检查")
    print("=" * 120)
    
    # 项目6
    if len(prs.slides) > 5:
        slide = prs.slides[5]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                print(f"\n项目6 - 应用场景: {len(table.rows)}行 × {len(table.columns)}列")
                print(f"当前案例数: {len(table.rows) - 1}个")
                for i in range(1, min(4, len(table.rows))):
                    print(f"  案例{i}: {table.cell(i, 0).text} - {table.cell(i, 1).text}")
                break
    
    # 项目7
    if len(prs.slides) > 6:
        slide = prs.slides[6]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                print(f"\n项目7 - 案例分析: {len(table.rows)}行 × {len(table.columns)}列")
                print(f"当前案例数: {len(table.rows) - 1}个")
                print(f"特性列数: {len(table.columns)}个")
                for i in range(1, min(4, len(table.rows))):
                    print(f"  案例{i}: {table.cell(i, 0).text} - {table.cell(i, 1).text}")
                break
    
    # 项目8
    if len(prs.slides) > 7:
        slide = prs.slides[7]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                print(f"\n项目8 - 技术对比: {len(table.rows)}行 × {len(table.columns)}列")
                print(f"对比维度数: {len(table.rows) - 1}个")
                print(f"对比协议数: {len(table.columns) - 1}个")
                for i in range(1, min(4, len(table.rows))):
                    print(f"  维度{i}: {table.cell(i, 0).text}")
                break
    
    print("\n" + "=" * 120)
    print("📌 需要扩展的内容:")
    print("=" * 120)
    print("项目6: 应用场景从10个扩展到15个（添加法律、媒体、交通、游戏、科研）")
    print("项目7: 案例从6个扩展到10个（添加更多行业案例）")
    print("项目8: 对比维度从9个扩展到15个（添加更多技术特性）")

if __name__ == "__main__":
    check_projects()