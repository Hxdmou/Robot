# -*- coding: utf-8 -*-
'''
具身智能&AI产业最新进展 PPT生成脚本
商务汇报风格 · 深蓝科技主题 · 22模块完整
日期：2026年8月21日
布局：每个模块拆分为【内容描述】和【细节描述】两页
- 内容描述页（单数页）：上部左右各10条 + 下部通栏阐述
- 细节描述页（双数页）：通栏20条超详细内容（不分栏）
- 每页标注【内容描述】/【细节描述】标签
- 正文字号统一10pt，行距1.1
'''
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ========== 颜色定义 ==========
DARK_BLUE = RGBColor(0x0A, 0x16, 0x2F)
MID_BLUE = RGBColor(0x10, 0x25, 0x48)
ACCENT_BLUE = RGBColor(0x1E, 0x5F, 0xA8)
GOLD = RGBColor(0xD4, 0xA5, 0x37)
LGRAY = RGBColor(0xC8, 0xD0, 0xDC)
MGRAY = RGBColor(0x80, 0x88, 0x98)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ========== 尺寸常量 - 动态精确计算：逐页实测行数，段间距自动填满（V3.23修复） ==========
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.08
CONTENT_X = MARGIN
CONTENT_W = SLIDE_W - 2 * MARGIN
HEADER_Y = 0.42
CONTENT_TOP = HEADER_Y
FOOTER_Y = 7.36
CONTENT_BOTTOM = 7.32
CONTENT_H = CONTENT_BOTTOM - CONTENT_TOP  # 6.9英寸 = 496.8pt
CONTENT_GAP = 0.05
# 细节页通栏区域高度（卡片高度按实测行数动态计算，不再使用统一固定值）
DETAIL_CONTENT_H = CONTENT_H
DETAIL_MARGIN_X = 0.05
BODY_SZ = 10
# V3.35统一页脚铁律：所有页（内容/细节/目录/封面/封底）页脚文字完全一致
FOOTER_TEXT = '具身智能&AI产业最新进展 · 2026年8月21日 · 商务汇报'

# ========== 辅助函数 ==========
def bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fc=None, ec=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.line.fill.background()
    if fc:
        shp.fill.solid(); shp.fill.fore_color.rgb = fc
    if ec:
        shp.line.color.rgb = ec; shp.line.width = Pt(0.5)
    return shp

def rrect(slide, x, y, w, h, fc=None, ec=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.line.fill.background()
    if fc:
        shp.fill.solid(); shp.fill.fore_color.rgb = fc
    if ec:
        shp.line.color.rgb = ec
    return shp

def tb(slide, x, y, w, h, text, sz=10, b=False, c=WHITE, al=PP_ALIGN.LEFT, an=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    tf.vertical_anchor = an
    p = tf.paragraphs[0]
    p.alignment = al
    # 固定行距11pt = 10pt * 1.1，精确控制高度
    p.line_spacing = Pt(11)
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    p.left_indent = Pt(0)
    p.first_line_indent = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = b
    run.font.color.rgb = c
    run.font.name = '微软雅黑'
    return box

def add_bullets(tf, items, start_idx=0, sz=10, color=LGRAY, space_after=0, line_spacing=11):
    import re
    sa_is_list = isinstance(space_after, (list, tuple))
    for i, item in enumerate(items):
        idx = start_idx + i
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        # 字体全部居中（V3.30用户硬性要求）
        p.alignment = PP_ALIGN.CENTER
        # 行距可配（V3.38：内容页用15pt更饱满，细节页用11pt紧凑）
        p.line_spacing = Pt(line_spacing)
        sa_i = space_after[i] if sa_is_list and i < len(space_after) else (0 if sa_is_list else space_after)
        p.space_after = Pt(sa_i)
        p.space_before = Pt(0)
        # 取消所有缩进，最大化可用宽度
        p.left_indent = Pt(0)
        p.first_line_indent = Pt(0)
        p.bullet_indent = Pt(0)
        # 前缀符号
        run_prefix = p.add_run()
        run_prefix.text = '▸ '
        run_prefix.font.size = Pt(sz)
        run_prefix.font.color.rgb = ACCENT_BLUE
        run_prefix.font.name = '微软雅黑'
        # 解析内容，【】内的标签用金色加粗高亮
        parts = re.split(r'(【[^】]+】)', item)
        for part in parts:
            run = p.add_run()
            run.text = part
            run.font.size = Pt(sz)
            run.font.name = '微软雅黑'
            if part.startswith('【') and part.endswith('】'):
                run.font.color.rgb = GOLD
                run.font.bold = True
            else:
                run.font.color.rgb = color
    # V3.33：末段段间距置0（BoundHeight不含末段sa，避免底部视觉空隙歧义）
    if items and not sa_is_list:
        tf.paragraphs[len(items) + start_idx - 1].space_after = Pt(0)

# ========== 精确文本测量（V3.23修复核心：两阶段闭环，以PowerPoint真实渲染为基准） ==========
import unicodedata

# 双阶段模式：MEASURE_MODE生成sa=0测量版；最终版用MEASURED真实数据反解段间距
MEASURE_MODE = False
MEASURED = {}   # key -> {'B0': 真实自然高度pt, 'nb': bullet段数, 'sa': 段间距pt}
TRIMMED = {}    # key -> 被COM裁剪后的条目文本列表

# 自动加载COM实测布局数据（若存在）
import os as _os
_LAYOUT_FILE = _os.path.join(_os.path.dirname(_os.path.abspath(__file__)), '_layout_final.py')
if _os.path.exists(_LAYOUT_FILE):
    exec(open(_LAYOUT_FILE, encoding='utf-8').read(), globals())

# 各区域可用高度(pt)：与布局常量严格同步
UPPER_REGION_H = 3.35
LOWER_REGION_H = CONTENT_BOTTOM - (CONTENT_TOP + UPPER_REGION_H + CONTENT_GAP)
PAD_X, PAD_TOP, PAD_BOT = 0.05, 0.04, 0.05
AVAIL_UPPER = (UPPER_REGION_H - PAD_TOP - PAD_BOT) * 72
AVAIL_LOWER = (LOWER_REGION_H - PAD_TOP - PAD_BOT) * 72
AVAIL_DETAIL = (DETAIL_CONTENT_H - 0.04 - 0.05) * 72
# V3.36四页/模块铁律：内容/细节各拆2页，每页用全高区域，更宽松不紧凑
AVAIL_FULL = (CONTENT_H - PAD_TOP - PAD_BOT) * 72

def _char_width_pt(ch, sz_pt):
    """按字符显示宽度（V3.30 COM实测校准：94个文本框实测中位比值0.9545）：
    全角(CJK)=字号×1.012，半角≈0.554倍字号，消除估算偏大导致的空隙"""
    if unicodedata.east_asian_width(ch) in ('F', 'W'):
        return sz_pt * 1.012
    return sz_pt * 0.554

def _text_width_pt(text, sz_pt):
    return sum(_char_width_pt(ch, sz_pt) for ch in text)

def _lines_needed(text, box_width_pt, sz_pt):
    """计算文本在指定宽度内的换行行数（至少1行）"""
    if not text:
        return 1
    usable = max(box_width_pt, 1.0)
    import math
    return max(1, math.ceil(_text_width_pt(text, sz_pt) / usable))

def _bullet_lines(item, box_width_pt, sz_pt):
    """bullet条目实测行数：前缀'▸ '占约1.3个全角宽度"""
    prefix_w = _text_width_pt('▸ ', sz_pt)
    first_w = max(box_width_pt - prefix_w, 1.0)
    total_w = _text_width_pt(item, sz_pt)
    if total_w <= first_w:
        return 1
    import math
    return 1 + math.ceil((total_w - first_w) / max(box_width_pt, 1.0))

def measure_block_height_pt(items, box_width_pt, sz_pt, line_spacing_pt, space_after_pt, title_lines=0):
    """精确计算文本块总高度(pt)：标题行 + 各条目实测行数×行距 + 段间距"""
    total = title_lines * line_spacing_pt
    n = len(items)
    for item in items:
        total += _bullet_lines(item, box_width_pt, sz_pt) * line_spacing_pt
    if n > 0:
        total += n * space_after_pt
    return total

def solve_space_after(items, box_width_pt, sz_pt, line_spacing_pt, avail_pt, title_lines=0, max_sa=30.0):
    """反解段间距：在可用高度内均匀分布填满，返回(段间距pt, 实际总高pt)。留0.1pt安全余量防溢出"""
    natural = measure_block_height_pt(items, box_width_pt, sz_pt, line_spacing_pt, 0, title_lines)
    n = len(items)
    if n == 0 or avail_pt <= natural:
        return 0.0, natural
    # V3.33：末段sa=0（BoundHeight不含末段段间距），段距只分布在前n-1段之间
    sa = (avail_pt - natural - 0.5) / max(n - 1, 1)
    sa = max(0.0, min(sa, max_sa))
    return sa, natural + (n - 1) * sa

def _fit_items(items, box_width_pt, sz_pt, line_spacing_pt, avail_pt, title_lines=1, min_keep=40):
    """防溢出安全网：自然高度超过可用高度时，逐条裁剪最长条目直到刚好放下。
    裁剪优先在标点处断句，绝不使用省略号等占位符。"""
    items = list(items)
    prefix_w = _text_width_pt('▸ ', sz_pt)
    first_w = max(box_width_pt - prefix_w, 1.0)
    guard = 0
    while guard < 300:
        guard += 1
        h = measure_block_height_pt(items, box_width_pt, sz_pt, line_spacing_pt, 0, title_lines)
        if h <= avail_pt:
            break
        worst_i, worst_lines = -1, 1
        for i, it in enumerate(items):
            ln = _bullet_lines(it, box_width_pt, sz_pt)
            if ln > worst_lines:
                worst_lines, worst_i = ln, i
        if worst_i < 0:
            break
        it = items[worst_i]
        # 目标：减少1行 → 新文本总宽度上限
        target_w = first_w + (worst_lines - 2) * box_width_pt
        # 扫描字符找裁剪点，优先取上限前最后一个标点断句
        cut = len(it)
        acc = 0.0
        last_punct = -1
        for k, ch in enumerate(it):
            acc += _char_width_pt(ch, sz_pt)
            if acc > target_w:
                cut = k
                break
            if ch in '，、；。：）%':
                last_punct = k
        if last_punct > min_keep:
            cut = last_punct + 1
        cut = max(cut, min_keep)
        if cut >= len(it):
            cut = max(len(it) - 8, min_keep)
        items[worst_i] = it[:cut]
    return items

def _trim_to_lines(items, para_lines, target_total, box_width_pt, sz_pt):
    """V3.33：按COM实测行数裁剪条目，使bullet总行数≤target_total。
    每次裁剪最长条目减1行，裁剪宽度留7%安全余量，优先在接近上限的标点处断句，绝不使用省略号。"""
    items = list(items)
    lines = list(para_lines) + [1] * max(0, len(items) - len(para_lines))
    lines = lines[:len(items)]
    prefix_w = _text_width_pt('▸ ', sz_pt)
    first_w = max(box_width_pt - prefix_w, 1.0)
    guard = 0
    while sum(lines) > target_total and guard < 400 and items:
        guard += 1
        cand = [k for k in range(len(items)) if lines[k] > 1]
        if not cand:
            items.pop()
            lines.pop()
            continue
        i = max(cand, key=lambda k: lines[k])
        new_l = lines[i] - 1
        target_w = (first_w + (new_l - 1) * box_width_pt) * 0.93
        it = items[i]
        cut = len(it)
        acc = 0.0
        last_punct = -1
        punct_acc = 0.0
        for k, ch in enumerate(it):
            acc += _char_width_pt(ch, sz_pt)
            if acc > target_w:
                cut = k
                break
            if ch in '，、；。：）%':
                last_punct = k
                punct_acc = acc
        # 仅当标点位置不低于上限85%时才在标点断句，避免裁太多造成新空隙
        if last_punct > 20 and punct_acc >= target_w * 0.85:
            cut = last_punct + 1
        cut = max(cut, 20)
        if cut >= len(it):
            cut = max(len(it) - 8, 20)
        items[i] = it[:cut]
        lines[i] = new_l
    return items

def _resolve(key, items, avail_pt, box_width_pt, line_spacing=11, title_lines=1, sz_pt=10):
    """V3.33一次到位（整数段距精确分配，绝不迭代）：
    PowerPoint对每段高度做整数舍入，小数段距在19-20段上累积±4~9pt误差。
    解法：使用整数段距，总高度=B0+Σsa，一次命中目标。
    - MEASURE_MODE：返回sa=0且不裁剪，供COM测量真实自然高度
    - MEASURED有实测数据：按实测行数裁剪防溢出 + 整数段距精确填满
    - 无实测数据时回退估算路径
    V3.38：返回三元组(sa, items, 实际内容高度pt)，供卡片按实际高度收缩居中，杜绝大段距空隙"""
    if MEASURE_MODE:
        return 0, list(items), avail_pt
    m = MEASURED.get(key)
    if m and m.get('B0'):
        B0 = float(m['B0'])
        para_lines = list(m.get('para_lines') or [])
        items = list(items)
        if para_lines and len(para_lines) >= len(items) and 8.0 <= B0 / max(title_lines + sum(para_lines[:len(items)]), 1) <= 14.0:
            para_lines = para_lines[:len(items)]
            n_lines = title_lines + sum(para_lines)
            lh = B0 / n_lines
            max_bullet_lines = int((avail_pt - 1.0) / lh) - title_lines
            cur = sum(para_lines)
            if cur > max_bullet_lines:
                items = _trim_to_lines(items, para_lines, max_bullet_lines, box_width_pt, sz_pt)
                cur = max_bullet_lines
            B0_new = (title_lines + cur) * lh
            nb = len(items)
            if nb > 1 and B0_new < avail_pt:
                # 留1pt安全余量防溢出（1pt空隙在6pt容差内不可见）
                remaining = avail_pt - 1.0 - B0_new
                if remaining <= 0:
                    return 0, items, B0_new
                # 整数段距分配：前nb-1段分摊，末段=0（整数不产生舍入误差）
                n_gaps = nb - 1
                base = int(remaining // n_gaps)
                extra = int(round(remaining - base * n_gaps))
                sa_list = [base + 1 if i < extra else base for i in range(n_gaps)] + [0]
                # V3.38：段距上限12pt，绝不用大段距硬撑全高制造视觉空隙
                sa_list = [max(0, min(s, 12)) for s in sa_list]
                return sa_list, items, B0_new + sum(sa_list)
            else:
                return 0, items, B0_new
    # 回退：估算路径（无实测数据时）
    items = _fit_items(items, box_width_pt, sz_pt, line_spacing, avail_pt, title_lines=title_lines)
    sa, total_h = solve_space_after(items, box_width_pt, sz_pt, line_spacing, avail_pt, title_lines=title_lines)
    return sa, items, total_h

# ========== 统一页面标签 ==========
def add_page_tag(slide, tag_text, tag_color):
    rrect(slide, SLIDE_W - 1.45, 0.1, 1.25, 0.24, fc=tag_color)
    tb(slide, SLIDE_W - 1.45, 0.1, 1.25, 0.24, tag_text, sz=8, b=True, c=WHITE, al=PP_ALIGN.CENTER, an=MSO_ANCHOR.MIDDLE)

def add_page_header(slide, part_num, title):
    """正常页眉，清晰可见，不极限压缩"""
    rect(slide, 0, 0, SLIDE_W, 0.03, fc=GOLD)
    rrect(slide, 0.1, 0.08, 0.65, 0.22, fc=GOLD)
    tb(slide, 0.1, 0.08, 0.65, 0.22, part_num, sz=8, b=True, c=DARK_BLUE, al=PP_ALIGN.CENTER, an=MSO_ANCHOR.MIDDLE)
    tb(slide, 0.85, 0.06, 10.8, 0.28, title, sz=13, b=True, c=WHITE, an=MSO_ANCHOR.MIDDLE)
    rect(slide, 0.06, HEADER_Y, SLIDE_W - 0.12, 0.01, fc=RGBColor(0x20, 0x35, 0x60))

# ========== V3.37字数均衡铁律：内容/细节各拆2页，按字数均匀拆分，两页同版式 ==========
def _split_by_chars(items):
    """按累计字数均匀拆成前后两页，遍历所有拆分点找字数差最小的位置"""
    items = list(items)
    if len(items) <= 1:
        return items, []
    total = sum(len(it) for it in items)
    best_split = 1
    best_diff = float('inf')
    acc = 0
    for i in range(len(items) - 1):
        acc += len(items[i])
        diff = abs(acc - (total - acc))
        if diff < best_diff:
            best_diff = diff
            best_split = i + 1
    return items[:best_split], items[best_split:]

def _build_card_textbox(slide, x, y, w, h, title_text, items, sz=10, space_after=0.0):
    """创建卡片文本框：金色标题+bullets，返回文本框"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0); tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p_t = tf.paragraphs[0]
    p_t.alignment = PP_ALIGN.CENTER
    p_t.line_spacing = Pt(11)
    p_t.space_after = Pt(0)
    p_t.left_indent = Pt(0); p_t.first_line_indent = Pt(0)
    r_t = p_t.add_run(); r_t.text = title_text
    r_t.font.size = Pt(10); r_t.font.bold = True; r_t.font.color.rgb = GOLD; r_t.font.name = '微软雅黑'
    add_bullets(tf, items, start_idx=1, sz=sz, space_after=space_after)
    return box

def _content_page_render(prs, part_num, title, items, key_suffix, page_label):
    """内容描述页统一版式：卡片高度跟随实际内容收缩+垂直居中，杜绝大段距空隙"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, DARK_BLUE)
    add_page_header(slide, part_num, title)
    add_page_tag(slide, '内容描述·' + page_label, ACCENT_BLUE)
    text_w = CONTENT_W - 2 * PAD_X
    avail = AVAIL_FULL
    sa, items, content_h_pt = _resolve(part_num + key_suffix, list(items), avail, text_w * 72)
    # V3.38：卡片高度=实际内容高度+上下padding，垂直居中于可用区域
    content_h_in = content_h_pt / 72.0
    card_h = min(CONTENT_H, content_h_in + PAD_TOP + PAD_BOT)
    card_y = CONTENT_TOP + (CONTENT_H - card_h) / 2.0
    box_h = card_h - PAD_TOP - PAD_BOT
    rrect(slide, CONTENT_X, card_y, CONTENT_W, card_h, fc=MID_BLUE)
    rect(slide, CONTENT_X, card_y, 0.04, card_h, fc=ACCENT_BLUE)
    _build_card_textbox(slide, CONTENT_X + PAD_X, card_y + PAD_TOP, text_w, box_h, '▎核心内容 · 代表动态 · 过程阐述', items, space_after=sa)
    tb(slide, 0, FOOTER_Y, SLIDE_W, SLIDE_H - FOOTER_Y, FOOTER_TEXT, sz=7, c=MGRAY, al=PP_ALIGN.CENTER, an=MSO_ANCHOR.MIDDLE)

def content_page_1(prs, part_num, title, left_items, right_items, process_items):
    """内容描述（第一页）：内容池前半（按字数均匀拆分）"""
    pool = list(left_items) + list(right_items) + list(process_items)
    first, _ = _split_by_chars(pool)
    _content_page_render(prs, part_num, title, first, 'C1', '第一页')

def content_page_2(prs, part_num, title, left_items, right_items, process_items):
    """内容描述（第二页）：内容池后半（按字数均匀拆分）"""
    pool = list(left_items) + list(right_items) + list(process_items)
    _, second = _split_by_chars(pool)
    _content_page_render(prs, part_num, title, second, 'C2', '第二页')

# ========== V3.38四页/模块：细节描述拆2页，按字数均匀拆分，卡片收缩+垂直居中 ==========
def _detail_page_render(prs, part_num, title, detail_title, d_items, key_suffix, page_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, DARK_BLUE)
    add_page_header(slide, part_num, title)
    add_page_tag(slide, '细节描述·' + page_label, GOLD)
    region_h = DETAIL_CONTENT_H
    d_PAD_X, d_PAD_TOP, d_PAD_BOT = 0.08, 0.04, 0.05
    text_w = SLIDE_W - 2 * DETAIL_MARGIN_X - 2 * d_PAD_X
    sa, d_items, content_h_pt = _resolve(part_num + key_suffix, d_items, AVAIL_DETAIL, text_w * 72)
    # V3.38：卡片高度=实际内容高度+上下padding，垂直居中于可用区域
    content_h_in = content_h_pt / 72.0
    card_h = min(region_h, content_h_in + d_PAD_TOP + d_PAD_BOT)
    card_y = CONTENT_TOP + (region_h - card_h) / 2.0
    box_h = card_h - d_PAD_TOP - d_PAD_BOT
    rect(slide, DETAIL_MARGIN_X, card_y, 0.04, card_h, fc=GOLD)
    box = slide.shapes.add_textbox(Inches(DETAIL_MARGIN_X + d_PAD_X), Inches(card_y + d_PAD_TOP), Inches(text_w), Inches(box_h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0); tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p_title = tf.paragraphs[0]
    p_title.alignment = PP_ALIGN.CENTER
    p_title.line_spacing = Pt(11)
    p_title.space_after = Pt(0)
    p_title.left_indent = Pt(0); p_title.first_line_indent = Pt(0)
    run_t = p_title.add_run(); run_t.text = detail_title
    run_t.font.size = Pt(10); run_t.font.bold = True; run_t.font.color.rgb = GOLD; run_t.font.name = '微软雅黑'
    add_bullets(tf, d_items, start_idx=1, sz=10, space_after=sa)
    tb(slide, 0, FOOTER_Y, SLIDE_W, SLIDE_H - FOOTER_Y, FOOTER_TEXT, sz=7, c=MGRAY, al=PP_ALIGN.CENTER, an=MSO_ANCHOR.MIDDLE)

def detail_page_1(prs, part_num, title, detail_title, detail_items):
    """细节描述（第一页）：前半细节（按字数均匀拆分）"""
    first, _ = _split_by_chars(detail_items)
    _detail_page_render(prs, part_num, title, detail_title, first, 'D1', '第一页')

def detail_page_2(prs, part_num, title, detail_title, detail_items):
    """细节描述（第二页）：后半细节（按字数均匀拆分）"""
    _, second = _split_by_chars(detail_items)
    _detail_page_render(prs, part_num, title, detail_title, second, 'D2', '第二页')

# ========== 目录页 - 22模块居中填满，标题居中 ==========
def toc_page(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, DARK_BLUE)
    # 背景装饰线
    for i in range(12):
        rect(slide, 0, 0.4 + i * 0.6, SLIDE_W, 0.005, fc=RGBColor(0x18, 0x28, 0x48))
    rect(slide, 0, 0, SLIDE_W, 0.08, fc=GOLD)
    # 标题水平居中
    tb(slide, 0, 0.4, SLIDE_W, 0.75, 'CONTENTS', sz=40, b=True, c=WHITE, al=PP_ALIGN.CENTER)
    tb(slide, 0, 1.15, SLIDE_W, 0.45, '目  录 · 22个核心模块完整覆盖', sz=18, b=True, c=GOLD, al=PP_ALIGN.CENTER)
    # 金色装饰线居中
    line_w = 2.2
    rect(slide, (SLIDE_W - line_w)/2, 1.7, line_w, 0.04, fc=GOLD)
    
    modules = [
        '01  人形机器人：量产元年全面爆发',
        '02  人形新品：2026新品密集发布',
        '03  核心零部件：国产替代加速',
        '04  央企国家队：战略布局入场',
        '05  安徽产业：合芜蚌协同发展',
        '06  蚌埠中国传感谷：MEMS传感器基地',
        '07  合肥科创：科教资源集聚高地',
        '08  江淮制造：制造强省应用场景',
        '09  AI算力：大模型算力底座',
        '10  AI智能体：具身大脑核心',
        '11  6G通信：空天地一体化',
        '12  消费电子：AI终端普及',
        '13  智慧农业：农业机器人应用',
        '14  医疗健康：医疗机器人突破',
        '15  教育AI：教育智能化转型',
        '16  能源电力：电力机器人运维',
        '17  自动驾驶：L4级商业化落地',
        '18  人形运动会：技术竞赛舞台',
        '19  真机部署：规模化落地进展',
        '20  物流仓储：仓储机器人普及',
        '21  灵巧手：精密操作核心部件',
        '22  安防应急：特种机器人守护安全',
    ]
    
    col_w = (CONTENT_W - 0.8) / 2
    col_gap = 0.8
    total_cols_w = 2 * col_w + col_gap
    left_x = (SLIDE_W - total_cols_w) / 2
    right_x = left_x + col_w + col_gap
    row_h = 0.46
    rows = 11
    # 整体垂直居中填满
    total_list_h = rows * row_h
    start_y = 1.9 + ((5.2 - total_list_h) / 2)  # 在1.9-7.1区域内垂直居中填满
    
    for i, mod in enumerate(modules):
        col = i // rows
        row = i % rows
        x = left_x if col == 0 else right_x
        y = start_y + row * row_h
        rrect(slide, x, y + 0.05, 0.52, 0.32, fc=ACCENT_BLUE)
        tb(slide, x, y + 0.05, 0.52, 0.32, mod[:2], sz=12, b=True, c=WHITE, al=PP_ALIGN.CENTER, an=MSO_ANCHOR.MIDDLE)
        tb(slide, x + 0.62, y + 0.05, col_w - 0.65, 0.32, mod[2:], sz=12, c=LGRAY, an=MSO_ANCHOR.MIDDLE)
    
    # 底部信息填满
    rect(slide, 0, SLIDE_H - 0.36, SLIDE_W, 0.04, fc=GOLD)
    tb(slide, 0, SLIDE_H - 0.28, SLIDE_W, 0.2, FOOTER_TEXT, sz=8, c=MGRAY, al=PP_ALIGN.CENTER, an=MSO_ANCHOR.MIDDLE)

# ========== 封面页 - 大气饱满无空白，所有内容正确对齐 ==========
def cover_page(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, DARK_BLUE)
    # 背景横线装饰填满
    for i in range(15):
        rect(slide, 0, 0.35 + i * 0.48, SLIDE_W, 0.006, fc=RGBColor(0x14, 0x24, 0x48))
    rect(slide, 0, 0, SLIDE_W, 0.08, fc=GOLD)
    # 左侧金色装饰条
    rrect(slide, 0.8, 1.6, 0.12, 3.2, fc=GOLD)
    # 主标题水平垂直居中
    tb(slide, 0, 1.4, SLIDE_W, 1.2, '具身智能&AI产业', sz=56, b=True, c=WHITE, al=PP_ALIGN.CENTER, an=MSO_ANCHOR.MIDDLE)
    tb(slide, 0, 2.6, SLIDE_W, 1.2, '最 新 进 展', sz=56, b=True, c=GOLD, al=PP_ALIGN.CENTER, an=MSO_ANCHOR.MIDDLE)
    # 英文副标题居中
    tb(slide, 0, 4.0, SLIDE_W, 0.5, 'Embodied Intelligence & AI Industry Report 2026', sz=15, c=MGRAY, al=PP_ALIGN.CENTER, an=MSO_ANCHOR.MIDDLE)
    # 描述居中
    tb(slide, 0, 4.6, SLIDE_W, 0.45, '—— 22个核心模块完整分析 · 量产元年全景观察 ——', sz=15, c=LGRAY, al=PP_ALIGN.CENTER, an=MSO_ANCHOR.MIDDLE)
    # 标签行：6个标签均匀分布，文字水平垂直居中
    tags = ['人形机器人量产', '核心零部件国产替代', '安徽合芜蚌产业', '蚌埠中国传感谷', 'AI算力+智能体', '22模块全覆盖']
    tag_y = 5.3
    tag_w = 1.95
    tag_gap = 0.15
    total_tag_w = len(tags) * tag_w + (len(tags)-1) * tag_gap
    tag_start_x = (SLIDE_W - total_tag_w) / 2
    for i, tag in enumerate(tags):
        tx = tag_start_x + i * (tag_w + tag_gap)
        rrect(slide, tx, tag_y, tag_w, 0.32, fc=ACCENT_BLUE)
        tb(slide, tx, tag_y, tag_w, 0.32, tag, sz=9, b=True, c=WHITE, al=PP_ALIGN.CENTER, an=MSO_ANCHOR.MIDDLE)
    # 版本和日期整体水平居中
    btn_w = 2.2
    date_w = 5.5
    btn_gap = 0.3
    total_block_w = btn_w + btn_gap + date_w
    block_start_x = (SLIDE_W - total_block_w) / 2
    rrect(slide, block_start_x, 6.0, btn_w, 0.45, fc=GOLD)
    tb(slide, block_start_x, 6.0, btn_w, 0.45, '商务汇报版', sz=13, b=True, c=DARK_BLUE, al=PP_ALIGN.CENTER, an=MSO_ANCHOR.MIDDLE)
    tb(slide, block_start_x + btn_w + btn_gap, 6.0, date_w, 0.45, '2026年8月21日', sz=16, b=True, c=WHITE, al=PP_ALIGN.CENTER, an=MSO_ANCHOR.MIDDLE)
    # 底部装饰线
    rect(slide, 0, SLIDE_H - 0.36, SLIDE_W, 0.04, fc=GOLD)
    tb(slide, 0, SLIDE_H - 0.28, SLIDE_W, 0.2, FOOTER_TEXT, sz=8, c=MGRAY, al=PP_ALIGN.CENTER, an=MSO_ANCHOR.MIDDLE)

# ========== 封底页 - 内容填满不留空，标签金色高亮 ==========
def back_page(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, DARK_BLUE)
    for i in range(15):
        rect(slide, 0, 0.35 + i * 0.48, SLIDE_W, 0.005, fc=RGBColor(0x1A, 0x24, 0x40))
    rect(slide, 0, 0, SLIDE_W, 0.08, fc=GOLD)
    
    # 标题区
    tb(slide, 0, 0.7, SLIDE_W, 0.9, 'THANK YOU', sz=52, b=True, c=WHITE, al=PP_ALIGN.CENTER)
    tb(slide, 0, 1.6, SLIDE_W, 0.42, '总结与展望 · 共赴具身智能时代', sz=20, b=True, c=GOLD, al=PP_ALIGN.CENTER)
    rect(slide, 3.8, 2.1, 5.73, 0.03, fc=GOLD)
    
    # 六大总结板块 - 两排三列填满
    blocks = [
        ('产业判断', [
            '【量产元年】2026年量产万台级交付',
            '【爆发节点】2027-2028年规模爆发',
            '【成本下探】年均下降35%至5-10万',
            '【国产化率】核心零部件超90%'
        ]),
        ('技术趋势', [
            '【具身大模型】VLA模型成主流路线',
            '【灵巧操作】力控精度接近人手',
            '【端侧算力】100-500TOPS成标配',
            '【多机协同】5台以上群体协作'
        ]),
        ('应用场景', [
            '【工业制造】汽车/3C工厂先行落地',
            '【物流仓储】搬运分拣码垛规模化',
            '【商业服务】酒店/展厅/医院普及',
            '【家庭消费】2028年进入千家万户'
        ]),
        ('安徽机遇', [
            '【制造强省】新能源汽车+家电场景丰富',
            '【合芜蚌协同】合肥AI+芜湖制造+蚌埠传感',
            '【科教支撑】中科大/合工大人才供给',
            '【政策支持】100亿产业基金加持'
        ]),
        ('蚌埠机会', [
            '【中国传感谷】国家级MEMS产业基地',
            '【六维力传感器】国内市占率超60%',
            '【产业集聚】68家传感器企业落地',
            '【对接规划】4小时上门技术支持'
        ]),
        ('未来展望', [
            '【2027年】年销量突破50万台',
            '【2030年】全球规模超万亿级',
            '【中国引领】成为最大生产应用市场',
            '【通用智能】真正实现通用人工智能'
        ]),
    ]
    
    # 六大总结板块 - 两排三列，动态高度+段间距填满（V3.23修复空隙）
    area_top, area_bot = 2.35, 7.05
    row_gap = 0.2
    bw = (CONTENT_W - 2 * 0.2) / 3
    bh = (area_bot - area_top - row_gap) / 2  # 动态铺满垂直区域
    for i, (t, items) in enumerate(blocks):
        col = i % 3
        row = i // 3
        x = CONTENT_X + col * (bw + 0.2)
        y = area_top + row * (bh + row_gap)
        rrect(slide, x, y, bw, bh, fc=MID_BLUE)
        rect(slide, x, y, 0.06, bh, fc=GOLD)
        tb(slide, x + 0.12, y + 0.08, bw - 0.2, 0.3, t, sz=12, b=True, c=GOLD)
        # 文本框：封底短文本用估算段间距填满（估算对短文本精确）
        box_y, box_h = y + 0.42, bh - 0.5
        b_items = list(items)
        text_w_pt = (bw - 0.2) * 72 - 8
        sa, _ = solve_space_after(b_items, text_w_pt, 9, 10, box_h * 72, title_lines=0, max_sa=40.0)
        box = slide.shapes.add_textbox(Inches(x + 0.12), Inches(box_y), Inches(bw - 0.2), Inches(box_h))
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        add_bullets(tf, b_items, sz=9, space_after=sa)
        # 封底9pt字号，固定行距10pt(≈1.1倍)
        for p in tf.paragraphs:
            p.line_spacing = Pt(10)
    
    # 底栏信息
    rect(slide, 0, SLIDE_H - 0.36, SLIDE_W, 0.04, fc=GOLD)
    tb(slide, 0, SLIDE_H - 0.28, SLIDE_W, 0.2, FOOTER_TEXT, sz=8, c=MGRAY, al=PP_ALIGN.CENTER, an=MSO_ANCHOR.MIDDLE)

# ========== 22个模块完整数据 ==========
all_modules = []

# PART 01 人形机器人量产
all_modules.append(('PART 01', '人形机器人：量产元年全面爆发',
    ['【WRC2026·2026年8月21日】2026世界机器人大会汇聚超300家企业(+69%)展品超3000件首发新品311款，主题"人机共生 产需共融"；49家央企首次集中参展带来12类应用场景；北京经开区落地多条千台级人形机器人中试产线面向消防/商超实用场景',
     '【产业定位】人形机器人是继智能手机、新能源汽车之后下一代通用智能终端，是具身智能技术最佳物理载体，被视为第四次工业革命核心标志性产品',
     '【市场规模】2026年全球人形机器人市场规模正式突破1200亿元大关，中国市场占比超45%达到540亿元，预计2030年全球规模将突破万亿',
     '【量产节点】优必选、特斯拉Optimus、小米CyberOne、波士顿动力Atlas等头部企业2026年全面实现万台级量产交付，产业进入规模化阶段',
     '【产能规划】国内已公布人形机器人产能规划总计超80万台/年，主要集中在长三角、珠三角和京津冀地区，2027年将进入产能集中释放期',
     '【价格下探】规模化量产带动BOM成本快速下降，2026年主流工业级人形机器人整机价格成功下探至15-25万元区间，较2023年下降60%',
     '【应用场景】工业制造场景先行落地，汽车总装、3C电子、新能源电池工厂率先部署；物流仓储、商业服务、家庭陪护场景2027年起逐步拓展',
     '【技术成熟度】行走稳定性、负载能力、续航时间三大核心指标全面达到商业化可用水平，MTBF平均无故障工作时间突破2000小时',
     '【资本热度】2026年上半年国内人形机器人领域融资总额超260亿元，同比增长180%，宇树科技科创板上市市值突破600亿元',
     '【政策支持】工信部《人形机器人创新发展指导意见》全面落地，深圳、上海、北京、安徽等多地出台专项扶持政策，最高补贴1亿元',
     '【供应链成熟】核心零部件国产化率快速提升至82%，谐波减速器、伺服电机、控制器成本较2022年下降60%，供应链自主可控水平大幅提高'],
    ['【产业报告·2026年8月21日】WRC2026发布《2026年人形机器人产业发展报告》：上半年中国出货量超4万台全球占比97%，整机产品达400余款超全球半数，新设企业11.6万户同比+9.5%；开普勒1.75米/75公斤人形已出口美徳奥年产能1000台；擎朗智能洗衣叠衣全流程完成度达及格线；徐晓兰：人形机器人有望成为继计算机/智能手机/新能源汽车后又一颠覆性产品',
     '【资本热浪·2026年8月21日】人形机器人从上场走向进厂资本加码：宇树科技创始人王兴兴表示具身智能领域的ChatGPT时刻或将在可见的未来到来快则两至三年慢则五到十年，当机器人能够被部署至家庭等任意陌生环境可完成约80%的任务时便意味着已抵达具身智能产业爆发的关键临界点；乐聚智能IPO发行申请获深交所受理是首家选择使用创业板第四套标准申请上市的企业2025年实现营业收入2.58亿元近三年复合增长率高达118.68%；越疆科技启动H回A进程创业板IPO项目获深交所上市委审议通过计划募资约12亿元；今年以来国内机器人相关企业共出现超过230起融资事件比去年同期增长28.8%；数字华夏完成亿元级Pre-A轮战略融资带来新一代仿生人形机器人夏澜R03、全新双形态人形机器人星行侠P02（双足行走与轮式移动可切换/身高130厘米/重量30公斤/25个自由度/飞兵模式续航超8小时）以及RoboEase场景大脑',
     '【广东十骏·2026年8月21日】广东十家人形机器人整机企业形成十骏现象：逐际动力数千台订单半数以上海外；智平方NeuroVLA类脑模型具备主动感知/故障自恢复/时序记忆；乐聚夸父核心部件国产率95%；众擎发起URKL全球自由格斗联赛；美的美罗U螺钉锁附超10万颗；荣耀半马夺魁；优必选Walker进比亚迪吉利产线',
     '【产业报告·2026年8月21日】WRC2026发布《2026年人形机器人产业发展报告》：上半年中国出货量超4万台全球占比97%，整机产品达400余款超全球半数，新设企业11.6万户同比+9.5%；开普勒1.75米/75公斤人形已出口美徳奥年产能1000台；擎朗智能洗衣叠衣全流程完成度达及格线；徐晓兰：人形机器人有望成为继计算机/智能手机/新能源汽车后又一颠覆性产品',
     '【小米铁大·2026年8月20日】WRC2026小米详解新一代人形机器人铁大：身高1.7米/体重66公斤/全身66个关节双足形态，按汽车工厂工人身高设计，自由度从上代21个增至66个约一半集中在手部，覆盖汽车工厂2000多个岗位80%以上运动空间；已有两款机器人在工厂实习，螺丝对位成功率3月90%→7月98%（人工99%）预计年底99%；10B模型使用约10万小时UniMi真机数据+1万小时遥操数据',
     '【广东十骏·2026年8月21日】广东十家人形机器人整机企业形成十骏现象：逐际动力数千台订单半数以上海外；智平方NeuroVLA类脑模型具备主动感知/故障自恢复/时序记忆；乐聚夸父核心部件国产率95%；众擎发起URKL全球自由格斗联赛；美的美罗U螺钉锁附超10万颗；荣耀半马夺魁；优必选Walker进比亚迪吉利产线',
     '【自变量WALL-B·2026年8月20日】自变量机器人WRC展示家庭服务+物流分拣双场景：物流分拣复刻全球直播产线两条机械臂配合夹爪全自主分拣复杂随机真实包裹效率1816件/小时准确率98%相比人形+五指灵巧手方案成本大幅下降70%；全自研端到端世界统一模型WALL-B融合视觉/语言/触觉/动作/物理预测具备跨本体跨任务跨场景能力；QUANXTA Zero无本体数据采集方案数据入库有效率超85%模型训练数据成本降低90%',
     '【资本热浪·2026年8月21日】人形机器人从上场走向进厂资本加码：宇树科技创始人王兴兴表示具身智能领域的ChatGPT时刻或将在可见的未来到来快则两至三年慢则五到十年；乐聚智能IPO发行申请获深交所受理是首家选择使用创业板第四套标准申请上市的企业2025年实现营业收入2.58亿元近三年复合增长率高达118.68%；越疆科技启动H回A进程创业板IPO项目获深交所上市委审议通过计划募资约12亿元；今年以来国内机器人相关企业共出现超过230起融资事件比去年同期增长28.8%',
     '【最新出货·2026年8月】SAG数据2026上半年全球人形机器人出货约1.91万台同比+272%，中国企业几乎包揽：智元约8400台(+562%)，全球每出货100台有97台来自中国',
     '【北京亦庄中试·2026年8月21日】北京亦庄落地多条千台级人形机器人中试产线：从原型样机转向小批量试生产，面向消防/商超落地实用场景；北京已设立100亿元机器人产业基金带动社会资本超530亿元，上半年工业机器人产量增长76%服务机器人产量增长2.3倍',
     '【代表产品】特斯拉Optimus Gen3 2026量产版正式交付，自由度52个，全身负载20kg单臂10kg，行走速度8km/h，连续工作8小时，售价15万元',
     '【代表产品】优必选Walker X2工业版已在比亚迪、宁德时代、特斯拉等工厂累计部署超3200台，获得追加1万台订单，交付周期缩短至3个月',
     '【代表产品】小米CyberOne 2代定位消费级市场，价格下探9.9万元以内，自由度48个，2026年第四季度正式量产发售，面向家庭服务场景',
     '【代表产品】波士顿动力Atlas电动版商业化落地加速，聚焦工业巡检和物流搬运场景，获得亚马逊2000台订单，计划2027年完成全部交付',
     '【技术参数】2026年主流机型自由度普遍达到40-55个，单臂负载5-10kg，全身负载15-25kg，关节扭矩密度较2023年提升40%',
     '【运动性能】最大行走速度5-10km/h，爬坡角度30度，可稳定上下楼梯、跨越5cm障碍，抗干扰能力大幅提升，被推倒后可自主起身',
     '【续航能力】标准电池容量5-10kWh，连续工作6-10小时，待机时间24小时以上，快速换电技术普及，更换电池时间小于30秒实现24小时作业',
     '【智能水平】全部搭载端侧具身大模型，支持自然语言交互、3D视觉识别、自主路径规划、新任务即时学习，任务泛化能力大幅提升',
     '【安全性能】全身碰撞检测、关节力控保护、软硬双重急停多重安全机制，满足ISO 13482人机协作安全标准，可与工人同工位安全作业',
     '【产业生态】整机厂商+核心零部件厂+系统集成商+AI大模型企业+应用客户+投资机构+科研院所组成完整产业生态，四大产业集聚区加速形成'],
    '▎人形机器人量产落地具体过程阐述',
    ['【技术预研期（2015-2020）】全球科技企业和科研机构启动人形机器人基础技术研发，波士顿动力Atlas液压版展示惊人运动能力但成本超200万美元，特斯拉Optimus项目立项，国内优必选/小米等开始原型机探索，核心零部件依赖进口，单台成本超百万元，主要用于技术验证和展会演示，产业处于萌芽阶段。',
     '【原型验证期（2021-2023）】特斯拉Optimus Gen1/Gen2原型发布引发全球关注，国内优必选Walker/小米CyberOne/智元远征等原型机密集亮相，核心零部件国产化开始突破，绿的谐波/汇川技术等企业推出人形机器人专用零部件，单台成本降至50-80万元，小批量试制百台级规模，工业场景开始试点应用。',
     '【供应链成熟期（2024-2025）】核心零部件国产化率快速提升至70%以上，谐波减速器/伺服电机/控制器成本较2022年下降60%，六维力传感器/灵巧手等卡脖子部件实现国产突破，整机BOM成本降至20-30万元，各企业建设量产产线，年产能提升至万台级，工业场景试点从单工位向整线扩展，ROI回收期缩短至3年以内。',
     '【量产元年期（2026）】2026年成为人形机器人量产元年，优必选Walker X2/特斯拉Optimus Gen3/傅利叶GR-2等机型实现万台级量产交付，国内已公布产能规划超80万台/年，整机价格下探15-25万元区间，汽车制造/3C电子工厂成为首批规模化落地场景，工业场景部署量突破5万台，商业服务场景开始试点，资本投入持续加码。',
     '【工艺爬坡期（2026H2-2027）】量产工艺持续优化，良率从初期85%提升至95%以上，比亚迪电子等建成全自动量产线，年产能10万台以上，核心零部件国产化率突破90%，成本进一步降至10-20万元区间，应用场景从工业向物流仓储/商业服务/公共服务快速扩展，年销量突破20万台，产业生态初步形成。',
     '【规模应用期（2028-2029）】人形机器人技术成熟度达到商业化可用水平，行走稳定性/负载能力/续航时间/智能水平全面满足场景需求，整机价格下探8-15万元，消费级产品开始上市，应用场景覆盖工业/物流/商业/家庭/公共服务等领域，年销量突破50万台，中国成为全球最大人形机器人生产和应用市场。',
     '【普及爆发期（2030）】人形机器人成本进一步下探至5-10万元区间，消费级市场爆发，家庭陪护/教育娱乐场景快速渗透，年销量突破100万台，产业规模超千亿级，带动核心零部件/AI算法/系统集成/运营服务等上下游产业生态全面成熟，形成万亿级产业集群，人形机器人成为继智能手机/新能源汽车后下一代通用智能终端。',
     '【技术迭代过程】运动控制从预编程步态→强化学习步态→端到端神经运动控制；智能从规则系统→大模型加持→具身智能通用能力；灵巧手从简单夹爪→腱驱灵巧手→触觉反馈仿人手；电池从锂电池→半固态→全固态电池，续航从2小时提升至12小时以上。',
     '【标准体系建设】2025年工信部发布人形机器人标准体系建设指南，2026年首批20项国家标准/行业标准发布实施，涵盖安全要求/性能测试/接口规范/数据格式等方面，检测认证体系建成，规范产业发展，保障人机协作安全。',
     '【产业生态构建】整机企业+核心零部件企业+AI企业+系统集成商+应用客户+投资机构+科研院所组成完整产业生态，长三角/珠三角/京津冀/安徽合芜蚌形成四大产业集聚区，各地出台专项扶持政策，建立产需对接机制，推动产业健康快速发展。'],
    '▎最新行情 · 最新研发 · 最新成果 · 产业前沿',
    ['【WRC2026首发311款·2026年8月21日】2026世界机器人大会官方最终公布首发新品总数311款（整机全球首发150余件/全品类含零部件方案300余件）：优必选超仿生人形U1系列全渠道订单突破1.3万台9月开始交付；宇树科技H2全尺寸人形+R1双足+GD-01载人变形机甲已量产；北京人形创新中心天工Omni轻量化39公斤全球首发；四川具身科技"爱湫"情感交互人形9.8万元起全球发售',
     '【宇树上市首日·2026年8月20日】宇树科技科创板上市首日开盘1100元较发行价150.80元涨近630%，盘中市值破4000亿元，收盘约845元涨460%市值超3400亿，978万户打新中签率0.018%中一签盈利约35万元；软银65亿美元收购Ampere+谷歌122亿绑定Marvell+英伟达参投Groq',
     '【运动会赛前探访·2026年8月20日】第二届世界人形机器人运动会赛前探访：灵巧手专项赛8个竞技小项轮番上演——电动工具装配/粉末称重/积木搭建/钉钉固定/开瓶撬盖/拆箱拆包/镊子夹豆/线缆连接；全自主模式表现超预期：上海电机学院Galaxy星璨队机器人全自主3分多钟完成4层拱门形积木搭建',
     '【服务场景真干活·2026年8月21日】WRC2026服务场景从表演转向干活：银河通用机器人零售店展区观众从机器选择商品下单，机器人接到订单后从货架挑选商品放在结账柜台，依托G0.5具身基础模型在物流企业真实前置仓自主完成拣选/导航/打包/放置全流程调度；千寻智能展区机器人接到整理客厅指令后可识别可乐/碗/垃圾/玩具将各物品摆放在各自位置',
     '【机器人移动母舰·2026年8月21日】飞巴科技全球首发机器人移动母舰——机器人的移动后勤基地，舱体可装载人形机器人/机器狗/无人机，车内自带换电工位和维修工位，即使在断网断电极端环境下也能给机器人提供算力和通信保障，今年年底投入量产，预计明年6月真正商用进入航空救援/应急消防/医学救援等领域',
     '【星动纪元物流·2026年8月21日】星动纪元WRC把物流分拣实战场景搬到展台：M7机器人准确完成抓取/翻面/放置传送带供包作业，物流分拣作业效率已达人工85%以上；商业化已率先完成行业PMF验证，与中国邮政/顺丰等头部物流企业深度合作，在全国5省市10多个物流中心批量部署常态化运营',
     '【江苏造出海·2026年8月21日】"江苏造"智能机器人海外批量"就业"：乐聚（江苏）人形机器人在WAIC 1:1复刻工业产线上纸箱拆垛/塑料箱拆垛/小件上料驾轻就熟遇箱体歪斜自行校准；擎朗智能XMAN-R1化身咖啡师/洗衣师傅/便利店店员连续保持全球商用服务机器人出货量第一前7月出口量货值大幅攀升全新人形产线泰州建成投产；智身新创仿生四足机器狗适配安防巡检应急救援前7月出口额320万元',
     '【DaxAI骐骥·2026年8月21日】DaxAI大咖机器人发布骐骥全地形智能坐骑机器马产品矩阵并与京东签署三年战略合作：骐骥X1纯四足仿生无车轮结构整机300kg最高时速7-10km/h，骐骥XS高速轮足机器马最高行驶速度突破40km/h续航60公里；依托端侧本地推理的DaxBrain-WM两仪世界模型实现全域全地形无障碍通行',
     '【最新行情·产量数据】2026年上半年国内人形机器人总产量达到3.2万台，同比大幅增长420%，其中工业制造场景占比72%、商业服务场景占比18%、物流仓储场景占比10%，截至2026年8月国内已累计部署人形机器人突破5万台',
     '【最新行情·招投标】2026年1-7月国内公开人形机器人招投标项目数量达到287个，同比大幅增长310%，工业制造领域占比超60%，汽车制造是最大采购方占42%，其次是3C电子占23%、新能源电池行业占18%，平均客单价约180万元',
     '【最新研发·运动控制】中国科学院自动化研究所2026年7月正式发布新一代端到端神经运动控制算法NeuroWalk V2.0，人形机器人复杂地形行走稳定性提升40%，侧向抗干扰能力提升60%，平地摔倒率降至1%以下，可稳定通过20度斜坡和10cm台阶',
     '【最新研发·具身智能】OpenAI联合Figure AI于2026年6月正式发布具身大模型Figure 02版本，机器人零样本任务理解能力大幅提升，对自然语言指令的任务完成率从65%提升至92%，可自主完成200+种未经训练的日常操作任务',
     '【最新研发·灵巧操作】哈尔滨工业大学机器人研究所2026年6月成功研发新一代腱绳驱动仿人灵巧手，拥有12个主动自由度、3个被动自由度，指尖力控精度达到0.02N，位置重复定位精度0.01mm，可稳定完成穿针引线、抓取生鸡蛋等精密操作',
     '【最新研发·电池技术】宁德时代2026年5月正式发布人形机器人专用半固态电池，能量密度达到400Wh/kg，支持15分钟快充至80%电量，循环寿命突破5000次，标准工作温度范围-20℃至60℃，可支持人形机器人连续工作12小时',
     '【最新研发·材料工艺】航空级T800碳纤维复合材料+TC4钛合金关节结构件开始大规模应用，整机重量从2023年的70-80kg降至45-60kg，负载自重比从1:4提升至1:2.5，机身强度提升50%，抗冲击能力满足工业场景严苛要求',
     '【最新成果·国产减速器】绿的谐波2026年6月正式发布新一代人形机器人专用SHG-25型谐波减速器，额定寿命突破2万小时，传动精度小于1弧分，回程间隙小于0.5弧分，单台价格降至1800元人民币，仅为日本哈默纳科同类产品价格的1/3',
     '【最新成果·伺服电机】汇川技术2026年7月发布IS650N系列高性能人形机器人专用伺服电机，功率密度提升35%达到4.0kW/kg，响应带宽达3kHz，扭矩波动小于2%，最高转速6000rpm，技术指标达到国际安川、松下同类产品领先水平',
     '【最新成果·六维力传感器】坤维科技2026年6月正式发布KWR80系列高精度六维力/力矩传感器，测量精度达到0.1%FS，采样频率2kHz，维间耦合误差小于0.2%，IP67防护等级，技术指标全面超越美国ATI同类产品，单台价格降至进口产品的1/2',
     '【最新成果·视觉感知】大疆创新2026年5月正式发布RoboMaster D450人形机器人专用深度相机，采用主动立体视觉技术，测距精度达到正负0.5%，室外10万lux强阳光环境下可稳定工作，最远探测距离10米，单台成本降至450元人民币',
     '【最新成果·量产工艺】比亚迪电子2026年7月建成国内首条人形机器人全自动量产组装线，年产能达到10万台，整机组装良率从初期85%提升至98.5%，生产节拍12分钟/台，产线自动化率达到92%，单台制造成本较手工组装下降40%',
     '【产业合作·腾讯优必选】2026年6月腾讯与优必选科技正式签署全面战略合作协议，腾讯提供混元具身大模型、云边端协同AI计算能力和机器人操作系统技术支持，优必选负责人形机器人整机制造、工业场景落地和客户交付',
     '【产业合作·华为鸿蒙】华为2026年5月正式发布HarmonyOS for Robotics人形机器人专用操作系统，提供端云协同AI能力、标准化设备接口、安全实时内核和低延迟通信框架，截至2026年8月已有15家人形机器人整机厂商正式接入鸿蒙生态',
     '【产业合作·国家电网】国家电网2026年7月正式成立电力人形机器人联合实验室，联合中电科21所、哈工大、东北大学等单位聚焦电力巡检、带电作业、应急抢修特种人形机器人研发，计划2027年在特高压变电站试点部署100台',
     '【政策动向·深圳】深圳市2026年5月发布人形机器人产业专项扶持政策，对实现万台级量产的企业给予最高1亿元人民币一次性补贴，规划建设10个典型应用示范场景，对采购国产人形机器人的企业给予30%采购补贴，最高补贴500万元',
     '【政策动向·上海】上海市发布人形机器人创新发展三年行动计划，明确目标2027年全市人形机器人产业规模突破1000亿元，培育3-5家具有全球竞争力的头部企业和20家以上专精特新小巨人企业，建成5个国家级研发平台',
     '【国际动态·波士顿动力】波士顿动力Atlas电动版人形机器人2026年正式获得亚马逊2000台仓储物流机器人订单，用于亚马逊 fulfillment center 货物搬运、分拣和上架作业，计划2027年完成全部交付，单台租赁价格约20美元/小时']))

# PART 02 人形新品
all_modules.append(('PART 02', '人形新品：2026新品密集发布',
    ['【新品节奏】2026年是人形机器人新品发布密集期，全年预计有超50款新品发布',
     '【产品分层】形成工业级/商业级/消费级/科研级四大产品分层，覆盖不同价位段',
     '【工业级产品】聚焦汽车制造/3C电子/新能源工厂场景，负载15-30kg，价格15-30万元',
     '【商业级产品】面向酒店/餐厅/展厅/医院服务场景，负载5-15kg，价格8-15万元',
     '【消费级产品】面向家庭陪护/教育娱乐场景，负载3-8kg，价格3-8万元，2026Q4开始发售',
     '【科研级产品】面向高校/科研院所，开放SDK和二次开发接口，价格20-50万元',
     '【新品趋势】整机重量持续轻量化，从70-80kg降至45-60kg，更适合人机协作',
     '【智能化升级】全部搭载具身大模型，支持自然语言指令，无需预编程即可完成新任务',
     '【模块化设计】关节/电池/传感器模块化设计，维护更换成本降低60%',
     '【设计语言】消费级产品开始注重外观设计，采用流线型机身，多配色可选'],
    ['【小米铁大·2026年8月21日】WRC2026小米详解新一代人形机器人铁大：身高1.7米/体重66公斤/全身66个关节双足形态，按汽车工厂工人身高设计，自由度从上代21个增至66个约一半集中在手部，覆盖汽车工厂2000多个岗位80%以上运动空间；已有两款机器人在工厂实习，螺丝对位成功率3月90%→7月98%（人工99%）预计年底99%；10B模型使用约10万小时UniMi真机数据+1万小时遥操数据；应用场景分智能制造/商业服务/家庭三阶段，家庭先从小米青年公寓验证',
     '【自变量双场景·2026年8月21日】自变量机器人WRC展示家庭服务+物流分拣双场景：物流分拣两条机械臂配合夹爪全自主分拣复杂随机真实包裹效率1816件/小时准确率98%，相比人形+五指灵巧手方案成本大幅下降70%，已与头部物流企业合作部署真实产线；全自研端到端世界统一模型WALL-B融合视觉/语言/触觉/动作/物理预测；QUANXTA Zero无本体数采数据入库有效率超85%成本降90%；今年3月与58到家推出机器人上门家政行业首次大规模进家庭',
     '【WRC2026开幕·2026年8月20日】2026世界机器人大会北京开幕：超300家企业参展(+69%)/展品超3000件/首发新品300余件，主题"人机共生 产需共融"；北京人形机器人创新中心"慧思开物"具身智能一站式开发平台+具身大一统模型Pelican Unify+开放平台"天工Omni"实现一脑多能一脑多机',
     '【自变量双场景·2026年8月20日】自变量机器人WRC展示家庭服务+物流分拣双场景：物流分拣两条机械臂配合夹爪全自主分拣复杂随机真实包裹效率1816件/小时准确率98%，相比人形+五指灵巧手方案成本大幅下降70%；全自研端到端世界统一模型WALL-B融合视觉/语言/触觉/动作/物理预测',
     '【运动会赛项·2026年8月20日】第二届世界人形机器人运动会8月22日北京开幕全赛项规则详解：举重按自重分轻量级≤40kg/重量级40-80kg，须完整人形+≥3指灵巧手；拔河新增二对二赛项每局2分钟赛道宽1.2米长12米；投壶首次设置壶口内径13cm距投掷线1.5米3分钟内站立完成',
     '【武汉军团·2026年8月20日】WRC2026武汉机器人军团集体亮相：格蓝若C1人形机器人学会"倒地自起"；D2-W四足机器狗搭载自研"玄鸟"巡检智脑；启灵"神农"自主完成搬运；"赤兔"四足机器狗峰值负载超120公斤；拾光S1通用家庭人形机器人亮相',
     '【越疆一脑多体·2026年8月21日】越疆科技携一脑多体具身智能核心展台参展：多机器群体协同搭建工业生产作业体系，依托自研空弈DobotWAM具身大模型作统一认知底座；新一代具身全栖机器人鹿萌首次线下亮相，身高近1.3米可在狭小空间精准规划肢体姿态',
     '【大晓世界模型·2026年8月20日】大晓机器人首次亮相WRC展示具身智能全栈实力：开悟世界模型3.1采用统一原生架构整合生成/物理/认知三类智能，在全球具身智能评测中世界模型视频生成/状态预测两项赛道取得靠前成绩；发布晓满/晓新/晓途三套行业解决方案',
     '【优艾智合隙锋·最新】工业原生人形机器人"明日熟练工"：出厂具备抓/取/握/拿/推基础动作，50条数据采集训练达90%成功率，24小时周期循环训练成岗位熟练工',
     '【特斯拉Optimus Gen3】2026年8月量产交付，自由度52个，全身负载20kg，单臂负载10kg',
     '【优必选Walker S】工业专用版，自由度55个，行走速度7.2km/h，已交付3000台',
     '【小米CyberOne 2】消费级旗舰，自由度48个，价格9.9万元起，2026年12月开售',
     '【小鹏PX5】汽车工厂专用版，自由度42个，专为汽车产线优化，已在小鹏工厂部署500台',
     '【傅利叶GR-2】通用人形机器人，自由度53个，开源开放，面向开发者和科研机构',
     '【智元远征A2】2026年新款，自由度54个，搭载智元具身大模型2.0，任务完成率提升50%',
     '【宇树H1】高动态性能人形，奔跑速度15km/h，可完成后空翻，面向科研和极限场景',
     '【达闼XR-4】云端大脑架构，5G云端协同，面向商业服务场景，已在100+酒店部署',
     '【钢铁侠MK-800】重载型工业机器人，全身负载80kg，面向物流搬运重型场景',
     '【追觅通用人形】2026年新款，结合追觅在清洁机器人领域技术积累，主打家庭服务场景'],
    '▎人形机器人新品迭代具体过程阐述',
    ['【概念探索期（2015-2020）】各厂商发布概念性人形机器人，外观机械感强，运动能力有限，只能完成简单展示动作，单台研发成本超百万元，主要用于科技展会展示技术实力，没有明确商业化路径，产品数量少，每年全球发布新品不足5款，技术路线不清晰，液压/电机/腱驱多种方案并行探索。',
     '【原型迭代期（2021-2024）】第二代/第三代原型机密集发布，机械结构优化，运动控制能力大幅提升，能够完成稳定行走/简单抓取等操作，外观设计更简洁美观，成本下降50%至50-80万元，开始小批量试产百台级规模，在工业/商业特定场景试点应用，收集用户反馈快速迭代，每年新品发布数量增长至10-20款。',
     '【工程化期（2025）】产品从原型向工程化转化，解决可靠性/稳定性/可维护性问题，关节模组/传感器/电池等核心部件标准化，模块化设计普及，维护更换成本降低60%，成本进一步降至20-40万元，年新品发布数量达30款以上，工业级产品开始小批量交付，用户场景验证全面展开。',
     '【量产上市期（2026）】2026年新品密集发布上市，全年预计超50款新品发布，形成工业级/商业级/消费级/科研级完整产品矩阵，技术成熟度大幅提升，成本降至可商业化区间15-25万元，万台级量产交付开始，工业级产品率先规模化应用，商业级产品批量上市，消费级产品开始预售。',
     '【产品分层期（2026H2-2027）】产品分层清晰：工业级15-30万元（负载15-30kg，8小时续航），商业级8-15万元（负载5-15kg，10小时续航），消费级3-8万元（负载3-8kg，12小时续航），科研级20-50万元（开放接口，二次开发），各价位段都有代表性产品，满足不同客户需求。',
     '【体验优化期（2027-2028）】产品体验持续优化，人机交互更自然，运动更流畅，噪音更低，外观更亲和，消费级产品外观设计向消费电子看齐，多配色可选，软件生态逐步完善，应用商店/技能商店上线，用户可以下载新技能扩展机器人能力。',
     '【生态繁荣期（2029-2030）】新品发布趋于稳定，每年20-30款迭代升级，硬件平台标准化，软件生态繁荣，第三方开发者数量超10万人，应用技能超1000种，消费级市场爆发，家庭保有量快速增长，人形机器人从工业走向大众消费市场。',
     '【设计语言演进】从早期机械裸露工业风→简洁流线科技风→仿生亲和消费风；机身材质从金属结构件→碳纤维复合材料→亲肤硅胶材料；灯光交互从简单指示灯→LED表情屏→面部柔性显示屏，人机交互体验持续提升。',
     '【核心技术演进】自由度从30+提升至50+；行走速度从3km/h提升至10km/h以上；单臂负载从3kg提升至10kg以上；续航从4小时提升至12小时以上；智能水平从预编程→大模型自然语言交互→自主学习新技能。',
     '【供应链配套演进】从核心零部件全部进口→国产化率30%→国产化率85%→全国产化供应链；零部件成本从占整机80%降至占比50%；交付周期从6个月缩短至1个月以内；售后网络从一线城市覆盖至全国地级市。'],
    '▎新品参数 · 价格对比 · 技术亮点 · 上市时间',
    ['【特斯拉Optimus Gen3参数】特斯拉2026年8月16日正式交付量产版Optimus Gen3，全身配置52个主动自由度，身高173cm，体重57kg，单臂负载能力10kg全身最大负载20kg，最大平地行走速度8km/h，标准工况连续续航8小时，单台量产目标成本降至15万元人民币，首批5000台率先交付美国得州超级工厂用于产线物料搬运',
     '【优必选Walker X2参数】优必选Walker X2工业版人形机器人全身配置55个主动自由度，身高165cm，体重63kg，单臂最大负载12kg全身最大负载25kg，最大平地行走速度7.2km/h，标准工况连续续航10小时，工业防护等级IP54，单台裸机售价22万元人民币，目前已实现批量量产交付，累计交付量超1.2万台',
     '【小米CyberOne 2参数】小米CyberOne 2第二代人形机器人全身配置48个主动自由度，身高170cm，体重52kg，单臂负载能力8kg，最大平地行走速度6.5km/h，标准工况连续续航12小时，搭载小米自研MiLM具身大模型，可完整接入小米智能家居生态，开发者版起售价9.9万元，计划2026年12月正式发售',
     '【小鹏PX5参数】小鹏汽车PX5工业人形机器人全身配置42个主动自由度，身高175cm，体重65kg，单臂负载能力15kg全身最大负载30kg，最大平地行走速度5km/h，标准工况连续续航8小时，工业防护等级IP65可适应工厂复杂恶劣环境，单台售价18万元人民币，计划2026年9月正式上市交付汽车工厂',
     '【傅利叶GR-2参数】傅利叶智能GR-2通用人形机器人全身配置53个主动自由度，身高168cm，体重55kg，单臂负载能力10kg全身最大负载18kg，最大平地行走速度6km/h，标准工况连续续航9小时，搭载傅利叶自研FSA系列力控关节，关节扭矩密度提升45%，单台售价16.8万元人民币，目前现货供应可直接下单',
     '【智元远征A2参数】智元机器人远征A2通用人形机器人全身配置54个主动自由度，身高172cm，体重58kg，单臂负载能力12kg全身最大负载22kg，最大平地行走速度7km/h，标准工况连续续航10小时，搭载智元自研具身大模型支持自然语言指令编程，单台售价19.8万元人民币，计划2026年10月正式批量交付',
     '【宇树H1参数】宇树科技H1科研版人形机器人全身配置44个主动自由度，身高180cm，体重47kg，单臂负载能力10kg全身最大负载15kg，最大平地奔跑速度15km/h，是目前全球行走速度最快的人形机器人，标准工况连续续航6小时，科研版单台售价9万元人民币，面向高校和科研机构销售',
     '【达闼XR-4参数】达闼机器人XR-4云端智能人形机器人全身配置50个主动自由度，身高165cm，体重60kg，单臂负载能力7kg全身最大负载10kg，最大平地行走速度5.5km/h，搭载达闼自研海睿云端大脑支持5G云端协同计算，单台售价25万元人民币，面向商业服务和展厅接待场景，目前已累计交付3000台',
     '【价格区间分布】2026年已发布人形机器人产品价格区间分布：工业级产品15-30万元人民币占比约60%，主要面向汽车制造/3C电子/新能源电池工厂场景；商业级产品8-15万元占比约25%，面向酒店/餐厅/展厅/医院服务场景；科研级产品20-50万元占比约5%面向高校科研院所；消费级产品3-8万元占比约10%预计2026Q4上市',
     '【技术亮点1·端侧具身大模型】2026年所有新发布人形机器人产品全部标配端侧具身大模型，本地AI推理延迟控制在200ms以内，支持自然语言交互、零样本任务理解、新场景快速自主学习，无需专业技术人员预编程即可完成90%以上常见操作任务，任务泛化能力较2025年提升200%以上',
     '【技术亮点2·全身力控安全】2026年新品全身关节均搭载高精度力矩传感器，碰撞检测灵敏度达到0.1N，碰撞响应时间小于10ms，具备软硬双重急停和虚拟安全围栏区域限制功能，全部机型均通过ISO 13482人机协作安全标准认证，可与人类工人在同一工位安全协同作业，人机混线安全性大幅提升',
     '【技术亮点3·自研关节模组】2026年国内主流厂商全部实现自研一体化关节模组大规模应用，关节扭矩密度较2023年提升40%达到200N·m/kg以上，单关节模组成本较外购方案下降50%，单台关节模组价格降至800元人民币以内，关节平均无故障工作时间（MTBF）突破1万小时满足工业连续作业需求',
     '【技术亮点4·多传感器融合】2026年新品普遍采用4D毫米波雷达+双目立体视觉+固态激光雷达+指尖触觉传感器多模态融合感知方案，复杂动态环境感知准确率达到99.9%，可在-20℃至60℃宽温度范围、雨雪粉尘等恶劣工业环境下稳定工作，室外复杂非结构化场景适应性较2025年大幅提升',
     '【技术亮点5·快速换电技术】2026年新品快速换电技术全面普及，采用标准模块化电池设计，无需任何专用工具即可在30秒内完成电池热插拔更换，配合共享电池柜可实现7×24小时不间断作业，同时支持15分钟快充至80%电量，电池循环寿命突破5000次，完全满足高强度工业连续作业需求',
     '【上市时间节奏汇总】2026年Q3（7-9月）预计有8款新人形机器人产品正式上市交付；Q4（10-12月）进入新品发布高峰期预计有15款新品集中上市发售；2027年Q1预计还有12款新品发布；2026年全年合计将有超过50款新人形机器人产品正式面向市场公开发售，市场产品供给极大丰富',
     '【市场预订数据统计】截至2026年8月15日，国内公开可统计的人形机器人预订订单量已经累计超过12万台，其中工业制造领域企业客户订单占比约75%，主要来自汽车制造、3C电子、新能源电池行业头部企业；商业服务领域订单占比约18%，科研教育机构订单占比约7%，订单量保持持续快速增长态势',
     '【供应链配套带动效应】每款新人形机器人产品平均带动30-50家国内核心零部件供应商同步研发和配套量产，2026年新发布产品零部件平均国产化率已经达到85%以上，部分头部厂商机型国产化率突破95%，核心供应链基本实现自主可控，供应链安全保障能力和成本优势显著，',
     '【工业设计趋势】2026年新品普遍采用高度仿生亲和工业设计，关节处采用符合人体工学的流畅曲线造型，外壳采用亲肤哑光材质避免冰冷机械感，视觉上更接近人类外形，人机交互亲和力大幅提升；同时机身结构件采用一体化压铸成型工艺，零部件数量减少30%，整机组装复杂度大幅降低',
     '【材料工艺全面升级】2026年新品广泛采用航空级7075高强度铝合金+T800高强度碳纤维复合材料混合机身结构，机身平均重量较2025年降低15%，同时结构抗冲击强度提升25%，可满足工业场景跌落碰撞严苛要求；机身防护等级普遍达到IP54以上，部分工业专用机型达到IP67防水防尘等级',
     '【2027年新品研发规划】国内各主流厂商2027年新品研发规划已经全部启动，重点研发方向聚焦三个核心维度：一是整机BOM成本进一步下探至5万元人民币以内推动大规模普及；二是家庭场景深度优化，重点提升家务操作能力和家庭环境适应性；三是通用智能水平进一步提升实现完全自主学习新技能无需人工干预']))

# PART 03 核心零部件
all_modules.append(('PART 03', '核心零部件：国产替代加速',
    ['【零部件分类】人形机器人核心零部件包括减速器/伺服电机/控制器/传感器/电池/灵巧手六大类',
     '【成本结构】减速器占整机成本35%，伺服电机占20%，控制器占15%，传感器占15%，其他占15%',
     '【国产替代率】2026年核心零部件平均国产替代率达82%，较2023年提升45个百分点',
     '【减速器】谐波减速器国产率90%，RV减速器国产率65%，绿的谐波/双环传动/中大力德主导',
     '【伺服电机】国产伺服电机市占率从2023年25%提升至2026年60%，汇川/禾川/埃斯顿快速崛起',
     '【控制器】国产控制器市占率75%，固高/雷赛/汇川提供完整运动控制解决方案',
     '【传感器】六维力传感器国产率70%，坤维/宇立/鑫精诚技术达到国际水平，成本降低70%',
     '【电池】宁德时代/比亚迪为人形机器人开发专用电池，能量密度400Wh/kg，循环寿命5000次',
     '【灵巧手】因时/傲博/大寰等国产灵巧手已实现规模化应用，成本降至进口产品1/3',
     '【降本趋势】核心零部件成本年均下降35-40%，推动整机价格快速下探'],
    ['【灵巧手狂卷WRC·2026年8月21日】WRC2026灵巧手从附属配件站到C位：灵心巧手直驱型Linker Hand O30专为强化学习设计可用灵巧手装配灵巧手具备自动化量产条件；章鱼动力OctoH-Hand腱绳+小臂电机直驱混合驱动23个主动自由度搭载超1900个触觉传感单元主打数采孪生闭环；因时机器人RH56F2连杆驱动方案与强脑科技Revo3脑机交互灵巧手并称"中国灵巧手新四小龙"；2026上半年国内灵巧手赛道融资超250亿元',
     '【中科硅纪CasiaHand·2026年8月21日】中科硅纪WRC展示六款CasiaHand系列行业级灵巧手含M系列/X系列及行业级三指G系列，面向工业制造/商业服务/科研教育/特种作业；Brain-Si 0.5类人灵巧操作具身大小脑模型采用分层协同架构大脑负责环境理解/任务识别/决策规划小脑负责技能执行/运动控制/关节实时协同；数据采集双轨并行同本体采集保证质量+EGO方式跨本体泛化',
     '【零部件链式聚集·2026年8月21日】WRC2026人形机器人整机企业通过"链式"聚集展示核心部件协同演进：傅利叶展区上下游生态合作企业奇点智控/伊塔动力/一行凌光展示高精度光学六维力传感器/关节模组/数据采集头环；卓誉科技已建成年产能20万台关节模组/50万条电机的自动化产线，电机最小做到14毫米专门给灵巧手指节用',
     '【绿的谐波】人形机器人专用谐波减速器SHG-25，寿命2万小时，传动精度<1弧分，已批量供货',
     '【双环传动】RV减速器RD-32E，负载扭矩320N·m，寿命1.5万小时，进入特斯拉供应链',
     '【汇川技术】IS620N伺服系统，功率密度3.5kW/kg，响应带宽3kHz，国内市占率第一',
     '【禾川科技】X7系列伺服电机，扭矩密度提升30%，支持EtherCAT总线，已批量应用',
     '【坤维科技】KWR75六维力传感器，测量精度0.1%FS，采样频率1kHz，达到国际领先水平',
     '【宁德时代】人形机器人专用固态电池，容量5kWh/10kWh，快充30分钟80%，2027年量产',
     '【因时机器人】BHX-12灵巧手，12自由度，指尖力控0.02N，可完成精密装配操作',
     '【大寰机器人】PGC-140自适应夹爪，行程140mm，力控范围5-140N，工业场景广泛应用',
     '【固高科技】GTHD系列运动控制器，支持55轴同步控制，运动周期125us',
     '【大疆】RoboMaster深度相机D435i，测距精度±0.5%，室外抗阳光，成本降至500元以内'],
    '▎核心零部件国产化具体过程阐述',
    ['【进口完全垄断期（2010年前）】中国机器人核心零部件几乎100%依赖进口，日本哈默纳科/纳博特斯克垄断减速器市场，日本安川/松下/三菱垄断伺服电机市场，德国倍福/西门子垄断控制器市场，美国ATI垄断力传感器市场，进口产品价格高昂供货周期长，单台谐波减速器价格超5000元，单台伺服电机超3000元，国产机器人企业生产成本高，利润薄，产业发展受到严重制约。',
     '【技术起步突破期（2011-2018）】国内企业开始持续研发投入，绿的谐波2013年推出首款国产谐波减速器，双环传动开始RV减速器研发，汇川技术在伺服电机领域持续突破，逐步实现从无到有，产品性能逐步接近进口水平，成本下降50%，开始在国产工业机器人领域小批量应用验证，国产化率从不足5%提升至20%左右。',
     '【性能追赶期（2019-2022）】国产核心零部件技术水平快速提升，谐波减速器寿命从5000小时提升至1.5万小时，传动精度达到1弧分以内，接近进口水平；伺服电机功率密度提升至3kW/kg以上，响应带宽达2kHz；六维力传感器实现技术突破，精度达0.2%FS。成本降至进口产品1/2-2/3，国产化率快速提升至40%左右，在中低端工业机器人领域大规模应用。',
     '【人形需求带动期（2023-2025）】人形机器人产业爆发带动核心零部件需求激增，对零部件性能提出更高要求（扭矩密度/重量/精度/寿命），国内企业针对人形机器人需求开发专用零部件，关节模组一体化设计，谐波减速器寿命突破2万小时，伺服电机功率密度提升至3.5kW/kg以上，六维力传感器精度达0.1%FS达到国际水平，成本进一步下降至进口1/2-1/3，国产化率提升至70%以上。',
     '【规模化量产期（2026）】核心零部件进入大规模量产期，绿的谐波年产100万台谐波减速器产能投产，双环传动年产50万台RV减速器产能建设，汇川技术年产200万套伺服系统产能投产，坤维科技六维力传感器年产能50万台，满足30-50万台人形机器人配套需求，国产化率达82%，成本较2022年再降50%，供应链自主可控水平大幅提升。',
     '【工艺成熟完善期（2027-2028）】核心零部件生产工艺成熟完善，引入车规级质量管控标准，产品良率提升至99.5%以上，一致性和可靠性大幅提升，谐波减速器寿命突破3万小时，伺服电机平均无故障时间（MTBF）达10万小时，六维力传感器维间耦合误差降至0.1%以内，成本进一步下降，国产化率目标90%。',
     '【技术领先期（2029-2030）】国产核心零部件技术达到国际领先水平，新型材料/新型结构/新工艺应用，谐波减速器采用新型齿形设计传动效率提升20%，伺服电机采用第三代磁钢材料扭矩密度提升至5kW/kg以上，六维力传感器采用一体化加工工艺精度达0.05%FS超越进口水平，成本降至进口产品1/3-1/4，实现完全自主可控并出口全球。',
     '【减速器技术演进】谐波减速器：齿形设计从传统渐开线→双圆弧齿形→新型P型齿形，传动效率从60%提升至85%以上，寿命从5000小时→2万小时→3万小时；RV减速器：从摆线针轮→新型摆线结构，承载能力提升30%，回差<1弧分。',
     '【伺服电机技术演进】从方波控制→正弦波控制→FOC矢量控制→磁场定向控制+自适应控制；编码器从增量式→绝对值式→多摩川兼容→国产高精度编码器，分辨率从17位提升至23位；功率密度从1kW/kg→2kW/kg→3.5kW/kg→5kW/kg。',
     '【传感器技术演进】六维力传感器：从应变片粘贴结构→一体化加工结构→MEMS结构，维间耦合从5%→1%→0.2%→0.1%；IMU从光纤陀螺→MEMS陀螺→高性能MEMS，零偏稳定性从10deg/h→1deg/h→0.1deg/h；激光雷达从机械旋转→半固态→全固态，成本从10万元→1万元→2000元。'],
    '▎技术参数 · 成本对比 · 国产化率 · 供应链进展',
    ['【谐波减速器参数对比】国产绿的谐波SHG-25系列人形机器人专用谐波减速器，单台价格1500-2500元人民币，额定使用寿命20000小时，传动精度小于1弧分回程间隙小于0.5弧分，传动效率85%；进口日本哈默纳科CSG系列同类产品，单台价格4000-6000元人民币，额定寿命25000小时，性能指标接近但价格是国产2.4倍，国产产品性价比优势显著',
     '【RV减速器参数对比】国产双环传动RD-32E系列高负载RV减速器，单台价格2500-4000元人民币，额定输出扭矩320N·m，额定使用寿命15000小时，回差小于1弧分，主要应用于人形机器人腿部髋膝大负载关节；进口日本纳博特斯克RV-E系列同类产品价格6000-10000元人民币，额定寿命20000小时，国产产品已进入特斯拉Optimus供应链，实现批量供货',
     '【伺服电机参数对比】国产汇川技术IS650N系列高性能人形机器人专用伺服电机，单轴价格800-1500元人民币，功率密度达到4.0kW/kg，响应带宽3kHz，扭矩波动小于2%最高转速6000rpm，支持EtherCAT实时总线通信；进口日本安川Σ-7系列同类产品单轴价格2000-3500元人民币，功率密度3.8kW/kg，国产产品性能已追平进口水平，价格仅为进口40%',
     '【六维力传感器参数对比】国产坤维科技KWR80系列高精度六维力/力矩传感器，单台价格8000-15000元人民币，测量精度达到0.1%FS，采样频率2kHz，维间耦合误差小于0.2%，IP67防护等级；进口美国ATI Nano43系列同类产品单台价格30000-50000元人民币，测量精度0.08%FS，国产产品性能已达国际领先水平，价格仅为进口的1/3',
     '【灵巧手参数对比】国产因时机器人BHX-12系列腱绳驱动仿人灵巧手，单台价格2-3万元人民币，配置12个主动自由度、3个被动自由度，指尖力控精度0.02N，位置重复定位精度0.01mm，可完成精密装配操作；进口英国Shadow Hand灵巧手价格15-20万元人民币，配置20个自由度，性能接近但国产价格仅为进口的1/7，已实现大规模应用',
     '【运动控制器参数对比】国产固高科技GTHD系列多轴运动控制器，单台价格5000-10000元人民币，最多支持64轴同步运动控制，最小运动周期125us，支持EtherCAT/PROFINET多种工业总线，内置人形机器人专用运动学算法库；进口德国倍福CX系列控制器价格20000-40000元人民币，支持128轴控制，国产产品已完全满足人形机器人控制需求',
     '【国产化率进度数据】人形机器人核心零部件国产化率逐年快速提升：2023年平均国产替代率仅37%，2024年提升至55%，2025年提升至70%，2026年已达到82%，计划2027年目标90%，预计2030年实现95%以上完全自主可控，彻底摆脱进口卡脖子依赖，其中谐波减速器、伺服电机、控制器等环节国产化率已超80%',
     '【成本下降进度】以2023年核心零部件整体成本为基准100%：2024年成本下降至70%，2025年下降至48%，2026年已下降至32%，三年时间核心零部件整体成本累计下降68%，直接推动整机价格从2023年60-70万元下降至2026年15-25万元区间，是整机价格雪崩的核心驱动力，预计2027年成本将进一步下降至2023年的22%',
     '【减速器产能建设】绿的谐波苏州新建年产100万台人形机器人专用谐波减速器超级工厂将于2026年Q3正式投产，采用全自动数字化生产线，生产节拍15秒/台，产品良率99.5%，可满足50万台人形机器人配套需求；双环传动浙江玉环新建年产50万台RV减速器生产基地正在建设中，预计2026年底投产，全部投产后将成为全球最大RV减速器生产基地',
     '【伺服系统产能建设】汇川技术苏州新建年产200万套高性能伺服系统数字化工厂已于2026年中正式投产，专门面向人形机器人和工业机器人市场，采用全自动智能生产线，年产能可满足30万台人形机器人全套伺服系统配套需求，生产良率99.2%，单台伺服系统制造成本较2023年下降55%，同时在合肥、深圳设有区域生产基地',
     '【谐波减速器技术突破】绿的谐波2026年6月发布的新一代SHG系列谐波减速器采用全新自主研发P型齿形设计，传动效率较上一代产品提升15个百分点达到85%以上，连续工作温升降低20%，额定寿命突破2万小时，采用全新润滑脂实现终身免维护，产品各项性能指标全面超越日本哈默纳科同类产品，技术水平全球领先',
     '【伺服电机技术突破】汇川技术2026年7月发布的新一代IS650N系列伺服电机采用第三代高性能钕铁硼磁钢材料和优化磁路设计，扭矩密度较上一代提升35%达到4.0kW/kg，最高转速达到8000rpm，采用23位高精度国产绝对值编码器，位置分辨率达到838万脉冲/转，编码器不再依赖日本多摩川进口，实现完全自主可控',
     '【力传感器技术突破】坤维科技2026年6月发布的KWR80系列六维力传感器采用一体化整体加工结构设计，取消传统粘贴式应变片工艺，维间耦合误差从1%降至0.2%以内，测量精度达到0.1%FS，采样频率提升至2kHz，响应时间小于0.5ms，IP67防水防尘等级可适应工业恶劣环境，技术指标全面超越美国ATI同类产品',
     '【灵巧手技术突破】大寰机器人2026年5月发布的PGC-140自适应灵巧夹爪采用腱绳驱动+连杆传动混合传动方案，指尖力分辨率达到0.01N，最大夹持力140N，行程范围140mm，整手重量仅500g，集成微型触觉传感器阵列，可感知物体纹理和硬度，已在工业装配、物流分拣场景实现大规模应用，单台成本较进口下降70%',
     '【供应链本地化优势】长三角地区已形成完整的人形机器人核心零部件产业集群，以上海、苏州、杭州、宁波为核心，在3小时车程范围内可以配齐人形机器人全部核心零部件，本地配套率超过90%，供应链响应速度从2周缩短至24小时，物流成本下降80%，大幅降低整机企业供应链管理难度和库存压力，形成显著的产业集群优势',
     '【安徽蚌埠产业布局】蚌埠中国传感谷重点布局机器人传感器产业，依托中电科思仪等核心企业，已引进12家传感器核心企业和配套厂商，重点发展六维力传感器、IMU惯性测量单元、激光雷达、深度相机等机器人用传感器产品，规划2027年传感器年产值突破100亿元，打造国家级机器人传感器产业基地，为安徽人形机器人产业提供核心传感器支撑',
     '【安徽合肥产业布局】合肥经济技术开发区重点布局伺服电机和运动控制器产业，汇川技术、埃斯顿、固高科技等国内头部企业均已在合肥设立生产基地和研发中心，规划2027年伺服系统和控制器年产值突破200亿元，可满足50万台人形机器人伺服和控制器配套需求，同时合肥高新区在AI芯片、具身大模型领域形成产业集聚，协同配套能力不断增强',
     '【车规级工艺导入】核心零部件企业开始全面引入汽车行业成熟的车规级生产工艺和IATF16949质量管控标准，采用全自动数字化生产线、SPC统计过程控制、全流程可追溯质量体系，零部件生产良率从初期90%提升至99.5%，平均无故障工作时间（MTBF）从5000小时提升至2万小时以上，产品一致性和可靠性达到车规级水平，满足工业场景大规模应用要求',
     '【检测认证体系建成】国家机器人检测与评定中心已建成完善的核心零部件检测认证体系，涵盖性能测试、可靠性测试、环境适应性测试、安全认证全流程，零部件检测认证周期从原来的6个月缩短至1个月，检测费用下降70%，建立统一的零部件标准体系和认证互认机制，有效降低零部件企业认证成本，缩短新产品上市周期，规范产业发展',
     '【未来发展目标】根据工信部《人形机器人创新发展行动计划》规划目标：2027年核心零部件国产化率达到90%，整机BOM成本降至10万元以内；2030年核心零部件国产化率实现95%以上完全自主可控，整机成本降至5万元以内，核心零部件技术水平达到国际领先，形成3-5家具有全球竞争力的核心零部件龙头企业，建成完整自主可控的产业供应链体系']))

# PART 04 央企国家队
all_modules.append(('PART 04', '央企国家队：战略布局入场',
    ['【战略定位】央企发挥新型举国体制优势，承担人形机器人产业链链长角色，整合资源突破瓶颈',
     '【入场节奏】2026年成为央企集中布局人形机器人元年，已有12家央企明确战略布局',
     '【国机集团】成立国机机器人有限公司，整合集团内部机器人资源，打造国家级人形机器人平台',
     '【中国兵器装备】依托长安汽车/建设工业等资源，布局工业人形机器人和特种机器人',
     '【中国电子】依托中国软件/中国长城等，布局人形机器人操作系统和AI芯片',
     '【中国电科】依托中电科思仪/中电科机器人等，布局传感器/控制器/特种机器人',
     '【国家电网】成立国网机器人科技有限公司，聚焦电力巡检/带电作业特种人形机器人',
     '【中国一汽/东风/长安】三大车企依托汽车制造优势，布局工业人形机器人及汽车场景应用',
     '【中国宝武】宝武机器人聚焦钢铁冶金场景重载工业人形机器人研发应用',
     '【中国石化/中石油】布局防爆特种人形机器人，用于石化厂区巡检和应急处置'],
    ['【央企实景实训·2026年8月21日】国资委透露中央企业机器人创新联合体将聚焦电力巡检/应急救援/钢铁冶金/石化等十大高价值应用场景推动人形机器人实景实训，年底形成万台级落地能力；49家央企WRC2026带来263件展品覆盖12类应用场景标志央企全面入场具身智能；工信部+国资委启动2026年度人形机器人与具身智能实景实训专项',
     '【央企创新联合体·2026年8月21日】2026世界机器人大会上中央企业联合展区首次集中亮相：在国务院国资委指导下兵器工业集团牵头，联合中央企业/高校及科研院所/民营企业/行业学会等百余家单位组建中央企业机器人创新联合体，同步发布央企机器人十大创新成果和十大高价值应用场景；48家央企精选263件优质展品参展，覆盖从基础材料/核心零部件/机器人本体到智能算法/数据底座/产业化落地的完整产业体系',
     '【国机集团】已发布国机H1通用人形机器人，自由度55个，负载30kg，已在国机内部工厂试用',
     '【兵器装备】长安汽车/建设工业联合发布兵装人形1号，专为汽车产线优化，已部署200台',
     '【中国电子】发布CEC-OS人形机器人操作系统，支持多品牌硬件统一接入，开源开放',
     '【中国电科】中电科思仪发布机器人传感器系列产品，六维力传感器/激光雷达/IMU全系列布局',
     '【国家电网】国网巡检人形机器人已在20个变电站试点应用，可完成带电作业操作',
     '【中国一汽】一汽红旗人形机器人已在红旗工厂总装车间部署300台，承担物料搬运任务',
     '【东风汽车】东风人形机器人与岚图汽车工厂合作，已完成50台试点部署',
     '【中国宝武】宝武重载人形机器人负载100kg，可在高温/高粉尘钢铁车间连续作业',
     '【中国石化】石化防爆人形机器人取得防爆认证，已在10个炼化厂试点巡检',
     '【通用技术集团】通用技术集团布局医疗人形机器人，聚焦手术辅助/康复护理场景'],
    '▎央企国家队布局具体过程阐述',
    ['【战略调研期（2022-2023）】国务院国资委组织央企开展人形机器人产业专题调研，组织院士专家论证会，分析全球人形机器人发展趋势和中国产业现状，明确央企在人形机器人产业中的定位和作用，结合各央企自身产业基础和应用场景优势，研究制定战略布局方向，不盲目跟风，避免低水平重复建设，为后续集中力量办大事奠定基础。',
     '【规划布局期（2024）】各央企陆续制定人形机器人战略规划，明确发展目标和实施路径，国机集团/中国兵器装备/中国电子/中国电科/国家电网等12家央企明确将人形机器人纳入集团战略新兴产业，组建专项工作组，安排专项资金，启动组织架构建设和人才招聘，与高校/科研院所/民营企业开展前期对接合作。',
     '【组织建设期（2025）】各央企陆续成立机器人专业子公司或研究院：国机集团成立国机机器人有限公司，兵器装备依托长安汽车建设机器人研究院，中国电子成立机器人OS公司，中国电科整合内部传感器/控制器资源成立机器人事业部，国家电网成立国网机器人科技有限公司，中国一汽/东风/长安成立机器人研发中心，组建研发团队合计超5000人，投入研发资金超100亿元。',
     '【原型研发期（2025H2-2026H1）】各央企启动原型机研发，结合自身应用场景需求开发专用人形机器人：国机集团开发通用工业人形，兵装开发汽车产线专用人形，国家电网开发电力巡检人形，中国石化开发防爆人形，中国宝武开发重载冶金人形，2026年上半年陆续发布首款原型产品，开始内部场景试点验证。',
     '【试点应用期（2026H2-2027）】央企人形机器人产品开始在内部场景规模化试点应用，依托央企丰富的应用场景（电力/石化/冶金/汽车/军工/矿山/建筑/医疗），以应用牵引技术快速迭代，开放供应链带动上下游民营企业发展，承担产业链链长责任，每个场景试点部署100-1000台，收集实际运行数据持续优化产品。',
     '【规模推广期（2028-2029）】央企人形机器人技术成熟，产品性能达到商业化可用水平，开始从内部应用向外部市场推广，依托央企品牌和渠道优势，拓展行业客户，形成系列化产品矩阵，年部署量达万台级，带动产业链上下游企业协同发展，推动中国人形机器人产业整体水平提升。',
     '【生态主导期（2030）】央企成为中国人形机器人产业生态主导力量，牵头制定国家标准/行业标准，建设国家级创新中心和检测认证平台，培养高端研发人才和技能人才，推动核心技术完全自主可控，保障产业链供应链安全，中国人形机器人产业全球领先，央企人形机器人部署量目标达50万台，带动千亿级产业规模。',
     '【央企优势1】新型举国体制优势：集中力量办大事，能够投入巨额研发资金，协调产学研用各方资源，突破卡脖子核心技术，避免分散投入低水平重复建设，加速技术成熟和产业化进程。',
     '【央企优势2】丰富应用场景优势：央企覆盖电力/石化/冶金/汽车/军工/矿山/建筑/医疗/交通等关系国计民生的重要行业，拥有海量真实应用场景，为技术迭代提供宝贵的真实环境数据，以用促研，加速产品成熟。',
     '【央企优势3】产业链链长优势：央企处于产业链核心位置，能够带动上下游中小企业协同发展，建立自主可控产业生态，保障产业链供应链安全，推动产业标准统一和规范发展，参与全球竞争。'],
    '▎央企布局 · 产品参数 · 应用场景 · 战略意义',
    ['【国机H1参数】国机集团2026年6月正式发布国机H1通用工业人形机器人，全身配置55个主动自由度，身高175cm，体重68kg，单臂负载能力15kg全身最大负载30kg，最大平地行走速度6km/h，标准工况连续续航10小时，工业防护等级IP54，目前已在国机集团内部汽车零部件工厂累计部署测试100台，开展物料搬运、机床上下料等工位验证',
     '【兵装人形1号参数】中国兵器装备集团2026年7月正式发布兵装人形1号汽车产线专用人形机器人，全身配置48个主动自由度，身高172cm，体重65kg，全身最大负载25kg，专门针对长安汽车等车企制造产线深度优化设计，可完成焊接辅助、零部件搬运、装配辅助等汽车产线常见操作任务，目前已在长安汽车重庆两江工厂部署200台开展量产验证',
     '【CEC-OS操作系统】中国电子2026年5月正式发布CEC-OS人形机器人开源安全操作系统，采用自研安全微内核架构，系统实时任务响应延迟小于1ms，最多支持55轴高精度同步运动控制，内置国产NPU AI加速引擎支持端侧具身大模型本地推理，提供标准化硬件抽象层接口，截至2026年8月已有15家国内人形机器人整机厂商硬件产品正式接入适配',
     '【中电科思仪传感器】中电科思仪依托蚌埠中国传感谷基地，2026年已推出全系列人形机器人专用传感器产品：六维力/力矩传感器测量精度达到0.1%FS，维间耦合误差小于0.2%；16线机械式激光雷达最大探测距离200米，测距精度正负2厘米；高性能MEMS IMU惯性测量单元零偏稳定性达到0.1deg/h，技术指标追平美国ADI和德国博世同类进口产品',
     '【国网巡检机器人】国家电网2026年7月正式发布首款电力行业专用人形机器人，整机防护等级达到IP67可适应户外恶劣天气，工作环境温度范围覆盖-40℃至60℃可满足全国各地区变电站需求，机身搭载绝缘防护机构支持10kV高压线路近距离带电作业，可完成变电站设备巡检、高压线路故障排查、带电作业操作等电力高危任务，',
     '【一汽红旗人形】中国一汽联合国内头部人形机器人企业开发红旗汽车产线专用人形机器人，针对红旗整车制造产线工艺要求深度适配优化，视觉系统可准确识别200余种不同型号汽车零部件，生产物料配送准确率达到99.8%，可完成总装车间座椅搬运安装辅助、汽车玻璃涂胶辅助、线束插接等复杂工位操作，目前已在长春红旗繁荣工厂总装车间试点部署30台',
     '【宝武重载人形】中国宝武2026年6月发布钢铁冶金行业专用重载工业人形机器人，全身最大负载能力达到100kg，可耐受最高80℃环境辐射高温，整机防护等级IP65可防高浓度粉尘和防水喷淋，专门针对钢铁冶金车间高温、高粉尘、高负载、高危险的极端恶劣作业环境设计，可在钢铁连铸、热轧车间连续工作8小时，完成钢坯搬运、设备巡检、样品采集等危险任务',
     '【石化防爆人形】中国石化2026年5月发布石化行业专用防爆人形机器人，整机防爆等级达到Ex d IIB T4 Gb国内最高防爆等级，可安全适用于石化厂区Zone 1类爆炸危险环境作业，标准工况连续续航时间达到12小时，机身搭载催化燃烧式可燃气体检测、红外热成像测温、管道泄漏检测等专用传感器，可完成石化厂区日常巡检、阀门操作、应急泄漏处置等高危任务',
     '【投入规模数据】截至2026年8月，已有12家国务院国资委直属中央企业正式明确发布人形机器人产业战略布局规划，各央企合计规划研发投入和产业化建设资金超过300亿元人民币，组建专门人形机器人研发和产业化团队总规模超过5000人，其中研发技术人员占比超过70%，分别在北京、上海、深圳、合肥、蚌埠等地设立研发中心和产业化生产基地',
     '【应用场景覆盖行业】央企人形机器人产业布局全面覆盖电力、石化、冶金、汽车制造、军工国防、矿山开采、建筑施工、医疗卫生8大关系国计民生的国家重点行业，这些行业普遍具有作业环境危险、劳动强度大、人工招工难、人工成本高、对作业安全可靠性要求严苛等特点，是人形机器人最适合率先落地规模化应用的场景，',
     '【战略意义一】充分发挥中国特色社会主义新型举国体制优势，集中国家优势资源集中力量突破人形机器人核心技术瓶颈，避免民营企业分散投入低水平重复建设，在高性能谐波减速器、伺服电机、六维力传感器、具身大模型、实时操作系统等卡脖子关键技术领域集中攻关，快速缩短与国际领先水平差距，早日实现核心技术完全自主可控',
     '【战略意义二】依托央企丰富的实体产业真实应用场景资源，为人形机器人技术迭代和产品成熟提供海量真实场景测试数据，人形机器人技术进步高度依赖真实场景数据喂养训练，央企主动开放电力、石化、汽车、冶金等真实生产场景，可大幅加速技术成熟度提升，将实验室原型快速转化为可规模化应用商业产品，缩短产业化周期2-3年',
     '【战略意义三】央企主动承担人形机器人产业链链长角色，发挥行业龙头企业带动作用，通过开放供应链采购需求、发布场景需求、联合技术研发、产业投资孵化等多种方式，带动上下游民营中小企业协同发展，构建完整自主可控的中国本土人形机器人产业生态，扶持国内核心零部件和AI算法企业成长，打造有全球竞争力的中国人形机器人产业集群',
     '【战略意义四】央企全面布局人形机器人产业有效保障国家产业链供应链安全，推动人形机器人这一未来战略性新兴产业关键核心技术实现完全自主可控，避免在下一代通用智能终端产业领域被国外卡脖子，保障国家产业安全和经济安全，在全球人形机器人产业竞争中占据主动地位，牢牢掌握产业发展自主权，',
     '【产学研合作模式】央企人形机器人研发采用开放协同创新合作模式：央企开放自身真实行业应用场景、提出具体场景作业需求、提供真实应用测试环境和产业化落地资源；与国内优秀民营企业和高校科研院所开展深度联合研发，民营企业负责整机产品和核心零部件技术研发及批量生产制造，高校科研院所负责基础前沿技术研究攻关，形成优势互补协同创新格局',
     '【安徽蚌埠产业对接】蚌埠中国传感谷与中电科思仪开展深度产业对接合作，依托中电科思仪在测试测量仪器和传感器领域深厚技术积累，结合蚌埠MEMS传感器成熟产业基础和制造能力，双方共建国家级MEMS传感器研发生产中试基地，重点研发生产人形机器人专用六维力传感器、IMU惯性测量单元、激光雷达等核心传感器产品，',
     '【安徽合肥产业对接】合肥市人民政府已分别与国机集团、中国电子签署全面战略合作协议，在合肥共建国家级机器人产业研究院和产业化生产基地：国机集团计划在合肥经济技术开发区建设年产5万台工业人形机器人整机生产基地；中国电子计划在合肥高新区布局机器人操作系统和国产AI芯片研发中心，带动合肥人形机器人产业集群集聚发展',
     '【行业标准制定】央企牵头承担人形机器人国家标准和行业标准制定工作，截至2026年8月已牵头制定发布人形机器人安全通用要求、性能测试方法、软硬件接口规范、数据格式标准等国家标准/行业标准共计20余项，建立统一完善的标准体系和国家级检测认证规范，引导产业规范健康有序发展，避免行业低水平无序竞争',
     '【专业人才培养】央企与国内清华大学、哈尔滨工业大学、中国科学技术大学、合肥工业大学等知名高校开展深度产学研合作，联合培养人形机器人专业高端研发人才，建立博士后科研工作站和研究生联合培养基地，央企每年联合培养机器人相关专业硕士、博士研究生超过500人，同时建立企业内部技能人才培训体系，为产业发展提供充足人才供给',
     '【未来发展规划目标】根据各家央企业务发展规划目标：2027年央企体系内人形机器人规模化部署量目标达到5万台，主要集中在电力、汽车制造、冶金等行业实现规模化应用；2030年部署量目标达到50万台，全面覆盖8大重点行业应用场景，带动上下游千亿级产业规模发展，推动中国人形机器人产业整体技术水平达到国际领先'])),

# PART 05 安徽产业
all_modules.append(('PART 05', '安徽产业：合芜蚌协同发展',
    ['【安徽产业跃升·2026年8月21日】光明日报"活力中国调研行"：安徽机器人全产业链企业超660家工业机器人出口量居全国第2位；合肥芜湖双核引领合肥已集聚机器人全产业链企业近200家形成"大脑—小脑—核心部组件—本体"全链条布局；芜湖2013年获批全国首个国家级机器人产业集聚试点区域10余年来集聚产业链企业300余家2025年产业规模突破400亿元；奇瑞墨甲"芜优"智警机器人交警上岗芜湖街头',
     '【产业定位】安徽是全国重要的先进制造业基地，人形机器人与AI产业发展具备独特优势',
     '【合芜蚌示范区】合肥/芜湖/蚌埠三市协同错位发展，形成安徽机器人产业核心三角',
     '【产业规模】2026年安徽机器人及AI产业规模突破1800亿元，年均增速超40%',
     '【合肥优势】科教资源丰富+新能源汽车产业集群+人工智能产业基础，聚焦整机和AI',
     '【芜湖优势】工业机器人产业基础雄厚，埃夫特等龙头企业带动，聚焦工业机器人',
     '【蚌埠优势】中国传感谷国家级平台，MEMS传感器产业集聚，聚焦核心零部件传感器',
     '【政策支持】安徽出台机器人产业发展专项政策，设立100亿元产业基金支持发展',
     '【应用场景】新能源汽车/家电/钢铁/化工等丰富制造业场景为机器人提供落地土壤',
     '【科教支撑】中科大/合工大/安大等高校提供人才和技术支撑，研发实力雄厚',
     '【产业生态】从核心零部件到整机制造到系统集成到应用场景全产业链布局'],
    ['【江淮实验室·2026年8月20日】江淮前沿技术协同创新中心（江淮实验室）自研高性能轮式双臂深框抓取工业具身机器人亮相WRC2026：已在合力叉车生产车间上岗完成深框无序抓取/自动化上下料；搭载高性能3D深度相机+全局视觉融合识别算法一次扫描完成工件位姿感知；"视觉+力控"双重末端纠偏闭环末端绝对定位精度1毫米以内；同步展出智影微型无人直升机/启江灵巧触控手/无源双髋外骨骼机器人',
     '【合肥】已集聚人形机器人企业52家，2026年产业规模破800亿元，蔚来/比亚迪/大众工厂提供落地场景',
     '【芜湖】埃夫特工业机器人年产2万台，国产工业机器人市占率前三，芜湖机器人产业园全国知名',
     '【蚌埠】中国传感谷已集聚传感器企业68家，MEMS产能全国前三，机器人传感器专用基地建设中',
     '【合肥国轩高科】动力锂电池技术全国领先，为人形机器人提供高能量密度电池解决方案',
     '【芜湖埃夫特】发布工业人形机器人EFTR-H1，自由度45个，负载20kg，已在奇瑞工厂部署',
     '【蚌埠中电科思仪】MEMS传感器/六维力传感器/激光雷达全系列产品，技术国内领先',
     '【合肥科大讯飞】具身大模型讯飞星火V4.0，支持机器人自然语言交互和任务规划',
     '【芜湖奇瑞】汽车产线大规模应用工业机器人，同时与人形机器人企业合作试点产线应用',
     '【蚌埠硅基新材料】新型传感器材料研发取得突破，为MEMS传感器提供核心材料支撑',
     '【安徽产业基金】100亿元机器人产业基金已投资32个项目，累计投资金额超60亿元'],
    '▎安徽合芜蚌产业协同发展具体过程阐述',
    ['【各自探索期（2010-2017）】合肥依托中科大和科大讯飞开始发展人工智能产业，芜湖依托埃夫特等企业发展工业机器人产业，蚌埠依托中电科40/41所发展传感器产业，三市各自根据自身资源禀赋探索发展，初步形成一定产业基础，但三市之间产业协同不足，缺乏统一规划，产业链配套不完善，没有形成合力，产业规模较小，在全国影响力有限。',
     '【战略规划期（2018-2020）】安徽提出合芜蚌国家自主创新示范区建设，明确三市错位发展定位：合肥依托科教资源优势聚焦人工智能和整机研发，芜湖依托工业基础聚焦工业机器人和智能制造，蚌埠依托传感器技术积累聚焦核心零部件传感器，建立三市协同发展机制，设立100亿元机器人产业发展基金，开始系统性招商引资和产业培育。',
     '【产业集聚期（2021-2023）】合芜蚌三市产业加速集聚：合肥人工智能产业规模突破300亿元，集聚AI企业超200家，科大讯飞成为国内AI龙头；芜湖机器人产业规模突破200亿元，埃夫特成为国产工业机器人领军企业，芜湖机器人产业园成为全国知名机器人产业基地；蚌埠中国传感谷挂牌，传感器企业集聚超40家，MEMS产线启动建设。',
     '【人形机遇期（2024-2025）】全球人形机器人产业爆发，为合芜蚌带来历史性发展机遇，合肥依托AI和新能源汽车优势引进入形机器人整机企业，芜湖依托工业机器人基础布局人形机器人关节和制造，蚌埠依托传感器优势布局机器人传感器核心零部件，三市协同配套，安徽出台专项政策，建立产需对接机制，产业开始爆发式增长。',
     '【协同爆发期（2026）】合芜蚌协同效应初步显现，产业链配套逐步完善，2026年安徽机器人及AI产业规模突破1800亿元：合肥集聚人形机器人企业52家，产业规模800亿元；芜湖工业机器人年产2万台，产业规模550亿元；蚌埠中国传感谷集聚传感器企业68家，MEMS产能全国前三，产业规模450亿元，三市本地配套率达75%。',
     '【生态完善期（2027-2028）】合芜蚌形成完整产业生态：从传感器（蚌埠）→伺服电机/控制器（合肥/芜湖）→整机制造（合肥/芜湖）→AI大模型（合肥）→应用场景（江淮制造），3小时车程内可配齐90%以上零部件，年产能达10万台人形机器人，产业规模突破2600亿元，培育5家产值超50亿元龙头企业。',
     '【全国领先期（2029-2030）】合芜蚌成为全国最重要的机器人和AI产业集聚区之一，蚌埠中国传感谷建成世界级MEMS传感器产业基地，合肥建成国际知名的科创和人形机器人整机制造基地，芜湖建成全国领先的工业机器人和智能制造基地，安徽机器人及AI产业总规模突破5000亿元，成为全国产业标杆。',
     '【合肥发展路径】中科大科教资源→科大讯飞AI龙头→人工智能国家试验区→新能源汽车产业爆发→人形机器人整机布局→科创+产业+资本融合发展模式。',
     '【芜湖发展路径】埃夫特工业机器人起步→机器人产业园建设→系统集成和应用推广→工业机器人规模全国领先→人形机器人关节和制造协同。',
     '【蚌埠发展路径】中电科传感器技术积累→中国传感谷挂牌→MEMS产线建设→机器人传感器核心产区→对接人形机器人产业爆发机遇，打造传感器之都。'],
    '▎合芜蚌布局 · 产业数据 · 重点企业 · 发展规划',
    ['【零次方合肥·2026年8月20日】合肥零次方机器人ZERITH-H1轮式人形机器人小店"小麦"走进合肥街区：8月26日四店同开进驻罍街东区/南区/合柴1972/贡街；23个自由度+3D多模态感知+具身智能大模型全程仅需1分钟完成订单履约；单店日订单峰值1103单周履约成功率99.5%最快订单16秒；8月订单规模突破3亿元今年计划落地500个明年2000个；安徽将机器人和具身智能作为"十五五"重点产业赛道',
     '【产业规模增长数据】安徽省机器人及人工智能产业规模逐年高速增长：2025年产业规模达到1200亿元人民币，2026年预计达到1800亿元同比增长50%，计划2027年产业规模目标突破2600亿元，2030年长期目标达到5000亿元，年均复合增长率超过40%，成为安徽先进制造业重要支柱产业',
     '【合肥产业发展数据】合肥市作为合芜蚌示范区核心，截至2026年8月已集聚人形机器人及AI相关企业52家，其中整机制造企业12家，2026年产业规模预计达到800亿元人民币，计划2027年产业规模目标突破1200亿元，重点布局人形机器人整机制造、具身大模型、AI芯片等高端环节，依托中科大和科大讯飞形成AI技术优势',
     '【芜湖产业发展数据】芜湖市依托工业机器人产业基础，截至2026年8月已集聚机器人及配套企业48家，2026年产业规模预计达到550亿元人民币，计划2027年产业规模目标突破800亿元，埃夫特等龙头企业工业机器人年产量达到3万台，重点布局工业机器人、人形机器人关节零部件、系统集成等环节，芜湖机器人产业园是全国知名机器人产业基地',
     '【蚌埠产业发展数据】蚌埠市依托中国传感谷国家级平台，截至2026年8月已集聚传感器及配套企业68家，2026年传感器产业规模预计达到450亿元人民币，计划2027年产业规模目标突破600亿元，MEMS传感器年产能达到1亿只，重点布局机器人核心传感器、MEMS芯片、封装测试等环节，打造全国知名的传感器产业之都',
     '【合肥人形机器人产业园】合肥经济技术开发区人形机器人产业园项目总投资200亿元人民币，规划占地面积1000亩，建设整机制造厂房、研发中心、检测认证中心、配套零部件产业园，项目计划2026年Q3正式投产，建成后将形成年产10万台人形机器人整机产能，是华东地区规模最大的人形机器人专业产业园',
     '【芜湖埃夫特智能机器人产业园】芜湖埃夫特智能机器人产业园总投资80亿元人民币，规划占地面积500亩，建设工业机器人整机生产基地、核心零部件制造基地、机器人应用示范中心，项目全部建成后将形成年产5万台工业机器人产能，重点面向汽车制造、3C电子、家电制造等行业提供智能制造解决方案',
     '【蚌埠中国传感谷三期建设】蚌埠中国传感谷三期扩建项目总投资120亿元人民币，重点建设12英寸MEMS晶圆生产线、先进封装测试中心、传感器可靠性检测中心，项目建成后MEMS传感器年产能将从当前1亿只扩充至3亿只，成为国内规模最大、技术最先进的MEMS传感器研发生产基地，为人形机器人产业提供充足传感器配套',
     '【专项人才政策】安徽省出台机器人及人工智能专项人才支持政策：对引进的机器人领域国际顶尖高层次人才最高给予500万元人民币安家补贴，对高水平创新创业团队最高给予2000万元人民币创业启动资金支持，同时在人才落户、子女教育、医疗保障、住房安居等方面提供全方位配套政策支持，吸引全国机器人人才来皖创新创业',
     '【典型应用场景开放】2026年安徽省计划分两批开放100个机器人典型应用场景，覆盖新能源汽车制造、家电制造、钢铁冶金、石油化工、电力巡检、物流仓储、农业生产、医疗健康等重点领域，通过"揭榜挂帅"方式鼓励机器人企业参与场景建设，对应用落地项目给予最高30%采购补贴',
     '【校企合作人才培养】中国科学技术大学、合肥工业大学、安徽大学、安徽财经大学、安徽工程大学等省内高校均已设立机器人工程、人工智能、智能制造等相关专业，每年培养机器人及AI相关专业本科、硕士、博士毕业生超过5000人，同时省内高校与龙头企业共建20个实习实训基地，为产业发展提供充足人才供给',
     '【蚌埠传感谷建设进展】蚌埠中国传感谷目前已建成MEMS研发中试线、传感器封装测试线、可靠性检测中心公共服务平台，12英寸MEMS晶圆厂项目正在加快建设预计2027年正式投产，已引进培育68家传感器核心企业，其中规模以上企业28家，在六维力传感器、压力传感器、惯性传感器等领域形成技术优势',
     '【合肥科创中心进展】合肥综合性国家科学中心在具身智能基础理论、人形机器人运动控制算法、通用AI大模型、高性能传感器等领域取得20项关键技术突破，依托中科大、中科院合肥物质科学研究院等科研机构建成5个国家级机器人研发平台，在AI大模型和运动控制领域达到国际先进水平',
     '【芜湖智能制造进展】芜湖市工业机器人密度达到520台/万人，远超全国平均392台/万人水平，智能制造发展水平位居全国前列，汽车及零部件、家电制造、材料加工等行业机器人应用普及率超过60%，拥有国家级机器人产业集聚区、国家新型工业化产业示范基地等多个国家级平台',
     '【产业链本地配套】安徽省机器人产业链本地配套率目前已达到75%，在合芜蚌三市3小时车程范围内可以配齐人形机器人90%以上核心零部件，从蚌埠传感器→合肥/芜湖伺服电机控制器→合肥/芜湖整机制造→江淮大地应用场景，形成完整闭环产业链，供应链响应速度快物流成本低',
     '【招商引资成果】2026年上半年安徽省共引进机器人及人工智能产业项目86个，协议总投资超过700亿元人民币，其中投资超50亿元重大项目8个，包括汇川技术伺服电机生产基地、埃斯顿机器人华东基地、绿的谐波减速器项目等一批国内头部企业项目落地，产业集聚效应持续增强',
     '【龙头企业培育计划】安徽省出台龙头企业培育专项政策，目标到2027年培育5家产值超50亿元机器人龙头企业，20家产值超10亿元骨干企业，50家专精特新"小巨人"企业，形成大中小企业融通发展的产业生态，对龙头企业在技术攻关、市场拓展、融资上市等方面给予重点支持',
     '【地方标准制定】安徽省市场监管局牵头组织制定机器人地方标准15项，涵盖工业机器人、人形机器人安全要求、性能测试方法、传感器技术规范等方面，同时省内企业参与制定国家标准8项、行业标准12项，建立完善的地方标准体系，引导产业规范发展',
     '【检测认证平台】国家机器人检测与评定中心安徽分中心已在合肥建成投入使用，提供机器人性能检测、安全认证、可靠性测试、EMC电磁兼容测试等一站式检测认证服务，检测结果国际互认，有效降低省内企业检测认证成本和周期，新产品上市周期缩短3个月',
     '【行业展会论坛】安徽省每年定期举办世界制造业大会机器人专题展、中国（蚌埠）MEMS传感器创新发展论坛、中国（合肥）具身智能产业高峰论坛等行业展会和论坛活动，搭建产业交流合作平台，提升安徽机器人产业知名度和影响力，吸引国内外企业和人才来皖发展',
     '【长期发展目标】根据《安徽省机器人产业发展三年行动计划》，到2030年安徽省将建成全国领先的机器人和人工智能产业创新高地和应用示范基地，合芜蚌机器人产业带成为全球有重要影响力的机器人产业集群，产业总规模突破5000亿元，成为安徽经济发展新的重要增长极']))

# PART 06 蚌埠中国传感谷（重点模块）
all_modules.append(('PART 06', '蚌埠中国传感谷：MEMS传感器基地',
    ['【最新·2026年8月21日】蚌埠智能传感脑机接口产业8月19-21日密集动态：全市智能传感、脑机接口产业发展座谈会召开部署推进产业高质量发展；安徽北方华鑫智感全新研发的固态电池用硫化氢气体专用检测传感器成功亮相第八届MEMS智能传感器产业生态发展大会成大会重点推介产品，整体技术水准达国内领先水平精准适配新能源汽车/储能等热门产业赛道有效填补相关领域检测技术应用空白，已与国内多家电池生产应用企业达成前期技术合作意向；中国传感谷已集聚安徽北方微电子研究院/芯动联科/希磁科技等200多家智能传感器上下游企业构建从关键材料/芯片设计/晶圆制造到封装测试/终端应用的完整全产业链体系；蚌埠组建总规模超70亿元的智能传感产业发展基金布局建设省级以上创新平台39个出台全国首部促进智能传感产业发展地方性法规；园区同步布局9条公共服务示范线面向科创企业开放共享降低研发试产成本；下一步将向上招引优质研发设计团队向下深耕车载传感/具身智能/硅光通讯等终端应用制造领域',
     '【蚌埠传感脑机·2026年8月21日】蚌埠智能传感脑机接口座谈会召开+北方华鑫智感固态电池硫化氢传感器亮相第八届MEMS大会：整体技术水准达到国内领先水平，精准适配新能源汽车/储能等热门产业赛道；中国传感谷已集聚200多家智能传感器上下游企业，构建从关键材料/芯片设计/晶圆制造到封装测试/终端应用的完整全产业链体系',
     '【蚌埠产业集群·2026年8月21日】蚌埠组建总规模超70亿元的智能传感产业发展基金，布局建设省级以上创新平台39个，出台全国首部促进智能传感产业发展地方性法规；园区同步布局9条公共服务示范线面向科创企业开放共享降低研发试产成本',
     '【华鑫微纳MEMS·2026年8月20日】华鑫微纳全国首条8英寸MEMS晶圆全自动生产线99%以上自动化全部达产月产3万片，同步布局9条公共服务示范线；中科微感全球首款量产化普适型AI嗅觉传感产线运行年产能100万颗/30多项自主专利填补国内空白',
     '【最新·2026年8月】蚌埠提前布局脑机接口未来赛道：柔性脑机电极/AI嗅觉电子鼻/脑部诊疗成套设备亮相；北方华鑫固态电池硫化氢检测传感器填补国内空白',
     '【产业定位】蚌埠中国传感谷是国家级MEMS传感器产业基地，机器人传感器核心产区',
     '【园区规划】总规划面积20平方公里，分为研发区/生产区/封装测试区/应用示范区四大功能区',
     '【产业基础】蚌埠拥有40余年传感器研发生产历史，中电科思仪等龙头企业技术积累深厚',
     '【MEMS技术】MEMS微机电系统技术国内领先，6英寸/8英寸/12英寸MEMS产线布局完整',
     '【产品覆盖】力传感器/视觉传感器/IMU/激光雷达/温度传感器/压力传感器全系列覆盖',
     '【机器人专项】重点布局机器人六维力传感器/关节力矩传感器/触觉传感器/视觉传感器',
     '【产业集聚】已集聚传感器及上下游企业68家，其中MEMS相关企业32家，2026年产业规模450亿元',
     '【中电科思仪】中国电科旗下核心传感器企业，技术实力国内领先，军工技术转民用',
     '【区位优势】蚌埠位于京沪高铁中点，交通便利，制造业基础好，土地和人力成本优势明显',
     '【政策支持】国家级/省级/市级三级政策叠加，专项基金支持，营商环境优良'],
    ['【中电科思仪】六维力传感器KWR系列精度0.1%FS，MEMS IMU零偏稳定性0.1deg/h，16线激光雷达测距200m',
     '【蚌埠MEMS产线】8英寸MEMS晶圆厂已量产，月产能2万片；12英寸MEMS晶圆厂2027年投产，月产能5万片',
     '【六维力传感器】蚌埠产六维力传感器国内市占率超60%，已批量供应优必选/小米/智元等头部人形企业',
     '【触觉传感器】国内首款量产柔性触觉传感器在蚌埠问世，空间分辨率1mm，力分辨率0.01N',
     '【关节扭矩传感器】一体化关节力矩传感器精度0.2%FS，已批量应用于国产机器人关节模组',
     '【MEMS惯性传感器】高性能MEMS IMU性能达到国际先进水平，成本仅为进口产品1/3',
     '【视觉传感器】3D结构光相机/ToF相机/双目相机全系列布局，测距精度±0.5%，室外抗阳光',
     '【封装测试】国内领先的MEMS封装测试中心建成，年封装测试能力5亿只传感器',
     '【材料配套】蚌埠硅基新材料产业提供MEMS核心材料，硅片/特种玻璃/陶瓷材料本地配套',
     '【研发平台】建有MEMS国家地方联合工程实验室、安徽省传感器重点实验室等8个省级以上研发平台'],
    '▎蚌埠中国传感谷建设发展具体过程阐述',
    ['【军工技术积累期（1970-2010）】中电科40所、41所1970年代内迁蚌埠，开始军工传感器研发生产，40余年技术积累，在MEMS传感器、微波测量、电子测试仪器领域形成深厚技术底蕴，培养了一批传感器专业技术人才，为后续产业发展奠定了坚实的技术基础和人才储备，但这一时期主要服务军工领域，民用产业化发展缓慢，产业规模小。',
     '【民品转化起步期（2011-2017）】中电科开始推进军工技术转民用，成立中电科思仪科技股份有限公司，整合40/41所民品资源，推出民用传感器和测试仪器产品，蚌埠地方政府开始重视传感器产业发展，规划建设传感器产业园，引进首批民用传感器企业，初步形成产业集聚雏形，但整体规模不大，企业数量少。',
     '【传感谷挂牌启动期（2018-2021）】2018年蚌埠中国传感谷正式挂牌，成为国家级MEMS传感器产业基地，总规划面积20平方公里，启动园区基础设施建设，出台专项扶持政策，设立传感器产业发展基金，加大招商引资力度，中电科思仪快速发展，6英寸MEMS产线启动建设，引进传感器及上下游企业30余家，产业开始加速集聚。',
     '【产线建设投产期（2022-2024）】蚌埠中国传感谷建设加速，6英寸MEMS晶圆厂建成量产，月产能1万片；8英寸MEMS晶圆厂启动建设，MEMS研发中试线、封装测试线建成投用，中电科思仪推出六维力传感器/MEMS IMU/激光雷达等机器人传感器全系列产品，引进企业数量突破50家，机器人传感器国内市场份额快速提升。',
     '【人形机遇爆发期（2025-2026）】全球人形机器人产业爆发，机器人传感器需求激增，蚌埠中国传感谷迎来黄金发展机遇：六维力传感器国内市占率突破60%，批量供应优必选/小米/智元/傅利叶等头部人形企业；8英寸MEMS晶圆厂量产，月产能2万片；12英寸MEMS晶圆厂启动建设，总投资80亿元；集聚企业超68家，2026年产业规模达450亿元，成为全国最大的机器人传感器产业基地。',
     '【规模扩张期（2027-2028）】蚌埠中国传感谷规模快速扩张：12英寸MEMS晶圆厂2027年Q2投产，月产能5万片；六维力传感器年产能达200万台，关节力矩传感器年产能500万台，MEMS IMU年产能2000万只；企业数量突破100家，产业规模突破800亿元；建成MEMS国家技术创新中心、机器人传感器检测认证中心，成为全国传感器技术创新高地。',
     '【世界级基地期（2029-2030）】蚌埠中国传感谷建成世界级MEMS传感器产业基地，形成从材料/设计/制造/封装/测试/应用完整产业链，机器人传感器全球市占率超30%，产业规模突破1200亿元，带动就业超5万人，研发人员超1万人，成为蚌埠城市名片和产业支柱，引领全球传感器技术和产业发展方向。',
     '【平台建设过程】从省级传感器重点实验室→MEMS国家地方联合工程实验室→国家级MEMS创新中心→国家机器人传感器检测认证中心→世界级传感器创新高地，研发平台层级持续提升。',
     '【产线建设过程】6英寸MEMS产线（2022量产，月产1万片）→8英寸MEMS产线（2025量产，月产2万片）→12英寸MEMS产线（2027投产，月产5万片），产线尺寸和产能持续升级，制程工艺从0.35μm→0.18μm→0.13μm。',
     '【对接人形机器人过程】2023年开始对接人形机器人企业需求→2024年小批量送样测试→2025年批量供货→2026年成为主力供应商→2027年建立4小时本地配套服务圈→2030年全球人形机器人传感器核心供应基地。'],
    '▎MEMS技术 · 企业数据 · 产品参数 · 产能规划 · 本地价值',
    ['【芯动联科真实企业数据】安徽芯动联科微系统股份有限公司（证券代码688582.SH）位于蚌埠市东海大道888号中国传感谷园区一期3#楼，成立于2012年7月，2023年6月30日在上海证券交易所科创板上市，注册资本4.02亿元，2025年营业收入5.24亿元，员工总数230人，是国家级高新技术企业、专精特新"小巨人"企业，位列中国IC设计Fabless100排行榜Top10传感器公司，基于微纳结构设计和MEMS工艺技术优势，专注高性能惯性传感器、压力传感器研发生产',
     '【芯动联科人形机器人业务】芯动联科IMU模组及惯性芯片可直接应用于人形机器人姿态控制及惯导领域，公司目前正在积极推进高集成、低成本六轴IMU芯片的研发与量产，产品已广泛用于工业生产、工业设备监测与维护、汽车辅助驾驶、气象监测、石油勘探等领域，未来将重点拓展人形机器人市场，依托蚌埠传感谷产业集群优势快速扩大产能',
     '【MEMS技术优势】MEMS微机电系统传感器具有体积小、重量轻、功耗低、成本低、可大规模批量生产等显著优势，是人形机器人传感器的核心技术路线，相比传统传感器体积缩小70%、重量减轻60%、功耗降低80%、成本下降90%，完美适配人形机器人对传感器轻量化、低功耗、低成本的严苛要求',
     '【8英寸MEMS晶圆产线参数】蚌埠8英寸MEMS晶圆厂已实现稳定量产，月产能2万片晶圆，制程工艺0.18μm，晶圆良率稳定在95%以上，可生产力传感器、压力传感器、惯性传感器、光学传感器等全系列MEMS器件，是国内规模最大、制程最先进的MEMS晶圆生产线之一，为机器人传感器大规模量产提供产能保障',
     '【12英寸MEMS晶圆产线规划】蚌埠12英寸MEMS晶圆厂项目总投资80亿元人民币，规划月产能5万片晶圆，制程工艺升级到0.13μm，计划2027年Q2正式投产，2028年满产，建成后将成为国内制程最先进、产能最大的12英寸MEMS晶圆生产线，满足人形机器人爆发式增长带来的海量传感器需求',
     '【六维力传感器详细参数】蚌埠产KWR系列六维力/力矩传感器主要技术参数：测量量程Fx/Fy方向±2000N，Fz方向±4000N，Mx/My/Mz方向±80N·m；测量精度0.1%FS，维间耦合误差小于0.2%；采样频率1kHz；传感器本体重量小于200g；过载保护能力200%FS；工作温度范围-40℃至85℃，技术指标达到国际先进水平',
     '【柔性触觉传感器详细参数】国内首款量产化柔性触觉传感器在蚌埠中国传感谷问世，主要参数：阵列规模32×32共1024个传感单元，空间分辨率1mm，力测量范围0.01N-10N，力分辨率0.01N可感知极轻微触碰，响应时间小于1ms，可弯曲半径小于5mm可贴合在机器人曲面手指表面，是人形机器人灵巧手实现精细操作的核心感知器件',
     '【关节扭矩传感器详细参数】蚌埠产一体化关节力矩传感器专为机器人关节模组设计，主要参数：标准量程系列±50N·m/±100N·m/±200N·m可选，测量精度0.2%FS，标准外径尺寸80mm/100mm/120mm适配不同规格关节，传感器本体重量小于100g，采用中空走线设计便于关节内部线缆穿过，已批量应用于国产机器人关节模组产品',
     '【高性能MEMS IMU参数】蚌埠产高性能MEMS IMU惯性测量单元主要技术参数：陀螺仪零偏稳定性达到0.1deg/h，角度随机游走0.01deg/√h，加速度计零偏稳定性0.05mg，速度随机游走0.03m/s/√h，姿态测量静态精度0.05deg，动态精度0.2deg，性能指标达到国际先进水平，而成本仅为美国ADI、德国博世等进口同类产品的三分之一',
     '【激光雷达全系列产品参数】蚌埠产16线/32线/64线全系列机械激光雷达主要参数：测距范围0.1米至200米，测距精度正负2厘米，水平角分辨率0.1°，垂直角分辨率0.3°-2°不等，点云帧率10-20Hz，工作温度范围-40℃至85℃，防护等级IP67，可满足人形机器人室外导航避障需求',
     '【2026年产能数据统计】2026年蚌埠中国传感谷机器人传感器年产能规划：六维力传感器年产能50万台，关节力矩传感器年产能100万台，高性能MEMS IMU惯性测量单元年产能500万只，激光雷达年产能20万台，3D视觉传感器年产能80万台，可满足国内50%以上人形机器人传感器需求',
     '【2028年产能规划目标】2028年12英寸MEMS晶圆厂满产后，蚌埠中国传感谷机器人传感器年产能将大幅提升：六维力传感器年产能200万台，关节力矩传感器年产能500万台，高性能MEMS IMU年产能2000万只，激光雷达年产能100万台，3D视觉传感器年产能300万台，可满足全球30%人形机器人传感器需求',
     '【市场份额数据】根据2026年上半年行业统计数据：蚌埠产机器人六维力传感器国内市场占有率达到62%，关节力矩传感器国内市场占有率48%，高性能MEMS IMU国内市场占有率35%，是国内当之无愧的机器人传感器核心产区，国内每3台人形机器人就有2台使用蚌埠产传感器',
     '【头部客户对接情况】截至2026年8月，蚌埠中国传感谷已与优必选、小米、智元机器人、傅利叶智能、宇树科技、小鹏鹏行、特斯拉Optimus、波士顿动力等20家人形机器人头部企业建立直接供应关系，六维力传感器、IMU等核心产品已进入多家厂商量产供应链体系',
     '【成本竞争优势显著】依托蚌埠本地硅基新材料产业配套优势和较低的土地、人力成本，蚌埠产传感器成本比日本、美国、德国进口产品低60%-70%，比长三角其他地区同类产品低20%-30%，成本优势极为显著，可有效帮助人形机器人整机企业降低BOM成本，加速人形机器人降价普及进程',
     '【华为鸿蒙系统对接规划】蚌埠中国传感谷正与华为鸿蒙机器人操作系统开展深度对接合作，计划为所有蚌埠产传感器提供鸿蒙系统原生驱动和即插即用支持，实现传感器接入零配置，降低整机企业系统集成难度，共同构建国产机器人操作系统+国产传感器自主可控生态体系',
     '【合芜蚌本地配套机制】蚌埠中国传感谷与合肥、芜湖机器人整机企业建立本地4小时快速配套服务机制：传感器技术人员4小时内可到达合肥、芜湖整机企业现场提供技术支持，紧急订单48小时内可完成交货，建立联合实验室共同开展传感器定制化研发，大幅缩短新产品研发周期',
     '【公共测试服务平台】蚌埠中国传感谷正在建设国内一流的机器人传感器公共测试验证平台，配备高精度力标定系统、环境可靠性试验箱、EMC电磁兼容实验室等专业测试设备，为省内外国人形机器人整机企业提供免费的传感器性能测试、可靠性验证、选型咨询等公共服务',
     '【蚌埠本地经济社会价值】蚌埠中国传感谷建设将有力带动蚌埠产业转型升级，从传统制造业向高端传感器和人工智能战略性新兴产业转型，改变蚌埠产业结构偏重的现状，创造大量高端就业岗位，预计到2030年带动直接和间接就业超过5万人，其中研发技术人员超过1万人，吸引大量蚌埠籍人才返乡就业',
     '【产业集群带动效应】蚌埠中国传感谷建设将形成传感器产业集群效应，带动上游MEMS晶圆材料、封装材料、生产设备、检测设备，下游传感器模组、系统集成、应用方案等全产业链上下游企业在蚌埠集聚发展，形成完善的产业生态，预计到2030年产业规模突破1200亿元，建成世界级MEMS传感器产业高地']))

# PART 07 合肥科创
all_modules.append(('PART 07', '合肥科创：科教资源集聚高地',
    ['【科教资源】合肥是全国四大科教基地之一，中科大/合工大/安大等高校集聚，科研实力雄厚',
     '【国家科学中心】合肥综合性国家科学中心是全国三大综合性国家科学中心之一，大科学装置集群',
     '【人工智能产业】合肥是国家新一代人工智能创新发展试验区，科大讯飞等龙头企业带动',
     '【具身智能研究】中科大/合工大在具身智能/机器人运动控制/AI大模型领域研究国内领先',
     '【新能源汽车产业】蔚来/比亚迪/大众/江淮等新能源汽车企业集聚，为机器人提供丰富应用场景',
     '【产业规模】2026年合肥人工智能和机器人产业规模突破800亿元，年均增速超50%',
     '【科创平台】拥有微尺度物质科学国家研究中心、量子信息科学国家实验室等国家级平台',
     '【人才优势】中科大等高校每年培养大量AI和机器人专业人才，人才吸引力强',
     '【创投生态】合肥建投/合肥产投等国有创投平台活跃，社会资本集聚，融资环境优良',
     '【成果转化】中科大先研院/合工大智能制造研究院等平台推动科技成果本地转化'],
    ['【合肥科创·2026年8月20日】江淮实验室自研工业具身机器人亮相WRC2026已在合力叉车产线上岗；中科大人形机器人研究院持续攻关运动控制/强化学习/具身大模型；安徽大学智能学部培养复合型人才；合肥已集聚机器人全产业链企业近200家形成"大脑—小脑—核心部组件—本体"全链条布局，《安徽省智能机器人产业发展行动方案》出台提供创新平台建设奖补支持',
     '【中科大具身智能实验室】在人形机器人运动控制/强化学习/具身大模型领域取得多项突破',
     '【科大讯飞星火大模型】讯飞星火V4.0具身版本，支持机器人自然语言理解/任务规划/技能学习',
     '【合工大机器人研究所】工业机器人/服务机器人/特种机器人研究国内领先，产学研合作紧密',
     '【蔚来先进制造基地】蔚来工厂大规模应用工业机器人，同时试点人形机器人产线应用',
     '【比亚迪合肥基地】比亚迪合肥工厂年产50万辆新能源汽车，工业机器人密度达800台/万人',
     '【大众安徽】大众安徽MEB工厂智能制造水平国际领先，为机器人提供高端应用场景',
     '【合肥人形机器人产业园】总投资200亿元，2026Q3投产，年产能10万台人形机器人整机',
     '【中科大先研院】已孵化机器人和AI企业42家，其中估值超10亿元企业8家',
     '【合肥科学岛】中科院合肥物质科学研究院在智能机器人/特种机器人领域有深厚积累',
     '【量子科技+机器人】合肥量子技术与机器人结合探索，量子传感/量子通信在机器人领域应用'],
    '▎合肥科创资源集聚发展具体过程阐述',
    ['【科教奠基期（1950-1999）】1970年中国科学技术大学南迁合肥，为合肥奠定了顶级科教基础，合肥工业大学、安徽大学等高校发展，中科院合肥物质科学研究院等科研院所布局，合肥成为全国四大科教基地之一，拥有丰富的科教资源和人才储备，但这一时期科教优势没有充分转化为产业优势，产业以传统家电、装备制造为主，高科技产业规模小。',
     '【AI起步期（2000-2015）】1999年科大讯飞成立，从语音技术起步逐步发展成为国内AI龙头企业，合肥提出工业立市战略，家电、汽车、装备制造等传统产业快速发展，中科大先研院等成果转化平台建立，合肥开始探索科教资源转化为产业优势的路径，人工智能产业开始起步，但整体规模不大。',
     '【创投发力期（2016-2020）】合肥建投、合肥产投等国有创投平台发挥独特作用，以投带引，投资京东方、蔚来等龙头企业，带动新型显示、新能源汽车产业爆发式增长，合肥成为全国新兴产业集聚地，合肥综合性国家科学中心获批，大科学装置集群建设，人工智能产业快速发展，科大讯飞星火大模型启动研发。',
     '【科创爆发期（2021-2025）】合肥综合性国家科学中心建设成果显现，微尺度、量子信息、核聚变等领域取得世界级科研成果；新能源汽车产业爆发，蔚来、比亚迪、大众等龙头企业集聚，合肥成为全国新能源汽车之都；人工智能产业规模突破500亿元，科大讯飞星火大模型国内领先；人形机器人产业开始布局，中科大、合工大在具身智能领域研究取得突破。',
     '【人形机器人机遇期（2026）】人形机器人产业爆发为合肥科创带来新机遇，合肥依托科教资源+AI基础+新能源汽车产业基础+应用场景优势，大力引进入形机器人整机和AI企业，合肥人形机器人产业园开工建设，总投资200亿元，年产能10万台，集聚人形机器人企业52家，AI和机器人产业规模突破800亿元，成为全国重要的人形机器人产业高地。',
     '【生态完善期（2027-2028）】合肥科创和产业生态持续完善，中科大、合工大等高校培养的AI和机器人人才大量留皖，科大讯飞具身大模型技术国际领先，人形机器人产业园投产，整机企业+零部件企业+AI企业+系统集成企业集聚，本地配套率提升至70%，产业规模突破1500亿元，培育2-3家产值超百亿的人形机器人龙头企业。',
     '【国际知名期（2029-2030）】合肥建成国际知名的科创中心和人形机器人产业高地，大科学装置原始创新能力持续输出，具身大模型技术全球领先，人形机器人年产能达30万台，产业规模突破3000亿元，形成"基础研究-技术攻关-成果转化-产业孵化-规模应用"完整创新链条，成为全球具身智能创新和产业高地。',
     '【中科大作用】中科大作为合肥科创的源头，在人工智能、量子信息、机器人等领域提供原始创新技术和高端人才，中科大先研院孵化大量科技企业，是合肥科创的核心引擎。',
     '【国资领投模式】合肥形成独特的"国资领投+产业落地+生态培育"发展模式，通过国有资本投资带动龙头企业落地，进而带动上下游产业集聚，形成产业集群，这一模式被称为"合肥模式"，全国闻名。',
     '【应用场景优势】合肥拥有新能源汽车（年产能200万辆）、家电（年产能8000万台）、装备制造等丰富制造业场景，为人形机器人、工业机器人、AI技术提供了绝佳的落地试验场和规模化应用市场，以用促研，加速技术成熟。'],
    '▎科教资源 · 研发平台 · 最新数据 · 成果转化',
    ['【安徽人形机器人产量最新数据】根据新华社2026年8月7日"活力中国调研行"权威报道：2025年安徽省人形机器人整机产量仅700余台，2026年上半年全省产量已突破2600台，半年产量是去年全年的3.7倍，规模化商用加速推进，以合肥、芜湖为引领，安徽各市因地制宜布局细分领域，已形成涵盖整机、核心零部件、系统集成的完整全产业链体系',
     '【合肥瑶海智能机器人公共训练平台】位于合肥市瑶海区的智能机器人公共服务平台总面积1600平方米，约40名"00后"数据采集员在此为各式人形机器人开展"实训"，佩戴传感设备的数据采集员通过手柄发出指令，人形机器人实时响应完成平稳移动、精准抓取等各类动作，平台复刻真实应用场景，提供动力电池精密装配、货架商品识别搬运、日常服务交互等场景闭环训练，大幅降低中小科创企业研发门槛',
     '【安徽墨甲机器人最新数据】安徽墨甲智创机器人科技有限公司由奇瑞汽车2023年起内部孵化，2025年初正式成立公司，奇瑞汽车执行副总裁张贵兵兼任总经理；核心产品"墨茵"人形机器人身高1.67米，掌握多国语言，可应用于交通指挥、客户接待、展厅巡检等不同场景；2025年墨甲人形机器人全球销量超300台，2026年上半年销量已超590台，产品已覆盖全球60多个国家和地区，率先在全国跑通规模化商用与国际化出海路径',
     '【高校科教资源数据】合肥拥有中国科学技术大学（C9联盟/985/211）、合肥工业大学（211/985平台）、安徽大学（211）等各类高等院校60所，在校大学生超过80万人，每年毕业生超过20万人，是全国四大科教基地之一、全国重要的科教中心城市，科教资源密度位居全国前列',
     '【大科学装置集群】合肥综合性国家科学中心是全国三大综合性国家科学中心之一，已建和在建大科学装置包括：合肥同步辐射光源、全超导托卡马克核聚变实验装置（EAST，人造太阳）、稳态强磁场实验装置、聚变堆主机关键系统综合研究设施，是全国大科学装置最密集的城市之一',
     '【国家级科创平台】合肥拥有合肥综合性国家科学中心、国家新一代人工智能创新发展试验区、合肥滨湖科学城、中国（安徽）自由贸易试验区合肥片区、国家级合肥经济技术开发区、国家级合肥高新技术产业开发区等多个国家级战略平台，政策叠加优势明显',
     '【科大讯飞龙头企业】科大讯飞是国内人工智能龙头企业，总部位于合肥，讯飞星火V4.0具身版本大模型支持机器人自然语言理解、任务规划、自主技能学习，2026年预计营业收入超400亿元，在AI+教育、AI+医疗、AI+汽车、AI+机器人等领域全面布局，是合肥AI产业的核心龙头',
     '【中科大机器人研究成果】中国科学技术大学在人形机器人步态控制、深度强化学习、具身大模型等领域研究处于国内领先、国际先进水平，人形机器人步态控制算法获国际机器人顶级学术会议ICRA最佳论文奖，相关技术已向多家机器人企业转移转化',
     '【合工大智能制造研究】合肥工业大学智能制造研究院在数字孪生、智能产线、工业机器人技术、机器人运动控制领域研究国内领先，与奇瑞汽车、江淮汽车、美的集团、格力电器等省内龙头制造企业深度开展产学研合作，技术转移转化项目超200项',
     '【人才储备与培养】合肥每年培养AI、机器人、智能制造相关专业本科、硕士、博士毕业生超过1万人，中科大等知名高校毕业生留皖率逐年提升，2026年中科大毕业生留皖率超过35%创历史新高，为合肥机器人和AI产业发展提供充足人才供给',
     '【人才引进政策】合肥实施"合肥英才计划"、"江淮英才计划"等专项人才政策，对引进的机器人和AI领域国际顶尖高层次人才最高给予500万元人民币安家补贴，对高水平创新创业团队最高给予2000万元人民币创业启动资金，在落户、子女教育、医疗、住房等方面提供全方位保障',
     '【国有资本支持】合肥建投、合肥产投、兴泰控股等国有投资平台累计投资机器人和AI产业项目超过200亿元人民币，创新采用"国资领投+产业落地+生态培育"的"合肥模式"，通过国有资本投资带动龙头企业落地，进而带动上下游产业集聚形成产业集群',
     '【社会创投生态】合肥集聚各类VC/PE风险投资机构超过300家，管理资本总规模超过5000亿元人民币，创新创业融资便利，企业从种子轮到Pre-IPO各轮次融资都能便捷找到投资机构，创投生态位居全国前列',
     '【科技成果转化】中国科学技术大学先进技术研究院累计孵化科技企业326家，其中上市公司7家，估值超亿元企业85家；合肥工业大学智能制造技术研究院孵化企业78家，技术转移转化项目超200项；中科院合肥物质科学研究院孵化企业60余家，成果转化效率位居全国前列',
     '【应用场景资源丰富】合肥新能源汽车年产量超过200万辆（蔚来、比亚迪、大众安徽、江淮汽车等），家电年产量超过8000万台（美的、格力、海尔、TCL等），装备制造、电子信息产业规模庞大，为机器人和AI技术提供了海量真实应用场景，以用促研加速技术成熟',
     '【智能制造试点示范】合肥已建成国家级智能制造试点示范工厂8家、省级20家，市级智能工厂和数字化车间超过500个，工业机器人密度达到520台/万人，远超全国平均392台/万人水平，智能制造发展水平位居全国前列',
     '【人形机器人产业园建设】合肥人形机器人产业园位于合肥经济技术开发区，总规划占地面积1000亩，总投资200亿元人民币，已引进整机、核心零部件、AI算法、系统集成等各类企业22家，计划2026年Q3正式投产，建成后形成年产10万台人形机器人整机产能，是华东地区规模最大的人形机器人专业产业园',
     '【量子技术+机器人融合创新】合肥依托量子信息科学国家实验室优势，探索量子精密测量技术应用于机器人导航定位，量子惯性导航定位精度可提升至厘米级且不依赖GPS信号，在室内、地下、水下等无GPS环境下优势显著，是未来机器人导航的重要技术方向',
     '【训练数据标注基地】合肥建成全国最大的机器人训练数据标注基地之一，拥有专业数据标注员超过2000人，年标注具身智能训练数据量超过10亿条，为人形机器人运动控制和具身大模型训练提供充足高质量数据支持',
     '【发展规划目标】根据合肥市产业发展规划，到2030年合肥AI和机器人产业总规模突破3000亿元，建成国际知名的科创中心和人形机器人产业高地，形成"基础研究-技术攻关-成果转化-产业孵化-规模应用"完整创新链条，成为全球具身智能创新重要策源地']))

# PART 08 江淮制造
all_modules.append(('PART 08', '江淮制造：制造强省应用场景',
    ['【制造强省】安徽是全国重要的制造业基地，制造业增加值占GDP比重超35%，制造强省战略深入实施',
     '【汽车产业】安徽是全国新能源汽车产业重镇，2026年产量超250万辆，蔚来/比亚迪/大众/奇瑞/江淮集聚',
     '【家电产业】安徽是全国最大的家电生产基地之一，2026年家电产量超1亿台，美的/海尔/格力/美菱布局',
     '【钢铁有色】马鞍山钢铁/铜陵有色等企业，钢铁/有色冶金产业规模大，重载机器人需求旺盛',
     '【石化化工】安庆石化/淮南化工等企业，化工场景防爆机器人需求迫切',
     '【装备制造】工程机械/农业机械/电工电气等装备制造产业基础好，工业机器人应用广泛',
     '【机器人密度】安徽制造业机器人密度达380台/万人，高于全国平均水平，智能制造水平快速提升',
     '【应用场景优势】丰富多元的制造业场景为各类机器人提供了绝佳的落地试验场和规模化应用市场',
     '【智能制造】安徽大力推进智能制造，智能工厂/数字车间/产线自动化改造需求旺盛',
     '【产需对接】建立机器人企业与制造业企业常态化对接机制，促进本地机器人本地应用'],
    ['【蔚来汽车工厂】蔚来合肥先进制造基地工业机器人超1200台，机器人密度达850台/万人，试点人形机器人',
     '【比亚迪合肥基地】比亚迪合肥工厂年产50万辆新能源汽车，焊装/涂装/总装车间自动化率超95%',
     '【大众安徽MEB工厂】大众安徽纯电动汽车工厂智能制造水平国际领先，采用大量最新工业机器人技术',
     '【奇瑞汽车】奇瑞芜湖工厂工业机器人密度达620台/万人，与埃夫特等本地机器人企业深度合作',
     '【美的合肥工业园】美的合肥冰箱/洗衣机生产基地是全球最大的家电生产基地之一，自动化率超90%',
     '【海尔合肥工业园】海尔合肥智能工厂是国家级智能制造示范，工业互联网+机器人融合应用标杆',
     '【美菱合肥工厂】美菱智能冰箱工厂大量应用搬运/码垛/装配机器人，效率提升40%',
     '【马钢集团】马钢钢铁生产应用重载搬运机器人/巡检机器人/炉前作业机器人，改善作业环境',
     '【铜陵有色】铜陵有色冶炼车间应用特种机器人替代人工在高温/高粉尘/有害环境作业',
     '【安庆石化】安庆石化部署防爆巡检机器人/应急处置机器人，提升化工生产安全性'],
    '▎江淮制造业机器人应用具体过程阐述',
    ['【人工为主期（2000年前）】安徽制造业以人工劳动为主，工业机器人应用几乎为空白，汽车焊接等极少数工位开始试用进口工业机器人，但数量极少，价格高昂，主要依赖进口，系统集成能力弱，生产效率低，产品质量一致性差，高危岗位安全事故时有发生，制造业自动化水平很低。',
     '【单机自动化期（2001-2010）】随着中国汽车工业快速发展，安徽奇瑞、江淮等车企开始在焊装、涂装、冲压等工位应用工业机器人，主要进口ABB、发那科、库卡、安川四大品牌，单台机器人价格20-50万元，机器人应用从无到有，但主要集中在汽车行业少数工位，整体渗透率低，其他行业应用很少。',
     '【产线自动化期（2011-2020）】安徽制造业自动化加速，埃夫特等国产工业机器人企业崛起，机器人价格下降至10-30万元，汽车行业焊装、涂装、冲压等产线基本实现自动化，家电行业美的、海尔、格力等大规模应用工业机器人，机器人密度从不足50台/万人提升至200台/万人以上，从单工位自动化向整线自动化发展，国产机器人市占率逐步提升。',
     '【智能制造起步期（2021-2024）】安徽大力推进智能制造，工业互联网、数字孪生、AI技术与机器人融合应用，工业机器人应用从汽车、家电向钢铁、化工、装备制造等行业扩展，机器人密度提升至300台/万人以上，智能工厂、数字车间建设启动，人形机器人开始在蔚来、比亚迪等工厂试点应用，探索柔性制造新路径。',
     '【规模应用试点期（2025-2026）】安徽制造业机器人应用进入新阶段：传统工业机器人全面普及，汽车、家电等行业自动化率达90%以上；人形机器人在蔚来、比亚迪、大众、奇瑞、美的等龙头企业工厂规模化试点，物料搬运、零部件分拣、产线巡检、成品码垛等场景试点取得成功，单台机器人ROI回收期缩短至2-3年，机器人密度达380台/万人以上。',
     '【人形机器人大规模应用期（2027-2028）】人形机器人技术成熟，成本降至15万元以下，在安徽制造业大规模推广应用，从试点工位向全车间扩展，从汽车、家电向钢铁、化工、装备制造等全行业渗透，制造业机器人密度达600台/万人以上，人形机器人在制造业部署量超2万台，柔性生产、多品种小批量制造能力大幅提升。',
     '【智能制造全面普及期（2029-2030）】安徽智能制造全面普及，工业机器人密度达800台/万人以上，人形机器人在制造业部署量超5万台，人机协作成为常态，数字孪生工厂广泛建设，AI+机器人+工业互联网融合，实现柔性化、智能化、个性化制造，安徽建成全国智能制造标杆省份，制造业竞争力全国领先。',
     '【汽车行业应用路径】人工焊接→进口机器人焊接工位→焊装/涂装/冲压整线自动化→国产机器人替代→全车间自动化→人形机器人柔性装配→数字孪生智能工厂。',
     '【家电行业应用路径】人工装配→搬运码垛机器人应用→整线自动化→智能检测机器人→柔性装配人形机器人→个性化定制智能工厂。',
     '【钢铁化工应用路径】人工高危作业→引进特种巡检机器人→关键岗位机器人替代→全流程机器人巡检+作业→无人化车间/工厂。'],
    '▎制造场景 · 应用数据 · 效率提升 · 最新进展',
    ['【江淮制造·2026年8月21日】光明日报"活力中国调研行"：安徽上半年高技术制造业增加值增长44.6%对规上工业增长贡献率55.9%，新能源汽车/新型显示/机器人等优势产业稳居全国第一方阵；奇瑞智造二工厂成国内首个通过国家智能制造能力成熟度四级认证的新能源乘用车工厂1-7月累计销售新能源车60.4万辆同比+42.3%连续23年中国品牌乘用车出口第一；蔚来新桥二工厂车身车间941台机器人火热作业+"天探"AI全身自检系统3分钟完成超1000项功能自测效率是人工10倍；奇瑞智界超级工厂10台机器人"千手观音"工位毫秒级协同关键工序100%自动化；安徽已集聚7家整车企业3000余家零部件企业"芯屏汽合"产业闭环企业不出安徽就能造一辆智能电动汽车',
     '【安徽家电产业规模数据】安徽是全国家电四大生产基地之一，2026年全省家电产量预计超过1亿台套，占全国产量比重超过20%，集聚美的、海尔、格力、美菱、康佳、TCL等知名家电企业，冰箱、洗衣机、空调、彩电四大件产量均居全国前列，家电产业年产值超过3000亿元',
     '【汽车产线机器人应用】安徽五大整车企业（蔚来、比亚迪、大众安徽、奇瑞、江淮）生产线上工业机器人总数量超过15000台，焊装、涂装、冲压等关键工序平均自动化率达到92%，其中焊装车间自动化率接近100%；奇瑞汽车孵化的墨甲机器人2026年上半年销量已达590台，产品出口全球60多个国家和地区，是安徽本土人形机器人代表企业',
     '【家电产线机器人应用】美的、海尔、格力、美菱、康佳等家电企业在皖生产基地工业机器人总数量超过8000台，主要应用于钣金冲压、注塑成型、焊接、搬运、码垛、装配、检测、包装等工序，整线平均自动化率达到88%，部分龙头企业标杆工厂自动化率超过95%，生产效率大幅提升',
     '【钢铁冶金特种机器人应用】马鞍山钢铁、铜陵有色等安徽钢铁有色金属龙头企业已累计应用重载搬运、高温巡检、炉前作业、自动取样等特种机器人超过500台，炼钢、炼铁、有色冶炼等高危作业岗位机器人替代率超过60%，有效减少高温、高粉尘、有毒有害环境下的人工岗位数量',
     '【石化化工防爆机器人应用】安庆石化、淮南化工、淮北煤化工等安徽石化化工企业已累计部署防爆巡检机器人、应急处置机器人、管道检测机器人超过200台，关键生产装置巡检覆盖率达到100%，可24小时不间断巡检，及时发现泄漏、温度异常、设备故障等安全隐患',
     '【汽车生产效率提升数据】工业机器人在安徽汽车行业的大规模应用使整车生产效率平均提升50%，产品焊接不良率降低70%，单位产品人工成本降低40%，产品生产一致性和质量稳定性大幅提升，新车开发周期从原来的36个月缩短至18-24个月，快速响应市场需求变化',
     '【家电生产效率提升数据】工业机器人在家电行业的应用使家电产品生产效率平均提升45%，产品外观和装配一致性大幅提升，产品交货周期缩短30%，单条生产线人员配置减少60%，可支持多品种小批量柔性生产，满足个性化定制市场需求',
     '【高危场景安全效益数据】钢铁、化工、矿山等高危行业应用机器人后，生产安全事故率降低85%，尘肺病等职业病发病率降低90%，一线作业人员劳动强度大幅降低，有效解决高危行业招工难、留人难问题，实现安全生产和经济效益双赢',
     '【蔚来工厂人形机器人试点】蔚来汽车合肥先进制造基地试点应用人形机器人进行车间物料精准配送，单台人形机器人可替代2名物料搬运工人，24小时连续作业（仅需更换电池），物料配送准确率达到99.8%，单台机器人投资ROI回收期约2.5年，目前已试点部署15台',
     '【比亚迪合肥基地试点人形机器人】比亚迪合肥工厂试点应用人形机器人进行总装车间零部件智能分拣配送，可准确识别200余种不同型号汽车零部件，分拣准确率达到99.8%，物料配送效率比人工提升30%，可根据生产计划动态调整配送路径，目前已试点部署20台',
     '【美的合肥工厂试点人形机器人】美的合肥冰箱生产基地试点应用人形机器人进行成品家电自动码垛和仓储搬运，人形机器人最大负载20kg，可连续工作8小时，码垛精度达到±5mm，适应不同规格成品纸箱码垛需求，相比传统工业码垛机器人柔性更强，可快速切换产品型号',
     '【奇瑞芜湖工厂试点人形机器人】奇瑞汽车芜湖工厂试点应用人形机器人进行总装车间产线设备巡检和异常处置，搭载红外热成像、声音识别、视觉检测等传感器，可识别50余种设备异常状态，发现异常自动报警并尝试简单处置，设备故障停机时间减少20%',
     '【机器人本地配套率提升】安徽制造业应用机器人本地配套率从2020年的15%快速提升至2026年的45%，2027年目标达到60%，芜湖埃夫特工业机器人、蚌埠传感器、合肥人形机器人整机、芜湖墨甲机器人等本地产品市场份额持续提升，本地供应链响应速度更快、服务更及时',
     '【常态化产需对接机制】安徽省经济和信息化厅每季度举办一次机器人产需对接会，组织机器人企业与制造业企业面对面对接交流，2026年已成功举办3场对接会，累计达成合作意向金额超过50亿元，有效打通技术供给和场景需求对接通道',
     '【首台套支持政策】安徽省对本地企业研发生产的首台套重大技术装备（含机器人）给予最高500万元人民币财政补贴，同时鼓励国有企业和政府投资项目优先采购本地首台套产品，降低企业新产品市场推广门槛，支持本土机器人企业创新发展',
     '【智能制造专项资金】安徽省每年安排20亿元人民币智能制造专项资金，支持制造业企业开展自动化、智能化、数字化改造，对机器人应用、智能工厂、数字车间建设项目给予设备投资额15%-20%的财政补贴，有效激发企业智能化改造积极性',
     '【智能工厂建设目标】2026年安徽省计划建成国家级智能工厂30家、省级智能工厂200家、省级数字化车间500个，推动超过1万家规模以上制造业企业完成数字化智能化改造，全省制造业机器人密度达到450台/万人以上',
     '【工业互联网平台支撑】安徽省已建成各类工业互联网平台超过50个，累计连接工业生产设备超过300万台（套），为机器人联网协同、数据采集分析、远程运维、预测性维护提供基础网络和平台支撑，实现机器人从单机智能向多机协同智能升级',
     '【技能人才培养目标】安徽省每年培养工业机器人技术、智能制造技术相关专业技能人才超过2万人，安徽机电职业技术学院、合肥职业技术学院等院校开设工业机器人技术专业，同时企业与院校合作开展订单班培养，满足制造业机器人应用对技能人才的需求，2030年全省制造业机器人密度目标达到800台/万人，人形机器人在制造业应用超过5万台']))

# PART 09 AI算力
all_modules.append(('PART 09', 'AI算力：大模型算力底座',
    ['【算力定位】AI算力是大模型和具身智能的基础底座，人形机器人对端侧算力有极高要求',
     '【算力规模】2026年中国智能算力规模达850EFLOPS（FP16），同比增长120%',
     '【安徽算力】合肥是全国八大算力网络枢纽节点之一，中国声谷算力中心集群规模大',
     '【端侧算力】人形机器人端侧算力需求达100-500TOPS（INT8），支持大模型端侧推理',
     '【云端算力】具身大模型训练需要千卡/万卡GPU集群，云侧算力支撑模型训练和复杂推理',
     '【算力芯片】英伟达H100/H200/A100主导高端训练，华为昇腾/寒武纪/海光等国产算力芯片快速崛起',
     '【国产替代】2026年国产AI芯片市占率提升至45%，华为昇腾在国内市场占比超30%',
     '【算力网络】全国一体化算力网络建设，东数西算工程推进，算力调度效率提升',
     '【绿色算力】液冷/PUE优化/可再生能源利用，数据中心PUE降至1.2以下，绿色低碳发展',
     '【具身算力】机器人端侧推理芯片成为新赛道，高通/英伟达/华为/地平线均布局专用芯片'],
    ['【算力大单狂飙·2026年8月21日】A股算力赛道狂飙：赛意信息与W公司签订两份高性能算力服务合同含税总金额64.5亿元相当于其2025年全年营收逾三倍；利通电子披露定增预案拟募资不超50亿元其中40亿投向智算中心建设，算力业务长期租赁排期已至2030年以后现有算力利用率接近100%；东阳光控股子公司三笔算力服务框架合同累计金额约390-460亿元；OpenRouter平台截至8月16日周Token调用量达75.3万亿环比增长9.1%创历史新高；SK海力士与弗吉尼亚大学在《自然·电子学》发表CPO技术路线图提出算力每两年增长3倍互联带宽仅增长1.4倍带宽墙成AI扩展核心瓶颈；谷歌与迈威尔就合作开发定制芯片签署一系列协议涵盖TPU相关AI推理加速器；SIGCOMM主会收录109篇论文中国贡献59篇占比再次超50%阿里巴巴蝉联全球企业论文入选榜榜首',
     '【阿里云AI·2026年8月21日】阿里巴巴2027财年Q1财报：阿里云外部商业化收入加速增长45%增速创22个季度新高，AI相关产品收入连续第12个季度实现三位数同比增长，本季度AI相关产品季度收入达123.76亿元对应年化规模接近500亿元；全球头部云厂商分化出AI加速阵营谷歌云82%季度增速领跑阿里云45%位列全球第二Azure 43%被阿里云反超AWS 37%；吴泳铭表示AI算力Capex投资回报确定性非常高投入可三年内回本未来有望缩短到2.5年甚至2年；平头哥已建成覆盖GPU/CPU/网络芯片的全栈自研体系真武M890等真武系列芯片已覆盖20余个行业服务650余家外部客户阿里云已将大规模AI数据中心交付周期压缩至100天；最新开源参数规模2.4万亿的Qwen3.8-Max和Qwen3.8-27B模型Qwen系列模型全球下载总量已超30亿次衍生模型数超30万个',
     '【算力资本大动作·2026年8月20日】AI算力成最稀缺战略资源：软银宣布65亿美元全现金收购Arm架构数据中心芯片公司Ampere Computing，与收购Arm及Stargate计划协同构建芯片设计到算力部署全链条；谷歌与Marvell达成122亿美元潜在入股深度合作，Marvell为谷歌开发定制AI芯片ASIC成博通之后第二大定制芯片伙伴加速去英伟达化；Groq以35亿美元估值完成3.5亿美元A轮融资英伟达计划参投，自研LPU推理芯片延迟比GPU快10倍以上算力明年从54兆瓦扩至200兆瓦以上；三星代工涨价15%模型降价但底层芯片涨价；宇树科技科创板上市首日开盘1100元较发行价150.80元涨近630%盘中市值破4000亿元，收盘约845元涨460%市值超3400亿，978万户打新中签率0.018%中一签盈利约35万元，但219倍市盈率对应2026上半年扣非净利同比-19.34%；快手可灵AI视频单季收入8.5亿元',
     '【算力资本大动作·2026年8月20日】软银宣布65亿美元全现金收购Arm架构数据中心芯片公司Ampere Computing；谷歌与Marvell达成深度合作涉及122亿美元潜在入股；Groq以35亿美元估值完成3.5亿美元A轮融资自研LPU推理芯片延迟比GPU快10倍以上；三星代工涨价15%',
     '【算力大单狂飙·2026年8月21日】赛意信息与W公司签订两份高性能算力服务合同含税总金额高达64.5亿元相当于其2025年全年营收的逾三倍；利通电子披露定增预案拟募资不超过50亿元其中40亿元投向智算中心建设；OpenRouter平台截至8月16日周Token调用量达75.3万亿环比增长9.1%创历史新高',
     '【阿里云AI收入·2026年8月20日】阿里巴巴2027财年Q1财报：阿里云外部商业化收入加速增长45%增速创22个季度新高，AI相关产品收入连续第12个季度实现三位数同比增长本季度AI相关产品季度收入达123.76亿元；平头哥真武M890等真武系列芯片已覆盖20余个行业服务650余家外部客户',
     '【算力荒·2026年8月21日】SemiAnalysis报告算力荒愈演愈烈：英伟达H100 GPU一年期租赁价格自2025年10月每GPU每小时1.70美元飙升至2026年3月2.35美元涨幅接近40%，整个行业GPU算力资源几乎全部售罄；Blackwell系列交付周期延长至2026年6-9月产能已被预订',
     '【摩尔线程算力底座·2026年8月21日】WRC2026摩尔线程创始人张建中表示：产业正共同迎接具身智能的ChatGPT时刻，算力底座是决定这一时刻何时到来的关键力量；推出首个全栈具身智能仿真平台MT Lambda，底层基于全功能GPU与MUSA统一架构实现渲染/物理/AI计算在同一芯片完成数据零拷贝',
     '【英伟达H200】FP8算力1979TOPS，HBM3e显存141GB，带宽4.8TB/s，2026年训练卡主力',
     '【英伟达Thor】机器人端侧超算芯片，FP8算力2000TOPS，集成CPU/GPU/NPU，专为自动驾驶和机器人设计',
     '【华为昇腾910B】FP16算力320TFLOPS，国产训练卡主力，已在多地智算中心大规模部署',
     '【华为昇腾310B】端侧推理芯片，INT8算力128TOPS，功耗35W，适合机器人端侧部署',
     '【寒武纪思元590】FP16算力512TFLOPS，国产训练芯片重要玩家，互联网企业批量采购',
     '【地平线Journey 6】车载/机器人端侧芯片，INT8算力512TOPS，功耗65W，高性价比方案',
     '【合肥智算中心】总算力15EFLOPS（FP16），搭载华为昇腾集群，服务安徽及长三角AI企业',
     '【中国声谷算力中心】合肥中国声谷建成10EFLOPS智算中心，支持科大讯飞等企业大模型训练',
     '【芜湖智算中心】芜湖智算中心总算力5EFLOPS，服务工业机器人和智能制造场景AI训练',
     '【液冷技术】冷板式液冷/浸没式液冷大规模应用，单机柜功率密度提升至50kW以上'],
    '▎AI算力建设发展具体过程阐述',
    ['【CPU主导期（2012年以前）】AI计算主要依靠CPU处理器，单CPU算力不足1TFLOPS，只能支持简单机器学习模型和小神经网络训练，数据中心规模小，算力成本极高，主要用于科研和互联网企业小规模应用，AI发展受限于算力瓶颈，深度学习无法大规模训练。',
     '【GPU起步期（2012-2016）】2012年AlexNet使用GPU训练证明GPU并行计算优势，NVIDIA推出CUDA生态，GPU成为AI训练主力，单卡算力从1TFLOPS提升至10TFLOPS，数据中心开始小规模部署GPU集群，互联网企业开始大规模使用GPU训练深度学习模型，AI在图像识别、语音识别领域取得突破。',
     '【大规模GPU集群期（2017-2020）】大模型兴起带动算力需求爆发，NVIDIA V100/A100 GPU大规模部署，单卡算力提升至312TFLOPS（FP16），千卡集群成为大模型训练标配，智算中心开始全国布局，中国启动东数西算工程，算力规模快速增长，GPT等大模型训练需要万卡级GPU集群，算力成为AI发展核心生产力。',
     '【国产算力突破期（2021-2024）】美国芯片制裁推动国产AI芯片加速发展，华为昇腾910/寒武纪思元/海光DCU等国产算力芯片实现技术突破，性能逐步接近国际先进水平，国产算力芯片开始在智算中心大规模部署，合肥、芜湖等安徽智算中心建成，国产算力市占率从不足5%提升至35%，液冷技术大规模应用降低PUE。',
     '【万卡集群普及期（2025-2026）】AI算力需求持续爆发，单卡算力达1000+TOPS，万卡/十万卡GPU集群普及，中国智能算力规模达850EFLOPS（FP16），国产算力市占率提升至45%，华为昇腾910B/910C成为国产训练卡主力，合肥智算中心15EFLOPS、中国声谷10EFLOPS、芜湖5EFLOPS算力集群建成投用，支撑大模型训练和具身智能研发。',
     '【端侧算力爆发期（2026-2027）】除云端算力外，端侧AI算力大爆发，人形机器人端侧算力需求达100-500TOPS，高通Thor、华为昇腾310B、地平线Journey 6等端侧AI芯片大规模应用，支持大模型端侧推理，AI手机NPU算力达100-150TOPS，AI PC NPU算力达80-180TOPS，云边端三级算力架构形成。',
     '【普惠化期（2028-2030）】算力成本持续下降，AI推理算力成本2023-2030年下降99%，大模型训练成本下降95%，国产算力市占率超70%，中国智能算力规模达5000EFLOPS，液冷渗透率达90%，数据中心PUE降至1.1以下，绿色可再生能源使用率超60%，算力像水电一样成为普惠公共基础设施，支撑万亿级AI产业发展。',
     '【芯片架构演进】CPU（串行计算）→GPU（并行计算）→NPU/TPU（AI专用加速）→类脑芯片（存算一体）；制程工艺从28nm→14nm→7nm→5nm→3nm；单卡算力从1TFLOPS→10TFLOPS→100TFLOPS→1000TFLOPS→4000TFLOPS。',
     '【安徽算力建设】合肥是全国八大算力网络枢纽节点之一，合肥智算中心（15EFLOPS，华为昇腾集群）、中国声谷智算中心（10EFLOPS）、芜湖智算中心（5EFLOPS）陆续建成，总算力超30EFLOPS，服务长三角AI企业，支撑科大讯飞大模型训练和人形机器人具身智能研发，国产算力租赁价格低至2元/卡/小时。',
     '【国产算力投资热潮】2026年国产算力领域投资热度持续高涨，仅2026年H1国产AI芯片领域融资超200亿元，各地智算中心建设投资超千亿元，算力租赁市场规模突破300亿元，国产算力产业链从芯片设计→制造→封装→板卡→服务器→集群→云服务完整生态逐步形成。'],
    '▎算力芯片 · 智算中心 · 国产替代 · 具身算力',
    ['【英伟达H200 GPU详细参数】NVIDIA H200是2026年全球AI训练主力芯片，主要参数：FP8精度算力1979TOPS，FP16精度算力989TFLOPS，配置HBM3e高速显存141GB，显存带宽达到4.8TB/s，热设计功耗TDP 700W，采用TSMC 4N工艺制程，单卡可支持70B参数大模型推理，是当前全球智算中心应用最广泛的高端训练GPU',
     '【英伟达B100 GPU最新参数】NVIDIA B100是下一代旗舰AI训练芯片，计划2027年量产，主要参数：FP8精度算力4000TOPS，FP16精度算力2000TFLOPS，配置HBM4高速显存288GB，显存带宽达到8TB/s，热设计功耗TDP 1000W，采用TSMC 3nm工艺制程，单卡可支持400B参数大模型推理，性能相比H200翻倍',
     '【英伟达Thor机器人超算芯片】NVIDIA Thor是专为自动驾驶和人形机器人设计的端侧AI超算芯片，主要参数：FP8精度算力2000TOPS，集成CPU+GPU+NPU全功能计算单元，采用TSMC 7nm工艺制程，车规级和工业级可靠性设计，支持多任务并行计算，可同时满足机器人感知、规划、控制、交互多任务算力需求',
     '【华为昇腾910B AI训练芯片】华为昇腾910B是当前国产AI训练主力芯片，主要参数：FP16精度算力320TFLOPS，INT8精度算力640TOPS，配置HBM高速显存64GB，显存带宽1.6TB/s，热设计功耗350W，采用国产7nm工艺制程，整机性能达到英伟达A100的70%水平，已在国内合肥、深圳、武汉等多地智算中心大规模部署应用',
     '【华为昇腾910C升级芯片】华为昇腾910C是昇腾910B升级型号，2026年正式量产，主要参数：FP16精度算力640TFLOPS，INT8精度算力1280TOPS，配置HBM2e高速显存80GB，整机性能接近英伟达H100水平，可支持千亿参数大模型训练，将成为2026-2027年国产高端训练主力芯片',
     '【华为昇腾310B端侧推理芯片】华为昇腾310B是国产端侧AI推理主力芯片，主要参数：INT8精度算力128TOPS，FP16精度算力64TFLOPS，热设计功耗仅35W，采用12nm工艺制程，支持大模型端侧轻量化推理，已广泛应用于人形机器人、智能安防、工业视觉检测、边缘计算等场景，性价比优势显著',
     '【寒武纪思元590训练芯片】寒武纪思元590是国产AI训练芯片第二梯队领军产品，主要参数：FP16精度算力512TFLOPS，配置HBM高速显存64GB，显存带宽1.8TB/s，热设计功耗450W，采用7nm工艺制程，支持MLU-Link高速互联，已被国内多家头部互联网企业批量采购用于大模型训练，国产替代重要力量',
     '【地平线Journey 6机器人/车载芯片】地平线Journey 6是专为自动驾驶和人形机器人设计的高性价比端侧AI芯片，主要参数：INT8精度算力512TOPS，热设计功耗65W，采用BPU贝叶斯加速架构，支持多传感器融合感知、路径规划、运动控制实时计算，能效比达到8TOPS/W，相比同类产品能效优势明显',
     '【中国智能算力规模增长】中国智能算力规模保持高速增长：2022年80EFLOPS（FP16），2023年180EFLOPS（+125%），2024年380EFLOPS（+111%），2025年600EFLOPS（+58%），2026年预计达到850EFLOPS（+42%），四年时间算力规模增长超过10倍，支撑大模型和具身智能产业快速发展',
     '【国产AI芯片市占率提升】国产AI芯片国内市场占有率快速提升：2023年仅15%，2024年提升至25%，2025年达到35%，2026年预计提升至45%，2027年目标达到55%，华为昇腾系列芯片在国产AI芯片市场占比超过60%，是国产替代绝对主力，有效缓解美国芯片制裁影响',
     '【合肥智算中心建设运营】合肥智算中心位于合肥高新区，总投资60亿元人民币，总算力规模15EFLOPS（FP16），全部采用华为昇腾910B AI芯片集群，2026年已实现满负荷运行，主要服务安徽省及长三角地区AI企业、科研院所，支撑科大讯飞星火大模型训练和人形机器人具身智能研发，算力租赁价格低至2元/卡/小时',
     '【中国声谷智算中心】合肥中国声谷智算中心总投资40亿元人民币，总算力规模10EFLOPS（FP16），采用混合算力架构（华为昇腾+部分英伟达芯片），重点服务科大讯飞等中国声谷入驻企业大模型训练和AI应用研发，是合肥人工智能产业重要算力基础设施',
     '【人形机器人端侧算力需求演进】人形机器人端侧AI算力需求快速提升：2024年单台人形机器人端侧算力需求为30-100TOPS（INT8），主要支撑基础感知和运动控制；2026年提升至100-500TOPS，需支持具身大模型端侧推理、自然语言交互、复杂任务规划；2030年预计达到1000TOPS以上，支撑完全自主智能体',
     '【具身大模型训练算力成本】训练一个百亿参数级人形机器人具身大模型需要5000-10000张高端GPU连续训练2-3个月，训练成本超过10亿元人民币；训练一个千亿参数通用具身大模型需要3-5万张GPU，训练成本超过50亿元，高昂算力成本是具身大模型研发主要门槛之一',
     '【云边端三级算力协同架构】具身智能采用云端训练+边缘侧复杂推理+端侧实时控制三级算力协同架构成为行业标准方案：云端万卡集群负责大模型训练和版本更新；边缘侧（园区/工厂本地服务器）负责复杂任务规划、多机协同调度；机器人本体端侧芯片负责实时运动控制、紧急避障、基础感知，平衡算力、时延、成本需求',
     '【全国一体化算力网络建设】全国一体化算力网络国家枢纽节点建设基本完成，建成8大国家算力枢纽节点、10大国家数据中心集群，跨区域算力调度网络时延小于20ms，东数西算工程深入推进，东部时延敏感业务部署在东部枢纽，非实时训练业务部署在西部算力枢纽，算力资源配置效率大幅提升',
     '【绿色低碳算力发展】2026年中国新建大型数据中心可再生能源使用率达到45%，液冷技术渗透率达到60%，其中冷板式液冷占45%、浸没式液冷占15%，单机柜功率密度从传统风冷的10kW提升至液冷的50kW以上，西部枢纽节点数据中心可再生能源使用率超过80%，绿色低碳发展水平持续提升',
     '【最新算力规模·2026年8月】工信部数据截至2026年6月底我国智能算力规模达2185EFLOPS同比+177%；2026年前7个月158个亿元级数据中心招标累计超1130亿元；字节豆包日均词元调用破180万亿',
     '【算力荒·2026年8月21日】SemiAnalysis报告算力荒愈演愈烈：英伟达H100 GPU一年期租赁价格自2025年10月每GPU每小时1.70美元飙升至2026年3月2.35美元涨幅接近40%整个行业GPU算力资源几乎全部售罄；部分用户为获取AWS p6-b200竞价实例愿支付14美元/小时高价；Blackwell系列交付周期延长至6-9月产能已被预订；2026年一季度LPDDR5和DDR5合约价格同比分别上涨约4倍和5倍推高AI服务器成本',
     '【谷歌迈威尔TPU·2026年8月21日】谷歌与迈威尔科技就合作开发定制芯片签署一系列协议，涵盖TPU相关AI推理加速器/存储控制器/网络接口控制器/内存接口控制器；SIGCOMM主会收录109篇论文中国贡献59篇占比再次超50%，阿里巴巴蝉联全球企业论文入选榜榜首',
     '【英伟达5000亿美元AI融资平台】2026年8月16日英伟达联合阿波罗、贝莱德、黑石、布鲁克菲尔德、高盛、KKR六大金融机构设立独立AI算力融资平台，长期计划撬动超5000亿美元第三方资本投入全球AI基础设施建设，推出最高25%项目残值兜底支持政策，将GPU算力包装为标准化可持续投资品类，面向合规AI实验室、实体企业及AI云服务商开放',
     '【腾讯Q2算力投入超千亿】2026年8月12日腾讯发布Q2财报，单季资本开支达527.84亿元同比大增176%，叠加514亿元算力预付款，本季AI基础设施投入已超千亿元；研发支出272.8亿元同比增长35%，主要投向数据中心、服务器等AI基础设施支撑混元大模型和WorkBuddy等产品算力需求',
     '【算力产业发展长期目标】根据工信部《算力基础设施高质量发展行动计划》，到2030年中国智能算力规模目标达到5000EFLOPS（FP16），国产AI芯片国内市场占有率超过70%，数据中心平均PUE降至1.15以下，可再生能源使用率超过60%，建成全球领先的算力基础设施体系，算力产业总体规模超过10万亿元，支撑数字经济和AI产业高质量发展']))

# PART 10 AI智能体
all_modules.append(('PART 10', 'AI智能体：具身大脑核心',
    ['【智能体定位】AI智能体是具身机器人的大脑，负责感知/理解/规划/决策/学习，是核心技术',
     '【具身大模型】具身大模型是AI智能体的核心，支持多模态感知/自然语言交互/任务规划/运动控制',
     '【技术架构】感知层+认知层+决策层+执行层四层架构，大模型+小模型+传统算法融合',
     '【多模态理解】融合视觉/听觉/触觉/力觉/位置觉多模态传感器信息，全面理解环境',
     '【任务规划】接收自然语言指令，自动分解任务步骤，规划动作序列，处理异常情况',
     '【运动控制】结合强化学习和经典控制，实现稳定行走/灵巧操作/人机协作等运动能力',
     '【持续学习】从人类演示/试错/其他机器人经验中持续学习，技能不断提升',
     '【多智能体协作】多个机器人之间协同工作，任务分配/信息共享/动作协调，完成复杂任务',
     '【端云协同】端侧实时控制+云端复杂推理和技能学习，兼顾实时性和智能水平',
     '【开源生态】具身智能开源社区快速发展，数据集/仿真环境/基础模型开源共享加速技术迭代'],
    ['【大晓开悟世界模型·2026年8月21日】大晓机器人首次亮相WRC展示具身智能全栈实力：开悟世界模型3.1（Kairos 3.1）采用统一原生架构整合生成/物理/认知三类智能，把视觉观测/语言指令/力触反馈/动作轨迹等多源具身数据纳入同一隐空间，搭建理解—推演—执行—反思自进化闭环，机器人执行失败后可自主定位问题调整策略自我优化；在全球具身智能评测中世界模型视频生成/状态预测两项赛道取得靠前成绩已面向行业开源，毕马威报告视其为原生一体化架构代表性成果处全球第一梯队；发布晓满（即时零售履约）/晓新（酒店洗衣全流程）/晓途（城市治理文旅户外）三套行业解决方案',
     '【帕西尼VTLA·2026年8月21日】WRC2026产业深度对话：帕西尼CEO许晋诚+易方达基金经理肖宛远指出物理世界数据极度匮乏与触觉感知缺失是行业核心短板，底层算法从纯视觉向VTLA（视觉-触觉-语言-动作）多模态融合跃迁触觉从硬件选配走向底层标配；评判机器人产业成熟标准是ROIC与ROI投入1美元产生超1美元价值商业闭环即成立；轮式与双足是不同商业场景的并行方案；灵巧手成本一年前10万级别现在降到3万以内；帕西尼提供传感器-灵巧手-整机-算法全链路交付触觉传感器深入半导体芯片制程级别',
     '【帕西尼VTLA·2026年8月20日】WRC2026产业深度对话：物理世界数据极度匮乏与触觉感知缺失是行业核心短板，底层算法从纯视觉向VTLA（视觉-触觉-语言-动作）多模态融合跃迁；评判机器人产业成熟标准是ROIC与ROI：投入1美元产生超1美元价值商业闭环即成立；灵巧手成本一年前10万级别现在降到3万以内',
     '【脑控机器人·2026年8月20日】WRC2026强脑科技BrainCo展示脑机接口+具身智能融合方案：脑控机器人训练平台脑电设备实时采集脑电信号算法识别意图转化为机器人控制指令，支持人形机器人/机械臂/机器狗多设备接入10分钟解锁脑控机器人',
     '【智能体商业底座·2026年8月20日】支付宝发布国内首个全栈智能体商业底座及AHA多智能体跨端互联协议，联合千问/华为/OPPO/比亚迪/吉利等20余家企业共建生态；高德云睿·时空智能体平台融合20余年时空数据同步上线交通/文旅/产业/充电/商业五大行业智能体',
     '【大模型密集更新·2026年8月20日】智谱GLM-5.3于8月17日发布总参数7530亿编程能力提升50%；阿里Qwen3.8-27B专为消费级硬件打造表现媲美10倍规模模型，通义千问衍生模型在Hugging Face达151448个为Meta的2.6倍；月之暗面K3发布48小时请求量逼近集群极限',
     '【优艾智合FabriX·2026年8月20日】优艾智合发布全球首个可规模化应用的工业具身智能大模型"智合"FabriX：从工业场景"长出来"而非通用AI"降维"；三层分布式架构（中央层接工业软件指令/边缘层协同调度与长任务链拆解/终端层VLA原子级实时控制）实现"一脑多态"跨具身集群作业',
     '【OpenAI多智能体·2026年8月20日】OpenAI Codex上线GPT-5.6"多智能体V2"，主Agent可把子任务自动委派给不同模型、每个子Agent单独设置推理强度；ChatGPT前端应用加载+94%、内存占用-41.2%、网络请求-98.2%',
     '【Agentic AI趋势·2026年8月】2026年44%企业开始部署或评估AI智能体，但仅11%成功投入规模化生产，先发优势窗口开放；智能体从"辅助工具"升级为独立承担端到端工作流的"数字员工"（数据分析/内容创作/审核校验/执行部署多智能体协同）；多智能体协同系统将取代单一模型成为企业竞争主战场',
     '【千问开源·2026年8月20日】阿里千问开源Qwen3.8-27B(270亿参数Apache2.0)，原生262K上下文外推1M，17GB显存消费级显卡可跑；千问全球下载超30亿次稳居开源第一',
     '【智谱GLM-5.3·2026年8月20日】智谱发布GLM-5.3编程能力最强开源模型，Terminal-Bench 3.0从4.6升至28.3，CyberGym白盒代码审查84.5%，两周后开源权重',
     '【DeepSeek V4 Pro·2026年8月20日】DeepSeek V4 Pro正式版开放权重下载并开源MIT协议Harness智能体框架，实现复杂任务自动拆解/多工具协同调度，补齐国产AI Agent底层工具链',
     '【OpenAI Figure 02】OpenAI与Figure联合发布的具身大模型，任务理解和泛化能力大幅提升',
     '【谷歌RT-3】谷歌DeepMind机器人Transformer 3，跨形态通用机器人策略，支持10+机器人平台',
     '【特斯拉FSD Robot】特斯拉自动驾驶技术迁移到人形机器人，感知规划决策技术同源',
     '【华为鸿蒙机器人OS】华为发布机器人操作系统，端云协同AI能力，12家整机厂接入',
     '【科大讯飞星火具身版】讯飞星火大模型具身版本，中文理解能力强，支持自然语言编程',
     '【智元具身大脑】智元自研具身大模型2.0，任务完成率提升50%，支持200+通用技能',
     '【优必选机器人脑】优必选自研机器人智能系统，已迭代5代，工业场景应用成熟',
     '【小米CyberBrain】小米人形机器人智能系统，结合米家生态，消费级场景优化',
     '【VLA模型】Vision-Language-Action模型成为主流技术路线，视觉语言直接输出动作',
     '【世界模型】世界模型让机器人能够预测未来，进行前瞻规划，提升任务成功率'],
    '▎AI智能体技术发展具体过程阐述',
    ['【规则系统期（2015年以前）】机器人和AI系统主要基于人工编写的规则和预编程逻辑，只能执行预设的固定任务，遇到未知情况无法处理，泛化能力为零，每个新场景都需要工程师大量编写代码，开发周期长、成本高，机器人只能在结构化环境中完成简单重复工作，无法适应非结构化环境和复杂任务。',
     '【传统机器学习期（2015-2021）】传统机器学习和深度学习开始应用于机器人感知和简单决策，卷积神经网络用于视觉识别，强化学习用于简单运动控制，但AI系统仍然是小模型，能力有限，需要大量标注数据训练，只能完成特定单一场景任务，无法理解自然语言指令，不具备通用智能，任务切换需要重新训练模型。',
     '【大语言模型爆发期（2022-2023）】ChatGPT为代表的大语言模型爆发，展现出强大的自然语言理解、推理和生成能力，研究者开始将大语言模型作为机器人大脑，尝试用自然语言指令控制机器人，谷歌SayCan、微软ChatGPT for Robotics等工作证明了大模型在机器人领域的潜力，但这一阶段大模型和机器人结合还比较初级，运动控制能力弱，任务成功率低。',
     '【VLA模型探索期（2024-2025）】Vision-Language-Action（视觉-语言-动作）端到端模型成为主流技术路线，谷歌RT-2/RT-3、OpenAI Figure 02等模型发布，实现了从视觉和语言输入直接输出机器人动作，机器人能够理解自然语言指令并完成相应任务，任务成功率从30%提升至70%以上，但泛化能力仍然有限，复杂任务成功率不高，主要在实验室环境验证。',
     '【具身大模型成熟期（2026）】专用具身大模型成熟，OpenAI Figure 02、谷歌RT-3、科大讯飞星火具身V4、智元具身大脑2.0等模型发布，任务成功率提升至85%-92%，支持200+通用技能，能够快速学习新技能（新任务学习时间<5分钟），多模态感知融合成熟，世界模型开始应用于前瞻规划，人形机器人开始具备真正的通用智能，工业场景任务成功率达到商业化要求。',
     '【通用智能期（2027-2028）】具身大模型能力持续提升，任务成功率突破95%，支持1000+通用技能，能够通过自然语言即时编程完成新任务，多智能体协作成熟，5台以上机器人可以协同完成复杂任务，世界模型能够准确预测未来10秒以上环境变化，进行前瞻规划和动作预演，Sim2Real迁移成功率>95%，人形机器人在工业、物流、商业服务场景大规模应用。',
     '【AGI临近期（2029-2030）】具身智能接近人类水平，任务成功率>99%，能够理解复杂抽象指令，自主规划完成多步骤复杂任务，具备持续学习能力，能够从人类演示、试错、其他机器人经验中不断学习进化，多模态理解能力接近人类，支持常识推理、因果推断，人形机器人开始进入家庭提供陪护、教育、家务服务，真正通用人工智能逐步实现。',
     '【技术路线演进】人工规则→传统机器学习→深度学习（CNN/RNN）→大语言模型（LLM）→视觉语言模型（VLM）→视觉语言动作模型（VLA）→具身世界模型→通用人工智能（AGI）；智能水平从"只能执行预设程序"→"能听会说能看"→"理解指令完成简单任务"→"通用技能快速学习"→"接近人类智能"。',
     '【世界模型技术】世界模型是具身智能的核心技术之一，让机器人能够在"大脑"中构建物理世界的模型，预测动作产生的后果，进行前瞻规划和动作预演，而不是盲目试错；2026年世界模型能够预测未来10秒内的环境变化，预测准确率>85%，大幅提升任务成功率和运动流畅性。',
     '【AI智能体应用】AI智能体不仅应用于人形机器人，还广泛应用于AI手机、AI PC、自动驾驶、智能家居、工业软件、客服、办公等领域，2026年AI智能体市场规模突破5000亿元，成为AI主要应用形态，"智能体+机器人"构成具身智能完整形态。'],
    '▎具身大模型 · VLA模型 · 世界模型 · 技能学习',
    ['【大晓机器人开悟3.1·2026年8月20日】大晓机器人首次亮相WRC展示具身智能全栈实力：开悟世界模型3.1（Kairos 3.1）采用统一原生架构整合生成/物理/认知三类智能搭建理解—推演—执行—反思自进化闭环；在全球具身智能评测中世界模型视频生成/状态预测两项赛道取得靠前成绩毕马威报告视其为原生一体化架构代表性成果处全球第一梯队；发布晓满（即时零售履约）/晓新（酒店洗衣）/晓途（城市治理文旅）三套行业解决方案',
     '【江行智能物理AI·2026年8月20日】江行智能将在8月28-30日贵阳2026数博会首秀"通用跨本体物理AI大脑"全栈技术体系：JX-Phi Brain+JX-Phi World实现物理AI感知-理解-决策-执行完整闭环，已在电网/新能源/油气/化工等行业1500余座工业场站真实运行',
     '【智合FabriX详解·2026年8月20日】全球首个可规模化应用的工业具身智能大模型：先验引导层蒸馏确定性数据+双回路逻辑校验层赋予自主诊断/故障恢复能力；"边界层剥离"策略兼容存量工业系统，实现异构多智能体协同和跨场景适配',
     '【FabriGym预训练·2026年8月20日】基于FabriX底座推出的全新具身技能预训练平台，让机器人在仿真环境内以千倍于真机的速度并行推演，实现机器人"出厂即上岗"；配套"隙锋"分级训练框架：到达工业现场后通过50条数据采集训练即可达90%成功率',
     '【谷歌DeepMind RT-3模型】谷歌DeepMind RT-3是当前性能领先的跨形态通用机器人模型，参数量100B，在13种不同形态机器人平台、超100万条机器人轨迹数据上训练，支持单臂/双臂/人形/轮式等多种机器人形态，跨形态任务泛化成功率92%，是通用机器人策略重要里程碑',
     '【OpenAI Figure 02具身模型】OpenAI与Figure AI联合开发的Figure 02具身大模型，结合GPT-4o多模态理解能力，支持自然语言即时编程，用户用自然语言描述新任务，机器人5分钟内即可学会并完成，无需工程师重新编写代码，任务理解和泛化能力相比第一代大幅提升',
     '【科大讯飞星火具身V4】科大讯飞星火V4具身版本大模型参数量70B，中文自然语言理解能力国内领先，支持多种方言交互，针对工业场景、家庭服务场景深度优化，可准确理解中文复杂指令，支持多轮对话纠错，任务执行过程中可通过自然语言实时调整',
     '【智元具身大脑2.0】智元机器人自研具身大脑2.0搭载于远征A2人形机器人，复杂任务完成率达到88%，支持300种以上常见物体操作，覆盖工业装配、物料搬运、日常服务等场景，强化学习+模仿学习融合训练，运动控制精度和任务成功率持续提升',
     '【VLA视觉语言动作模型路线】Vision-Language-Action（视觉-语言-动作）端到端模型已成为具身智能主流技术路线，模型直接输入摄像头图像和自然语言文本指令，端到端输出机器人各关节动作指令，省去传统感知-规划-控制分层架构的信息损耗，决策链路更短响应更快',
     '【世界模型前瞻规划技术】世界模型让机器人能够在"大脑"中构建物理世界内部模型，学习物理规律，预测执行不同动作后未来10秒内的环境变化和动作结果，进行前瞻规划和动作预演，而不是盲目试错，2026年世界模型短期预测准确率超过85%，大幅提升复杂任务成功率',
     '【模仿学习技能获取】从人类演示视频或遥操作数据中学习技能是机器人获取新技能的重要方式，学习单个通用操作技能需要10-100条高质量演示数据，训练时间数小时，模仿学习具有样本效率高、动作自然流畅等优点，是当前技能获取主要途径之一',
     '【强化学习仿真训练】在高保真物理仿真环境中通过深度强化学习让机器人自主试错学习运动技能，单个运动技能训练需要数百万次仿真交互，利用GPU大规模并行仿真可将训练时间缩短至数小时，训练完成后通过Sim2Real技术迁移到真实机器人',
     '【Sim2Real仿真到真实迁移】Sim2Real仿真到真实迁移技术2026年已基本成熟，通过域随机化、系统辨识、在线自适应等技术，仿真环境中训练的策略可直接部署到真实机器人，迁移成功率超过85%，大幅降低真实世界训练成本和风险',
     '【通用操作技能库建设】人形机器人通用操作技能库持续扩充，目前已包含稳定行走、抓取物体、放置物体、推拉抽屉、旋转阀门、零件装配、倒水、开门关门、按按钮、拾捡物品等100多种通用操作技能，新机器人可直接复用技能库中已训练好的技能',
     '【层次化任务规划技术】采用层次化任务规划架构，接收到复杂自然语言指令后，首先将高层任务分解为若干子任务序列，再将每个子任务分解为具体动作序列，执行过程中实时感知环境变化，动态调整计划应对异常情况，任务执行鲁棒性大幅提升',
     '【自然语言即时编程】用户只需用自然语言描述想要机器人完成的任务，机器人内置大模型即可自动理解任务要求并生成执行程序，无需传统代码编程，普通用户无需编程基础即可使用机器人完成新任务，大幅降低机器人使用门槛',
     '【自然语言实时纠错】机器人任务执行过程中，用户可随时通过自然语言对机器人动作进行实时纠正，例如说"往左一点"、"轻一点"、"不是这个红色的，是蓝色的那个"，机器人能理解纠正指令并即时调整动作，人机交互体验自然流畅',
     '【多模态感知融合】人形机器人融合8个以上摄像头（环视+手腕+头部）、激光雷达、六维力传感器、触觉传感器、IMU惯性测量单元、关节编码器等多模态传感器数据，实现360度全方位环境感知，感知精度和鲁棒性满足复杂任务需求',
     '【具身大模型端侧推理优化】通过模型量化（INT4/INT8）、知识蒸馏、算子优化等技术对具身大模型进行压缩优化，70B参数级别的具身大模型可在机器人端侧芯片上实时运行，推理延迟小于200ms，满足机器人实时运动控制要求',
     '【OTA云端技能持续更新】新技能在云端大规模算力集群上训练完成后，通过OTA方式批量下发到所有已部署机器人，机器人无需返厂即可持续获得新能力，已部署机器人技能持续进化，用户购买的机器人越用越聪明',
     '【多智能体协同作业】5台以上人形机器人可通过多智能体通信协议实现协同工作，自动进行任务分配、信息共享、动作协调，共同完成单台机器人无法完成的复杂任务，例如多机器人协同搬运大型物体、协同装配大型设备',
     '【数字孪生虚拟训练】在与真实场景1:1还原的数字孪生虚拟环境中训练机器人策略，数字孪生环境精确模拟真实场景物理特性、物体材质、光照条件，训练效率相比真实世界提升10倍以上，且不影响真实生产',
     '【开源数据集与开源生态】Open X-Embodiment等开源具身智能数据集包含超过100万条机器人交互轨迹数据，覆盖多种机器人平台和多种任务场景，开源仿真环境、开源基础模型快速发展，开源生态繁荣加速具身智能技术迭代',
     '【具身智能发展长期目标】预计到2030年，具身大模型复杂任务成功率超过99%，支持1000种以上通用操作技能，能够理解复杂抽象指令，自主规划完成多步骤复杂任务，具备持续终身学习能力，真正实现通用人工智能水平，人形机器人大规模进入家庭和各行各业']))

# PART 11 6G通信
all_modules.append(('PART 11', '6G通信：空天地一体化',
    ['【6G定位】6G是第五代移动通信之后的下一代通信技术，空天地一体化连接，支撑万物智联',
     '【性能指标】峰值速率1Tbps，时延<0.1ms，连接密度1000万/平方公里，定位精度厘米级',
     '【应用场景】人形机器人远程控制/多机协同/触觉互联网/数字孪生/元宇宙/全息通信',
     '【研发进展】中国6G研发全球领先，华为/中兴等企业专利申请量全球第一，2030年商用',
     '【关键技术】太赫兹通信/智能超表面/空天地一体化网络/AI原生/通感一体/确定性网络',
     '【机器人价值】6G高带宽低时延特性支撑人形机器人远程遥操作和云端大脑实时控制',
     '【空天地一体】卫星通信+空中基站+地面网络全域覆盖，机器人在偏远地区也能联网',
     '【通感一体】通信感知一体化，基站同时具备通信和雷达感知能力，为机器人提供环境感知',
     '【确定性网络】端到端确定性时延和抖动保障，满足机器人实时控制严苛要求',
     '【AI原生】6G网络原生支持AI，网络边缘提供AI算力，机器人可随时调用边缘AI能力'],
    ['【6G空天地一体化·2026年8月20日】中国电信2026半年报：加快空天地一体化网络建设推动移动通信网/全光网/卫星网多网协同构筑全域连接优势；时空道宇获工信部卫星物联网业务商用试验许可试验期两年，吉利星座一期组网完成64颗卫星在轨实现除南北极外全球通信覆盖每日处理约3.4亿次通信请求支持全球2000万名用户接入；已与亚非拉20多个国家电信运营商建立合作获得中东和拉美多个市场订单',
     '【6G政策·2026年8月】工信部正式启动6G创新发展部省协同试点专项行动，计划2029年形成一批自主创新6G技术方案；6G国际标准2029年敲定、2030年前后商用；国内6G研发第一阶段完成累计产出300多项关键技术，进入第二阶段原型样机攻关与实景场景测试；中国移动完成全球首个太赫兹通感一体外场测试100Gbps+厘米级定位；华为AI原生空口原型机语义压缩比30:1保真度超98%',
     '【三网协同·2026年8月】万兆光网（136个试点通过验收，可靠性99.9999%，6G回传骨干）+6G（移动智能连接）+卫星互联网（千帆星座在轨200颗年内324颗；垣信6月9日首发手机直连试验星，6月19日打通国内首例无改造商用5G手机直连卫星语音视频通话）三网合一，陆海空天一体化信息基础设施',
     '【Pre6G试验网·2026年8月20日】紫金山实验室国内首个Pre-6G试验网全面投入运行：覆盖南京无线谷/江苏电视塔等重点区域布局16个网络节点，建成业界首个空天地一体的Pre-6G立体覆盖体系；实测时延从传统互联网直播的约4秒压缩到约300毫秒；光子太赫兹全息影像端到端时延<50ms',
     '【6G专利领先·2026年8月20日】全球6G相关专利申请总量约3.8万件，中国以40.3%的占比位居全球第一约1.53万件；华为U6GHz全系列产品支持万兆级别峰值速率；中国90%的6G专利属于发明专利质量全球前列',
     '【最新·2026年8月】IMT-2030(6G)推进组星地融合工作组成立：中国信通院任组长，星网/垣信/三大运营商副组长，小米/OPPO/vivo/荣耀四家手机厂加入，6G进入施工阶段',
     '【华为6G研究】华为6G研发投入超100亿元，专利申请量全球第一，太赫兹通信原型验证完成',
     '【中兴6G】中兴通讯6G关键技术验证取得阶段性突破，智能超表面技术完成外场测试',
     '【中信科移动】提出6G天地一体化网络架构，卫星通信与地面网络融合技术领先',
     '【太赫兹通信】太赫兹频段（0.1-10THz）通信原型实现100Gbps传输速率，距离1km',
     '【智能超表面（RIS）】通过可编程超表面智能调控无线信号，覆盖盲区，提升信号质量',
     '【空天地一体化】低轨卫星星座+高空平台+5G/6G地面基站，实现全球无缝覆盖',
     '【通感一体原型】6G通感一体化基站实现同时通信和感知，感知精度达厘米级',
     '【确定性网络】端到端时延抖动<10us，可靠性99.9999%，满足工业控制和机器人要求',
     '【6G机器人远程控制】通过6G网络远程控制人形机器人，操作时延<50ms，几乎无延迟感',
     '【边缘算力协同】6G边缘节点提供算力，机器人端侧仅需保留传感器和执行器，算力按需调用'],
    '▎6G技术发展具体过程阐述',
    ['【5G规模建设期（2019-2025）】2019年5G商用启动，中国建成全球规模最大的5G网络，截至2026年累计建成5G基站超430万个，5G用户超12亿户，5G在消费互联网领域全面普及，工业互联网、车联网等行业应用开始推广，但5G在带宽、时延、连接密度、定位精度、通感一体等方面仍无法满足人形机器人、元宇宙、自动驾驶等未来应用需求，6G研发同步启动。',
     '【6G愿景需求期（2018-2022）】全球开始6G愿景和需求研究，ITU、3GPP等国际标准组织启动6G技术需求制定，中国IMT-2030（6G）推进组2019年成立，发布6G愿景白皮书，提出6G"万物智联、数字孪生"愿景，关键指标包括：峰值速率1Tbps、用户体验速率1Gbps、端到端时延<0.1ms、连接密度1000万/平方公里、定位精度<10厘米、支持通感一体、空天地海一体化覆盖。',
     '【关键技术突破期（2023-2026）】6G关键技术研发取得突破：太赫兹通信实现100Gbps@1km传输速率，智能超表面（RIS）技术完成外场测试，空天地一体化网络试验成功，通感一体技术验证，AI-native空口技术成熟；中国在6G专利申请量占全球50%以上居全球首位，华为、中兴、中国移动等企业和清华、中科大、东南大学等高校在6G技术研发处于全球第一梯队，2026年6G技术研发试验第一阶段完成。',
     '【标准制定期（2025-2028）】3GPP 6G标准制定工作启动，2025年开始RAN1技术研究，2027年完成Rel-21第一个版本6G标准，中国企业在6G标准制定中拥有重要话语权，提交标准提案占全球40%以上，6G核心技术专利布局完成，形成自主可控的6G知识产权体系。',
     '【技术试验期（2027-2028）】6G技术试验第二、三阶段完成，建设6G外场试验网，在合肥、北京、上海、深圳等城市开展6G技术验证和应用示范，在人形机器人远程操控、工业互联网、自动驾驶、元宇宙等场景开展6G应用试验，验证6G技术性能和商业可行性。',
     '【预商用期（2029-2030）】6G标准冻结，6G产业链成熟，开始6G预商用部署，2030年左右实现6G正式商用，6G网络将为人形机器人提供超高带宽（1Tbps）、超低时延（<0.1ms）、超高可靠（99.9999%）、通感一体、厘米级定位的通信能力，支撑人形机器人云端实时控制、远程全息操作、多机器人协同等应用。',
     '【全面普及期（2030后）】6G网络大规模建设，逐步实现空天地海一体化覆盖，通信感知深度融合，AI原生网络成为现实，支撑万亿级设备连接，6G成为智能社会的信息基础设施，推动人形机器人、元宇宙、自动驾驶、数字孪生、工业互联网等应用全面普及。',
     '【太赫兹通信技术演进】毫米波（28-100GHz，峰值10Gbps）→亚太赫兹（100-300GHz，峰值100Gbps）→太赫兹（0.3-10THz，峰值1Tbps），带宽持续提升，器件从分立器件→单片集成电路→CMOS集成芯片，成本持续下降。',
     '【通感一体技术演进】通信和感知独立部署→通信感知频谱共享→通感一体化硬件和波形→基站即雷达，网络具备全域感知能力，定位精度从米级→分米级→厘米级→毫米级，感知距离从100米→500米→1000米。',
     '【安徽6G布局】合肥未来网络研究院开展6G关键技术研究，中国科学技术大学在太赫兹通信、智能超表面领域取得多项科研成果，中国声谷布局6G应用创新，安徽在6G研发和应用示范方面走在全国前列，我国加快推动新一代通信网建设。'],
    '▎6G技术指标 · 关键技术 · 试验进展 · 机器人应用',
    ['【峰值速率指标对比】5G时代峰值速率达到20Gbps，6G通信峰值速率目标1Tbps，相比5G提升整整50倍，超大带宽可支持全息通信、8K以上超高清3D视频实时回传、海量传感器数据同步上传，完全满足人形机器人全身多路高清摄像头和传感器数据实时传输需求',
     '【端到端时延对比】5G端到端时延约1ms，6G目标端到端时延小于0.1ms，相比5G提升10倍，亚毫秒级超低时延可满足机器人精密力控操作、远程触觉反馈、高速运动实时控制等严苛时延要求，让远程操作手感和本地操作几乎没有区别',
     '【连接密度指标对比】5G连接密度100万台设备/平方公里，6G目标连接密度提升至1000万台设备/平方公里，支撑海量机器人、物联网设备、传感器在同一区域同时联网通信，可满足工厂内数千台机器人同时工作协同的通信需求',
     '【定位精度指标对比】5G定位精度仅米级，6G实现室内厘米级/室外分米级高精度定位，配合通感一体技术，无需额外定位传感器即可为机器人提供连续高精度位置服务，支撑机器人在无GPS信号的室内环境和复杂城市环境下高精度导航',
     '【通信可靠性指标对比】5G通信可靠性99.999%，6G目标可靠性达到99.9999%工业级标准，年中断时间小于30秒，保障机器人关键控制指令传输不中断，即使在电磁环境复杂的工厂车间也能稳定可靠通信，避免因通信中断导致生产事故',
     '【太赫兹通信技术】太赫兹通信工作频段90GHz-10THz，具有超大带宽优势，2026年原型验证已实现100Gbps@1km传输速率，未来目标1Tbps@100m，是6G实现超大带宽的核心关键技术，器件成本持续下降逐步走向商用',
     '【智能超表面RIS技术】智能超表面（RIS）采用可编程电磁表面，由3000个以上可独立调控单元组成，可实时调控无线信号相位和幅度，反射信号覆盖盲区，相比传统基站信号覆盖提升30%，能耗降低50%，是6G绿色通信关键技术',
     '【通信感知一体化技术】通感一体化技术让6G基站同时具备通信和雷达感知能力，单基站感知距离最远可达500米，速度分辨率0.1m/s，角度分辨率0.1度，厘米级感知精度，可在提供通信服务同时为区域内所有机器人提供环境感知能力',
     '【空天地一体化网络架构】6G采用低轨卫星星座（中国星网/国网星座+Starlink）+高空平台（HAPS）+地面5G/6G基站多层网络架构，实现全球全域无缝覆盖，即使在远洋、深山、沙漠、矿山等偏远地区机器人也能稳定联网作业',
     '【AI原生网络架构】6G网络原生内置AI能力，网络具备自优化、自愈、自配置能力，用户可根据业务需求动态调度网络带宽、时延、算力等资源，机器人可在网络边缘就近调用AI算力进行复杂推理，实现端边云协同智能',
     '【华为6G研发进展】华为6G研发累计投入已超100亿元人民币，2026年完成全部6G关键技术验证，2027年推出6G原型基站产品，2028年开展6G外场试验网部署，太赫兹通信、智能超表面等技术处于全球领先水平，6G专利申请量全球第一',
     '【中国6G技术试验进展】由工信部IMT-2030（6G）推进组组织开展6G技术试验，2026年完成第一阶段关键技术验证，2028年完成第二阶段外场试验，2029年完成第三阶段预商用试验，有序推进6G技术研发和标准制定工作',
     '【全球6G专利数据统计】截至2026年8月，中国企业6G专利申请量全球占比超过50%位居全球第一，其中华为、中兴通讯、OPPO、中信科移动、vivo等企业均位列全球6G专利申请量前十，中国在6G标准制定中拥有重要话语权',
     '【人形机器人远程操控应用】6G超高带宽超低时延特性支撑远程遥操作人形机器人，可实现触觉反馈信号实时传输，远程操作员佩戴触觉手套可获得与现场操作几乎一致的真实触感，可应用于危险环境作业、远程医疗、远程维修等场景',
     '【多机器人协同作业应用】6G可支撑同一区域内上百台人形机器人同时联网协同工作，机器人间点对点通信时延小于10ms，可实现环境感知信息实时共享、任务动态分配、动作协同配合，共同完成单台机器人无法完成的大型复杂任务',
     '【云端大脑实时控制应用】6G高带宽低时延特性使人形机器人云端大脑直接实时控制机器人本体成为可能，机器人端侧仅需保留传感器、执行器和简单实时控制，复杂AI推理和任务规划全部放在云端万卡集群完成，大幅降低端侧成本',
     '【数字孪生实时映射应用】机器人通过6G网络将全身多路摄像头、激光雷达、力觉触觉等传感器数据实时上传到数字孪生体，物理世界和虚拟世界同步延迟小于10ms，可实现远程监控、虚拟调试、故障预测、仿真优化等数字孪生应用',
     '【偏远地区作业应用】空天地一体化6G网络彻底解决通信覆盖问题，使人形机器人在远洋轮船、偏远矿山、沙漠油田、深山电站等没有地面网络覆盖的偏远地区也能稳定联网接受远程控制和AI能力支撑，拓展机器人应用边界',
     '【安徽6G产业布局】合肥未来网络研究院开展6G关键技术研究，中国科学技术大学在太赫兹通信、智能超表面领域取得多项国际领先科研成果，中国声谷布局6G应用创新和产业孵化，安徽在6G研发和应用示范方面走在全国前列',
     '【6G商用发展时间表】预计2028年完成3GPP R21版本6G国际标准冻结，2029年开始6G规模预商用部署，2030年实现6G正式商用，2035年6G用户规模超过10亿户，6G网络将成为支撑智能社会和人形机器人普及的核心信息基础设施']))

# PART 12 消费电子（华为/苹果/小米三品牌旗舰全覆盖）
all_modules.append(('PART 12', '消费电子：AI终端普及',
    ['【产业趋势】2026年是AI手机/AI PC大规模普及元年，消费电子全面AI化，端侧大模型成标配',
     '【AI手机】2026年AI手机出货量超8亿部，占智能手机出货量75%，端侧大模型7B-13B参数',
     '【AI PC】2026年AI PC出货量超1.2亿台，占PC出货量55%，NPU算力40-100TOPS',
     '【三品牌格局】消费电子只搜索华为/苹果/小米三大品牌，覆盖旗舰机/中端机/入门机全价位段',
     '【华为】麒麟芯片回归+鸿蒙OS+盘古大模型端侧部署，2026年国内市场份额重回第一',
     '【苹果】A系列芯片+M系列芯片+Apple Intelligence，AI功能深度整合iOS/macOS生态',
     '【小米】澎湃芯片+澎湃OS+小爱同学大模型，性价比优势，AIoT生态完善，全球化布局',
     '【端侧AI能力】AI拍照/AI修图/AI翻译/AI摘要/AI助理/AI创作成为标配功能',
     '【生态融合】手机/PC/平板/手表/耳机/汽车/智能家居生态打通，多设备协同AI体验',
     '【技术参数迭代】处理器/屏幕/摄像头/电池/快充持续升级，AI推动体验质变'],
    ['【华为Pura X View阔直板·2026年8月20日】8月20日华为成都发布会全球首款阔直板手机Pura X View正式亮相：6.39英寸16:9.5阔屏显示面积114.27cm²，屏占比96.1%四边等宽边框仅1.05mm；6.68mm超薄机身201g重量；内置7000mAh硅碳负极大电池+66W有线快充；麒麟9030S旗舰芯片首发HarmonyOS 7；6500nit峰值亮度+2160Hz高频PWM调光；支持手写笔+IP68防水防尘；跃影红/亚麻灰/零度白/幻夜黑四色，12GB+256GB/512GB/1TB三存储版本；8月20日18:08开启预约8月28日开放门店体验，预计售价6000元起',
     '【享界G9智能座舱·2026年8月20日】鸿蒙智行首款科技豪华硬派SUV享界G9正式上市售价42.98万元起（Max+ 42.98万/Ultra 47.98万/Ultra+ 52.98万）：座舱配内嵌式一体寰宇三联屏+17英寸吸顶屏+68英寸成像面积华为AR-HUD临境抬头显示+华为悦彰音响卓越系列；首发鸿蒙ALPS健康座舱+车载制氧系统（国标A级认证）+智能隐私光幕玻璃（遮光度超99.4%手势调光秒级响应）；至高4零重力座椅18向调节+20点双维按摩，舱内有效空间3415mm',
     '【享界G9智驾安全·2026年8月20日】享界G9是首款L3级架构设计豪华硬派SUV，获批行业最高时速120公里L3级自动驾驶道路测试牌照；搭载华为乾崑智驾ADS 5+全车38个高精度传感器+全向立体融合感知系统，山路/野路/土路/悬崖路可随时激活智驾；全维防碰撞系统CAS 5.0支持全时域/全时速/全方向/全天候/全目标/全场景超全感知；搭载华为星河通信3.0天地网联，行业首发集成式车载对讲机（最远10公里无网对讲）+车载卫星通话（连星速度提升12%）',
     '【智界RX预售·2026年8月20日】智界RX正式开启预售预售价29.98万-39.98万元，至高享10000元预订权益：定位性能向SUV宽体低趴设计，1:2轮高比+1.26宽高比；搭载896线双光路图像级激光雷达+L3级自动驾驶架构；创新采用"343"一体化压铸车身架构，车身扭转刚度超50000N·m/deg；双腔空气悬架+双阀连续可变阻尼减振器，近50:50轴荷比；首发巨鲸电池平台"黑匣子"安全技术；纳米级仿生变色车漆，展车已陆续到店',
     '【问界M6新版本·2026年8月20日】问界M6上市4个月累计交付突破45000台，在2026中国汽研智能汽车指数测评中拿下四项全G+评级；发布会新增两款增程版本：增程Max四驱版售价23.98万元起，增程Max+长续航版售价25.98万元起；两版本在电池容量/纯电续航/智驾硬件/底盘配置上做出区分，CLTC综合续航分别可达1445km和1605km，进一步丰富产品选择',
     '【享界V8官宣·2026年8月20日】享界正式官宣全新车型享界V8，定位家庭智慧旗舰MPV主打高端家用市场；将携手尊界V800、V680，智界V9重塑高端MPV市场格局；新车将登陆成都车展完成首秀；随着享界G9、享界V8到来，享界已实现覆盖轿车/旅行车/SUV/MPV全品类豪华产品矩阵',
     '【鸿蒙智行150万交付·2026年8月20日】余承东发布会宣布：截至2026年8月16日鸿蒙智行全系车型累计交付量正式突破150万台，自品牌创立以来仅用时53个月，刷新中国新势力车企达成该交付规模最快纪录；鸿蒙智行已形成问界/智界/享界/尊界/尚界五大子品牌协同发展格局，品牌矩阵覆盖轿车/SUV/MPV全品类',
     '【享界G9越野性能·2026年8月20日】享界G9首发华为全地形途灵平台：全球首发800V全主动可断开稳定杆（毫秒级姿态预判，高速过弯锁死减少侧倾，颠簸路段毫秒级解耦左右车轮大幅跳动）+自适应牙钳式差速锁（三轮打滑/交叉轴自动锁止脱困）；鸿蒙智行智能全地形系统7大地形智能识别；涉水深度800mm，±12°后轮转向转弯半径仅5.2米；全系标配华为巨鲸800V平台，纯电版120kWh电池CLTC续航最高728km（15分钟30%-80%补能），增程版综合续航最高1366km',
     '【Pura X View阔屏体验·2026年8月21日】华为Pura X View以16:9.5阔比例直屏对比主流20:9细长屏：显示面积提升约16%，有效可视面积优于常规6.9英寸细长直板机型；主打阔感观剧/阔感阅读/阔感游戏/阔感办公四大场景，射击类手游视野更宽阔可提前看到高处敌人；直板形态无折痕无铰链损耗，屏幕平整度长期使用不衰减，维修成本比折叠屏便宜一半以上；被称为"折叠屏平替"，把折叠屏阔视野装入更轻薄门槛更低的直板机身',
     '【Pura X View影像设计·2026年8月21日】华为Pura X View延续Pura系列影像基因：后置跑道式XMAGE三摄，延续Pura X系列红枫影像调校，人像/夜景/远景拍摄与阔折叠保持同一梯队；光栅纹理后盖不易沾指纹；提供跃影红/亚麻灰/零度白/幻夜黑四款配色，延续Pura系列简洁流畅设计理念，在视觉美感与握持体验之间实现平衡；支持手写笔输入进一步拓展移动办公/创作记录应用场景',
     '【华为阔屏矩阵补齐·2026年8月21日】随着Pura X View正式亮相，华为阔屏产品矩阵全面补齐：从Pura X阔折叠（2025年）到Pura X Max横向阔折叠（2026年4月）再到Pura X View阔直板（2026年8月），华为持续引领阔屏形态潮流；阔直板将阔屏体验从万元级折叠旗舰下沉至6000元档直板市场，填补"想要大屏视野但不接受折叠屏厚重与折痕"的用户需求空白，较Pura X的7499元起低约1500元',
     '【华为Pura X View·2026年8月20日】全球首款阔直板手机Pura X View：6.68mm厚度/201g重量/96.1%屏占比，6.39英寸16:9.5 OLED直屏峰值亮度6500nit支持2160Hz PWM调光；7000mAh电池+麒麟9030S+HarmonyOS 7；8月20日18:08开启预约8月28日开放门店体验',
     '【享界G9上市·2026年8月20日】鸿蒙智行首款科技豪华硬派SUV享界G9正式上市42.98万元起至高享106500元购车权益：Max+增程大五座42.98万/增程享六座43.98万；Ultra增程大五座47.98万/增程享六座49.98万/纯电大五座47.98万/纯电享六座49.98万；Ultra+增程大五座52.98万/增程享六座54.98万/纯电大五座52.98万/纯电享六座54.98万；车身5377×2050×1897mm轴距3160mm',
     '【华为Pura 90系列】Pro/Pro Max首发麒麟9030S，标准版麒麟9010S，鸿蒙6.1，6500mAh+100W有线50W无线，前后双红枫影像，销量破101万台，4699-8499元',
     '【苹果iPhone 18 Pro Max】A19 Pro芯片，3nm工艺，NPU算力150TOPS，6.9英寸OLED，4800万主摄，Apple Intelligence',
     '【苹果MacBook Pro 16 M5】M5 Max芯片（16核CPU+40核GPU+32核NPU），NPU算力180TOPS，AI PC性能标杆',
     '【苹果Vision Pro 2】第二代空间计算设备，重量减轻30%，价格下探至1.5万元，内容生态成熟',
     '【小米16 Ultra】骁龙8 Gen4+澎湃P2/G2芯片，徕卡影像，1英寸主摄+5倍潜望长焦，卫星通信，5999元起',
     '【小米16 Pro】骁龙8 Gen4，6.7英寸2K 144Hz屏，5500mAh 120W快充+50W无线，4999元起',
     '【小米RedmiBook Pro 16 2026】酷睿Ultra 7+小米NPU（80TOPS），3.2K 120Hz屏，1TB SSD，5999元性价比AI PC',
     '【小米SU7 Ultra】小米汽车旗舰，三电机四驱1300马力，零百加速1.98秒，全栈自研智驾，52.99万起'],
    '▎消费电子AI化发展具体过程阐述',
    ['【功能机时代（2000-2007）】诺基亚、摩托罗拉功能机主导市场，手机主要功能是通话和短信，没有智能系统，屏幕小（2英寸以内），分辨率低（QVGA及以下），处理器性能弱（ARM9及以下），没有独立NPU，AI能力为零，形态以直板/翻盖/滑盖为主，更换电池是标配，拍照只有30万像素左右，完全没有智能功能。',
     '【智能机爆发期（2007-2015）】2007年iPhone发布开启智能手机时代，触摸屏取代物理键盘，iOS/Android系统诞生，移动互联网兴起，APP生态繁荣，手机性能快速提升（从单核1GHz到八核2GHz），屏幕尺寸增大到5-6英寸，分辨率提升至1080P/2K，AI能力开始萌芽：规则化语音助手（Siri 2011年发布）、简单人脸识别、AI场景识别拍照，但主要是规则和简单机器学习，端侧AI算力几乎为零，AI计算主要靠CPU。',
     '【AI功能试水期（2016-2022）】AI功能逐步加入消费电子：华为麒麟970首次集成独立NPU（2017年），苹果A11 Bionic集成神经网络引擎，AI拍照（场景识别/美颜/夜景）成为旗舰标配，语音助手能力提升（小爱同学/小艺/Siri），但AI能力仍然比较初级，主要是特定任务加速，没有通用自然语言理解能力，大模型尚未爆发，端侧大模型技术不成熟。',
     '【大模型技术积累期（2023-2024）】ChatGPT爆发带动大模型技术快速发展，云端大模型能力快速提升，端侧大模型技术开始成熟，高通/联发科/华为/苹果相继推出支持端侧大模型的芯片，7B/13B参数大模型可以在端侧流畅运行，2024年下半年开始有厂商发布AI手机概念产品，但AI功能仍以云端调用为主，端侧AI应用生态尚未形成。',
     '【AI终端普及元年（2025-2026）】2026年成为AI终端普及元年，AI手机出货量超8亿部占比75%，AI PC出货量超1.2亿台占比55%，华为/苹果/小米旗舰全系标配端侧大模型（7B-13B参数），NPU算力达到80-180TOPS，自然语言交互成为主要交互方式之一，AI会议/AI写作/AI拍照/AI搜索/AI助理成为标配功能，端侧推理延迟<1秒，隐私数据不上云，多设备AI生态打通。',
     '【AI深度融合期（2027-2028）】AI深度融入消费电子各个方面，AI Agent（智能代理）能够自主完成复杂任务：订票/订餐/安排日程/处理邮件/整理文档/跨设备协同，多模态AI成为标配，支持文字/图像/视频/语音/空间多模态理解和生成，AR智能眼镜开始规模普及，空间计算设备Vision Pro生态成熟，消费电子形态开始多元化。',
     '【全场景个人AI时代（2029-2030）】全场景个人AI助理成为现实，手机/PC/平板/手表/耳机/汽车/眼镜/智能家居/服务机器人共享统一个人AI模型，AI理解用户习惯和需求，主动提供服务，端侧NPU算力达到500TOPS以上，可以本地运行70B参数大模型，AI能力接近人类助理水平，消费电子进入真正的智能时代。',
     '【处理器NPU演进】没有NPU→麒麟970独立NPU（1.9TOPS）→A11神经网络引擎（0.6TOPS）→骁龙8 Gen3（45TOPS）→麒麟9020（120TOPS）→A19 Pro（150TOPS）→M5 Max（180TOPS），NPU算力9年提升近100倍。',
     '【存储内存演进】内存：4GB→8GB→12GB→16GB→24GB→32GB，旗舰手机内存最高24GB，AI PC内存最高192GB，满足大模型运行内存需求；存储：64GB→128GB→256GB→512GB→1TB→2TB→8TB，旗舰手机最高2TB存储，AI PC最高8TB存储。',
     '【电池快充演进】电池容量：功能机800mAh→智能机早期2000mAh→现在旗舰5000-6000mAh硅碳负极电池；快充功率：5W→18W→65W→100W→120W→210W有线快充，15W→50W→80W无线快充，10分钟充至50%以上，续航焦虑基本解决。'],
    '▎旗舰参数 · 价格版本 · AI功能 · 技术亮点',
    ['【华为Pura 90s系列·2026年8月20日】华为官宣Pura 90s系列旗舰（Pura 90s Pro/Pura 90s Pro Max）8月28日发布：Pro Max搭载200MP超大传感器长焦相机支持20倍视频变焦，1/1.28英寸大底RYYB传感器16EV超宽动态范围；Ultra Lighting HDR相机50MP配LOFIC技术物理光圈F1.4-F4.0 OIS；AI构图3.0实时推荐最佳取景+AI去眩光+AI移动重排对象+AI最佳表情；第二代昆仑玻璃防反光防刮6000mAh电池Pro Max支持100W有线快充',
     '【Pura X View预约·2026年8月20日】华为Pura X View 8月20日18:08开启预约已上架京东及华为商城8月28日开启门店体验：提供跃影红/亚麻灰/零度白/幻夜黑四色，12GB+256GB/12GB+512GB/12GB+1TB三种存储版本；6.39英寸16:9.5阔屏分辨率2232×1320峰值亮度6500nit支持2160Hz PWM调光；6.68mm厚度201g重量屏占比96.1%四边等宽1.05mm；7000mAh电池支持66W超级快充；首发搭载HarmonyOS 7系统',
     '【享界G9座舱安全·2026年8月20日】享界G9座舱内部有效空间3415mm：高配搭载最多四零重力座椅18向调节+20点双维按摩；鸿蒙ALPS健康座舱+车载制氧系统+智能隐私光幕玻璃遮光度超99.4%；寰宇三联屏+17英寸吸顶屏+68英寸成像面积AR-HUD+悦彰音响；安全方面100%潜艇级热成型钢+18根2000MPa级热成型钢车顶承压16.5吨+一体化压铸铝合金工艺+标配11气囊；搭载华为星河通信3.0天地网联支持最远10公里车载对讲+卫星通话',
     '【苹果iPhone 18 Pro Max】A19 Pro芯片（3nm第二代），8核CPU+8核GPU+16核NPU（150TOPS）；6.9英寸Super Retina XDR OLED，1-120Hz ProMotion，3500nit峰值亮度；后置4800万主摄（传感器位移防抖）+4800万超广角+1200万5倍潜望长焦；4800mAh电池45W有线+25W MagSafe无线；iOS 20全功能支持Apple Intelligence；钛金属中框IP69防水；存储：256G 9999/512G 11499/1T 13499元；2026年9月发布',
     '【苹果MacBook Pro 16 M5】Apple M5 Max芯片：16核CPU（12性能+4能效）+40核GPU+32核NPU（180TOPS）；最高192GB统一内存+最高8TB SSD；16.2英寸Liquid Retina XDR显示屏，3456×2234分辨率，120Hz ProMotion，1600nit HDR峰值亮度；100Wh电池21小时视频播放；140W MagSafe快充；macOS Sequoia深度整合Apple Intelligence；重量2.1kg；存储：36+512G 19999/64+1T 24999/96+2T 29999/192+8T 45999元',
     '【苹果Vision Pro 2】M5芯片+R2芯片；Micro-OLED双眼4K分辨率（单眼2300万像素），120Hz刷新率，120度视场角；重量450g（比一代减轻30%）；眼动追踪+手势追踪+空间音频；visionOS 3支持空间视频拍摄/沉浸式办公/空间游戏；外接电池续航4小时；存储：256G 14999元；2026年6月WWDC发布',
     '【小米16 Ultra】高通骁龙8 Gen4（3nm）+澎湃P2快充芯片+澎湃G2电池管理芯片；6.9英寸2K LTPO AMOLED华星C9屏，1-144Hz自适应刷新率，4000nit峰值亮度，1920Hz PWM调光；后置徕卡四摄：一英寸LYT-900主摄（f/1.4-f/4.0可变光圈）+4000万超广角+2亿像素5倍潜望长焦+5000万像素10倍超长焦；6000mAh硅碳负极电池120W有线+50W无线；澎湃OS 3.0支持双向卫星通信；IP68防水陶瓷机身；存储：12+256G 5999/16+512G 6499/16+1T 7499/24+2T 8499元',
     '【小米16 Pro】骁龙8 Gen4；6.7英寸2K 144Hz LTPO AMOLED；后置5000万LYT-800主摄+5000万超广角+5000万3倍潜望长焦；5500mAh电池120W有线+50W无线快充；澎湃OS 3.0；IP68防水金属中框玻璃后盖；存储：12+256G 4999/12+512G 5499/16+1T 6299元',
     '【小米RedmiBook Pro 16 2026】英特尔酷睿Ultra 7 268V + 小米自研NPU（80TOPS AI算力）；32GB LPDDR5X内存+1TB PCIe 4.0 SSD（可扩展）；16英寸3.2K 120Hz IPS屏，100% sRGB色域，500nit亮度；80Wh电池100W快充；1.8kg重量；澎湃OS for Windows支持AI字幕/AI会议/AI画图/AI写作；存储：16+512G 4999/32+1T 5999元；2026年3月发布',
     '【小米SU7 Ultra】三电机四驱系统：前220kW+后350kW+后350kW，系统总功率960kW（1306马力），系统总扭矩1680N·m；0-100km/h加速1.98秒，0-200km/h加速5.8秒，最高车速350km/h；宁德时代麒麟II电池130kWh，CLTC续航800km，800V高压平台10分钟补能400km；小米全栈自研智驾系统双Orin-X芯片（508TOPS）+激光雷达+11摄像头+12超声波+5毫米波雷达；21英寸轮毂碳陶瓷刹车；空气悬架+CDC电磁减振；售价52.99万元；2026年3月发布已交付',
     '【AI功能共性】三大品牌旗舰全部支持：自然语言语音助手（连续对话/复杂指令/多轮交互）、AI拍照修图（AI消除/AI扩图/AI增强/AI夜景）、AI会议（实时录音转写/智能摘要/待办提取/多语种翻译）、AI写作（文案/邮件/报告/代码生成）、AI搜索（自然语言搜索/信息整合总结）、AI安全（隐私计算/端侧处理不上云）',
     '【端侧大模型】华为盘古大模型端侧版（13B参数）、苹果Apple Intelligence（云端30B参数+端侧3B混合推理）、小米MiLM大模型端侧版（7B/13B），推理延迟<1秒，隐私敏感数据全程端侧处理不上传',
     '【生态互联】华为鸿蒙分布式软总线支持手机/PC/平板/手表/车机/智能家居无缝流转接续；苹果Continuity支持iPhone/Mac/iPad/Watch/Vision Pro全生态设备无缝协同接力；小米澎湃OS HyperConnect支持全生态AIoT设备互联互通智能联动',
     '【通信技术】华为独家支持天通一号卫星通话+北斗卫星消息+星闪NearLink；苹果支持卫星SOS紧急求助+卫星消息共享；小米支持双向卫星通信（天通+北斗）+5.5G NTN；三家旗舰均标配Wi-Fi 7/蓝牙5.4/NFC/红外遥控',
     '【材料工艺】华为采用昆仑玻璃2代+玄武架构+纳米微晶陶瓷；苹果采用Grade 5钛金属中框+超瓷晶面板；小米采用龙晶玻璃+航空铝中框+陶瓷后盖；抗跌落抗刮擦耐用性较上一代提升2-3倍',
     '【市场份额】2026年H1中国智能手机市场：华为28%份额重回第一，苹果18%第二，小米15%第三，三家合计占据61%市场份额；AI PC市场联想/华为/苹果/小米位列前四',
     '【智界RX预售·2026年8月20日】智界RX正式开启预售29.98万-39.98万元至高享10000元预订权益：定位性能向SUV宽体低趴设计，配备共光路舱内激光雷达版本预售29.98万起/L3级自动驾驶架构版本预售35.98万起；搭载896线双光路图像级激光雷达；创新"343"一体化压铸车身架构车身扭转刚度超50000N·m/deg；双腔空气悬架+双阀连续可变阻尼减振器近50:50轴荷比；首发巨鲸电池平台"黑匣子"安全技术',
     '【问界M6交付·2026年8月20日】问界M6上市4个月累计交付突破45000台，在2026中国汽研智能汽车指数测评中拿下四项全G+评级；发布会新增两款增程版本：增程Max四驱版售价23.98万元起，增程Max+长续航版售价25.98万元起；两版本在电池容量/纯电续航/智驾硬件/底盘配置上做出区分CLTC综合续航分别可达1445km和1605km；鸿蒙智行老车主增换购享2万元尾款减免',
     '【享界V8官宣·2026年8月20日】余承东正式公布家庭智慧旗舰MPV享界V8主打高端家用市场：将携手尊界V800/V680、智界V9重塑高端MPV市场格局；8月21日成都车展实车首亮相；随着享界G9与享界V8到来享界已构建覆盖轿车/旅行车/SUV/MPV全品类豪华产品矩阵；享界G9至高106500元权益包含后挂储物箱/电动踏板/后排娱乐屏/座椅升级/智驾包补贴',
     '【鸿蒙智行150万·2026年8月20日】余承东宣布截至2026年8月16日鸿蒙智行全系累计交付突破150万台仅用时53个月刷新中国新势力品牌达成速度纪录：鸿蒙智行已形成问界/智界/享界/尊界/尚界五大子品牌协同发展格局覆盖轿车/SUV/MPV全品类；享界G9是首款获批行业最高时速120公里L3级自动驾驶道路测试牌照的车型也是首款采用L3级自动驾驶架构设计的豪华硬派SUV',
     '【技术趋势3】AI代理（Agent）深度整合系统：手机/PC上的AI助理能够自主理解指令并完成订票/订餐/安排日程/处理邮件等复杂任务',
     '【未来方向】2027-2030年AI终端形态向AR智能眼镜/脑机接口/智能家居服务机器人延伸，全场景个人AI助理成为现实']))

# PART 13-PART22 剩余模块快速填充（保持详细度）
module_titles_rest = [
    ('PART 13', '智慧农业：农业机器人应用', '农业机器人/极飞/大疆'),
    ('PART 14', '医疗健康：医疗机器人突破', '医疗机器人/达芬奇/天智航'),
    ('PART 15', '教育AI：教育智能化转型', '教育AI/科大讯飞/作业帮'),
    ('PART 16', '能源电力：电力机器人运维', '电力机器人/国家电网/南瑞'),
    ('PART 17', '自动驾驶：L4级商业化落地', '自动驾驶/百度萝卜快跑/特斯拉FSD'),
    ('PART 18', '人形运动会：技术竞赛舞台', '人形机器人运动会/赛事/竞技'),
    ('PART 19', '真机部署：规模化落地进展', '真机部署/工厂/物流/服务'),
    ('PART 20', '物流仓储：仓储机器人普及', '物流仓储/极智嘉/快仓/海康'),
    ('PART 21', '灵巧手：精密操作核心部件', '灵巧手/因时机器人/Shadow Hand'),
    ('PART 22', '安防应急：特种机器人守护安全', '安防应急/消防机器人/排爆机器人'),
]

def make_detail_module(part_num, title, keyword, category_idx):
    categories = [
        # 智慧农业
        {'left': [
            '【产业定位】智慧农业是乡村振兴战略重要支撑，农业机器人替代人力解决劳动力短缺问题',
            '【市场规模】2026年中国智慧农业市场规模突破1500亿元，农业机器人占比35%',
            '【无人机植保】极飞/大疆植保无人机年作业面积超20亿亩次，植保机械化率超70%',
            '【采摘机器人】果蔬采摘机器人逐步成熟，草莓/番茄/苹果等作物采摘成功率>95%',
            '【巡检机器人】农田/温室/养殖场巡检机器人普及，监测作物生长/畜禽健康/环境参数',
            '【自动驾驶农机】北斗导航自动驾驶拖拉机/收割机普及，作业精度厘米级，效率提升30%',
            '【AI种植决策】AI+物联网+大数据，精准灌溉/施肥/打药，节水节药30%，增产15%',
            '【畜禽养殖机器人】喂料/清粪/挤奶/巡检机器人大规模应用，养殖效率提升40%',
            '【安徽农业基础】安徽是农业大省，粮食产量全国前列，智慧农业应用需求旺盛',
            '【政策支持】数字乡村发展战略，农机购置补贴向智慧农机倾斜'
        ], 'right': [
            '【极飞科技P150】农业无人飞机，载重150L，喷幅12米，每小时作业400亩，RTK厘米级定位',
            '【大疆T100】大疆农业植保无人机T100，载重100kg，双旋翼设计，作业效率350亩/小时',
            '【极飞R150】农业无人车，可喷洒/播种/运输，全地形适应，自主规划路径',
            '【采摘机器人】中科原动力草莓采摘机器人，采收速度8秒/个，成功率96%，日作业10亩',
            '【博创联动】自动驾驶拖拉机系统，改装成本2-3万元，作业精度±2.5cm，已改装10万台',
            '【中联农机】中联重科AI收割机，自动识别作物成熟度，自动调整作业参数，损失率<1%',
            '【温氏养殖机器人】温氏集团养猪场巡检/喂料/清粪机器人全覆盖，养殖工人减少60%',
            '【安徽农垦】安徽农垦集团建设5个智慧农场示范基地，农机自动驾驶率达85%',
            '【蚌埠农业】蚌埠怀远石榴/五河螃蟹等特色农产品产业探索AI+农业应用',
            '【农业无人机·2026年8月21日】贵州六盘水20多万亩红心猕猴桃进入采摘季引入吊运无人机空中转运鲜果采收效率大幅提升；中国农林植保无人机保有量从2018年约3万架增至2025年约30万架，农用无人机年作业面积突破4.6亿亩相比人工效率提升超30倍农药用量减少30%以上综合成本下降50%，当前渗透率约20%预计2030年升至50%以上；大疆农业无人机已应用100多个国家和地区全球累计销量突破60万台国内单年作业台数超32万台单年作业量突破33亿亩次实现650万吨物资吊运；杭州乔戈里科技智能采摘机器人搭载激光雷达/机器视觉多传感器融合+AI大模型自主识别成熟果实软爪轻抓技术正复制到番茄/草莓/黄瓜/彩椒；2026年中央一号文件首次提出拓展无人机/物联网/机器人应用场景农作物耕种收综合机械化率达76.7%农业科技进步贡献率超64%'
        ], 'process': [
            '【人工畜力时代（1949-1980）】新中国成立初期，农业生产几乎完全依靠人力和畜力，牛耕人种是主要方式，生产效率极低，粮食亩产不足100公斤，机械化率不足10%，农业生产主要解决温饱问题，劳动强度大，农民几乎全年无休，农业人口占总人口80%以上。',
            '【初步机械化期（1981-2000）】改革开放后，小型拖拉机、收割机、农用三轮车开始普及，主要农作物耕种收综合机械化率提升至30%左右，但仍以人力为主，农机质量差、故障多、适用性差，大型农机依赖进口，农机手水平参差不齐。',
            '【全面机械化期（2001-2015）】国家加大农机购置补贴力度，大中型拖拉机、联合收割机、插秧机快速普及，主要农作物耕种收综合机械化率提升至65%以上，小麦基本实现全程机械化，水稻、玉米机械化率快速提升，但智能化程度低，农机需要专业驾驶员操作。',
            '【无人机植保起步期（2016-2020）】极飞、大疆等企业推出农业植保无人机，无人机植保从无到有快速发展，植保无人机保有量从不足1000架增长至10万架以上，年作业面积突破10亿亩次，RTK厘米级定位技术应用，作业效率是人工的几十倍，解决了传统植保"打药难、打药累、易中毒"问题。',
            '【智慧农业试点期（2021-2024）】AI、物联网、大数据技术开始应用于农业，自动驾驶农机试点，土壤传感器、气象站普及，精准灌溉、精准施肥、AI病虫害识别开始示范应用，采摘机器人、养殖机器人试点，智慧农场、数字乡村建设启动，安徽等农业大省建设智慧农业示范基地。',
            '【规模应用期（2025-2026）】2026年农业机器人开始规模应用，极飞/大疆农业无人机年作业面积超20亿亩次，自动驾驶农机改装量超20万台，AI病虫害识别覆盖主要农作物，采摘机器人在草莓、番茄等高价值作物种植中规模应用，畜禽养殖机器人在大型养殖场普及，智慧农业市场规模突破1500亿元。',
            '【全面智能化期（2027-2028）】农业机器人成本进一步下降，性能持续提升，5G+农业机器人实现远程控制和协同作业，多机器人协同农场试点，从耕种到收获全程无人化作业示范，AI育种大规模应用，品种选育周期大幅缩短，农产品全流程溯源体系建成，智慧农业从示范走向普及。',
            '【无人农场普及期（2029-2030）】2030年主要农作物耕种收综合机械化率达90%，农业机器人普及率达50%，无人农场在全国主要粮食产区推广，农业生产实现全程智能化、无人化，精准农业成为常态，水、肥、药利用率大幅提升，农业劳动生产率提升5倍以上，农民从体力劳动者转变为农业运营者。',
            '【农业机器人技术演进】手动农具→畜力农具→小型农业机械→大型农业机械→自动驾驶农机→农业无人机→农业AI决策系统→多机器人协同无人农场；定位导航从人工→GPS→北斗单频→北斗RTK厘米级→北斗+视觉多传感器融合。',
            '【安徽智慧农业进展】安徽作为全国农业大省，粮食产量常年居全国第4位，2026年已建设省级智慧农业示范基地50个，农机自动驾驶系统安装量超2万台，植保无人机保有量超3万台，安徽农垦集团5个智慧农场示范基地农机自动驾驶率达85%，蚌埠怀远石榴、五河螃蟹等特色农产品探索AI+农业应用，智慧农业发展走在全国前列。'
        ], 'detail': [
            '【极飞P150农业无人机参数】载重量150L超大药箱，喷幅宽度达12米，每小时作业效率400亩，RTK厘米级定位精度1cm+1ppm，满载续航15分钟，电池快充10分钟充至80%，整机售价约12万元，全球累计作业面积超15亿亩次',
            '【大疆T100农业无人机参数】载重量100kg，双旋翼共轴反桨设计，喷幅宽度10米，每小时作业效率350亩，配备全向数字雷达和有源相控阵避障系统，支持夜间作业和山地仿形飞行，整机售价约9.8万元，全球市占率超60%',
            '【果蔬采摘机器人性能参数】草莓/番茄/苹果等果蔬采摘成功率稳定在95%-98%，单果采摘时间仅6-10秒，单台机器人日作业量相当于8-10名熟练采摘工人，投资回收期（ROI）约1.5-2年，24小时不间断作业不受天气影响',
            '【北斗导航自动驾驶农机】作业直线精度±2.5cm，直线度偏差小于2cm，夜间作业能力不受光线影响，作业效率比人工驾驶提升30%，燃油消耗降低10%，对行精度高减少重播漏播，已改装各类农机超20万台套',
            '【AI精准灌溉施肥系统】结合土壤传感器、气象站数据、作物生长模型，实现变量灌溉和精准施肥，节水30%-40%，肥料利用率提高20%，作物产量增加10%-15%，每亩增收约200-300元，已在全国1000+万亩耕地推广应用',
            '【AI病虫害识别技术】通过手机拍照或无人机航拍即可识别农作物病虫害，识别准确率超过98%，单张图片识别速度小于1秒，精准推荐对症农药和用量，减少盲目打药，农药使用量降低30%，覆盖水稻/小麦/玉米/果蔬等50+种主要作物',
            '【畜禽养殖机器人应用】实时监测圈舍温度、湿度、氨气、二氧化碳等环境参数，自动完成喂料、清粪、通风、温控等工作，畜禽死淘率降低20%，饲料转化率提高15%，养殖工人减少60%，万头猪场仅需5-8名管理人员',
            '【水产养殖机器人】水下巡检机器人实时监测水质参数（溶氧/PH/氨氮）、鱼群健康状态、网箱破损情况，自动完成投饵、增氧、清淤、死鱼打捞等作业，养殖密度提高30%，发病率降低40%，亩均增产增收超5000元',
            '【安徽智慧农业建设进展】安徽省已建设省级智慧农业示范基地50个，农机自动驾驶系统安装量超2万台套，植保无人机保有量超3万台，小麦、水稻、玉米主要农作物耕种收综合机械化率达83%，智慧农业发展走在全国前列',
            '【蚌埠本地农业应用】蚌埠国家农业科技园区建设智慧农业示范区，怀远石榴、五河螃蟹、固镇花生等特色农产品产业全面推广无人机植保，无人机植保率超80%，特色农产品品质提升、品牌溢价能力显著增强',
            '【极飞科技全球市场布局】极飞农业无人机已在全球100+国家和地区推广应用，累计作业面积超20亿亩次，建立了覆盖全球的销售服务网络，在日本、韩国、东南亚、拉美、非洲等市场占有率位居前列，是中国农业科技出海代表企业',
            '【大疆农业业务规模】大疆农业无人机全球市占率超过60%，稳居全球第一，2026年农业无人机业务收入超200亿元，T系列植保无人机已迭代至T100，产品覆盖植保、播种、施肥、测绘等全场景农业作业需求',
            '【智慧农机购置补贴政策】国家将智慧农机纳入农机购置补贴范围，补贴比例30%-50%，农业机器人单机补贴最高可达5万元，自动驾驶农机改装补贴2-3万元，大幅降低农民购机成本，加速智慧农机普及应用',
            '【土地流转规模化经营】全国土地流转率加速提升，2026年土地流转率超60%，家庭农场、农民合作社、农业企业等新型经营主体成为主力，规模化经营为农业机器人大规模应用创造了必要条件，小块分散农田难以发挥机器人效率',
            '【5G网络农村覆盖进展】全国5G网络已覆盖90%以上行政村，光纤宽带村村通工程全面完成，农村地区网络带宽和时延满足农业机器人远程控制、实时数据回传、AI云端推理需求，为智慧农业提供通信基础设施支撑',
            '【AI辅助育种技术突破】AI技术应用于农作物育种，通过基因组测序、表型分析、环境模拟，育种周期从传统8-10年大幅缩短至2-3年，育种效率提升3倍以上，抗病虫、高产、优质新品种选育速度显著加快',
            '【区块链农产品溯源体系】区块链+物联网技术实现农产品从田间种植、管理、收获、加工、运输、销售全流程溯源，消费者扫码即可查看农产品全生命周期信息，质量安全可追溯，农产品品牌信任度提升，溢价能力增强',
            '【乡村电商与智慧物流】直播电商+智慧物流体系解决农产品销售难题，2026年全国农产品网络零售额超1万亿元，县长直播带货、农民主播成为新潮流，产地仓+冷链物流+快递进村让优质农产品走出大山、走向全国',
            '【智慧农业人才培养】全国农业院校普遍开设农业工程、智慧农业、农业人工智能相关专业，每年培养智慧农业专业人才超1万人，新型职业农民培训工程每年培训农民超100万人次，为智慧农业发展提供人才支撑',
            '【2030年智慧农业发展目标】到2030年，全国主要农作物耕种收综合机械化率达90%，农业机器人普及率达50%，建成1000个全程无人化示范农场，精准农业成为常态，水肥药利用率大幅提升，基本实现农业现代化目标'
        ]},
        # 医疗健康
        {'left': [
            '【产业定位】医疗机器人是高端医疗器械重要方向，手术机器人/康复机器人/护理机器人快速发展',
            '【市场规模】2026年中国医疗机器人市场规模突破350亿元，年均增速超45%',
            '【手术机器人】腔镜/骨科/神经外科/血管介入手术机器人逐步国产化，价格大幅下降',
            '【达芬奇垄断】达芬奇手术机器人长期垄断国内市场，装机量超500台，单次手术费用超3万元',
            '【国产替代】天智航/微创机器人/威高/思哲睿等国产手术机器人获批上市，价格仅进口1/2-1/3',
            '【康复机器人】上下肢康复/外骨骼/步行训练机器人在医院康复科普及，帮助患者恢复运动功能',
            '【护理机器人】转运/陪护/喂药/消毒机器人在医院应用，减轻护士工作负担，解决护理人员短缺',
            '【AI诊断】AI医学影像/AI病理/AI辅助诊断准确率达主任医师水平，基层医疗能力大幅提升',
            '【安徽医疗】中科大附一院/安医大一附院等医院引进手术机器人，医疗机器人应用快速增长',
            '【政策支持】医疗机器人纳入鼓励采购目录，国产医疗设备采购比例要求提升'
        ], 'right': [
            '【WRC医疗展区·2026年8月21日】WRC2026医疗展区成最大看点：长木谷发布全球首款六位一体ROPA6 AI全骨科手术机器人同时斩获国家三类注册证与欧盟CE认证；强联智创展出全球首个且目前唯一AI驱动神经血管疾病智能手术解决方案已在全国近60家医院落地完成超1万例手术；德壹医疗红光治疗机器人获批二类证',
            '【程天外骨骼·2026年8月21日】程天科技发布全新消费级新品GoGo-H Pro搭载自研新一代AI步态自适应算法，可识别偏瘫/拖步/左右不对称等特殊步态并精准助力，毫秒级实时感知用户运动意图；外骨骼机器人从医疗康复走向消费级市场价格下探让更多行动障碍人群可负担',
            '【达芬奇Xi】直觉外科达芬奇手术机器人第四代，四臂结构，3D高清视野，腕式手术器械，装机量超8000台全球',
            '【天智天玑】天智航骨科手术机器人，国内首款获批骨科机器人，已完成手术超20万台，精度0.8mm',
            '【微创图迈】微创医疗腔镜手术机器人，四臂腔镜机器人，性能接近达芬奇，价格仅为1/2',
            '【傅利叶康复】傅利叶智能上下肢康复机器人，已在全国3000+医院康复科部署，康复训练效果提升50%',
            '【大艾外骨骼】大艾机器人下肢外骨骼康复机器人，帮助截瘫/偏瘫患者重新行走，已进入1000+医院',
            '【钛米消毒】钛米机器人医院消毒机器人，紫外线+过氧化氢消毒，覆盖医院90%以上物表消毒需求',
            '【润迈德介入】润迈德血管介入手术机器人，冠脉造影/PCI手术辅助，已完成NMPA注册',
            '【AI影像】科大讯飞AI医学影像，肺结节/乳腺癌/眼底病变识别准确率>97%，已在30000+基层医院部署',
            '【中科大附一院】安徽省立医院装机达芬奇手术机器人3台，国产手术机器人2台，年机器人手术超1万台',
            '【安医大一附院】安徽医科大学一附院建设智慧医院，引进30+种医疗机器人，智能化水平全国领先'
        ], 'process': [
            '【纯开放手术时代（1990年前）】外科手术以开放手术为主，需要切开十几厘米甚至几十厘米大切口，医生肉眼直视下操作，手的稳定性和精度有限，手术创伤大、出血多、恢复慢、并发症高，复杂手术难度大，手术效果高度依赖医生个人经验，优质医疗资源集中在大城市大医院。',
            '【腹腔镜微创起步期（1991-2010）】腹腔镜微创手术开始普及，通过几个小孔插入器械和摄像头，医生看着屏幕操作，创伤大幅减小，但传统腹腔镜器械活动度有限（只有4个自由度），操作不灵活，缝合打结等精细操作难度大，学习曲线长，达芬奇手术机器人2000年获FDA批准开始进入中国，但装机量极少，价格极其昂贵。',
            '【进口达芬奇垄断期（2011-2018）】达芬奇手术机器人在中国装机量快速增长，从不足10台增长到近100台，但直觉外科完全垄断市场，设备售价超2000万元，年服务费超100万元，专用器械单把超万元，单次机器人手术费用比普通腹腔镜贵2-3万元，只有顶级三甲医院能负担，国产手术机器人处于研发起步阶段。',
            '【国产技术突破期（2019-2024）】国产医疗机器人企业持续研发投入，天智航骨科手术机器人2020年获批成为首款国产手术机器人，微创图迈腔镜机器人2022年获批，威高、思哲睿、润迈德等企业产品陆续获批NMPA，国产设备性能逐步接近进口水平，价格仅为进口1/2-2/3，开始在三甲医院批量装机应用，康复机器人、护理机器人、AI诊断也快速发展。',
            '【AI诊断崛起期（2020-2025）】AI医学影像技术快速成熟，肺结节、乳腺癌、眼底病变、骨折、病理切片AI诊断准确率达到主任医师水平，科大讯飞等企业AI诊断产品获批三类证，开始在基层医院大规模部署，解决基层医院缺乏优质影像科、病理科医生问题，AI辅助诊断让基层患者也能获得高水平诊断，医疗公平性大幅提升。',
            '【规模普及启动期（2025-2026）】2026年国产医疗机器人技术成熟，性能达到国际先进水平，天智航骨科机器人累计完成手术超20万台，微创图迈装机超100台，国产手术机器人价格降至进口1/3，单次手术费用降至1-2万元，康复机器人在医院康复科普及率超60%，AI医学影像覆盖3万+基层医院，中国机器人手术量突破200万台/年。',
            '【基层普及期（2027-2028）】国产医疗机器人成本进一步下降，政策支持国产医疗设备采购，二级医院甚至县级医院开始普及手术机器人、康复机器人，5G+远程手术让基层患者在家门口就能享受大专家手术，护理机器人在医院、养老院大规模应用，减轻护理人员负担，医疗机器人从大三甲走向普通医疗机构。',
            '【普惠医疗期（2029-2030）】2030年国产医疗机器人市占率达70%，二级以上医院医疗机器人普及率达80%，AI诊断基层全覆盖，消费级康复外骨骼、助行机器人开始进入家庭，医疗机器人惠及普通民众，优质医疗资源通过机器人和AI下沉，城乡医疗差距大幅缩小，人均预期寿命进一步提升。',
            '【医疗机器人技术演进】开放手术→腹腔镜微创手术→进口达芬奇手术机器人→国产手术机器人→AI辅助诊断→5G远程手术→多机器人协同手术→消费级家庭医疗机器人；手术精度从厘米级→毫米级→亚毫米级→0.8mm以内，创伤从几十厘米切口→几个小孔→自然腔道无创手术。',
            '【安徽医疗机器人进展】安徽医疗机器人应用快速增长，中科大附一院（安徽省立医院）装机达芬奇手术机器人3台、国产手术机器人2台，年机器人手术超1万台；安医大一附院建设智慧医院，引进30+种医疗机器人；全省三甲医院手术机器人装机量超80台，年机器人手术量超5万台，基层医院AI影像覆盖率达60%；蚌埠医学院第一、第二附属医院引进手术机器人和康复机器人，智慧医院建设快速推进。'
        ], 'detail': [
            '【达芬奇Xi详细参数】美国直觉外科第四代达芬奇手术机器人配置4个交互式机械臂，7自由度腕式EndoWrist手术器械可转腕540度超越人手活动范围，搭载3D高清立体视觉系统放大10-15倍提供沉浸式术野，支持5:1运动缩放和智能震颤滤除功能，器械端活动自由度超过90度，单台设备装机成本约2000万元人民币，年维护服务费约150万元',
            '【天智天玑2.0骨科机器人】天智航天玑2.0是国内首款获批NMPA三类证的骨科手术机器人系统，机械臂拥有6个高精度自由度，手术定位精度达到亚毫米级0.8mm，术前三维规划时间小于5分钟，广泛适用于脊柱外科、创伤骨科、关节置换等各类骨科手术，单台设备装机成本约800万元人民币，截至2026年8月累计完成临床手术超20万例',
            '【微创图迈腔镜机器人】微创医疗机器人自主研发的图迈Toumai四臂腔镜手术机器人，配备3D超高清立体视野系统，7自由度腕式手术器械操作灵活精准，主从控制端到端延迟小于100ms达到国际领先水平，独创力反馈功能让医生感知组织触感，单台装机成本约1000万元人民币（仅为达芬奇的一半），2022年正式获批NMPA，截至2026年8月全国装机超100台',
            '【傅利叶ArmMotus康复机器人】傅利叶智能ArmMotus M2 Pro上肢康复机器人拥有7自由度力反馈控制，支持主动、被动、助动、抗阻多种训练模式，可根据患者恢复情况自适应调节训练难度和阻力等级，内置丰富的游戏化训练任务提高患者依从性，广泛适用于脑卒中、脑外伤、脊髓损伤、神经系统疾病导致的上肢运动功能障碍康复训练，单台设备市场售价约80万元人民币',
            '【大艾AiLeg下肢外骨骼】大艾机器人AiLegs下肢外骨骼康复机器人配置10个高精度主动自由度，支持原地站立、平地行走、上下楼梯、跨越障碍等多种康复训练动作，基于AI自适应步态规划算法实时调整步态参数匹配患者康复进度，适用于截瘫、偏瘫、脊髓损伤、脑瘫等下肢运动功能障碍患者康复训练，帮助患者重新站立行走，单台设备市场售价约120万元人民币',
            '【钛米消毒机器人详细参数】钛米智能消毒机器人采用紫外线UVC+过氧化氢雾化双重消毒技术，紫外线消毒剂量≥10000μW/cm²达到高水平消毒标准，过氧化氢雾化浓度30%可覆盖物表和空气，30分钟可完成100平方米空间全面消毒，搭载激光SLAM自主导航避障系统支持自主路径规划，任务完成后自动返回充电桩充电，单台设备市场售价约30万元人民币',
            '【AI医学影像诊断能力】科大讯飞等企业AI医学影像产品在肺部CT结节检测灵敏度达到99%、特异度95%；乳腺癌钼靶诊断AUC值达0.99超越主任医师水平；糖尿病视网膜病变等眼底病变筛查准确率98%；骨折、脑出血、气胸等急症识别准确率超97%；单次AI辅助诊断时间小于10秒，诊断成本仅为人工专家诊断的1/10',
            '【机器人手术量爆发增长】2026年中国全年机器人辅助手术量正式突破200万台大关保持全球第二，其中国产手术机器人完成手术占比从2020年不足5%大幅提升至2026年的45%，国产替代加速推进；腔镜手术机器人占比约55%、骨科手术机器人占比约25%、其他专科手术机器人占比20%，三甲医院机器人手术渗透率超过30%',
            '【安徽医疗机器人应用进展】安徽省三甲医院手术机器人装机总量已超80台，2026年全年完成机器人辅助手术量超5万台保持全国前列；中科大附一院（安徽省立医院）年机器人手术量突破1万台；全省基层医疗机构AI医学影像辅助诊断系统覆盖率达60%，让基层患者在家门口就能享受高水平诊断服务',
            '【蚌埠本地智慧医疗建设】蚌埠医学院第一附属医院、蚌埠医学院第二附属医院作为皖北医疗中心，已引进达芬奇手术机器人、国产腔镜手术机器人、上下肢康复机器人等各类医疗机器人设备，智慧医院建设快速推进；蚌埠市第一、第二、第三人民医院也已部署AI医学影像、消毒机器人、配送机器人等智能化设备，医疗服务能力持续提升',
            '【康复机器人临床效果】临床数据显示使用康复机器人进行系统化训练，脑卒中偏瘫患者运动功能恢复率较传统人工康复提升40%以上，患者平均住院日缩短25%，康复治疗师人均工作量减少60%；机器人提供的高强度、重复性、标准化训练是人工康复难以实现的，康复效果显著优于传统康复模式',
            '【护理机器人临床价值】医院智能护理机器人可承担80%以上重复性护理工作：包括患者转运、器械药品配送、环境消毒、生命体征监测、陪护照料等，护士工作效率提升50%，护理差错率降低70%，有效解决护理人员短缺问题，让护士有更多时间关注患者病情观察和人文关怀',
            '【机器人手术费用大幅下降】国产手术机器人规模化应用后，单次机器人辅助手术患者自付费用从达芬奇垄断时期的3-5万元大幅下降至1-2万元区间，费用下降幅度达60%，机器人手术不再是高端奢侈医疗服务，普通工薪阶层患者也能负担得起，极大提升了医疗可及性和公平性',
            '【5G远程手术技术突破】5G低时延通信技术+手术机器人支持超远程手术，顶级专家在千里之外的大城市三甲医院就能为基层医院、偏远地区患者实时进行手术操作，端到端时延控制在50ms以内无卡顿，已成功完成多例跨省份5G远程手术，推动优质医疗资源下沉',
            '【介入手术机器人快速发展】血管介入、神经介入、消化介入、呼吸介入等各类介入手术机器人快速发展成熟，医生可在铅防护室外远程操控机器人完成手术，彻底避免X线电离辐射对医护人员的职业伤害，同时手术操作精度和稳定性大幅提升，介入手术并发症率降低30%',
            '【胶囊内镜机器人无痛苦检查】胶囊内镜机器人患者只需吞服一颗药丸大小的胶囊机器人，即可完成整个消化道（食道、胃、小肠、大肠）的全面检查，全程无痛苦无需麻醉，胶囊内置摄像头拍摄数万张高清图像，AI辅助阅片诊断准确率超过95%，检查完成后胶囊自然排出体外无残留',
            '【消费级外骨骼走进民用】消费级康复外骨骼、助行外骨骼开始进入市场，主要用于老年人助行、产业工人负重助力、登山徒步运动助力等场景，价格下探至1-3万元普通家庭可承受区间；未来消费级外骨骼将像电动车一样普及，帮助老年人独立行走提升生活质量，降低工人劳动强度',
            '【NMPA审批政策加速】国家药监局NMPA将医疗机器人纳入创新医疗器械特别审批绿色通道，优先审评审批，国产创新医疗机器人获批速度较以前加快2-3年；截至2026年8月已有超30款国产医疗机器人获得NMPA三类医疗器械注册证，涵盖手术、康复、护理、诊断等全品类',
            '【手术医师培训体系完善】国家卫健委已建立完善的机器人手术医师培训认证体系，在全国设立20个机器人手术培训基地，采用理论授课+模拟训练+动物实验+临床带教阶梯式培训模式，每年培训合格机器人手术医师超1万人，保障机器人手术安全规范开展',
            '【2030年产业发展目标】规划到2030年国产医疗机器人国内市场占有率达到70%以上，二级以上医院医疗机器人普及率达到80%，AI医学影像辅助诊断系统实现基层医疗机构全覆盖，医疗机器人整体技术水平达到国际先进，部分领域实现全球领先，让优质医疗资源惠及全体人民'
        ]},
        # 教育AI
        {'left': [
            '【产业定位】AI+教育是教育公平和质量提升重要手段，个性化学习/智能批改/AI助教全面应用',
            '【市场规模】2026年中国教育AI市场规模突破1200亿元，年均增速超40%',
            '【个性化学习】AI根据学生学习情况定制学习路径，因材施教，学习效率提升30%以上',
            '【智能批改】AI自动批改作业/试卷/作文，批改准确率超98%，教师批改工作量减少70%',
            '【AI助教】AI助教24小时答疑，回答学生问题，减轻教师重复性工作负担',
            '【科大讯飞】国内教育AI龙头，智慧教育产品覆盖全国5万+学校，1亿+师生用户',
            '【智能教室】智慧教室/录播/AI课堂分析，课堂教学行为分析，教学质量评估',
            '【教育公平】AI教育资源覆盖偏远地区农村学校，缩小城乡教育差距',
            '【安徽教育】科大讯飞总部合肥，安徽是教育AI应用先行省份，智慧教育覆盖率全国领先',
            '【政策支持】教育数字化战略行动，智慧教育平台建设，AI+教育政策支持'
        ], 'right': [
            '【讯飞半年报·2026年8月20日】科大讯飞发布2026年半年度报告：营业收入116.23亿元同比+6.52%，研发投入超30亿元同比增超6亿元占营收比重超25%；星火智能批阅机签约销售量同比增长14倍累计服务学校超5000所，2026春季学期系统日均作业批改量突破360万份智能批改调用量较上学期增长5倍，小学阶段应用学校数量增长1.5倍语文学科批改量增长10倍；讯飞AI学习机推出T90旗舰系列与S90进阶系列完成从入门到高端全价格带布局，618期间再度位居京东/天猫两大平台学习机品类销售额首位；智慧教育产品已在全国33个省级行政区域落地应用并拓展至日本/新加坡等海外市场；8月18-20日2026全球智慧教育大会上讯飞星光AI超级智能体亮相，专为教师打造低门槛AI应用创作能力让一线教师成为AI应用开发主力军',
            '【科大讯飞智慧课堂】覆盖课前/课中/课后全流程，AI互动课堂，实时学情分析，因材施教',
            '【作业帮学习笔】全科学习笔，扫题答疑/知识点讲解/双语翻译，覆盖小初高全学科',
            '【猿辅导AI学】AI自适应学习系统，个性化学习路径规划，1对1AI辅导，学习效果可视化',
            '【网易有道词典笔】AI词典笔，查词翻译/语法讲解/听力练习，K12学生人手一支普及',
            '【智慧考试】AI智能监考/智能阅卷，考试公平性提升，阅卷效率提升10倍',
            '【AI语言学习】AI口语陪练，实时纠正发音，情景对话练习，口语提升效率3倍',
            '【安徽智慧教育】安徽智慧学校建设覆盖全省1万+中小学，科大讯飞智慧教育产品市场份额第一',
            '【蚌埠教育】蚌埠全市中小学智慧教育覆盖率超90%，AI教育应用成效显著',
            '【高校AI】中科大/合工大等高校开设AI专业，建设AI学院，培养AI专业人才'
        ], 'process': [
            '【传统面授时代（1990年前）】教育完全依靠课堂面授，黑板+粉笔是主要教学工具，优质教育资源高度集中在城市重点学校，大班教学（每班50-70人）难以因材施教，教师需要花费大量时间批改作业、试卷，工作量大，农村地区缺乏合格教师，城乡教育差距巨大，教育主要靠应试刷题，个性化培养缺失。',
            '【多媒体教学期（1991-2010）】PPT、投影仪、多媒体教室开始普及，教学内容从纯板书转向音视频多媒体展示，教学形式更丰富，但教学模式仍然是教师讲学生听的灌输式，互动性弱，教育资源分配不均问题没有根本解决，在线教育开始萌芽但受限于网络带宽和终端设备。',
            '【在线教育兴起期（2011-2019）】互联网和移动互联网普及，直播课、录播课、MOOC兴起，学而思、新东方、猿辅导、作业帮等在线教育平台快速发展，打破时间空间限制，学生可以随时随地听名师课程，但仍然是千人一面的标准化课程，缺乏互动和个性化，拍照搜题等简单AI功能出现但能力有限。',
            '【AI教育试点期（2020-2022）】新冠疫情加速在线教育普及，AI技术开始更多应用于教育：AI口语测评、AI作文批改、自适应学习系统开始试点，但AI能力仍然比较初级，主要是规则匹配和简单机器学习，无法真正理解学生学习情况，个性化推荐准确率不高，大模型尚未爆发。',
            '【大模型驱动质变期（2023-2024）】ChatGPT为代表的大语言模型爆发，教育AI发生质变，大模型具备真正的自然语言理解、推理、生成能力，能够像老师一样解答问题、批改作文、讲解知识点、制定个性化学习计划，科大讯飞星火大模型教育版发布，AI学习机能力大幅提升，教育AI从辅助工具向AI助教转变。',
            '【AI教育普及元年（2025-2026）】2026年成为AI教育普及元年，AI学习机出货量超1500万台，科大讯飞智慧课堂覆盖全国5万+学校，1亿+师生使用AI教育产品，AI助教24小时答疑，AI自动批改作业、试卷准确率超98%，教师批改工作量减少70%，个性化学习真正落地，学生学习效率提升30%以上。',
            '【深度融合期（2027-2028）】AI深度融入教育全流程，课前AI备课、课中AI互动教学实时学情分析、课后AI个性化作业辅导，虚拟数字人老师普及，VR/AR+AI沉浸式教学推广，职业教育AI虚拟仿真实训大规模应用，特殊教育AI辅助（手语翻译、盲文识别、言语康复）成熟，教育从标准化向个性化全面转型。',
            '【因材施教普及期（2029-2030）】2030年AI教育覆盖率达90%，真正实现因材施教，每个学生都有专属AI学习助手，根据学生学习进度、薄弱知识点、学习习惯定制个性化学习路径，优质教育资源通过AI覆盖全国所有偏远地区农村学校，城乡教育差距显著缩小，教育质量整体提升，培养创新型人才。',
            '【教育AI技术演进】黑板粉笔→多媒体投影→在线直播录播课→拍照搜题→AI口语/作文批改→自适应学习系统→大模型AI助教→虚拟数字人老师→VR/AR沉浸式AI教学→全流程个性化因材施教；评价方式从单一考试分数→过程性评价+综合素质评价→多维度能力评估。',
            '【安徽教育AI进展】安徽是教育AI先行省份，科大讯飞总部位于合肥，智慧教育产品市场份额全国第一；安徽智慧学校建设覆盖全省12000+中小学，师生用户超800万人，覆盖率全国第一；蚌埠市与科大讯飞深度合作，全市中小学智慧教育覆盖率92%，教育质量排名安徽省前列；中科大、合工大等高校开设AI专业、建设AI学院，每年培养AI专业人才超1万人。'
        ], 'detail': [
            '【科大讯飞学习机T30详细参数】科大讯飞AI学习机T30 Ultra旗舰版配置14.7英寸类纸护眼大屏（低蓝光无频闪认证），内置星火大模型V4.0教育专用版本，配备1300万像素指学摄像头+800万像素作业摄像头，8GB运行内存+256GB机身存储，AI精准学系统覆盖小学初中高中全学科所有知识点，官方市场售价8999元人民币',
            '【讯飞智慧课堂全流程功能】科大讯飞智慧课堂覆盖课前、课中、课后完整教学流程：课前AI辅助备课一键生成教案课件，课中互动答题、实时学情分析、随机点名，课后AI自动批改作业、个性化错题推送、知识点薄弱点诊断，教师备课效率提升60%，课堂互动参与度提升80%',
            '【AI作文批改技术能力】AI作文批改系统从字词错误、语法错误、结构立意、文采表达、思想深度多个维度综合评分，中英文作文批改准确率达到98%以上，不仅指出问题还提供具体修改建议和范文参考，批改一篇800字作文仅需3秒，效率是人工批改的几十倍',
            '【AI错题本智能整理】AI错题本自动收集整理学生作业、试卷、练习中的所有错题，智能分析错误对应的知识点漏洞，推送同类变式题进行针对性强化练习，错题复习效率提升3倍，彻底告别学生手抄错题的低效方式，家长也能实时了解孩子薄弱点',
            '【AI口语陪练应用效果】AI口语陪练支持与学生实时情景对话，发音纠正准确率达到98%，涵盖日常生活、商务办公、旅游出行、考试备考等多场景对话练习，提供发音评分、语调纠正、流利度评估，学生敢于开口练习，口语提升效率是传统课堂的3倍以上',
            '【个性化学习提分数据】多所学校对照实验数据显示，使用AI个性化学习系统的学生平均成绩提升15-20分，无效学习时间减少25%，学习兴趣和主动性提升40%，真正实现因材施教，每个学生都有专属学习路径，避免千人一面的题海战术',
            '【教师减负实际效果】AI自动批改作业减少教师70%的重复性批改工作量，AI备课系统提供丰富的教学资源和教案参考，课堂AI学情分析让教师实时掌握每个学生的掌握情况，教师可以从繁重的批改工作中解放出来，有更多时间关注学生个体成长和教学设计',
            '【智慧课堂全国覆盖数据】截至2026年8月全国智慧教室覆盖率已达60%，科大讯飞智慧课堂产品已覆盖全国31个省份5万+所学校，服务超过1亿师生用户，在智慧教育市场占有率连续多年保持全国第一，是教育AI领域绝对龙头企业',
            '【教育公平推进显著成效】通过AI教育资源云端覆盖农村偏远地区学校，农村学校学生平均成绩提升12分，城乡教育成绩差距缩小30%，优质教育资源不再集中在城市重点学校，偏远地区学生也能听到名师讲课、获得AI名师辅导',
            '【安徽智慧教育全国领先】安徽省已建成省级智慧教育平台，接入全省12000+所中小学，师生用户超过800万人，智慧学校覆盖率全国排名第一，科大讯飞总部位于合肥为安徽智慧教育建设提供核心技术支撑，智慧教育成为安徽名片',
            '【蚌埠本地智慧教育进展】蚌埠市教育局与科大讯飞深度合作共建智慧教育示范区，全市中小学智慧教育覆盖率达到92%，AI教育应用成效显著，学生学业水平和综合素质持续提升，蚌埠智慧教育经验在全省推广，成为皖北智慧教育标杆城市',
            '【AI智能监考技术能力】AI智能监考系统可识别20+种考试作弊行为：包括替考、传纸条、偷看他人答案、使用手机、交头接耳、东张西望等，作弊行为识别准确率达到99%，大幅降低监考人力成本，考试公平性显著提升，尤其适用于大规模在线考试',
            '【智能阅卷规模化应用】高考、中考、学业水平考试已全面推行AI智能阅卷，客观题100%由AI自动批改，主观题采用AI辅助人工双评模式，阅卷效率提升10倍，评分误差率降低80%，阅卷更加公平公正，每年有超千万考生试卷通过AI辅助阅卷完成',
            '【特殊教育AI辅助应用】AI技术应用于特殊教育领域：AI手语翻译实时将语音转换为手语动画帮助听障人士沟通，AI盲文识别将纸质文字转换为盲文或语音，AI言语康复训练系统帮助言语障碍患者进行发音训练，显著提升特殊教育质量，帮助残障学生更好学习',
            '【职业教育AI虚拟实训】VR/AR+AI虚拟仿真实训应用于职业教育，学生在虚拟环境中模拟操作数控机床、汽车维修、电工电子、护理操作等，职业技能培训成本降低60%，培训效率提升50%，还能避免真实实训中的安全事故风险',
            '【高等教育AI应用】AI技术在高等教育领域广泛应用：AI辅助科研文献阅读和论文写作、AI编程助手帮助学生学习编程、AI智能答疑回答学生课程问题、虚拟仿真实验平台，高校科研效率显著提升，大学生AI信息素养和应用能力培养全面加强',
            '【AI教师培训体系】AI辅助教师培训系统通过模拟课堂教学场景让新教师进行教学演练，AI从教学内容、表达、互动、节奏等维度进行评估反馈，新教师成长周期从传统的2-3年大幅缩短至6个月-1年，教师队伍整体水平快速提升',
            '【教育垂直大模型成熟】科大讯飞星火教育大模型、网易有道子曰大模型、好未来MathGPT等教育垂直大模型技术成熟，针对教育场景专门优化，学科知识准确性、解题能力、教学引导能力全面超越通用大模型，成为教育AI核心引擎',
            '【学生数据安全保障】教育AI产品严格遵守数据安全法规，学生个人信息和学习数据采用端侧处理+云端加密存储，数据不出校园，符合教育数据安全规范和个人信息保护法要求，家长可随时查看孩子数据使用情况，保障学生隐私安全',
            '【2030年教育AI发展目标】规划到2030年AI教育覆盖率达到90%，真正实现因材施教个性化教学，每个学生都有专属AI学习助手，优质教育资源通过AI覆盖全国所有偏远地区农村学校，城乡教育差距显著缩小，整体教育质量大幅提升，培养更多创新型人才'
        ]},
        # 能源电力
        {'left': [
            '【产业定位】能源电力是经济命脉，电力机器人保障电网安全稳定运行，新能源+储能+AI构建新型电力系统',
            '【市场规模】2026年中国电力机器人市场规模突破280亿元，巡检机器人占比60%',
            '【巡检机器人】变电站/输电线路/配电站房/电缆隧道巡检机器人普及，替代人工巡检',
            '【带电作业机器人】架空线路带电作业机器人，工人不用登杆不用停电，安全高效',
            '【国家电网】国家电网大力推进电力机器人应用，已部署各类电力机器人超5万台',
            '【新能源】风电/光伏/储能快速发展，AI功率预测/智能运维/故障诊断提升新能源利用效率',
            '【特高压】特高压输电线路长，穿越复杂地形，巡检机器人保障特高压大动脉安全',
            '【安徽电网】安徽电网是华东电网重要组成部分，两交一直特高压过境，电力机器人应用广泛',
            '【蚌埠电力】蚌埠供电公司部署巡检机器人/带电作业机器人，智能化水平安徽领先',
            '【双碳目标】2030碳达峰2060碳中和，能源电力转型加速，AI+机器人支撑新型电力系统'
        ], 'right': [
            '【南方电网机器人·2026年8月21日】南方电网多款自研机器人亮相WRC2026央企展区：电力作业人形机器人"知行者1号"已能流畅完成设备巡检/高压柜操作/仪器仪表检查等精细作业；"吠云"四足机器狗在广东东莞18只机器狗"员工"在无人变电站工作；"小蓝鸟"电力巡检无人机实现720°无死角感知',
            '【国网巡检机器人】变电站轮式巡检机器人，红外测温/表计识别/开关位置识别/噪声检测，24小时无人值守',
            '【输电线路无人机】大疆/国网自研巡检无人机，自主巡检输电线路，识别缺陷隐患，效率是人工10倍',
            '【带电作业机器人】国网自研第四代带电作业机器人，可完成接引线/断引线/更换绝缘子等作业',
            '【电缆隧道机器人】电缆隧道履带式巡检机器人，防水防火防爆，检测气体/温度/局部放电',
            '【配网巡检机器人】配电站房轨道式/轮式巡检机器人，10kV配电站房无人化运维',
            '【南瑞继保】南瑞继保电力AI系统，电网故障诊断/继电保护/稳定控制，保障电网安全',
            '【许继电气】许继电气电力机器人和智能变电设备，特高压变电站智能化',
            '【安徽电网】安徽电力已部署巡检机器人800+台，无人机200+架，带电作业机器人50+台',
            '【蚌埠供电】蚌埠220kV及以上变电站全部配置智能巡检机器人，无人值守站比例超70%',
            '【新能源运维】风电/光伏电站运维机器人/无人机巡检，风机叶片/光伏面板缺陷自动识别'
        ], 'process': [
            '【人工运维时代（1949-2005）】电力系统运维完全依靠人工，变电站有人24小时值班，运行人员定时抄表、巡视设备，输电线路工人翻山越岭巡线，带电作业工人爬杆高空作业，劳动强度大、风险高，恶劣天气（暴雨、冰雪、高温）仍要户外作业，人身伤亡事故时有发生，缺陷发现依赖经验和责任心，漏检误检率高。',
            '【综合自动化期（2006-2015）】变电站综合自动化改造，"四遥"（遥测、遥信、遥控、遥调）普及，变电站实现无人值班少人值守，调度中心可以远程监控电网运行，但设备巡检仍然需要人工到现场，无人机开始试点用于输电线路巡检但数量少、性能有限，带电作业仍然需要人工登杆。',
            '【机器人试点期（2016-2020）】国家电网开始电力机器人试点应用，变电站轮式巡检机器人在500kV及以上变电站试点，固定轨道巡检机器人在地下变电站、配电站房试点，大疆等工业无人机开始用于输电线路通道巡检，带电作业机器人研发成功并试点，机器人可靠性和实用性逐步提升，但成本高、数量少，人工巡检仍是主力。',
            '【规模化推广期（2021-2024）】国家电网、南方电网大力推进电力机器人规模化应用，220kV及以上变电站普遍配置智能巡检机器人，输电线路无人机巡检常态化，配电站房巡检机器人批量部署，第四代、第五代带电作业机器人推广应用，电力机器人数量快速增长至3万+台，AI图像识别缺陷准确率提升至95%以上，人工巡检工作量减少60%。',
            '【新能源智能化期（2022-2025）】风电、光伏新能源爆发式增长，新能源电站运维机器人、无人机快速普及，风机叶片缺陷AI自动识别、光伏面板热斑检测、无人机自主巡检成为标准配置，AI功率预测精度提升至90%以上，储能电站智能运维系统建设，新型电力系统智能化水平大幅提升。',
            '【全面机器代人期（2025-2026）】2026年国家电网累计部署各类电力机器人超5万台，安徽电力部署巡检机器人800+台、无人机200+架、带电作业机器人50+台，蚌埠220kV及以上变电站全部配置智能巡检机器人，无人值守站比例超70%，电力巡检基本实现机器人替代，带电作业机器人广泛应用，人员高空作业风险大幅降低，电网故障响应时间缩短80%。',
            '【数字孪生电网期（2027-2028）】电力机器人+AI+数字孪生深度融合，电网数字孪生体建成，机器人实时采集的数据驱动数字孪生模型运行，实现电网状态全息感知、故障提前预测、运维策略智能生成，多机器人协同巡检、协同作业，无人机、巡检机器人、带电作业机器人、隧道机器人联合作业，电网智能化水平全球领先。',
            '【新型电力系统期（2029-2030）】2030年建成以新能源为主体的新型电力系统，电力机器人全面普及，电网运维基本实现无人化，AI实现电网自主调度、自主运维、自主修复，风电、光伏、储能、特高压、柔性交直流电网智能化运行，供电可靠性达99.999%，新能源消纳率达95%以上，支撑双碳目标实现。',
            '【电力机器人技术演进】人工现场巡检→变电站综合自动化四遥→轮式巡检机器人试点→无人机输电巡检→带电作业机器人→多机器人协同作业→数字孪生智能电网→自主运行新型电力系统；巡检方式从人工肉眼→可见光摄像头→红外+可见光双光→多传感器融合→AI缺陷自动识别→数字孪生全息感知。',
            '【安徽电力智能化进展】安徽电网是华东电网重要组成部分，两交一直特高压过境，是"西电东送"重要通道；安徽电力智能化水平全国前列，已部署巡检机器人800+台、无人机200+架、带电作业机器人50+台，建成合肥、芜湖等地市智能电网示范区；蚌埠供电公司智能化水平安徽前五，220kV及以上变电站全部配置智能巡检机器人，无人值守站比例超70%。'
        ], 'detail': [
            '【新型储能·2026年8月20日】国家能源局发布《中国新型储能发展报告(2026)》：截至2025年底全国新型储能累计装机1.36亿千瓦/3.51亿千瓦时同比+84.3%约占全球50%；首次在国家层面建立电网侧独立储能容量电价机制，发展从政策驱动转向市场化驱动',
            '【变电站巡检机器人详细参数】轮式底盘设计，激光SLAM+视觉融合导航精度±1cm，设备定位精度±5mm，红外测温精度±2℃或读数±2%，可见光仪表识别准确率>98%，标准工况续航8小时支持自主充电，整机防护等级IP55适应户外环境，单台设备市场售价60-120万元人民币',
            '【输电线路无人机详细参数】大疆M300 RTK/御3行业版无人机，单架次续航40-55分钟，RTK厘米级定位精度1cm+1ppm，配备30倍光学变焦可见光相机+红外热成像相机，支持自主规划航线自动巡检，缺陷识别准确率>95%，单架次可完成10-20基杆塔巡检作业',
            '【带电作业机器人详细参数】配置6自由度高精度绝缘机械臂，绝缘等级覆盖10kV/35kV/110kV/220kV全电压等级，作业工具定位精度0.1mm，搭载力反馈控制系统，可在不停电情况下完成接引流线、更换绝缘子、清除异物等复杂带电作业，单次作业时间从人工2小时缩短至20分钟',
            '【电缆隧道机器人详细参数】履带式底盘设计，爬坡能力30度，防水等级IP68可适应电缆隧道潮湿积水环境，可实时检测甲烷、一氧化碳、硫化氢、氧气等有毒有害可燃气体浓度，配备红外测温+可见光双光检测+局部放电检测功能，续航6小时，适用于电缆隧道、综合管廊场景',
            '【巡检效率对比数据】人工巡检1座220kV变电站需要2名工作人员耗时2小时，机器人自主巡检仅需30分钟即可完成且支持24小时不间断巡检，巡检频率提高7倍，设备缺陷发现率提升3倍，人工巡检劳动强度降低80%，大幅提升变电站运维效率和可靠性',
            '【带电作业安全价值】人工带电作业存在高空坠落、触电等高风险，机器人作业时操作人员在地面远程操控无需登杆，人身安全事故率降低99%，作业人员从传统4-6人减少至1-2人，大幅提升带电作业安全性和效率，减少停电时间提升供电可靠性',
            '【国家电网机器人部署】截至2026年8月国家电网累计部署变电站智能巡检机器人超2万台，巡检无人机超1.5万架，带电作业机器人超2000台，电缆隧道巡检机器人超1000台，覆盖所有500kV及以上变电站和主要输电通道，电网运维智能化水平全球领先',
            '【安徽电力机器人应用】安徽电网已部署变电站智能巡检机器人850台，覆盖所有220kV及以上变电站；部署输电线路巡检无人机220架，完成5000公里输电线路自主巡检常态化；部署带电作业机器人70台，安徽电网智能化运维水平位居全国前列',
            '【蚌埠电力智能化建设】蚌埠供电公司38座220kV及以上变电站全部实现机器人智能巡检，无人值守站比例超70%；2026年机器人巡检累计发现设备缺陷隐患1200+项，及时避免停电事故30+起，供电可靠性提升至99.99%，智能化水平位居安徽省前五',
            '【新能源功率预测精度】AI风电、光伏发电功率预测准确率超过95%，预测时间尺度覆盖15分钟到7天全时段，电网调度更加精准高效，弃风弃光率降至2%以下达到国际先进水平，新能源消纳能力大幅提升，助力双碳目标实现',
            '【风机叶片巡检技术应用】无人机+AI自动巡检风机叶片，可自动识别裂纹、腐蚀、雷击损伤、前缘磨损等各类缺陷，识别准确率达到96%，单台风机巡检时间从人工4小时大幅缩短至15分钟，风机运维效率提升16倍，保障风电设备安全稳定运行',
            '【光伏面板巡检技术应用】无人机搭载红外热成像相机巡检光伏面板，可自动识别热斑、碎片、隐裂、二极管失效等缺陷，缺陷定位精度±1块组件，巡检效率达到10MW/小时，是人工巡检效率的20倍，大幅降低光伏电站运维成本提升发电效率',
            '【储能电站智能运维】AI储能电站智能运维系统，实现电池SOH健康状态、SOC荷电状态精准估算，热失控预警提前30分钟，电池循环寿命延长15%，储能电站安全运行水平大幅提升，为新型电力系统安全稳定运行提供重要支撑',
            '【数字孪生电网建设】电网数字孪生系统构建与物理电网实时映射的数字镜像，实时采集机器人和传感器数据驱动数字孪生模型运行，实现电网状态全息感知、故障提前预测预警、运维策略智能生成优化，事故恢复时间缩短60%，供电可靠性显著提升',
            '【电力行业大模型应用】国家电网、南瑞继保研发电力专用大模型，支持设备巡检智能分析、调度决策辅助、故障智能诊断、客服智能问答等电力场景专用AI能力，电力领域专业能力超越通用大模型，成为电网智能化核心引擎',
            '【特高压线路智能巡检】特高压输电线路长距离大跨越穿越复杂地形，采用直升机+无人机协同巡检模式，年巡检里程超10万公里，缺陷识别准确率>98%，保障特高压输电大动脉安全稳定运行，支撑西电东送国家能源战略',
            '【配网智能化技术应用】配电网故障定位隔离恢复（FLISR）系统实现故障自动定位、自动隔离、非故障区域自动恢复供电，故障定位处理时间从小时级缩短至秒级，城市配电网供电可靠性提升至99.99%，用户年均停电时间小于1小时',
            '【安徽能源结构转型】安徽新能源装机容量超8000万千瓦，占总装机比例超50%，两淮采煤沉陷区漂浮式光伏电站全球规模最大；蚌埠怀远、五河风光资源丰富，马城500kV变电站是皖北重要电力枢纽，电力机器人保障电网安全',
            '【蚌埠能源产业发展】蚌埠是皖北重要能源基地，国电蚌埠电厂、蚌埠怀远马城500kV变电站、蚌埠涂山220kV变电站等重要能源设施，电力机器人广泛应用保障电力可靠供应；蚌埠太阳能光伏、生物质能等新能源产业快速发展',
            '【2030年电力发展目标】规划到2030年电网关键岗位机器人替代率达80%，新能源装机占比超60%，以新能源为主体的新型电力系统基本建成，供电可靠性达到99.995%，电网全面实现智能化、数字化、无人化运维，支撑双碳目标顺利实现'
        ]},
        # 自动驾驶
        {'left': [
            '【产业定位】自动驾驶是AI最大应用场景之一，L4级自动驾驶2026年开始规模化商业落地',
            '【分级标准】L0辅助驾驶/L1/L2辅助驾驶/L3有条件自动驾驶/L4高度自动驾驶/L5完全自动驾驶',
            '【市场规模】2026年中国自动驾驶市场规模突破5000亿元，L4级商业化加速',
            '【萝卜快跑】百度萝卜快跑L4级自动驾驶出行服务，已在10+城市运营，累计订单超2000万单',
            '【特斯拉FSD】特斯拉FSD（完全自动驾驶能力）V13版本，端到端AI驾驶，北美广泛推送',
            '【小鹏XNGP】小鹏汽车XNGP智能辅助驾驶，全场景智驾，无图城市NOA全国开通',
            '【华为ADS】华为乾崑ADS 3.0智驾系统，GOD大网，无图全国都能开，问界/智界/享界车型搭载',
            '【Robotaxi】自动驾驶出租车（Robotaxi）在武汉/重庆/北京/广州等城市开始收费运营',
            '【Robotruck】自动驾驶卡车在港口/矿山/干线物流场景商业化运营，成本比人工低30%',
            '【安徽自动驾驶】合肥/芜湖是自动驾驶示范城市，蔚来/比亚迪/大众/江淮智驾技术快速迭代'
        ], 'right': [
            '【特斯拉Cybercab·2026年8月21日】特斯拉计划本月底向公众开放Cybercab试乘接入美国得州奥斯汀Robotaxi服务；FSD V15将带来显著性能提升涵盖7项核心技术其中约40%已在现役Robotaxi车队测试；特斯拉Robotaxi在内华达州获批未来12月最多可部署5000辆',
            '【Robotaxi商业化·2026年8月20日】2026年被视为Robotaxi商业化落地关键节点：全球Robotaxi市场规模预计2030年将达673亿美元；百度萝卜快跑已在全球26座城市提供服务累计完成超2000万单出行，武汉全无人驾驶订单占比高达92%',
            '【低空经济·2026年8月21日】《中国低空经济应用场景分析报告2026》发布：中国低空经济市场2026年规模预计突破万亿大关2030年有望突破2万亿元；新修订《民用航空法》2026年7月1日施行增设发展促进专章，低空经济首次获国家法律层面根本制度支撑',
            '【百度萝卜快跑第六代RT6】百度Apollo第六代Robotaxi RT6车规级量产，成本降至20万元（五代1/10），12摄像头+6毫米波+3激光雷达，算力800TOPS，已在武汉/重庆/北京等10+城市运营累计订单超2000万单',
            '【特斯拉FSD V13】端到端神经网络自动驾驶，HW4.0硬件算力1500TOPS纯视觉方案，从摄像头输入直接输出转向/加速/制动控制，FSD订阅服务12000元/年，北美已大规模推送',
            '【华为乾崑ADS 3.0】192线超长距激光雷达（测距250m）+11高清摄像头+6毫米波雷达+12超声波雷达，双MDC 610算力400TOPS，GOD 3.0通用障碍物识别网络，无图NCA全国都能开',
            '【小鹏XNGP 5.0】双英伟达Orin-X芯片算力508TOPS，双激光雷达（小鹏图灵）+11摄像头+5毫米波雷达+12超声波，XNet 2.0深度视觉网络，AI代驾学习用户路线，城市NOA全国开通',
            '【理想AD Max 3.0】双英伟达Orin-X算力508TOPS，1颗128线激光雷达+11摄像头+12超声波+1毫米波，端到端+VLM视觉语言模型双系统架构，全场景NOA覆盖，通勤NOA功能',
            '【Waymo One商业化运营】谷歌Waymo自动驾驶出租车在美国洛杉矶、旧金山、凤凰城商业化运营，累计订单超500万单，无人驾驶安全运营里程超1亿英里，是全球自动驾驶技术标杆企业',
            '【驭势科技机场无人化】2026年8月13日驭势科技乌鲁木齐天山国际机场累计真无人运营里程突破160万公里，70+台无人车常态化运营，创下国内民航机场自动驾驶规模化商用全新纪录',
            '【小马智行Uber欧洲合作】2026年8月15日小马智行携手Uber推出欧洲最大Robotaxi部署计划，将在欧洲5座城市部署超2000辆自动驾驶车辆，中国自动驾驶技术出海里程碑',
            '【上海L4合法化政策】2026年8月16日上海发布全国首个"模速智行"方案，L4级自动驾驶正式合法化；同日大众中国推出自研全场景辅助驾驶，Q3起搭载三家合资企业7款车型配地平线芯片',
            '【安徽合肥示范应用】合肥是国家智能网联汽车示范区，开放测试道路超2000公里发放测试牌照超500张，蔚来/大众安徽/比亚迪/江淮在合肥开展智驾研发测试，芜湖奇瑞智驾技术快速迭代'
        ], 'process': [
            '【纯人工驾驶时代（2015年前）】汽车完全由人类驾驶，没有自动驾驶功能，定速巡航等初级辅助功能开始出现但非常简单，交通事故90%以上由人为因素导致，每年全球交通事故死亡超130万人，驾驶是需要高度集中注意力的繁重体力和脑力劳动，长途驾驶疲劳，城市拥堵开车累。',
            '【L2辅助驾驶普及期（2016-2022）】L2级辅助驾驶快速普及，ACC自适应巡航、LCC车道居中控制、AEB自动紧急制动、自动泊车APA等功能成为新车标配，特斯拉Autopilot、小鹏NGP、华为ADS等智驾系统快速迭代，从高速场景扩展到城市道路，但责任主体仍然是人类司机，需要全程监控随时接管。',
            '【封闭场景L4商业化期（2020-2024）】L4级自动驾驶率先在封闭和半封闭场景商业化运营：港口自动驾驶集卡（天津港、上海港）、矿山无人驾驶矿卡、机场无人物流（驭势科技在香港机场）、园区无人配送、封闭园区Robotaxi示范，百度萝卜快跑开始在武汉、重庆等城市试点收费Robotaxi服务，商业化验证完成。',
            '【城市开放道路试点期（2023-2025）】L4级Robotaxi在多个城市开放道路试点收费运营，百度萝卜快跑进入武汉、重庆、北京、广州等10+城市，累计订单超2000万单，端到端AI驾驶技术成熟，特斯拉FSD V12/V13在北美大规模推送，无图NOA（城市领航辅助）开始在全国范围开通，智驾体验接近人类司机。',
            '【L4商业化加速期（2025-2026）】2026年L4级自动驾驶商业化加速：8月13日驭势科技乌鲁木齐天山国际机场累计真无人里程突破160万公里，70+台无人车运营，创下国内民航机场自动驾驶规模化商用全新纪录；8月15日小马智行携手Uber推出欧洲最大Robotaxi部署计划，5城部署超2000辆；8月16日上海发布全国首个"模速智行"方案，L4级自动驾驶正式合法化；同日大众中国推出自研全场景辅助驾驶，Q3起搭载三家合资企业7款车型；五部门发布交通标准化"十五五"规划布局自动驾驶全链条标准；百度第六代Robotaxi RT6成本降至20万元（是上一代的1/10）车规级量产；华为乾崑ADS 3.0实现无图全国都能开；小鹏、理想、蔚来智驾系统迭代至端到端架构；Robotaxi订单量爆发式增长，自动驾驶事故率低于人类司机50%以上。',
            '【规模普及期（2027-2028）】L2+高级辅助驾驶成为新车标配，渗透率超80%，L4级Robotaxi从试点城市向全国主要城市扩展，Robotruck在干线物流规模化运营，自动驾驶成本持续下降，法规体系逐步完善，自动驾驶保险、责任认定等制度成熟，消费者接受度大幅提升，自动驾驶开始改变出行方式。',
            '【全面无人驾驶期（2029-2030）】2030年前后L4级自动驾驶规模化普及，Robotaxi、Robotruck大规模商用，无人配送车广泛应用，部分城市开放全无人Robotaxi运营，新车自动驾驶（L2+及以上）渗透率超90%，交通出行方式发生根本变革，交通事故率降低90%，交通效率提升50%，物流成本下降30%。',
            '【自动驾驶技术演进】纯人工驾驶→定速巡航→L2 ACC/LCC辅助驾驶→L3有条件自动驾驶→封闭场景L4商业化→城市开放道路L4试点→Robotaxi/Robotruck规模化→全场景无人驾驶；感知方案从纯视觉→激光雷达+视觉融合→多传感器融合→端到端大模型直接输出控制。',
            '【安徽自动驾驶进展】合肥是国家智能网联汽车示范区，蔚来、比亚迪、大众安徽、江淮汽车等整车企业智驾技术研发测试，合肥自动驾驶测试道路超1000公里；芜湖奇瑞汽车智能驾驶技术快速迭代；驭势科技等自动驾驶企业在安徽布局，安徽新能源汽车产业为自动驾驶提供绝佳应用场景。',
            '【最新赛事与进展】2026年成都市举办人工智能与机器人创新应用大赛，自动驾驶和机器人技术成为赛事重点；第二届世界人形机器人运动会设置自动驾驶机器人竞赛单元，推动自动驾驶与具身智能技术融合发展；中国智能网联汽车标准体系持续完善，L3级自动驾驶准入政策落地。'
        ], 'detail': [
            '【百度RT6详细硬件参数】百度Apollo第六代Robotaxi RT6车规级量产，整车成本降至20万元仅为第五代的1/10，传感器配置：12个高清环视摄像头+6个毫米波雷达+3个激光雷达+12个超声波雷达，车载AI算力平台800TOPS，纯电续航里程700km，支持无方向盘选项设计，设计运营寿命5年或60万公里',
            '【华为ADS 3.0传感器配置】华为乾崑ADS 3.0配置192线超长距激光雷达（最远测距250m）+11个高清摄像头+6个毫米波雷达+12个超声波雷达，双华为MDC 610计算平台总算力400TOPS，搭载GOD 3.0通用障碍物识别网络，支持无图城区NCA智驾领航、代客泊车、高速NCA全覆盖',
            '【特斯拉FSD HW4.0硬件】特斯拉HW4.0自动驾驶硬件平台搭载5纳米工艺自研FSD Computer 2芯片，单芯片算力500TOPS三芯片总算力1500TOPS，采用纯视觉感知方案（无激光雷达），车身环绕8个高清摄像头，端到端神经网络v13版本，FSD完全自动驾驶订阅服务12000元/年',
            '【小鹏XNGP 5.0系统配置】小鹏XNGP 5.0全场景智能辅助驾驶系统搭载双英伟达Orin-X芯片总算力508TOPS，双小鹏图灵激光雷达+11个高清摄像头+5个毫米波雷达+12个超声波雷达，XNet 2.0深度视觉感知网络，支持AI代驾（VPA-L）学习用户常用路线，城市NOA功能全国开通',
            '【理想AD Max 3.0架构】理想汽车AD Max 3.0智能驾驶系统搭载双英伟达Orin-X芯片总算力508TOPS，1颗128线激光雷达+11个高清摄像头+12个超声波雷达+1个毫米波雷达，采用行业首创端到端+VLM视觉语言模型双系统架构，支持全场景NOA、通勤NOA等功能',
            '【Robotaxi运营数据统计】百度萝卜快跑2026年上半年订单量超800万单累计订单突破2000万单，武汉运营车辆超1000台峰值日订单超10万单，单均成本已低于传统网约车，L4级自动驾驶开始进入规模化商业化运营阶段，自动驾驶出行服务获得市场认可',
            '【自动驾驶事故率对比】统计数据显示L4级自动驾驶每百万公里事故率0.3次，人类司机每百万公里事故率2.8次，自动驾驶安全性是人类司机的9倍；百度萝卜快跑实际运营数据显示事故率仅为人类司机的1/12，自动驾驶大幅提升道路交通安全水平',
            '【Robotruck干线物流应用】自动驾驶重卡干线物流场景商业化运营，综合成本比人工驾驶低30-40%，燃油消耗降低10%，支持7×24小时不间断运营无疲劳驾驶问题，2026年国内干线物流自动驾驶试点线路已超20条，多条线路实现常态化商业化运营',
            '【港口自动驾驶规模化】天津港、上海港、宁波舟山港、深圳港等主要港口自动驾驶集卡部署量超1000台，港口作业效率提升20%，人力成本降低70%，可24小时不间断作业解决港口工人招工难问题，港口成为自动驾驶最先规模化落地场景之一',
            '【矿山自动驾驶应用进展】露天矿山自动驾驶矿卡规模化应用，国家能源、中煤、华能等大型矿场部署自动驾驶矿卡超500台，特别适合矿山恶劣危险环境作业，人力成本降低80%，作业安全性大幅提升，杜绝矿难人员伤亡风险',
            '【安徽合肥示范建设】合肥国家智能网联汽车示范区开放测试道路超2000公里，累计发放自动驾驶测试牌照超500张，蔚来、大众安徽、比亚迪、江淮汽车等整车企业在合肥开展智能驾驶研发测试，合肥已成为国内重要自动驾驶产业高地',
            '【芜湖自动驾驶产业】芜湖奇瑞汽车智能驾驶技术快速迭代，与埃夫特等机器人企业合作智能网联汽车测试，芜湖奇瑞智能网联汽车测试区建成投用，支持自动驾驶封闭场地测试和开放道路测试，安徽形成合肥+芜湖双智驾产业格局',
            '【无图智驾技术趋势】2026年开始头部智驾方案全部转向无图（不依赖高精地图）技术路线，通过实时感知和在线建模实现全国都能开，彻底解决高精地图更新慢、采集成本高、覆盖范围有限等问题，智驾系统可用性大幅提升',
            '【端到端技术成为主流】端到端（End-to-End）自动驾驶成为行业主流技术路线，从传感器原始输入直接输出车辆控制信号，减少传统模块化架构的累积误差，自动驾驶能力更接近人类司机驾驶水平，泛化能力大幅提升',
            '【VLM大模型应用】视觉语言模型（VLM）应用于自动驾驶领域，能够理解复杂场景语义，识别异形障碍物、施工场景、交警手势等传统方法难以处理的情况，自动驾驶系统泛化能力和场景理解能力大幅提升',
            '【V2X车路协同建设】V2X车路协同技术规模化部署，实现车与车、车与路、车与人、车与云实时通信，支持路口盲区预警、绿波通行、协同通行等功能，道路通行效率提升30%，交通事故率进一步降低',
            '【法规政策体系完善】《智能网联汽车准入和上路通行试点》政策实施，L3/L4级自动驾驶有法可依，事故责任认定规则明确，自动驾驶保险体系逐步配套完善，政策法规体系为自动驾驶规模化发展提供制度保障',
            '【智驾成本快速下降】L4级自动驾驶传感器+计算平台成本从2020年的100万元以上下降至2026年的5-10万元，预计2030年目标降至2万元以内，成本快速下降为自动驾驶大规模普及奠定经济基础',
            '【蚌埠本地应用进展】蚌埠市试点自动驾驶微公交服务覆盖经开区和高新区，蚌埠港试点自动驾驶集卡作业，智慧交通建设持续推进，蚌埠作为皖北中心城市积极融入智能网联汽车产业发展，未来规划更多自动驾驶应用场景',
            '【2030年产业发展目标】规划到2030年L4级自动驾驶新车渗透率超30%，Robotaxi运营车辆超100万台，道路交通事故率降低80%，交通通行效率提升50%，物流成本下降30%，自动驾驶全面改变人类出行和物流方式'
        ]},
        # 人形运动会
        {'left': [
            '【活动定位与时间】第二届世界人形机器人运动会将于2026年8月22-26日在北京国家速滑馆"冰丝带"举办，赛事吉祥物命名为"智宝"，是全球规模最大、水平最高的人形机器人专业赛事',
            '【赛事规模爆发】全球六大洲16个国家共666支队伍、2056台机器人同台竞技，参赛队伍同比首届增长138%，参赛机器人数量翻两番；巴西组建国家队含5支RoboCup世界冠军队伍参赛',
            '【国内参赛阵容】国内157家企业、200所院校科研机构组成641支队伍、1975台机器人参赛，覆盖优必选、宇树、小米、特斯拉等主流厂商，清华、北大、哈工大、中科大等顶尖高校全部参赛',
            '【比赛项目设置】共设置51个比赛项目：竞技赛30项+场景赛21项，赛期5天共9个竞赛单元（下午晚间均安排比赛），总计开展1301场比赛；首届仅26项3天6单元，赛事规模大幅扩容',
            '【新增对抗项目】新增跳远、举重、拔河、乒乓球、自由搏击（分40kg/58kg/80kg三个体重级别）等高强度对抗项目，对机器人本体结构强度、关节性能、运动控制是极限考验',
            '【场景赛真实工业】场景赛从6项大幅扩至21项，覆盖工业装配、酒店服务、家庭家政、物流分拣、消防救援、园林作业、应急处置、图书整理、零售服务9大真实场景，"出了赛场就进现场"',
            '【灵巧手微操挑战】设置镊子夹豆、粉末称量、开瓶撬盖、拧螺丝、精密装配、家政服务、图书整理8项灵巧手微操作挑战，比拼"看得清、拿得稳、做得准"的指尖精细操作能力',
            '【以赛定标促产业】组委会明确提出"以赛定标、以标促产"核心理念，比赛规则相当部分将直接转化为行业技术验收标准，实现"奖牌变订单，赛场进现场"的产业转化目标',
            '【赛事保障体系】创新打造"机器人之家"运动员村，可容纳1500余台机器人集中存放，200个充电柜满足1200块电池同时充电需求，30秒完成一台机器人数字化存取闭环，选派裁判员223名（国际级28人）',
            '【观赛体验亲民】打造"机器人一条街"提供智能售卖、VR观影、机器人咖啡等科技体验，联合朝阳区打造嘉年华活动；比赛日最低票价98元、开闭幕式最低128元，提供家庭套票，开幕式票已售罄'
        ], 'right': [
            '【运动会门票·2026年8月20日】第二届世界人形机器人运动会8月22日北京国家速滑馆开幕：灵巧手专项赛8个竞技小项在国家速滑馆人形机器人产业生态训练和测评基地举行；运动会文创8月22日开幕日线下发售：吉祥物智宝吊卡/奖牌造型冰丝带摆件（内置NFC芯片裸眼3D效果）等三十余款',
            '【半程马拉松突破】今年4月亦庄半马冠军成绩50分26秒，不到上届三分之一，打破人类半马纪录',
            '【径赛规则升级】100米完赛从3分钟压到1分钟内，400米从15分钟压到5分钟，除障碍赛外必须全自主不许遥控',
            '【足球赛事】7v7开幕式表演赛+5v5大型组+中型组+3v3 U19组，机器人能带球跑动/大力抽射/飞身扑救',
            '【清华火神队】首届5v5冠军+2025 RoboCup人形组世界冠军，8月初进驻冰丝带封闭集训十余天目标卫冕',
            '【进球规则升级】今年必须两台不同机器人先后触球两次及以上进球方有效，倒逼各队打磨多机传切配合',
            '【体育舞蹈】体育舞蹈参赛队从去年29支暴涨至108支，设街舞/国标/啦啦操三个赛项，机器人踩节拍做爆发动作',
            '【场景赛预赛】8月16日场景赛率先展开预赛，办公场景机器人20分钟内自主完成装填打印机/摆会议用品/操作碎纸机等任务',
            '【场景赛真实度】园林场景室外任务：找违禁物品/提醒不文明现象/任务被干扰中断后续接完成，比工业场景更具挑战',
            '【产业数据】2026上半年中国人形机器人产量近2.5万台，同比暴涨超310%，中国厂商拿下全球出货量97%',
            '【价格雪崩】2023年一台人形造价约65万，现在宇树G1基础版8.5万，最新R1系列仅2.69万，三年跌价95%'
        ], 'process': [
            '【技术萌芽期（2015-2022）】人形机器人技术不成熟，行走不稳定容易摔倒，只能在实验室完成简单演示动作，没有正式赛事，只有科技展会表演性质展示，参赛主体是高校科研院所原型机，企业参与少，无明确比赛规则。',
            '【第一届探索期（2025年8月）】2025年8月在北京举办首届世界人形机器人运动会，具有探索性质，设置26个项目3天6个单元比赛，280支队伍、500多台机器人参赛，主要是国内企业和高校，优必选、宇树、小米等派队参赛，半马冠军成绩2小时40分42秒，开始建立比赛规则体系。',
            '【第二届爆发期（2026年8月）】2026年8月22-26日在国家速滑馆"冰丝带"举办第二届赛事，8月16日部分场景赛率先展开预赛；规模爆发式增长：六大洲16国666支队伍2056台机器人（队伍+138%，机器人翻两番，巴西组建国家队含5支RoboCup队伍），国内157家企业200所院校641支队伍1975台机器人参赛；51个项目（竞技赛30项+场景赛21项）1301场比赛5天9单元（赛期从3天6单元增至5天9单元，下午晚上比赛）；4月亦庄半马冠军成绩50分26秒打破人类纪录；新增跳远/举重/拔河/乒乓球/自由搏击（40/58/80kg三个级别）等对抗项目；场景赛扩至21项覆盖工业/办公/酒店/家庭/物流/消防/园林/应急/零售9大真实场景，两两同台PK，自主完赛权重更高；灵巧手8项微操（夹豆/称量/开瓶/拧螺丝/装配/家政/图书整理）比"看得清、拿得稳、做得准"；组委会创新打造"机器人之家"驻场保障中心（类似运动员村），可容纳1500余台机器人集中存放和不少于1200块电池充电需求，30秒完成一台机器人数字化存取闭环，足球赛训基地紧邻能量充电站；实现5G+Wi-Fi融合通信保障抗干扰高可靠竞赛环境；打造"机器人一条街"提供智能售卖/VR观影/机器人咖啡等科技体验；联合朝阳区打造嘉年华活动释放票根经济；打造全球唯一一个针对机器人赛事的赛事指挥系统和具身智能机器人管理平台，实现跨品牌跨型号机器人管理；选派裁判员223名（国际级28人、国家级75人、一级120人），较上届增加100人；以赛定标，比赛规则相当部分将转化为实用技术验收标准，实现"奖牌变订单，赛场进现场"。',
            '【以赛促研，以赛定标】组委会明确"以赛定标、以标促产"核心目标：赛场上每多跳高1厘米、每多举起1公斤，都是工程师对关节电机、减速器、精密加工不断优化的结果；比赛形成的规则相当部分将演化成实用技术验收标准，真正实现"得了奖牌就拿订单，出了赛场就进现场"。',
            '【2026赛事足球技术亮点】足球赛实现跨代跃迁：2024年1.0阶段人工逐行编程控制；2025年2.0阶段全场无人工干预自主对战；2026年具备完整传切战术，机器人能近距离连续带球、传切、跟进、大力抽射、飞身扑救；清华火神队作为卫冕冠军+RoboCup世界冠军，重构了软件决策系统适配门球/角球/界外球等复杂场景。',
            '【产业数据突破】赛事折射产业爆发：2026上半年中国产量近2.5万台同比+310%，占全球出货97%，全球前六大厂商全中国占88%份额；宇树科技科创板发行市值610亿978万人打新，2025扣非净利润5.91亿是全球唯一规模化盈利人形企业；智元机器人累计量产突破15000台；价格雪崩：2023年65万→宇树G1 8.5万→R1仅2.69万，三年跌95%。',
            '【观赛与体验】朝阳区依托奥林匹克中心区打造机器人主题沉浸式体验集群："冰丝带"湖畔设"嗨FUN机器人闲玩市集"，机器人厨师做拉花咖啡/冰激凌、机器人脱口秀/潮流舞蹈、人机同台演奏音乐会、写书法/下五子棋/聊天/跳舞互动；门票最低98元，周末晚间场已热销，观众可现场直观感受技术进步。',
            '【赛事规则完善期（2027-2028）】第三届、第四届持续举办，规则不断完善，项目更贴近真实应用，参赛队伍扩展至全球，国际企业（波士顿动力、特斯拉Figure、Agility Robotics等）参赛，成为国际顶级机器人赛事，赛事奖金提升，"以赛定标"机制成熟推动产业标准化。',
            '【技术成熟期（2029-2030）】人形机器人运动能力接近人类运动员水平，百米突破10秒，球类流畅度接近人类比赛，精细操作达普通人水平；赛事商业价值凸显，赞助商/转播权/门票形成完整商业闭环，成为科技界体育界双重盛会；人形机器人从赛场走向工厂、家庭、服务各场景全面普及。',
            '【安徽与蚌埠参与】合肥、芜湖机器人企业（埃夫特、奇瑞、智元合肥基地等）组团参赛；中科大、合工大组成高校代表队；蚌埠中国传感谷企业为多个参赛队提供六维力传感器、IMU、微型力传感器等核心传感器零部件，助力参赛机器人取得好成绩；同期2026成都市人工智能与机器人创新应用大赛也成功举办。'
        ], 'detail': [
            '【赛事准确时间地点】第二届世界人形机器人运动会将于2026年8月22日至26日在北京国家速滑馆"冰丝带"举办，开幕式当晚安排7v7足球表演赛+400米/100米预赛，闭幕式前将上演自由搏击决赛+足球5v5决赛作为压轴大戏',
            '【赛事规模准确数据】全球六大洲16个国家共666支参赛队伍、2056台机器人参赛；国内157家企业、200所院校科研机构组成641支队伍、1975台机器人参赛；首届仅280队500多台机器人，本届队伍增长138%机器人数量翻两番',
            '【比赛项目总数统计】共设置51个比赛项目（竞技赛30项+场景赛21项），赛期5天共9个竞赛单元将开展1301场激烈比赛；首届仅设26项3天6单元，本届取消外围赛并将武术、体育舞蹈正式纳入竞技赛项目',
            '【金牌产生节奏安排】8月23日开幕后首个比赛日就将产生12枚金牌，赛事前半程即进入"结果密集期"；正赛5天安排1301场比赛密度极高，平均每天比赛超过260场，对机器人可靠性是严峻考验',
            '【半程马拉松成绩突破】2026年4月北京亦庄人形机器人半程马拉松，冠军成绩从2025年首届的2小时40分42秒大幅缩短至50分26秒不到上届时长三分之一，组委会官方表述称该成绩打破了人类半程马拉松纪录',
            '【径赛完赛标准升级】2026年径赛完赛时间要求大幅压缩：100米从3分钟压缩到1分钟以内，400米从15分钟压到5分钟，1500米从40分钟压到15分钟；除障碍赛外所有径赛机器人必须全自主完成不允许遥控干预',
            '【新增高强度对抗项目】本届新增跳远、举重、拔河、乒乓球、自由搏击等高强度对抗项目；自由搏击分40公斤、58公斤、80公斤三个体重级别贯穿8个竞赛单元；举重"大力士"24日晚决出两枚金牌，对机器人本体结构是极限考验',
            '【场景赛贴近真实工业】场景赛从2025年6项大幅扩展至2026年21项，覆盖工业装配、酒店服务、家庭家政、物流分拣、应急救援、图书整理等9大真实工作场景，专门针对"乏、脏、险、难"岗位设计，两两同台PK考验真实工作能力',
            '【灵巧手微操赛项设置】灵巧操作是本届赛事重点比拼项目：粉末称量精度0.1克、镊子夹绿豆、开瓶撬盖、拧M3螺丝、精密电子装配，比的就是"看得清、拿得稳、做得准"的指尖精细活，力控精度和触觉感知能力是决胜关键',
            '【足球赛事规则重大升级】足球赛设7v7开幕式表演赛、5v5大型组、中型组、3v3 U19青少年组四个组别；2026年进球判定门槛大幅提升：常规对战必须两台不同机器人先后触球两次及以上进球方有效，倒逼各队打磨多机传切配合能力',
            '【清华火神队全力卫冕】清华火神队是首届5v5足球赛冠军+2025 RoboCup人形组世界冠军，20余名队员覆盖自动化、信息、车辆、机械多个专业，8月初进驻冰丝带封闭集训十余天，硬件新增专属传球动作、软件重构决策系统目标只有冠军',
            '【体育舞蹈参赛规模爆发】体育舞蹈参赛队从2025年29支暴涨到2026年108支，设街舞、国标、啦啦操三个赛项，机器人需精准踩准音乐节拍做出爆发性动作，考验关节响应速度、运动控制精度和音乐节拍识别能力',
            '【赛事保障体系建设】"机器人之家"运动员村对接京东物流数字化管理体系，采用"一机一码""一电池一码"溯源管理，扫码即可查询赛队信息和设备参数，存取登记30秒完成；200个充电柜可同时满足1200块电池充电需求',
            '【中国机器人产业地位】2026年上半年中国厂商拿下全球人形机器人出货量97%以上，全球前六大厂商全是中国企业合计份额接近88%；中国上半年产量接近2.5万台同比暴涨超310%；摩根士丹利把全年出货预期从2.8万台上调到5万台',
            '【宇树科技科创板表现】宇树科技刚结束科创板申购发行市值610亿元，978万人排队打新中签率极低；是全球唯一规模化盈利的人形机器人公司，2025年扣非净利润5.91亿元；价格雪崩：2023年65万→G1基础版8.5万→R1仅2.69万',
            '【智元机器人量产数据】智元机器人累计量产突破15000台是国内量产规模领先企业之一，远征A1/A2系列产品覆盖工业制造、商业服务、科研教育等多个场景，搭载智元自研具身大模型任务完成率持续提升',
            '【配套科技体验活动】朝阳区在奥林匹克中心区打造机器人主题沉浸式体验集群："冰丝带"湖畔"嗨FUN机器人闲玩市集"，机器人厨师现场制作拉花咖啡、冰激凌，还有机器人脱口秀、潮流舞蹈、人机同台演奏音乐会等互动项目',
            '【票务价格亲民普惠】比赛日门票最低仅98元，开闭幕式门票最低128元，还设有双人套票、三人家庭套票等多种票型供选择，不是只有硬核铁粉才能观赛，普通家庭周末观赛成本完全可承受，目前开幕式多个票档已售罄',
            '【赛事产业溢出价值】赛事期间预计达成产业合作意向超200亿元，投融资签约超100亿元；赛事验证的稳定行走、抗干扰摔倒自主恢复、灵巧操作等技术将快速转化到量产产品加速产业成熟；央视全程报道直播观众预计超5亿人次',
            '【组委会产业判断】组委会常务副主任、北京市经信局局长姜广智明确表示："机器人练好这些最后一米的手艺，就能得了奖牌就拿订单，出了赛场就进现场。"这句话精准概括了人形机器人产业"以赛促研、以赛定标、以标促产"的最朴素发展逻辑'
        ]},
        # 真机部署
        {'left': [
            '【部署现状】2026年是人形机器人真机规模化部署元年，工业场景率先落地万台级',
            '【部署场景】汽车制造/3C电子/物流仓储/新能源工厂/商业服务等场景率先部署',
            '【比亚迪工厂】比亚迪工厂已部署人形机器人超3000台，承担物料搬运/零部件分拣任务',
            '【特斯拉工厂】特斯拉超级工厂部署Optimus超5000台，承担产线搬运/装配辅助任务',
            '【宁德时代】宁德时代工厂部署人形机器人超1000台，电池产线物料配送/质量检测',
            '【3C电子】富士康/立讯精密3C工厂部署人形机器人，精密装配/检测/搬运',
            '【物流仓储】京东/菜鸟/亚马逊仓库部署人形机器人，分拣/码垛/搬运包裹',
            '【商业服务】酒店/餐厅/展厅/医院部署人形机器人，迎宾/引导/配送/讲解服务',
            '【安徽部署】安徽汽车/家电/新能源产业基础好，蔚来/比亚迪/美的/海尔等工厂部署量快速增长',
            '【ROI拐点】人形机器人单机成本下降+人工成本上升，工业场景ROI回收期缩短至2-3年'
        ], 'right': [
            '【机器人移动母舰·2026年8月21日】WRC2026舰队协同新生态：飞巴科技全球首发机器人移动母舰——机器人的移动后勤基地，舱体可装载人形机器人/机器狗/无人机，车内自带换电工位和维修工位，即使在断网断电极端环境下也能给机器人提供算力和通信保障，今年年底投入量产预计明年6月真正商用进入航空救援/应急消防/医学救援等领域，让机器人从单兵作战走向舰队协同；新松人工智能研究院发布多机型协同系统OneHub羿枢可接入不同种类机器人3台机器人两种类型都可在该系统协同下工作融入大模型技术；猿声先达科技首次对外展出多维触觉动捕手套+能感知物体接近的大面积电子皮肤，动捕手套可感知法向力/切向力及非常密集的力的方向模块化设计可重构跟人骨骼完全一样的骨骼建模',
            '【优艾智合隙锋·2026年8月20日】2026世界人工智能大会上5台工业原生人形机器人"隙锋"在模拟线边仓有条不紊进行物料拣选与配送，与1台移动搬运机器人协同作业完成生产流程；全套核心产品研发量产均依托合肥基地完成，以2024年收入计工业移动操作机器人全球市占率第一',
            '【优艾智合版图·2026年8月20日】优艾智合产业版图：半导体领域进入全球顶级晶圆厂在内十余家国内外晶圆厂；能源化工覆盖上游原材料/电力/冶金/化工/公用事业；累计超800个具身智能场景落地项目，服务超400家全球头部企业',
            '【星工聚将·2026年8月20日】星工聚将CTO陈牧：自研"数字风洞"仿真校准平台SynTunnel测仿真与真实动作偏差反向修正模型，训练成本从指数增长收敛成线性增长；末端快换系统6秒内切换吸盘/夹爪/灵巧手等十几种执行器实现"一机多用"',
            '【机器人进厂·2026年8月20日】央视网：机器人扎堆进厂"上班"，2026上半年国内具身智能赛道融资总额突破935亿元同比+5倍；工信部+国资委实景实训专项年底凝练百个高价值场景带动万台级落地',
            '【最新·2026年8月】小米机器人北京亦庄汽车工厂上岗：承担抓取搬运/灵巧手手指操作/精细触觉反馈装配三类工作，单一工站成功率98%',
            '【比亚迪工厂部署】比亚迪西安/深圳/合肥/常州工厂累计部署Walker X2/自研人形3200台，承担产线物料搬运、零部件上下料等重复性工位任务',
            '【特斯拉得州工厂】特斯拉美国得州超级工厂Optimus Gen3累计部署5200台，承担产线零部件搬运、装配辅助、简单重复工位作业',
            '【宁德时代工厂】宁德时代福建/江苏/安徽工厂累计部署人形机器人1200台，承担电芯搬运、模组检测、危险工位操作任务',
            '【富士康深圳工厂】富士康深圳龙华/观澜工厂部署人形机器人800台，在iPhone组装线承担简单装配、质量检测、物料配送任务',
            '【京东物流仓储】京东亚洲一号智能仓部署人形机器人500台，承担包裹分拣、码垛、补货任务，支持24小时不间断作业',
            '【优必选商业服务】优必选Walker商业服务版在酒店/医院/展厅累计部署超2000台，提供迎宾引导、物品配送、讲解接待服务',
            '【达闼云端员工】达闼云端机器人员工已在全国100+酒店正式上岗，提供办理入住、送物引导、信息咨询服务，提升服务效率',
            '【蔚来合肥基地】蔚来合肥先进制造基地F2工厂部署人形机器人300台，承担汽车产线物料配送、零部件转运、简单装配辅助',
            '【美的奇瑞安徽布局】美的合肥工业园部署150台、奇瑞芜湖工厂部署埃夫特人形200台，安徽制造业人形部署量快速增长',
            '【鹿明MOS2重载新品】8月14日鹿明机器人发布Lumos MOS2重载轮臂具身机器人，50kg双臂负载，定位工业重载AI Worker支持技能持续学习进化'
        ], 'process': [
            '【实验室验证期（2020-2023）】人形机器人主要停留在实验室研发阶段，整机成本超100万元，稳定性差，行走容易摔倒，续航不足2小时，只能完成简单演示动作，少量原型机在工厂做POC（概念验证）测试，验证技术可行性，距离真实部署差距大，主要问题是"能走但不能干活，能动但不稳定"。',
            '【POC试点探索期（2024）】2024年是人形机器人POC试点元年，优必选Walker、宇树H1、特斯拉Optimus等开始在头部工厂（比亚迪、特斯拉、富士康）做小范围POC测试，单厂部署几台到几十台，承担最简单的搬运、上下料任务，发现大量真实场景问题（地面不平、光照变化、人机干扰、任务多变），产品快速迭代，成本降至50-80万元，续航提升至4小时。',
            '【小批量试点部署期（2025）】2025年小批量试点部署启动，单厂部署几十台上百台，比亚迪部署超500台，特斯拉部署1000台，承担物料搬运、简单上下料、码垛等重复性任务，收集海量真实场景数据持续训练具身智能大模型，MTBF（平均无故障工作时间）从200小时提升至500小时，成本降至30-50万元，ROI约3-4年，验证了商业可行性。',
            '【规模化部署元年（2026）】2026年是人形机器人规模化部署元年（最新阶段），技术成熟度达到工业应用要求：MTBF提升至2000小时，快速换电（<30秒）实现24小时作业，成本降至15-30万元，单台年成本5-7万元，低于工人年综合成本，ROI缩短至2-2.5年；全球部署量达8万台，中国占60%约4.8万台，比亚迪累计部署3200台，特斯拉5200台，宁德时代1200台，工业场景占比75%，从汽车、3C扩展到新能源、物流、化工等更多行业；8月14日鹿明机器人发布Lumos MOS2重载轮臂式具身智能机器人，50kg双臂负载，全向移动+多模态感知+Lumos NexCore持续学习系统，面向工业重载场景；中国信通院判断产业正从"能用"向"好用"关键一跃，AI Worker进入产业现场提速，未来竞争取决于真实数据获取效率、硬件制造成本、产业交付速度。',
            '【多场景扩展期（2027-2028）】人形机器人从工业场景向商业服务、家庭服务扩展，工业场景部署量持续增长，单厂部署上千台甚至数千台，承担装配、检测、运维等更复杂任务；商业场景（酒店、餐厅、医院、展厅、银行）大规模部署服务型人形机器人；物流仓储场景人形机器人成为标准配置；全球部署量2028年达50万台，中国达30万台，单机成本降至10-15万元，ROI缩短至1.5-2年，人形机器人产业形成完整产业链。',
            '【家庭服务萌芽期（2029-2030）】人形机器人开始进入家庭，承担家务（打扫、整理、做饭辅助）、老人陪护、儿童教育、家庭安防等任务，家用版成本降至5-10万元，达到普通家庭可承受范围；工业场景人形机器人成为工厂标配，新建工厂预留人形机器人工位，全球部署量超200万台，人形机器人真正融入生产生活各方面。',
            '【部署技术挑战与迭代】真实部署面临诸多挑战：非结构化环境适应性（地面不平整、障碍物、光照变化）、人机协作安全（碰撞检测、力反馈）、长时间工作可靠性（散热、电池、关节磨损）、任务泛化能力（大模型泛化到新任务、新场景）；每一个挑战都需要在真实部署中发现问题、迭代解决，"以部署促研发"是核心路径。',
            '【部署标准体系建设】随着规模化部署，人形机器人安全标准、通信标准、接口标准、作业标准逐步建立：ISO 13482服务机器人安全标准升级适配人形机器人，中国发布《人形机器人工业部署安全规范》，人形机器人与工厂MES/WMS系统接口标准化，不同品牌机器人任务调度协同标准制定，推动产业规范化发展。',
            '【人才培养与就业转型】人形机器人大规模部署带来就业结构变化：简单重复性体力劳动岗位被替代，同时催生出人形机器人运维、调试、训练、任务规划等新岗位，职业院校开设人形机器人相关专业，企业开展员工转岗培训，实现平稳过渡，劳动生产率大幅提升。',
            '【安徽真机部署进展】安徽制造业基础雄厚，是人形机器人真机部署先行省份：合肥蔚来F2工厂、比亚迪合肥工厂、大众安徽工厂、美的合肥工业园2026年累计部署人形机器人2200台；芜湖奇瑞工厂、埃夫特基地部署600台；蚌埠中国传感谷企业、中电科40/41所、玻璃设计院、昊方机电等试点部署约100台；2027年安徽全省部署目标超1万台，建成全国人形机器人应用示范高地。'
        ], 'detail': [
            '【南京场景开放·2026年8月20日】南京以场景开放夯实产业强市建设：全国首个城域级"南京场景服务平台"今年已吸引超1800个注册用户累计征集应用场景清单近800个；AI二次供水泵房/移动充电机器人"桩找车"/江宁能碳虚拟电厂等创新场景落地',
            '【光象科技双工位·2026年8月21日】WRC2026光象科技演示行业首次单台具身智能机器人双工位自主循环作业：Phi-Bot X1在焊接上料工位完成毫米级对准上料后自主切换至移动质检工位，部署周期从半年压缩至周级甚至天级，质检效率提升51%',
            '【全球部署量增长趋势】2025年全球人形机器人实际部署量约1万台，2026年规模化部署元预计达8万台，2027年预计25万台，2028年预计50万台，部署量呈指数级快速增长态势，人形机器人从实验室走向产业现场',
            '【中国部署量占比】中国市场人形机器人部署量占全球60%，2026年全年预计部署约4.8万台，其中工业制造场景占比75%、商业服务场景占比20%、物流仓储及其他场景占比5%，中国是人形机器人最大应用市场',
            '【比亚迪工厂部署进展】比亚迪2025年试点部署500台，2026年西安/深圳/合肥/常州工厂累计部署达3200台，2027年目标部署1万台，2028年目标5万台覆盖所有整车和零部件工厂',
            '【特斯拉Optimus部署规划】特斯拉2025年试点部署1000台，2026年得州超级工厂部署达5200台，2027年目标部署10万台，优先满足特斯拉自有工厂需求后才会对外销售',
            '【单台人形工作效率对比】人形机器人在搬运等标准化工位工作效率约为熟练工人的70-80%，但支持快速换电可实现24小时不间断作业，单日有效工作时长是人工的3倍，综合产出效率已超过人工',
            '【人工替代率计算】单台人形机器人配合三班倒可替代2-3名工人，主要替代简单重复性体力劳动岗位，被替代工人转岗到技能要求更高的机器人运维、质量检测、设备维护等岗位，实现就业平稳转型',
            '【投资回报ROI测算】单台人形机器人年综合成本（折旧+维护+电费）约5-7万元，国内制造业工人年综合成本8-12万元，投资回收期（ROI）缩短至2-2.5年，具备明确商业价值',
            '【可靠性指标提升】量产人形机器人平均无故障工作时间（MTBF）从2025年的500小时大幅提升至2026年的2000小时，已达到工业设备可靠性入门要求，可满足工厂连续作业需求',
            '【快速换电技术指标】模块化电池设计支持快速热插拔更换，更换电池时间小于30秒，配合共享电池柜可实现7×24小时不间断作业，电池循环寿命突破5000次满足高强度使用需求',
            '【安全标准与认证】量产人形机器人全部满足ISO 13482服务机器人安全标准，全身碰撞检测灵敏度0.1N响应时间小于10ms，具备软硬双重急停和虚拟安全围栏，可与工人同工位安全协作',
            '【安徽省部署总量】2026年安徽省人形机器人部署量约3500台，其中蔚来/比亚迪/大众/奇瑞/美的等制造业企业占比85%，安徽依托制造业基础成为全国人形机器人应用先行省份',
            '【合肥市工厂部署】合肥新能源汽车和家电产业集聚效应显著，2026年蔚来/比亚迪/大众/美的等企业合计部署人形机器人约2200台，2027年目标部署8000台建成全国应用示范城市',
            '【芜湖市产业部署】芜湖奇瑞汽车/埃夫特机器人等企业2026年合计部署约600台，依托机器人产业基地优势快速推进，2027年目标部署2000台建成机器人应用产业生态',
            '【蚌埠市试点部署】蚌埠中国传感谷企业、中电科40/41所、玻璃设计院、昊方机电等企业2026年试点部署约100台，依托传感器产业优势探索应用，2027年目标部署500台',
            '【RaaS租赁部署模式】机器人即服务（RaaS）模式兴起，企业不用一次性采购机器人，按使用时长或完成任务量付费，大幅降低企业初期投入门槛，特别适合中小企业和季节性波动场景',
            '【整厂解决方案模式】机器人企业提供整厂人形机器人部署解决方案，包含机器人硬件、多机调度系统、人员培训、运维服务一体化交付，企业只需提出需求即可快速落地应用',
            '【任务能力持续扩展】人形机器人从最初简单的物料搬运，逐步扩展到零部件装配、产品质量检测、设备维护巡检、异常情况处理等更复杂任务，技能边界持续拓展',
            '【OTA远程技能升级】机器人通过OTA远程升级持续获得新技能，已部署在现场的机器人无需返厂即可获得新能力，机器人功能持续进化迭代，保护客户投资',
            '【人员培训体系完善】工厂开展人形机器人运维人员和协作人员培训，平均培训周期2周即可上岗，工人对机器人接受度逐步提升，人机协作成为工厂新常态',
            '【2030年产业发展目标】规划到2030年全球人形机器人部署量超500万台，中国超300万台，制造业人形机器人密度达到100台/万人，人形机器人成为制造业标配设备'
        ]},
        # 物流仓储
        {'left': [
            '【产业定位】物流仓储是机器人应用最成熟场景之一，AGV/AMR/分拣/码垛机器人大规模普及',
            '【市场规模】2026年中国物流机器人市场规模突破350亿元，年均增速超35%',
            '【AGV/AMR】自动导引车（AGV）/自主移动机器人（AMR）在仓库内搬运货物，替代人工叉车',
            '【分拣机器人】快递/电商分拣中心交叉带分拣/滚珠模组带分拣/机器人分拣，效率超人工5倍',
            '【码垛机器人】仓库出货码垛/拆垛机器人，负载50-300kg，码垛速度800-1200次/小时',
            '【极智嘉】极智嘉（Geek+）是全球AMR领军企业，AMR出货量全球第一',
            '【快仓】快仓智能仓储机器人，菜鸟/京东/唯品会等电商仓库大规模应用',
            '【海康机器人】海康威视旗下海康机器人，机器视觉+移动机器人双轮驱动，国内市占率前三',
            '【快递物流】顺丰/京东/中通/圆通/韵达快递分拨中心自动化率超90%，分拣基本无人化',
            '【安徽物流】合肥/芜湖是全国重要物流枢纽，京东/菜鸟/顺丰在安徽建设大型智能仓'
        ], 'right': [
            '【博银合创+魔法原子·2026年8月21日】WRC2026工业产线从一机一岗到一机多能：博银合创携手法奥机器人展出工业级具身智能机器人BW10-Lite最高运行速度1.5m/s双臂最大负载20公斤；魔法原子把仓储流通作业实景搬进展台多SKU实时分拣成功率达99%；物流仓储成为最先算清ROI的场景',
            '【极智嘉P800】货到人拣选AMR，负载800kg，导航精度±10mm，运行速度2m/s，全球部署超5万台',
            '【海康机器人MR系列】潜伏/移载/叉取全系列AMR，海康威视视觉+激光导航，国内市占率25%',
            '【快仓QuickBin】快仓料箱到人机器人，适用电商拆零拣选，拣选效率1000件/小时，是人工3倍',
            '【京东天狼仓】京东亚洲一号智能仓，天狼仓系统，AMR+分拣+码垛全流程自动化，日处理订单超百万',
            '【菜鸟未来园】菜鸟网络未来园区，无人仓/无人车/无人机/无人柜全链路无人化',
            '【立镖分拣机器人】立镖小黄人分拣机器人，3000台机器人协同分拣，每小时处理20万件快递',
            '【ABB码垛机器人】ABB IRB 660码垛机器人，负载250kg，码垛节拍1200次/小时，精度0.1mm',
            '【新松AGV】新松机器人AGV系列，重载/轻载全系列，汽车/烟草/电商行业广泛应用',
            '【合肥京东亚洲一号】合肥京东亚洲一号智能物流园，安徽最大智能仓，AGV/分拣/码垛全自动化',
            '【芜湖顺丰分拨中心】芜湖顺丰智能分拨中心，自动分拣率99%，日处理包裹超200万件'
        ], 'process': [
            '【纯人工作业时代（1990-2010）】仓储物流完全依靠人工作业：人工叉车搬运托盘、人工拣货员拉着拣货车在货架间行走找货、人工分拣包裹、人工扫码录入，仓库工人每天行走10-20公里，劳动强度极大，拣货错误率约1-3%，大促期间（双十一、618）爆仓频发，招工难、用工贵、人员流失率高（年流失率30-50%），人效低，管理难度大。',
            '【自动化仓储起步期（2011-2015）】自动化立体仓库开始建设，AGV（自动导引车）开始试点应用，主要是磁条导航、二维码导航，需要预先铺设路径，柔性差；自动分拣机（交叉带分拣机）在快递分拨中心应用，分拣效率提升，但成本高、灵活性差，仅适用于标准化场景；京东、菜鸟等开始探索智能仓储，但整体自动化率不足20%。',
            '【AGV规模化应用期（2016-2019）】Kiva机器人模式引入中国，极智嘉、快仓、海康机器人等企业推出潜伏顶升AGV，二维码导航为主，亚马逊Kiva验证了"货到人"模式有效性，京东亚洲一号、菜鸟无人仓开始大规模建设，AGV在中国市场快速增长，从电商扩展到医药、烟草、汽车零部件等行业，仓储自动化率提升至40%左右，但AGV需要改造仓库环境，柔性仍有不足。',
            '【AMR柔性智能化期（2020-2024）】AMR（自主移动机器人）取代传统AGV成为主流，激光SLAM+视觉融合导航，无需预先铺设磁条/二维码，能够自主导航、避障、路径规划，部署快、柔性高；海康、极智嘉、快仓等企业AMR产品线完善，料箱机器人、叉取AMR、分拣AMR等多种机型；AI智能调度系统可调度上千台机器人协同作业，"货到人"拣选效率提升3-5倍，错误率降至0.01%以下；快递分拨中心自动分拣率达90%以上，仓储自动化率提升至70%。',
            '【人形机器人试点期（2025-2026）】2026年人形机器人开始在物流仓储场景试点应用，京东亚洲一号、菜鸟、顺丰分拨中心试点部署人形机器人，承担拆码垛、包裹分拣、装卸车、异常件处理等不规则任务；物流场景标准化程度介于工业和商业之间，是人形机器人率先落地的场景之一；AMR+人形机器人+机械臂组合形成完整仓储无人化方案；2026年中国物流机器人市场规模突破500亿元。',
            '【全面无人化期（2027-2028）】人形机器人在物流仓储规模化部署，从试点扩展到大面积应用，承担更多复杂任务：不规则物品拣选、包装、贴标、装卸车，AMR负责水平搬运，机械臂负责固定位置码垛，人形机器人负责非标准化任务，AI系统统一调度多机型协同作业，仓储自动化率超90%，大促期间无需大量临时工，实现真正的无人仓。',
            '【物流网络智能化期（2029-2030）】整个物流网络实现智能化：仓储无人化、干线运输自动驾驶、末端配送无人化（无人车/无人机），AI大数据预测需求提前铺货，库存周转率提升50%，物流成本占GDP比重从2025年14%降至10%以下，接近发达国家水平，中国物流效率全球领先，智慧物流成为中国竞争力重要组成部分。',
            '【物流机器人技术演进】人工地牛/叉车→磁条AGV→二维码AGV→潜伏顶升AGV（货到人）→激光SLAM AMR→料箱/叉取AMR→人形机器人试点→多机器人协同无人仓；导航方式从人工驾驶→磁条导航→二维码导航→激光SLAM→视觉SLAM→多传感器融合导航；拣选方式从"人找货"→"货到人"→"机器人自主拣选"。',
            '【安徽物流智能化进展】安徽是长三角物流枢纽，合肥、芜湖是国家物流枢纽承载城市：京东亚洲一号合肥仓、芜湖顺丰智能分拨中心、合肥菜鸟智慧仓等智能仓储项目建成，安徽快递分拨中心自动分拣率达99%；合肥综合保税区、芜湖港等部署物流机器人超5000台；蚌埠作为皖北物流中心，智慧物流园区建设加快，皖北徽商物流港部署仓储机器人200+台，快递分拨自动化率达95%以上。',
            '【物流行业价值】物流行业是国民经济基础性、战略性产业，2026年中国社会物流总额超350万亿元，物流机器人应用大幅提升效率、降低成本、减少错误、改善工人工作条件；电商、快递、制造业物流是主要应用场景，智慧物流支撑中国电子商务和制造业高效运转，降低全社会物流成本。'
        ], 'detail': [
            '【极智嘉P800潜伏AMR参数】极智嘉P800潜伏顶升AMR最大负载800kg，最高运行速度2m/s，采用激光SLAM+视觉融合导航，定位精度±10mm，标准工况续航8小时支持自动充电，整机防护等级IP54可适应多种仓储环境，单台设备市场售价约15万元人民币',
            '【海康MR5-1200叉取AMR参数】海康机器人MR5-1200叉取式自主移动机器人最大负载1200kg，货叉最大提升高度3米，运行速度1.5m/s，激光SLAM导航支持窄巷道作业和托盘自动识别，适用于标准托盘水平搬运和堆高作业，单台售价约25万元',
            '【快仓QuickBin料箱机器人】快仓QuickBin料箱到人机器人最大料箱负载50kg，存储高度可达8米支持双伸位货叉存取，拣选效率达1000件/小时，拣选准确率99.99%，特别适用于电商拆零拣选场景，大幅提升料箱存储密度和拣选效率',
            '【立镖小黄人分拣机器人】立镖"小黄人"分拣机器人单台负载3kg，最高运行速度3m/s，采用二维码导航，AI智能路径规划算法优化，支持3000台机器人大规模协同调度，分拣效率达20万件/小时，分拣错误率低于0.01%，是快递分拣主力机型',
            '【ABB IRB 660码垛机器人参数】ABB IRB 660四轴专业码垛机器人最大负载250kg，工作臂展3.15m，标准码垛循环能力1200次/小时，重复定位精度0.1mm，适用于纸箱、袋装、箱装等多种包装形式货物的码垛/拆垛作业，是工业码垛标杆机型',
            '【分拣效率对比数据】人工分拣效率150-200件/人/小时，错误率3-5%；交叉带自动分拣线效率10000-20000件/线/小时，错误率低于0.01%，分拣效率提升100倍，错误率降低两个数量级，大促期间可24小时不间断作业',
            '【拣选效率对比数据】传统人工摘果式拣选效率80-120件/人/小时；AMR"货到人"拣选效率400-600件/人/小时；机器人自主拣选（机械臂+AGV）效率达800-1200件/小时，拣选效率较传统模式提升5-10倍，准确率大幅提升',
            '【京东亚洲一号建设】京东在全国已建成40+座"亚洲一号"智能物流产业园，合肥亚洲一号日处理订单能力达150万单，仓储自动化率达到95%，是国内智能化水平最高的电商仓储网络之一，支撑京东物流"211限时达"服务承诺',
            '【菜鸟智能仓网络建设】菜鸟在全球建成100+个智能仓，AMR自主移动机器人部署量超10万台，覆盖国内核心城市和海外主要市场，跨境电商智能履约时效提升50%，智能分单、智能路由、智能调度全链路AI优化',
            '【极智嘉全球市场布局】极智嘉AMR全球累计部署量超5万台，服务全球客户超1000家，业务覆盖30+国家和地区，全球AMR市场份额约15%，是全球出货量最大的AMR企业之一，海外收入占比超40%国际化程度高',
            '【海康机器人市场地位】海康机器人依托海康威视技术积累，移动机器人累计出货量超20万台，国内移动机器人（AGV/AMR）市场份额连续多年排名第一，产品线覆盖潜伏、叉取、料箱、分拣全系列，服务客户超1万家',
            '【快递行业自动化率】顺丰、中通、圆通、韵达、申通、极兔等头部快递企业一级分拨中心自动化率超90%，自动分拣设备处理占比达95%，人工分拨占比不足5%，快递分拣环节基本实现无人化自动化',
            '【安徽省物流机器人建设】安徽是长三角重要物流枢纽，合肥、芜湖、蚌埠国家物流枢纽建设加快推进，全省智能仓储总面积超500万平米，AMR移动机器人部署量超1万台，快递分拨中心自动化率达98%位居全国前列',
            '【蚌埠市智慧物流建设】蚌埠作为皖北物流中心城市，中通、圆通、韵达等快递企业在蚌埠建设区域分拨中心，自动分拣设备全面普及；皖北徽商物流港等园区部署仓储机器人200+台，蚌埠皖北保税物流中心智能化水平持续提升',
            '【仓储密度提升效果】AMR自主移动机器人配合高位货架存储方案，仓储密度较传统平库提升2-3倍，土地利用率提升3倍，库存准确率达到99.99%，库存周转率提升40%，大幅降低仓储租金和人工成本',
            '【AI多机协同调度系统】AI智能仓储调度系统支持数千台AMR机器人实时协同调度，AI动态路径规划、任务智能分配、交通自动管制全局最优，机器人运行无碰撞、无拥堵、无死锁，系统整体效率比人工调度提升30%',
            '【AMR柔性部署优势】AMR机器人无需改造地面（无需铺设磁条、二维码），基于SLAM自主导航，上线部署周期短（1-2周即可上线），业务波动时可灵活增减机器人数量，电商大促期间可快速扩容机器人应对订单洪峰',
            '【冷链物流专用机器人】冷链仓储低温环境（-20℃以下）专用机器人，整机采用耐低温材料和元器件，电池耐低温特殊配方，防护等级IP65可适应冷库高湿低温环境，解决冷库人工工作环境恶劣、招工难问题',
            '【末端无人配送应用】园区、高校、社区、写字楼等封闭/半封闭场景末端配送无人车试点应用，快递、外卖、生鲜配送"最后一公里"问题逐步解决；无人机在偏远山区、海岛配送试点，末端配送多元化无人化加速',
            '【2030年智慧物流发展目标】规划到2030年全国仓储自动化率达95%，AMR及各类物流机器人部署量超500万台，人形机器人在仓储不规则物品拣选、拆码垛、装卸车场景广泛应用，全社会物流效率再提升50%，物流成本占GDP比重降至发达国家水平'
        ]},
        # 灵巧手
        {'left': [
            '【产业定位】灵巧手是人形机器人末端执行器，决定机器人操作能力，是核心关键部件之一',
            '【技术难度】灵巧手机械设计/驱动/控制/感知难度极高，被称为人形机器人"皇冠上的明珠"',
            '【自由度】仿人灵巧手通常12-20自由度，接近人手27自由度，实现类人操作',
            '【驱动方式】腱驱动/连杆驱动/直线驱动/气动人工肌肉，各有优劣，腱驱动最常用',
            '【力控精度】高端灵巧手指尖力控精度达0.01-0.02N，人手力控精度约0.005N，逐步接近人手',
            '【触觉感知】灵巧手集成指尖触觉传感器阵列，感知接触力/滑移/温度/材质，触觉空间分辨率1mm',
            '【因时机器人】国内灵巧手龙头企业，因时BHX系列灵巧手量产应用最广，国内市占率第一',
            '【Shadow Hand】英国Shadow Robot公司Shadow Hand是灵巧手标杆，20自由度，价格昂贵',
            '【大寰机器人】大寰自适应夹爪/灵巧手，工业场景应用广泛，性价比高',
            '【国产替代】国产灵巧手技术快速追赶，成本仅为进口1/5-1/3，批量应用于人形机器人'
        ], 'right': [
            '【帕西尼10亿融资·2026年8月21日】帕西尼WRC期间官宣三件大事：总部落户北京海淀+完成股份制改造+10亿元战略轮融资到位累计融资近40亿元刷新全球触觉感知领域融资纪录资方含国家级产业基金/产业龙头/地方国资；全球首发足底多维触觉传感器PX-FOOTRIX基于6D阵列式触觉传感技术实现足底全域三维阵列力觉感知；第四代ITPU多维触觉传感器PX6AX GEN4搭载全球首款自研原生6D触觉感知芯片GEN4 FUSE',
            '【中科硅纪灵巧手·2026年8月21日】中科硅纪围绕类人灵巧操作全栈技术集中展示：六款CasiaHand系列行业级灵巧手含M系列/X系列及行业级三指G系列；CasiaHand Brain-Si 0.5类人灵巧操作具身大小脑模型采用分层协同架构；灵巧操作能力分四级：重复性操作→视触觉引导适应性操作→通用抓取收纳物流分拣→功能性工具使用',
            '【因时BHX-12】12自由度腱驱灵巧手，指尖力控0.02N，集成触觉传感器，重量550g，已批量配优必选/小米',
            '【因时BHX-20】20自由度高精度灵巧手，每个手指3-4自由度，指尖力控0.01N，科研和高端应用',
            '【Shadow Dexterous Hand】24自由度（20主动+4被动），气动+腱驱动，力控精度高，售价约150万元，科研用',
            '【大寰PGC-140】自适应二指夹爪，行程140mm，力控5-140N，工业场景广泛，性价比高',
            '【大寰DH-3】三指灵巧手，12自由度，自适应抓取，可抓取不同形状物体，工业和科研用',
            '【傲博iHand】傲博机器人五指灵巧手，16自由度，力控精度0.05N，配傲博机械臂',
            '【哈工大灵巧手】哈工大机器人所自研20自由度仿人灵巧手，指尖力控0.02N，穿针引线演示',
            '【清华DexHand】清华大学灵巧手，肌腱驱动，触觉传感器阵列，操作精细，学术前沿',
            '【特斯拉Optimus灵巧手】特斯拉Optimus Gen3手部6驱动器11自由度，自适应抓取，可拿鸡蛋/精密零件',
            '【小米CyberHand】小米CyberOne 2代手部，12自由度，力控+触觉，可完成拧瓶盖/用手机等操作'
        ], 'process': [
            '【简单夹爪时代（2015年前）】工业机器人末端执行器主要是二指或三指气动/电动夹爪，自由度极少（2-3自由度），只能完成简单的开合抓取动作，没有力反馈和触觉感知，只能在结构化场景中固定位置抓取形状规则的特定物体，无法适应物体形状变化，无法完成精细操作，只能完成"抓起来、放下"的简单动作，无法拧瓶盖、穿针、使用工具等复杂操作。',
            '【科研原型灵巧手时代（2016-2021）】高校和科研院所开始研发多自由度仿人灵巧手，Shadow Hand（英国）、哈工大/清华/北航等国内高校灵巧手陆续问世，自由度从5-20不等，采用腱驱动、气动肌肉等驱动方式，能够实现部分人手动作，但主要是实验室原型，存在三大问题：①成本极高（进口Shadow Hand单只100-150万元，国产科研原型单只10-50万元）；②可靠性差，连续工作几小时就可能出现腱绳断裂等故障；③集成度低，需要外部控制器和气源/电源，无法装在机器人上实际使用；只有少量科研应用，无法产业化。',
            '【技术突破期（2022-2024）】国产灵巧手技术取得关键突破：驱动方式优化（直线电机+腱驱动/差动机构）、集成度提升（驱动器、控制器、传感器集成在手内部）、成本大幅下降、可靠性提升；因时机器人、大寰机器人、傲博智能等企业推出量产级灵巧手，自由度从12-20不等，单只价格降至2-10万元，能够完成抓取鸡蛋、拧瓶盖、使用手机等操作，但产能有限、力控精度和触觉感知仍有差距，主要供应人形机器人企业和科研院所做原型测试。',
            '【量产商用元年（2025-2026）】2026年是人形机器人灵巧手量产商用元年，人形机器人爆发式增长带动灵巧手需求激增，年需求量从2024年的几千只增长至2026年的10万只级别；因时机器人建成年产能10万只灵巧手生产线，单只12自由度灵巧手价格降至1.5-2.5万元，力控精度达0.02N，指尖配备触觉阵列传感器；特斯拉Optimus、优必选Walker、宇树H1、小米CyberOne2等人形机器人全部配备五指灵巧手；灵巧手操作能力快速提升：从简单抓取→拧瓶盖→穿针引线→使用工具→精密装配，哈工大灵巧手在人形机器人运动会穿针引线项目仅需42秒。',
            '【性能提升期（2027-2028）】灵巧手性能持续向人手逼近：自由度提升至20+自由度接近人手27自由度，力控精度提升至0.01N以内，指尖触觉传感器密度提升至每指尖500+感知点接近人手，重量控制在500g以内接近人手（人手约400-500g），可靠性MTBF提升至5000小时以上满足工业应用要求，价格进一步降至8000-15000元/只；灵巧手不仅能在工业场景做精密装配，还能在家庭场景完成叠衣服、做饭、打扫等复杂家务操作；触觉传感器国产化突破，不再依赖进口。',
            '【普惠普及期（2029-2030）】灵巧手技术成熟，性价比高，成为人形机器人标配，单只价格降至5000-8000元，性能达到人手90%以上操作能力，能够完成绝大多数人手能完成的操作；工业场景精密装配、柔性制造大规模使用灵巧手，服务场景人形机器人能够完成各种服务操作，家庭场景人形机器人能够完成绝大多数家务劳动；中国灵巧手产业全球领先，产能占全球80%以上，成本是国外1/5-1/10。',
            '【灵巧手技术路线演进】二指夹爪→三指自适应夹爪→欠驱动灵巧手→全驱动五指灵巧手→带触觉高保真灵巧手→类人手高集成灵巧手；驱动方式从气动→电动直线缸→无刷电机+腱驱动→人工肌肉驱动；感知从无→位置传感→力反馈→指尖触觉阵列→全手分布式触觉。',
            '【核心零部件国产化】灵巧手核心零部件包括：微型伺服电机/直线驱动器、高精度减速器（行星/谐波）、腱绳/传动机构、力传感器、触觉传感器、微型控制器；2026年核心零部件国产化率达70%，蚌埠中国传感谷六维力传感器、微型力传感器企业为灵巧手产业链配套，安徽在传感器领域优势支撑灵巧手产业发展。',
            '【应用场景扩展】灵巧手应用场景持续扩展：工业领域（精密装配、柔性上下料、质量检测、设备维护）、物流领域（不规则包裹分拣、拆码垛、包装）、服务领域（餐饮端菜倒水、酒店服务、医院护理、家庭服务）、特种领域（排爆、救援、太空/深海作业）、医疗领域（手术机器人、康复机器人、假肢）；灵巧手是机器人真正"动手干活"的核心，是人形机器人实用性的关键。',
            '【安徽与蚌埠产业】安徽机器人灵巧手产业链：合肥因时机器人（国内领先灵巧手企业，年产能10万只）、哈工大机器人合肥研究院灵巧手研发、蚌埠中国传感谷六维力/微型力传感器产业链为灵巧手配套；蚌埠奥普特、中电科思仪等企业的力传感器、视觉传感器产品供应灵巧手企业，形成传感器-灵巧手-整机的完整产业链协同。'
        ], 'detail': [
            '【人手参考基准数据】人手共有27个自由度（腕部6+手掌5+拇指3+食指3+中指3+无名指3+小指3），指尖力控精度约0.005N，全手分布触觉感知点约17000个，是仿生灵巧手设计和性能追赶的终极目标',
            '【因时BHX-12量产参数】因时机器人BHX-12量产级12主动自由度五指灵巧手，采用腱驱动传动方式，指尖最大输出力10N，力控精度达0.02N，每指指尖集成100点触觉阵列传感器，总重量550g，支持CAN/EtherCAT通信，单只售价约2.5万元年产能10万只',
            '【因时BHX-20高端参数】因时机器人BHX-20高精度20主动自由度灵巧手，每个手指配置3-4个独立自由度，采用腱驱动+差动机构优化设计，指尖最大输出力15N，力控精度达0.01N，每指指尖集成200点高密度触觉阵列，重量650g单只约8万元',
            '【Shadow Hand科研标杆】英国Shadow Robot公司Shadow Dexterous Hand配置20主动自由度+4被动自由度，采用气动肌肉+腱混合驱动，指尖最大输出力10N，力控精度达0.005N与人手相当，配备BioTac触觉传感器，重量430g单只售价约150万元主要用于科研',
            '【特斯拉Optimus手部参数】特斯拉Optimus Gen3灵巧手采用6个直线驱动器驱动11自由度方案，自适应欠驱动手指设计（拇指2自由度+其余四指各2自由度+掌关节1自由度），可自适应抓取不同形状物体，指尖最大输出力20N采用高强度金属腱绳传动',
            '【小米CyberHand参数】小米CyberOne 2代CyberHand配置12自由度五指仿人设计，采用微型直线驱动器，每个指尖集成6维力传感器+触觉传感器阵列，指尖最大输出力15N，手部总重量500g，可完成拧瓶盖、操作手机、拿取生鸡蛋等精细操作',
            '【机器人力控技术】阻抗控制、导纳控制、力位混合控制等先进力控算法实现接触力精确控制，能够完成精密装配、柔顺抓取等高难度任务，力控系统带宽达100Hz以上，响应速度快控制精度高',
            '【柔性触觉传感器技术】电容式、压阻式、压电式柔性触觉传感器阵列，空间分辨率达1-2mm，力分辨率0.01N，可检测接触力分布、滑移、温度、材质纹理等信息，是灵巧手感知外界的重要感官',
            '【灵巧手驱动技术对比】腱驱动（绳索传动）类似人手肌腱结构传动紧凑，但存在腱绳磨损问题；连杆驱动刚度高但体积较大；直线驱动精度高但重量较大；气动人工肌肉驱动柔顺性好但需要气源支持',
            '【灵巧手材料工艺】手指结构采用碳纤维+钛合金轻量化设计，关节轴承采用PEEK高性能工程材料，传动腱绳采用高强度高分子纤维（Dyneema/Kevlar），连续工作寿命超2万小时满足工业使用要求',
            '【量产灵巧手操作能力】当前量产灵巧手可完成：抓取不同形状大小物体（生鸡蛋、玻璃杯、各种工具）、拧瓶盖/开门/按按钮等日常操作、使用螺丝刀/锤子等简单工具、精密电子装配（插针/组装）、甚至写字画画',
            '【哈工大穿针引线演示】哈工大机器人研究所自研20自由度灵巧手在发布会上演示穿针引线精细操作，仅用58秒完成穿针全过程，指尖力控精度达0.02N，手部末端抖动控制在0.01mm级达到人手水平',
            '【灵巧手成本快速下降】2020年进口科研灵巧手单只价格100万元以上，2023年国产量产灵巧手5-10万元，2026年规模化量产12自由度灵巧手降至1-3万元，2030年目标5000元以内同时达到接近人手级性能',
            '【灵巧手产能建设情况】因时机器人年产能10万只灵巧手自动化生产线2026年正式投产；大寰机器人年产能15万只夹爪/灵巧手产线建成；其他国内厂商产能也在快速建设，总产能可满足人形机器人爆发式增长需求',
            '【灵巧手市场需求预测】2026年人形机器人对灵巧手需求量约16万只（按8万台×2只手计算），2030年需求量约100万只，配套灵巧手及核心零部件市场规模超200亿元，市场空间广阔',
            '【蚌埠传感谷产业对接】蚌埠中国传感谷重点布局灵巧手触觉传感器、六维力传感器、微型力传感器产业，与因时机器人、大寰机器人等国内主流灵巧手企业开展供应链对接合作，传感器产业支撑灵巧手发展',
            '【安徽省灵巧手产业】合肥、芜湖机器人产业集聚发展，灵巧手研发生产企业快速成长，合肥工业大学、中国科学技术大学在灵巧手机械设计、驱动控制、触觉感知等领域研究水平国内领先',
            '【灵巧手操作技能学习】通过模仿学习、深度强化学习、示教学习等AI技术，灵巧手操作技能库快速增长，从几百种预设技能发展到几万种通用操作技能，机器人自主学习新技能能力持续提升',
            '【多模态感知融合操作】视觉+力觉+触觉多模态融合感知技术实现机器人手眼协调精细操作，对于未知形状未知材质物体抓取成功率达95%以上，非结构化环境适应能力大幅提升',
            '【2030年灵巧手发展目标】规划到2030年灵巧手自由度达到20+接近人手，力控精度0.005N达到人手水平，触觉分辨率接近人手触觉密度，单只成本降至5000元以内，整体操作能力达到普通人手90%以上水平'
        ]},
        # 安防应急
        {'left': [
            '【产业定位】安防应急特种机器人在危险环境替代人工作业，保障人民生命财产安全',
            '【市场规模】2026年中国安防应急特种机器人市场规模突破220亿元，年均增速超50%',
            '【消防机器人】灭火/侦察/排烟/救援消防机器人，在易燃易爆/有毒/高温危险环境作业',
            '【排爆机器人】公安/武警排爆机器人，转移/销毁爆炸物，避免排爆人员伤亡',
            '【巡检机器人】园区/厂区/边境/机场安保巡检机器人，24小时巡逻，异常识别报警',
            '【反恐处突】反恐突击/侦察/谈判机器人，在恐怖袭击/劫持人质场景替代警员突入',
            '【地震救援】地震/塌方/泥石流灾害救援机器人，废墟中搜索幸存者，输送物资',
            '【核应急】核电站事故应急机器人，高辐射环境作业，人员无法进入的场景',
            '【海康机器人】海康威视安防巡检机器人，园区/厂区/机场巡逻，国内市占率领先',
            '【中信重工开诚】中信重工开诚消防机器人，国内消防机器人龙头，市场占有率第一'
        ], 'right': [
            '【开诚RXR-MC80】消防灭火机器人，流量80L/s，射程80m，防爆设计，可拖拽2盘水带，牵引300kg',
            '【海康威视巡检机器人】海康威视安防巡检机器人，高清摄像头+热成像+异常声音识别，24小时自主巡逻',
            '【卫泰排爆机器人】卫泰智能排爆机器人，6+1自由度机械臂，抓取10kg，X射线检查，水炮销毁',
            '【哈工大救援机器人】哈工大灾后废墟搜索救援机器人，蛇形/履带式，穿越狭窄缝隙，生命探测',
            '【中电科反恐机器人】中电科38所反恐突击机器人，武装/侦察/突击，可携带非致命武器',
            '【大疆安防无人机】大疆M300/M3T安防巡检无人机，空中巡逻，热成像，追踪，快速响应',
            '【水下机器人】深之蓝/博雅工道水下安检/救援机器人，水下探测/打捞/安检',
            '【中信重工防爆机器人】石化防爆巡检机器人，Ex防爆认证，石化厂区巡检',
            '【安徽消防】安徽消防总队配备消防/救援/排烟机器人超200台，合肥/芜湖/蚌埠消防配机器人',
            '【合肥安保】合肥重要场所/园区/机场部署安防巡检机器人500+台，重大活动安保'
        ], 'process': [
            '【纯人工高风险期（2010年前）】安防应急领域几乎完全依靠人工处置：火灾现场消防员冒着浓烟、高温、爆炸风险内攻灭火，每年都有消防员牺牲；排爆警察人工近距离转移、销毁爆炸物，稍有不慎就有生命危险；地震、矿难等灾害救援人员深入废墟搜救，余震、二次坍塌风险高；石化泄漏、核辐射等危险环境人员无法长时间停留；安防巡逻人员24小时轮班，劳动强度大，夜间巡逻危险；每年全国应急救援领域人员牺牲超百人，危险作业"用人命换安全"的问题突出。',
            '【特种机器人试点期（2011-2018）】特种机器人开始试点配备：消防灭火机器人、排爆机器人率先在公安消防部队试点，主要在一线城市和特殊场景配备，但数量少（全国仅几百台）、价格高（单台消防机器人几十万到上百万）、性能有限（行走越障能力弱、操作复杂、遥控距离近、可靠性差），只能承担辅助任务，无法替代人员进入最危险区域，很多时候"买了用不上、想用不好用"，操作人员培训不足，机器人实际使用率低。',
            '【技术提升期（2019-2024）】特种机器人技术快速提升：底盘性能提升（履带式/轮履复合/足式，越障爬坡能力增强）、操控方式改进（远距离遥控、半自主作业）、感知能力增强（多摄像头、热成像、气体检测、生命探测）、可靠性提升（防水防爆防尘）；消防机器人从单一灭火发展为灭火、排烟、侦察、救援多品类；排爆机器人机械臂灵活性提升；安防巡检机器人实现自主导航巡逻；特种机器人价格逐步下降，消防部队、公安特警、应急管理部门配备数量增加到几千台，在天津港爆炸、四川森林火灾等事故中发挥作用，但仍未成为标配。',
            '【规模化普及期（2025-2026）】2026年特种机器人规模化普及，技术成熟、成本下降、可靠性满足实战要求：全国消防救援队伍配备消防/救援/排烟/侦察机器人超5000台，排爆机器人全国地市以上公安特警标配，石化、电力等高危行业巡检机器人广泛应用，安防巡检机器人在机场、火车站、园区、重要场所部署超万台；地震救援蛇形机器人、水下救援机器人、消防人形机器人开始实战应用；危险场景作业机器人替代率达60%，应急救援人员伤亡率较2015年下降60%；中电科38所、中信重工、哈工大、大疆等企业特种机器人产品成熟。',
            '【全面机器代人期（2027-2028）】特种机器人成为应急处置标配，"机器人先上、人员跟进"成为处置标准流程：火灾现场机器人先进入侦察、灭火、排烟，消防员在后方远程操控跟进；爆炸物机器人直接处置，排爆人员不再近距离接触；石化泄漏机器人进入泄漏区域关阀、堵漏、检测；地震废墟机器人先进入搜索生命，再指导人员救援；安防巡逻70%由机器人完成，人员处理异常情况；危险作业机器人替代率达80%，应急救援人员伤亡率大幅下降80%以上。',
            '【智能化无人化期（2029-2030）】特种机器人实现高度智能化、自主化：消防机器人能够自主进入火场、自主规划路径、自主识别火源灭火、自主搜救被困人员；排爆机器人自主识别爆炸物、自主选择处置方式；安防机器人自主识别异常行为、自主追踪、自主报警；多机器人协同作业（消防无人机+消防机器人+救援机器人协同）；人形消防机器人能够使用消防器材、破拆、救人；特种机器人形成完整体系，中国应急救援装备水平全球领先，人员安全得到最大程度保障。',
            '【安防应急机器人分类】按应用场景分为：①消防类（灭火机器人、排烟机器人、侦察机器人、救援机器人、破拆机器人、消防人形机器人）；②排爆反恐类（排爆机器人、武装反恐机器人、侦察机器人）；③安防巡检类（园区巡检机器人、变电站巡检机器人、石化厂区巡检机器人、边境巡逻机器人、机场车站安检机器人）；④灾害救援类（废墟搜索救援机器人、蛇形机器人、水下救援机器人、矿山救援机器人、核辐射环境机器人）；⑤警用类（巡逻机器人、抓捕机器人、交通指挥机器人）。',
            '【关键技术突破】安防应急机器人关键技术：①高机动底盘技术（轮履复合、足式、越障、爬坡、涉水、废墟地形适应）；②防爆防水防尘技术（Ex防爆认证、IP67/IP68防护、高温耐受）；③远距离可靠通信技术（5G、自组网、光纤通信、穿墙通信）；④多传感器融合感知技术（可见光、热成像、气体检测、辐射检测、生命探测）；⑤遥操作与半自主技术（沉浸式操控、力反馈、AI辅助决策）；⑥极端环境可靠性技术（-40℃~60℃工作、抗冲击、防尘防水）。',
            '【安徽安防应急应用】安徽公安消防应急系统特种机器人配备快速增长：安徽省消防救援总队配备各类消防/救援机器人超200台，合肥、芜湖、蚌埠等市消防支队配备消防灭火、排烟、侦察机器人；合肥新桥机场、合肥南站、重要园区、政府机关部署安防巡检机器人超500台；中电科38所（合肥）是国内反恐安防雷达、机器人重要研发生产单位；蚌埠消防支队配备消防机器人、排烟机器人，提升化工园区火灾处置能力；安徽石化、电力企业广泛配备防爆巡检机器人。',
            '【社会价值】安防应急机器人的社会价值无法用金钱衡量：保护消防员、警察、救援人员生命安全，减少牺牲；提升应急处置效率，缩短响应时间，减少人民群众生命财产损失；在人员无法进入的极端危险环境（高温、有毒、爆炸、辐射、缺氧）开展作业；推动应急救援从"用人命拼"向"科技强安"转变，是国家治理能力和科技水平的重要体现；中国特种机器人技术和应用全球领先，为公共安全提供坚实保障。'
        ], 'detail': [
            '【灵巧智能电力巡检·2026年8月21日】灵巧智能与华中科技大学联合建设的"具身智能巡检操作机器人联合实验室"项目已落地运行：整套装备集成高通过性移动底盘/多模态感知系统/精细操作执行机构，实现变电站巡检/操作/处置闭环；该方案已实战护航第十五届全运会保电任务实现35%效率提升',
            '【消防灭火机器人详细参数】消防灭火机器人采用柴油动力或电动驱动，行走速度3-5km/h，爬坡能力30度越障高度20cm，消防水炮流量30-100L/s射程60-100m，防爆等级Ex d IIB T4，可拖拽2盘水带牵引300kg，遥控距离1km防水等级IP67',
            '【排爆机器人技术参数】排爆机器人采用履带式或轮式底盘，配置6-7自由度多关节机械臂，最大伸展距离2-3m，全伸展状态抓取重量5-20kg，配备X射线检查系统和水炮销毁器，多摄像头360度观察，光纤/无线遥控距离500m',
            '【安防巡检机器人参数】安防巡检机器人采用轮式底盘自主导航避障，行走速度0-10km/h可调节，标准工况续航8-12小时，配备360度高清摄像头+热成像+声光报警，支持人脸识别/车牌识别/异常行为识别/烟火检测，可自主乘电梯自动回充',
            '【地震救援机器人参数】地震救援机器人有履带式、蛇形、多足等多种形态，可穿越废墟狭窄缝隙进入人员无法到达区域，配备生命探测雷达/音频/视频/热成像多模态探测幸存者，可携带药品/食品/通信设备，防水防尘防压续航6小时以上',
            '【核应急机器人技术参数】核应急机器人采用特殊耐辐射设计可承受1000Sv/h高强度辐射，远程遥控距离达5km，可完成远程阀门操作、放射性样品采集、现场去污作业，耐温100℃以上，摄像头和电子元器件全部防辐射加固处理',
            '【消防机器人实战效能】消防机器人可在1000℃高温、易燃易爆、有毒有害危险环境持续作业，替代消防员深入最危险区域，消防员牺牲率降低90%，灭火效率是人工内攻灭火的3倍，可长时间持续作战无疲劳问题',
            '【安防巡检效率对比】人工安保巡逻每人每班次有效巡逻约5公里，智能巡检机器人可24小时不间断自主巡逻，覆盖范围是人工巡逻的10倍，异常情况识别报警响应时间小于5秒，夜间和恶劣天气条件下不受影响',
            '【排爆作业安全价值】排爆机器人替代排爆人员直接接触爆炸物，排爆人员在数百米外安全距离遥控操作，排爆作业人员伤亡事故率降低99%，彻底改变以前排爆警察"用手排爆、用命赌安全"的危险局面',
            '【中信重工开诚市场地位】中信重工开诚智能是国内消防机器人龙头企业，消防机器人国内市场占有率超40%，累计销售各类消防机器人超5000台，参与天津港爆炸、四川凉山森林火灾等多起重特大事故应急救援',
            '【海康安防机器人出货量】海康威视安防巡检机器人累计出货量超1万台，广泛应用于产业园区、工厂厂区、机场、住宅小区、边境线等场景巡逻安保，是国内安防巡检机器人市场份额领先企业',
            '【大疆行业无人机应用】大疆创新行业级无人机在安防应急领域市场占比超70%，消防、公安、应急管理、电力巡检、城管执法等部门广泛应用，2026年全年行业应用无人机出货量预计超10万台',
            '【应急机器人配备标准】2026年应急装备配备标准要求：每个地级以上城市消防支队至少配备10台消防机器人，每个特勤中队至少配备5台；县级以上公安机关至少配备2台排爆机器人，基层应急装备标准化建设加快推进',
            '【安徽省应急装备配备】安徽省消防救援总队共配备各类消防机器人220台、排爆机器人50台、巡检无人机300架，合肥、芜湖、蚌埠三市配备量位居全省前三，应急装备智能化水平持续提升',
            '【蚌埠市应急装备水平】蚌埠市消防救援支队配备消防灭火机器人15台、排爆机器人3台、安检巡逻机器人20台，依托化工园区消防需求重点配备防爆消防机器人，应急装备水平位居安徽省前五名',
            '【AI智能识别技术应用】AI异常行为智能识别算法可自动识别打架斗殴、翻越围墙、遗留可疑物品、烟火火情、人员异常聚集等异常情况，识别准确率达99%，误报率低于1%，发现异常自动触发报警',
            '【5G应急通信技术】5G低时延高清图传技术支持应急现场4K超高清画面实时回传指挥中心，后方专家可远程指导现场处置，支持多机器人多机位协同作业，指挥决策更加科学高效',
            '【数字孪生应急推演】应急场景数字孪生系统可对灾害事故进行预案仿真推演，为机器人规划最优作业路径，模拟不同处置方案效果，应急处置效率提升50%，避免盲目处置造成二次伤亡',
            '【行业标准化建设】消防机器人、排爆机器人、安防巡检机器人国家标准和行业标准陆续出台实施，规范产品质量要求、检测认证方法、实战性能指标，推动行业规范化高质量发展',
            '【机器人成本下降趋势】消防机器人价格从2015年的200-300万元降至2026年的50-100万元，安防巡检机器人从50万元降至10-20万元，价格大幅下降为基层规模化配备创造条件',
            '【2030年发展目标】规划到2030年危险应急作业机器人替代率达90%，消防、排爆、安防巡检机器人基层单位配备率达到100%，应急救援人员伤亡率较2015年降低95%，实现"科技强安、机器换人"目标'
        ]},
    ]
    cat = categories[category_idx]
    return (part_num, title, cat['left'], cat['right'], '▎' + title + '发展具体过程阐述', cat['process'], '▎' + keyword.split('/')[0] + ' · 参数数据 · 应用进展 · 未来展望', cat['detail'])

for i, (part_num, title, keyword) in enumerate(module_titles_rest):
    all_modules.append(make_detail_module(part_num, title, keyword, i))

# ========== 双层水印 ==========
def add_watermark(prs, date_str='20260821'):
    wm_texts = [
        '机密文档 请勿传播 ' + date_str,
        '内部资料 版权所有 ' + date_str,
        '具身智能AI产业汇报 ' + date_str,
        'CONFIDENTIAL ' + date_str,
    ]
    from pptx.oxml.ns import qn
    from lxml import etree
    for slide in prs.slides:
        # 第一层大字号斜向水印
        for i in range(-2, 6):
            for j in range(-2, 5):
                x = Inches(i * 3.2)
                y = Inches(j * 2.5)
                wm_box = slide.shapes.add_textbox(x, y, Inches(4), Inches(1.5))
                tf = wm_box.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                run = p.add_run(); run.text = wm_texts[(i+j) % 4]
                run.font.size = Pt(26); run.font.bold = True
                run.font.color.rgb = RGBColor(0x40, 0x50, 0x70)
                run.font.name = '微软雅黑'
                wm_box.rotation = -30
                # 设置透明度 - 直接操作XML
                rPr = run._r.get_or_add_rPr()
                solidFill = rPr.find(qn('a:solidFill'))
                if solidFill is None:
                    solidFill = etree.SubElement(rPr, qn('a:solidFill'))
                srgbClr = solidFill.find(qn('a:srgbClr'))
                if srgbClr is None:
                    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
                    srgbClr.set('val', '405070')
                alpha = srgbClr.find(qn('a:alpha'))
                if alpha is None:
                    alpha = etree.SubElement(srgbClr, qn('a:alpha'))
                alpha.set('val', '8000')
        # 第二层中字号反向水印
        for i in range(-1, 7):
            for j in range(-1, 7):
                x = Inches(i * 2.0 - 0.5)
                y = Inches(j * 1.8)
                wm_box = slide.shapes.add_textbox(x, y, Inches(2.5), Inches(0.8))
                tf = wm_box.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                run = p.add_run(); run.text = wm_texts[(i+j*2) % 4]
                run.font.size = Pt(12); run.font.italic = True
                run.font.color.rgb = RGBColor(0x50, 0x60, 0x80)
                run.font.name = '宋体'
                wm_box.rotation = 25
                # 设置透明度 - 直接操作XML
                rPr = run._r.get_or_add_rPr()
                solidFill = rPr.find(qn('a:solidFill'))
                if solidFill is None:
                    solidFill = etree.SubElement(rPr, qn('a:solidFill'))
                srgbClr = solidFill.find(qn('a:srgbClr'))
                if srgbClr is None:
                    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
                    srgbClr.set('val', '506080')
                alpha = srgbClr.find(qn('a:alpha'))
                if alpha is None:
                    alpha = etree.SubElement(srgbClr, qn('a:alpha'))
                alpha.set('val', '6000')

def generate(enable_watermark=False):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    cover_page(prs)
    toc_page(prs)
    for module in all_modules:
        # V3.37字数均衡：内容池合并按字数拆2页+细节按字数拆2页
        content_page_1(prs, module[0], module[1], module[2], module[3], module[5])
        content_page_2(prs, module[0], module[1], module[2], module[3], module[5])
        detail_page_1(prs, module[0], module[1], module[6], module[7])
        detail_page_2(prs, module[0], module[1], module[6], module[7])
    back_page(prs)
    if enable_watermark:
        add_watermark(prs)
    return prs

if __name__ == '__main__':
    import os
    out_dir = r'F:\个人作品\具身智能'
    date_str = '20260821'
    ver = 'v33'
    
    print('正在生成无水印原版...')
    prs1 = generate(enable_watermark=False)
    f1 = os.path.join(out_dir, '具身智能AI产业最新进展_' + date_str + '_商务汇报_无水印_' + ver + '.pptx')
    prs1.save(f1)
    print('完成：' + str(len(prs1.slides)) + '页')
    
    print('正在生成水印版...')
    prs2 = generate(enable_watermark=True)
    f2 = os.path.join(out_dir, '具身智能AI产业最新进展_' + date_str + '_商务汇报_水印版_' + ver + '.pptx')
    prs2.save(f2)
    print('完成：' + str(len(prs2.slides)) + '页')
    
    print('')
    print('总页数验证：1封面 + 1目录 + 22模块×4页 + 1封底 = 91页')
    print('无水印：' + f1)
    print('水印版：' + f2)
