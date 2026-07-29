
# ============================================================================
# 免责声明与AI使用规范
# ============================================================================
# 本文件仅供技术研究与学习交流使用，不得用于任何非法用途。
#
# AI使用规范：
#   1. 使用本文件相关内容时须遵守所在地法律法规及伦理准则
#   2. 不得用于侵犯他人合法权益、危害网络安全、破坏公共秩序的活动
#   3. 涉及自动化决策的场景须确保人工复核机制与可解释性
#   4. 处理个人信息时须符合数据保护相关法规要求
#
# 风险提示：
#   本文件内容按"现状"提供，不保证绝对准确无误。
#   使用者须自行评估风险，因使用本文件导致的任何损失由使用者承担。
# ============================================================================

from pptx import Presentation
import os

def analyze_ppt_colors(ppt_path):
    prs = Presentation(ppt_path)
    
    print(f"分析PPT文件: {ppt_path}\n")
    print("="*80)
    
    for i, slide in enumerate(prs.slides):
        print(f"\n【第 {i+1} 张幻灯片】")
        print("-"*60)
        
        # 分析背景
        try:
            if hasattr(slide.background, 'fill'):
                if slide.background.fill.type == 1:  # 纯色背景
                    bg_rgb = slide.background.fill.fore_color.rgb
                    print(f"背景色: {bg_rgb}")
                else:
                    print("背景: 渐变或图片")
        except:
            print("背景: 无法读取")
        
        # 分析表格
        table_count = 0
        for shape in slide.shapes:
            if shape.has_table:
                table_count += 1
                print(f"\n  表格 {table_count}:")
                table = shape.table
                
                # 分析表头
                if len(table.rows) > 0:
                    header_row = table.rows[0]
                    print(f"    行数: {len(table.rows)}, 列数: {len(table.columns)}")
                    
                    for j, cell in enumerate(header_row.cells[:3]):  # 只看前3列
                        # 单元格背景
                        try:
                            if cell.fill.type == 1:
                                cell_rgb = cell.fill.fore_color.rgb
                                print(f"    列 {j} 背景: {cell_rgb}")
                        except:
                            pass
                        
                        # 文字颜色
                        try:
                            if len(cell.text_frame.paragraphs) > 0:
                                para = cell.text_frame.paragraphs[0]
                                if len(para.runs) > 0:
                                    run = para.runs[0]
                                    if hasattr(run.font.color, 'rgb'):
                                        text_rgb = run.font.color.rgb
                                        print(f"    列 {j} 文字: {text_rgb}, 字号: {run.font.size}")
                        except:
                            pass
                break
    
    print("\n" + "="*80)
    print("分析完成！")

if __name__ == "__main__":
    ppt_path = EXTERNAL_PROJECT_DIR  # was: A2A_PROTOCOL_AI_AGENT_V14_蓝色字体.pptx
    if os.path.exists(ppt_path):
        analyze_ppt_colors(ppt_path)
    else:
        print(f"文件不存在: {ppt_path}")