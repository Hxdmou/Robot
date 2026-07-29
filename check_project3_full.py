
# ============================================================================
# 免责声明与AI使用规范
# ============================================================================
# 本文件仅供技术研究与学习交流使用，不得用于任何非法用途。
#
# AI使用规范：
#   1. 使用本文件相关内容时须遵守所在地法律法规及伦理准则
#   2. 不得用于侵犯他人合法权益、危害网络安全、破坏公共秩序的活动
#   3. 涉及自动化决策的场景须确保人工复核机制与可解释性
#   4. 处理个人信息时须符合数据保护相关法规要求
#
# 风险提示：
#   本文件内容按"现状"提供，不保证绝对准确无误。
#   使用者须自行评估风险，因使用本文件导致的任何损失由使用者承担。
# ============================================================================

from pptx import Presentation

def check_project3_display():
    ppt_path = EXTERNAL_PROJECT_DIR  # was: A2A_PROTOCOL_AI_AGENT_V14_终极完整版.pptx
    prs = Presentation(ppt_path)
    
    print("=" * 120)
    print("📌 项目3完整内容检查")
    print("=" * 120)
    
    if len(prs.slides) > 2:
        slide = prs.slides[2]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                print(f"\n表格尺寸: {len(table.rows)}行 × {len(table.columns)}列\n")
                
                # 显示每一行的完整内容
                for i in range(len(table.rows)):
                    row_data = []
                    for j in range(len(table.columns)):
                        cell_text = table.cell(i, j).text.strip()
                        row_data.append(cell_text)
                        if not cell_text:
                            print(f"⚠️ 警告: 第{i+1}行第{j+1}列为空！")
                    
                    # 显示整行内容
                    print(f"\n第{i+1}行:")
                    for j, col_name in enumerate(row_data):
                        print(f"  列{j+1}: {col_name}")
                
                # 检查是否有L10
                has_l10 = False
                for i in range(len(table.rows)):
                    if "L10" in table.cell(i, 0).text:
                        has_l10 = True
                        print(f"\n✅ 找到L10: {table.cell(i, 1).text}")
                        break
                
                if not has_l10:
                    print("\n❌ 未找到L10边缘层！")
                
                break

if __name__ == "__main__":
    check_project3_display()