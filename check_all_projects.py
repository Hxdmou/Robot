
# ============================================================================
# 免责声明与AI使用规范
# ============================================================================
# 本文件仅供技术研究与学习交流使用，不得用于任何非法用途。
#
# AI使用规范：
#   1. 使用本文件相关内容时须遵守所在地法律法规及伦理准则
#   2. 不得用于侵犯他人合法权益、危害网络安全、破坏公共秩序的活动
#   3. 涉及自动化决策的场景须确保人工复核机制与可解释性
#   4. 处理个人信息时须符合数据保护相关法规要求
#
# 风险提示：
#   本文件内容按"现状"提供，不保证绝对准确无误。
#   使用者须自行评估风险，因使用本文件导致的任何损失由使用者承担。
# ============================================================================

from pptx import Presentation

def check_all_projects():
    print("=" * 120)
    print("🔍 检查 docs/presentations/A2A_ENTERPRISE_PPT_V12.pptx 所有项目")
    print("=" * 120)
    
    ppt_path = EXTERNAL_PROJECT_DIR  # was: A2A_ENTERPRISE_PPT_V12.pptx
    prs = Presentation(ppt_path)
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n📄 幻灯片 {slide_idx + 1}:")
        
        # 获取幻灯片标题
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.text.startswith("项目"):
                    title = shape.text[:20]
                    break
        
        # 查找表格
        table_found = False
        for shape in slide.shapes:
            if shape.has_table:
                table_found = True
                table = shape.table
                print(f"   标题: {title}")
                print(f"   表格: {len(table.rows)}行 × {len(table.columns)}列")
                
                # 打印第一列内容
                first_col = []
                for i in range(min(5, len(table.rows))):
                    cell_text = table.cell(i, 0).text.strip()
                    first_col.append(cell_text)
                print(f"   第一列: {first_col}...")
                break
        
        if not table_found:
            print(f"   标题: {title}")
            print("   └─ 无表格")

if __name__ == "__main__":
    check_all_projects()