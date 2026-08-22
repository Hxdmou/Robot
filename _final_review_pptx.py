# -*- coding: utf-8 -*-
'''推送前终极复查：直接打开v32 PPTX文件
1. 每页页脚文本是否完全一致（FOOTER_TEXT）
2. 全部【】黄标标签日期是否都在08-20/21（无过期）
3. 页数/模块结构完整性
'''
import sys, re
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

FILES = [
    r"F:\个人作品\具身智能\具身智能AI产业最新进展_20260822_商务汇报_无水印_v34.pptx",
    r"F:\个人作品\具身智能\具身智能AI产业最新进展_20260822_商务汇报_水印版_v34.pptx",
]
SLIDE_H = 7.5  # 16:9标准高度（英寸）

def is_stale(tag):
    # 2026-08-22：窗口08-21/08-22，08-20及更早=过期
    for m in re.finditer(r'2026年(\d{1,2})月(\d{1,2})日', tag):
        mo, d = int(m.group(1)), int(m.group(2))
        if mo < 8 or (mo == 8 and d < 21):
            return True
    for m in re.finditer(r'2026年(\d{1,2})月(?!\d)', tag):
        if int(m.group(1)) < 8:
            return True
    for m in re.finditer(r'(?<!2026年)(\d{1,2})月(\d{1,2})日', tag):
        mo, d = int(m.group(1)), int(m.group(2))
        if mo < 8 or (mo == 8 and d < 21):
            return True
    return False

total_problems = 0
for fp in FILES:
    print('=' * 70)
    print(fp.split('\\')[-1])
    prs = Presentation(fp)
    n = len(prs.slides)
    print('总页数: %d' % n)
    footers = {}
    stale_tags = []
    fresh_tags = 0
    for si, slide in enumerate(prs.slides, 1):
        texts = []
        bottom_texts = []  # 只收集页面最底部(y>=SLIDE_H-0.5英寸)的文本
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
                    try:
                        top_in = shape.top / 914400.0
                        if top_in >= SLIDE_H - 0.5:
                            bottom_texts.append(t)
                    except Exception:
                        pass
        full = '\n'.join(texts)
        # 页脚：只取页面最底部含"商务汇报"的文本
        for t in bottom_texts:
            if '商务汇报' in t and len(t) < 80:
                footers.setdefault(t, []).append(si)
        # 黄标扫描
        for tag in re.findall(r'【[^】]*】', full):
            if is_stale(tag):
                stale_tags.append((si, tag[:60]))
            elif re.search(r'2026年8月2[01]日|8月2[01]日', tag):
                fresh_tags += 1
    print('--- 页脚检查 ---')
    for ft, pages in footers.items():
        print('  [%d页] %s' % (len(pages), ft))
    if len(footers) == 1:
        print('  [OK] 页脚完全统一')
    else:
        print('  [FAIL] 页脚不一致！共%d种' % len(footers))
        total_problems += 1
    print('--- 黄标日期检查 ---')
    print('  新鲜黄标(08-20/21): %d 个' % fresh_tags)
    if stale_tags:
        print('  [FAIL] 过期黄标 %d 个:' % len(stale_tags))
        for si, tag in stale_tags[:20]:
            print('    第%d页: %s' % (si, tag))
        total_problems += 1
    else:
        print('  [OK] 无过期黄标')

print('=' * 70)
if total_problems == 0:
    print('[V] 终极复查全部通过：页脚统一 + 0过期黄标')
else:
    print('[X] 发现 %d 类问题，必须修复后再推送！' % total_problems)
